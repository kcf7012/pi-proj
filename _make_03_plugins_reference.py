"""
簡報 2/4：Plugin 技術參考 (03-plugins-reference.pptx)
約 45 張
對應：03-plugins-reference.md
完整技術規格：元件、manifest、CLI、開發工具
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 45

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "Plugin 技術參考（完整規格）",
        "所有元件、manifest 欄位、CLI 指令、開發工具",
        tag="#03 · 技術參考"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "完整技術規格的權威參考",
        slide_num=2, total=TOTAL, source="03-plugins-reference.md"
    )

    objectives = [
        ("🧩", "7 種 Plugin 元件", "Skills / Agents / Hooks / MCP / LSP / Monitors / Themes"),
        ("📦", "Manifest 完整架構", "plugin.json 所有欄位詳解"),
        ("🌍", "環境變數", "CLAUDE_PLUGIN_ROOT / DATA / PROJECT_DIR"),
        ("💾", "快取與檔案解析", "Plugin 怎麼被複製與隔離"),
        ("📂", "完整目錄結構", "標準 plugin 應該長什麼樣"),
        ("🛠️", "CLI 完整指令", "init / install / uninstall / prune / list / tag"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(2.2)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_cols = 3
    grid_rows = 2
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.8)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_PRIMARY
        card.line.width = Pt(2)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.2), box_w, Inches(0.7),
            font_size=40, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.2), y + Inches(1.0), box_w - Inches(0.4), Inches(0.5),
            font_size=15, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.5), box_w - Inches(0.4), Inches(0.6),
            font_size=11, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Plugin 元件總覽", "7 種元件一次看懂")

    # ============================================================
    # Slide 4：7 種元件速覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin 的 7 種元件",
        "Skills / Agents / Hooks / MCP / LSP / Monitors / Themes",
        slide_num=4, total=TOTAL, source="03 § Plugin 元件參考"
    )

    components = [
        ("📚", "Skills", "skills/", "新增可呼叫的 skill，建立 `/name` 快捷方式"),
        ("🤖", "Agents", "agents/", "提供專門的 subagent"),
        ("⚡", "Hooks", "hooks/hooks.json", "事件處理程式，自動回應 Claude Code 事件"),
        ("🔌", "MCP Servers", ".mcp.json", "連接外部工具和服務"),
        ("🧠", "LSP Servers", ".lsp.json", "Language Server Protocol，即時程式碼智慧"),
        ("📡", "Monitors", "monitors/monitors.json", "背景 monitor，持續整個工作階段"),
        ("🎨", "Themes", "themes/", "顏色主題（實驗性）"),
    ]

    box_w = Inches(6.1)
    box_h = Inches(0.85)
    box_gap_x = Inches(0.13)
    box_gap_y = Inches(0.1)
    grid_start_x = Inches(0.5)
    grid_start_y = Inches(1.7)

    for i, (icon, name, path, desc) in enumerate(components):
        row = i // 2
        col = i % 2
        x = grid_start_x + col * (box_w + box_gap_x)
        y = grid_start_y + row * (box_h + box_gap_y)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(1.5)

        # icon
        h.add_text_block(
            slide, icon,
            x + Inches(0.2), y + Inches(0.1), Inches(0.7), Inches(0.65),
            font_size=24
        )
        # 名稱
        h.add_text_block(
            slide, name,
            x + Inches(0.9), y + Inches(0.1), Inches(1.7), Inches(0.3),
            font_size=14, bold=True
        )
        # 路徑
        h.add_text_block(
            slide, path,
            x + Inches(0.9), y + Inches(0.45), Inches(1.7), Inches(0.3),
            font_size=10, color=h.COLOR_GRAY_TXT, italic=True
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(2.7), y + Inches(0.2), box_w - Inches(2.9), Inches(0.5),
            font_size=11
        )

    h.add_callout(
        slide, "💡 布林 frontmatter 欄位（disable-model-invocation 等）接受 yes/no/on/off/1/0（v2.1.218+）",
        Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 5：Skills 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Skills 詳解",
        "Plugin 根目錄的 skills/ 目錄",
        slide_num=5, total=TOTAL, source="03 § Skills"
    )

    h.add_code_block(
        slide, """skills/
├── pdf-processor/
│   ├── SKILL.md
│   ├── reference.md (optional)
│   └── scripts/ (optional)
└── code-reviewer/
    └── SKILL.md""",
        Inches(0.7), Inches(1.7), Inches(6.0), Inches(2.2),
        font_size=13
    )

    h.add_bullet_list(
        slide, [
            "Skills 是包含 `SKILL.md` 的目錄；commands 是簡單的 markdown 檔案",
            "Skills 和 commands 在安裝 plugin 時**自動發現**",
            "Plugin 根目錄中的單一 `SKILL.md` 會被載入為**單一 skill**",
            "設定 frontmatter `name` 欄位以控制 skill 的叫用名稱",
        ],
        Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.5),
        font_size=12
    )

    h.add_text_block(
        slide, "📋 Plugin 安裝範圍",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["範圍", "設定檔", "使用案例"],
        [
            ["`user`", "`~/.claude/settings.json`", "個人 plugins 跨所有專案（預設）"],
            ["`project`", "`.claude/settings.json`", "團隊 plugins 透過版本控制共享"],
            ["`local`", "`.claude/settings.local.json`", "專案特定 plugins，gitignored"],
            ["`managed`", "Managed settings", "受管的 plugins（唯讀）"],
        ],
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(2.0),
        font_size=11
    )

    # ============================================================
    # Slide 6：Agents 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Agents 詳解",
        "agents/ 目錄中的 subagent 定義",
        slide_num=6, total=TOTAL, source="03 § Agents"
    )

    h.add_code_block(
        slide, """---
name: agent-name
description: 此 agent 的專長以及 Claude
  應何時叫用它
model: sonnet
effort: medium
maxTurns: 20
disallowedTools: Write, Edit
---

詳細的系統提示，描述 agent 的角色、
專業知識和行為。""",
        Inches(0.7), Inches(1.7), Inches(7.0), Inches(3.0),
        font_size=11
    )

    h.add_bullet_list(
        slide, [
            "支援的 frontmatter：`name`、`description`、`model`、`effort`、`maxTurns`",
            "`tools`、`disallowedTools`、`skills`、`memory`、`background`、`isolation`",
            "唯一有效的 `isolation` 值是 `\"worktree\"`",
            "Plugin agents 在 @-mention 下拉中以**範圍名稱**出現（`my-plugin:code-reviewer`）",
        ],
        Inches(8.0), Inches(1.7), Inches(4.8), Inches(3.0),
        font_size=12
    )

    h.add_callout(
        slide, "⚠️ 安全考量：plugin agents **不支援** `hooks`、`mcpServers` 和 `permissionMode`",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.5),
        icon="", font_size=13
    )

    h.add_text_block(
        slide, "完整詳細資訊見 05-subagents.md",
        Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.4),
        font_size=14, color=h.COLOR_GRAY_TXT, italic=True
    )

    # ============================================================
    # Slide 7：Hooks 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hooks 詳解",
        "hooks/hooks.json 或 plugin.json 內聯",
        slide_num=7, total=TOTAL, source="03 § Hooks"
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Write|Edit",
        "hooks": [
          {
            "type": "command",
            "command": "\\"${CLAUDE_PLUGIN_ROOT}\\"/scripts/format-code.sh"
          }
        ]
      }
    ]
  }
}""",
        Inches(0.7), Inches(1.7), Inches(7.5), Inches(2.8),
        font_size=10
    )

    h.add_text_block(
        slide, "5 種 Hook 類型",
        Inches(8.5), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`command`：執行 shell 命令或指令碼",
            "`http`：POST JSON 事件到 URL",
            "`mcp_tool`：呼叫 MCP server 上的工具",
            "`prompt`：用 LLM 評估提示",
            "`agent`：執行 agentic 驗證器",
        ],
        Inches(8.5), Inches(2.1), Inches(4.5), Inches(2.0),
        font_size=11
    )

    h.add_callout(
        slide, "💡 對 plugin 自己的 bundled MCP server 的 hooks 必須使用**範圍名稱**（`plugin::<server>`）",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5),
        icon="", font_size=12
    )

    h.add_text_block(
        slide, "完整事件表與生命週期見 06-hooks.md",
        Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.4),
        font_size=14, color=h.COLOR_GRAY_TXT, italic=True
    )

    # ============================================================
    # Slide 8：MCP Servers 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "MCP Servers 詳解",
        ".mcp.json 或 plugin.json 內聯",
        slide_num=8, total=TOTAL, source="03 § MCP servers"
    )

    h.add_code_block(
        slide, """{
  "mcpServers": {
    "plugin-database": {
      "command": "${CLAUDE_PLUGIN_ROOT}/servers/db-server",
      "args": ["--config", "${CLAUDE_PLUGIN_ROOT}/config.json"],
      "env": {
        "DB_PATH": "${CLAUDE_PLUGIN_ROOT}/data"
      }
    },
    "plugin-api-client": {
      "command": "npx",
      "args": ["@company/mcp-server", "--plugin-mode"]
    }
  }
}""",
        Inches(0.7), Inches(1.7), Inches(7.0), Inches(3.0),
        font_size=11
    )

    h.add_text_block(
        slide, "🔌 整合行為",
        Inches(8.0), Inches(1.7), Inches(4.8), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "啟用 plugin 時 Plugin MCP servers **自動啟動**",
            "Servers 在 Claude 工具組中顯示為標準 MCP 工具",
            "與 Claude 現有工具無縫整合",
            "可獨立於使用者 MCP servers 設定",
            "`/reload-plugins` 中途運行時，**保留配置未變更** servers 的即時連線",
        ],
        Inches(8.0), Inches(2.1), Inches(4.8), Inches(3.0),
        font_size=12
    )

    h.add_callout(
        slide, "💡 MCP 與 LSP server 都有自己合併 hook / 配置的規則，見對應章節",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 9：LSP Servers 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "LSP Servers 詳解",
        "Language Server Protocol，即時程式碼智慧",
        slide_num=9, total=TOTAL, source="03 § LSP servers"
    )

    h.add_code_block(
        slide, """{
  "go": {
    "command": "gopls",
    "args": ["serve"],
    "extensionToLanguage": {
      ".go": "go"
    }
  }
}""",
        Inches(0.7), Inches(1.7), Inches(6.0), Inches(2.2),
        font_size=13
    )

    h.add_text_block(
        slide, "🔍 LSP 整合提供",
        Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**即時診斷**：每次編輯後立即看到錯誤和警告",
            "**程式碼導航**：跳轉定義、找參考、懸停資訊",
            "**語言感知**：型別資訊和文件",
        ],
        Inches(7.0), Inches(2.1), Inches(5.8), Inches(1.8),
        font_size=12
    )

    h.add_text_block(
        slide, "🛠️ 設定欄位速查",
        Inches(0.5), Inches(4.1), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["欄位", "描述"],
        [
            ["`command` ✅ 必需", "LSP 二進位檔（必須在 PATH 中）"],
            ["`extensionToLanguage` ✅ 必需", "副檔名對應到語言識別碼"],
            ["`args` / `env`", "命令列引數與環境變數"],
            ["`transport`", "stdio（預設）或 socket"],
            ["`startupTimeout` / `shutdownTimeout`", "啟動/關閉逾時（毫秒）"],
            ["`restartOnCrash` / `maxRestarts`", "崩潰後自動重啟（v2.1.205+）"],
            ["`diagnostics`", "是否推診斷到 Claude context（預設 true）"],
        ],
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.4),
        font_size=11
    )

    h.add_callout(
        slide, "⚠️ 必須單獨安裝語言伺服器二進位檔；Plugin 只配置連接，不包含伺服器本體",
        Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 10：Monitors 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Monitors 詳解",
        "背景 monitor，每行 stdout 變 Claude 通知",
        slide_num=10, total=TOTAL, source="03 § Monitors"
    )

    h.add_code_block(
        slide, """[
  {
    "name": "deploy-status",
    "command": "\\"${CLAUDE_PLUGIN_ROOT}\\"/scripts/poll-deploy.sh",
    "description": "Deployment status changes"
  },
  {
    "name": "error-log",
    "command": "tail -F ./logs/error.log",
    "description": "Application error log",
    "when": "on-skill-invoke:debug"
  }
]""",
        Inches(0.7), Inches(1.7), Inches(7.0), Inches(2.7),
        font_size=10
    )

    h.add_text_block(
        slide, "📡 行為說明",
        Inches(8.0), Inches(1.7), Inches(4.8), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "Plugin 啟用時自動啟動 monitor",
            "每個 monitor 執行 shell 命令，**持續整個工作階段**",
            "每個 stdout 行傳遞給 Claude 作為通知",
            "僅在**互動式 CLI** 工作階段執行（host 沒 Monitor tool 時跳過）",
            "若要內聯，設 `plugin.json` 的 `experimental.monitors` 為相同陣列",
        ],
        Inches(8.0), Inches(2.1), Inches(4.8), Inches(2.5),
        font_size=12
    )

    h.add_text_block(
        slide, "🔧 必需 / 選用欄位",
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["欄位", "必需", "描述"],
        [
            ["`name`", "✅", "plugin 內唯一識別碼（防重複程序）"],
            ["`command`", "✅", "session 工作目錄中執行的 shell 命令"],
            ["`description`", "✅", "監視內容的摘要（任務面板顯示）"],
            ["`when`", "❌", "`always`（預設）或 `on-skill-invoke:<skill>`"],
        ],
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.5),
        font_size=11
    )

    h.add_callout(
        slide, "⚠️ Monitor `command` **無法參考 `${user_config.*}`**（不會接收 CLAUDE_PLUGIN_OPTION_<KEY> 環境變數）",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 11：Themes + Skills 目錄 Plugins
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Themes 與 Skills 目錄 Plugins",
        "兩個輕量但實用的特殊形態",
        slide_num=11, total=TOTAL, source="03 § Themes / Skills 目錄"
    )

    h.add_text_block(
        slide, "🎨 Themes（實驗性）",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """{
  "name": "Dracula",
  "base": "dark",
  "overrides": {
    "claude": "#bd93f9",
    "error": "#ff5555",
    "success": "#50fa7b"
  }
}""",
        Inches(0.5), Inches(2.1), Inches(6.0), Inches(2.3),
        font_size=11
    )

    h.add_text_block(
        slide, "📚 Skills 目錄 Plugins",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "任何 skills 目錄下含 `.claude-plugin/plugin.json` 的資料夾",
            "下一個 session 自動作為 `<name>@skills-dir` 載入",
            "**無需 marketplace** 和**無需安裝步驟**",
            "plugin 是**就地發現**的（不複製到快取）",
            "`~/.claude/skills/`：每個專案（僅自己）",
            "`<project>/.claude/skills/`：專案範圍（受信任後）",
        ],
        Inches(6.733), Inches(2.1), Inches(6.0), Inches(3.0),
        font_size=11
    )

    h.add_callout(
        slide, "💡 專案範圍的 `@skills-dir` plugins 只從你啟動 Claude Code 的目錄的 `.claude/skills/` 載入",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5),
        icon="", font_size=12
    )

    h.add_bullet_list(
        slide, [
            "SKILL.md 變更**立即生效**  /  其他元件（hooks/、.mcp.json、agents/）需 `/reload-plugins` 或重啟",
            "想停載 → 刪除資料夾或 `claude plugin disable my-tool@skills-dir`（無 uninstall 步驟）",
        ],
        Inches(0.7), Inches(6.1), Inches(12), Inches(0.8),
        font_size=12
    )

    # ============================================================
    # Slide 12：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "Plugin Manifest 完整架構", "plugin.json 所有欄位")

    # ============================================================
    # Slide 13：plugin.json 完整架構
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "plugin.json 完整架構",
        "Manifest 是**選用**的（但建議提供中繼資料）",
        slide_num=13, total=TOTAL, source="03 § Plugin manifest 完整架構"
    )

    h.add_code_block(
        slide, """{
  "name": "plugin-name",
  "displayName": "Plugin Name",
  "version": "1.2.0",
  "description": "Brief plugin description",
  "author": {
    "name": "Author Name",
    "email": "author@example.com",
    "url": "https://github.com/author"
  },
  "homepage": "https://docs.example.com/plugin",
  "repository": "https://github.com/author/plugin",
  "license": "MIT",
  "keywords": ["keyword1", "keyword2"],
  "skills": "./custom/skills/",
  "commands": ["./custom/commands/special.md"],
  "agents": ["./custom/agents/reviewer.md"],
  "hooks": "./config/hooks.json",
  "mcpServers": "./mcp-config.json",
  "outputStyles": "./styles/",
  "lspServers": "./.lsp.json",
  "experimental": {
    "themes": "./themes/",
    "monitors": "./monitors.json"
  },
  "dependencies": [
    "helper-lib",
    { "name": "secrets-vault", "version": "~2.1.0" }
  ]
}""",
        Inches(0.5), Inches(1.7), Inches(8.0), Inches(5.0),
        font_size=9
    )

    h.add_text_block(
        slide, "📌 關鍵事實",
        Inches(8.7), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**`name` 是唯一必需**（kebab-case，無空格）",
            "省略 manifest → 自動探索預設位置，從目錄名衍生名稱",
            "**無法識別的欄位會被忽略**（可用同一 manifest 給 VS Code / Cursor / npm）",
            "`claude plugin validate` 報告為**警告**，不是錯誤",
            "類型錯誤仍會失敗（不是警告就過）",
            "傳 `--strict` 在 CI 將警告視為錯誤",
        ],
        Inches(8.7), Inches(2.1), Inches(4.5), Inches(3.5),
        font_size=12
    )

    # ============================================================
    # Slide 14：中繼資料欄位
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "中繼資料欄位速查",
        "10 個中繼資料欄位",
        slide_num=14, total=TOTAL, source="03 § 中繼資料"
    )

    h.add_comparison_table(
        slide,
        ["欄位", "類型", "說明"],
        [
            ["`$schema`", "string", "JSON Schema URL（編輯器自動完成，載入時忽略）"],
            ["`displayName`", "string", "/plugin 選擇器顯示名稱（可含空格，v2.1.143+）"],
            ["`version`", "string", "語義版本；設定後 plugin 會被固定"],
            ["`description`", "string", "plugin 用途簡短說明"],
            ["`author`", "object", "作者資訊（name、email、url）"],
            ["`homepage` / `repository`", "string", "文件 URL / 原始程式碼 URL"],
            ["`license`", "string", "SPDX 識別碼（MIT、Apache-2.0…）"],
            ["`keywords`", "array", "探索標籤"],
            ["`defaultEnabled`", "boolean", "未設定時是否啟用（v2.1.154+）"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 defaultEnabled 優先級",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`defaultEnabled` 是「其他因素未決定時」的後備  /  **使用者的 `enabledPlugins` 設定優先**",
            "**相依性要求優先**：被另一個啟用 plugin 所需時，Claude Code 會寫入 `true`",
            "適用場景：新增成本或使用者應選擇加入的 plugin（連接外部服務）",
        ],
        Inches(0.7), Inches(5.7), Inches(12), Inches(1.4),
        font_size=12
    )

    # ============================================================
    # Slide 15：元件路徑欄位
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件路徑欄位速查",
        "9 個元件路徑 + 路徑行為規則",
        slide_num=15, total=TOTAL, source="03 § 元件路徑"
    )

    h.add_comparison_table(
        slide,
        ["欄位", "取代 / 新增", "用途"],
        [
            ["`skills`", "**新增**", "skill 目錄（預設 `skills/` 始終掃描）"],
            ["`commands`", "**取代**", "平面 `.md` skill 檔案或目錄"],
            ["`agents`", "**取代**", "agent 檔案路徑"],
            ["`hooks`", "自有合併", "hook 設定路徑或內聯設定"],
            ["`mcpServers`", "自有合併", "MCP 設定路徑或內聯設定"],
            ["`lspServers`", "自有合併", "LSP 設定"],
            ["`outputStyles`", "**取代**", "輸出樣式檔案/目錄"],
            ["`experimental.themes`", "**取代**", "色彩主題檔案/目錄"],
            ["`experimental.monitors`", "**取代**", "背景 monitor 設定"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=12
    )

    h.add_text_block(
        slide, "📌 路徑行為規則",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "所有路徑必須相對於 plugin 根目錄，並以 `./` 開頭  /  可指定多個路徑（陣列）",
            "來自自訂路徑的元件使用相同的命名和命名空間規則",
            "**hooks / MCP / LSP** 有自己的合併規則（見對應元件章節）",
        ],
        Inches(0.7), Inches(5.7), Inches(12), Inches(1.4),
        font_size=12
    )

    # ============================================================
    # Slide 16：使用者設定 userConfig
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "使用者設定：userConfig",
        "讓使用者在啟用時輸入值",
        slide_num=16, total=TOTAL, source="03 § userConfig"
    )

    h.add_code_block(
        slide, """{
  "userConfig": {
    "api_endpoint": {
      "type": "string",
      "title": "API endpoint",
      "description": "Your team's API endpoint"
    },
    "api_token": {
      "type": "string",
      "title": "API token",
      "description": "API authentication token",
      "sensitive": true
    }
  }
}""",
        Inches(0.7), Inches(1.7), Inches(7.0), Inches(3.0),
        font_size=11
    )

    h.add_text_block(
        slide, "🔧 支援的類型",
        Inches(8.0), Inches(1.7), Inches(4.8), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`string` / `number` / `boolean`",
            "`directory` / `file`",
            "欄位：title、description、default、sensitive、required",
            "string 類型可加 `multiple`（字串陣列）",
            "number 類型可加 `min` / `max`",
        ],
        Inches(8.0), Inches(2.1), Inches(4.8), Inches(2.5),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 儲存位置與替換",
        Inches(0.5), Inches(5.1), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**非敏感值**：`settings.json` 的 `pluginConfigs` 金鑰下",
            "**敏感值**：macOS Keychain 或 `~/.claude/.credentials.json`",
            "替換：`${user_config.KEY}` 用於 MCP/LSP server 設定與 hook 命令",
            "環境變數：`CLAUDE_PLUGIN_OPTION_<KEY>`（大寫）匯出到 hook 程序",
        ],
        Inches(0.7), Inches(5.5), Inches(12), Inches(1.5),
        font_size=12
    )

    h.add_callout(
        slide, "⚠️ 在 shell 中執行的 hook `command` **拒絕 `${user_config.*}`**（避免 shell 注入攻擊）",
        Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.3),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 17：環境變數速查
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "3 個關鍵環境變數",
        "Plugin 路徑與持久資料",
        slide_num=17, total=TOTAL, source="03 § 環境變數"
    )

    h.add_comparison_table(
        slide,
        ["變數", "解析為", "用途"],
        [
            ["`${CLAUDE_PLUGIN_ROOT}`", "plugin 安裝目錄的絕對路徑", "與 plugin 捆綁的指令碼、二進位檔、設定檔"],
            ["`${CLAUDE_PLUGIN_DATA}`", "持久目錄（首次參考時建立）", "已安裝的依賴（`node_modules`、Python venv）"],
            ["`${CLAUDE_PROJECT_DIR}`", "專案根目錄", "專案本地指令碼和設定檔"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.5),
        font_size=12
    )

    h.add_text_block(
        slide, "📌 替換發生在哪些欄位？",
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["Plugin 元件", "佔位符解析的欄位"],
        [
            ["Skill / agent 內容", "佔位符出現的任何地方"],
            ["Hook / monitor 命令", "佔位符出現的任何地方"],
            ["MCP `stdio` servers", "`command`、`args`、`env`"],
            ["MCP `http` / `sse` / `ws` servers", "`url`、`headers`、`headersHelper`"],
            ["LSP servers", "`command`、`args`、`env`、`workspaceFolder`"],
        ],
        Inches(0.5), Inches(4.7), Inches(12.333), Inches(2.0),
        font_size=11
    )

    h.add_callout(
        slide, "💡 `${CLAUDE_PLUGIN_DATA}` 解析為 `~/.claude/plugins/data/{id}/`（非 [a-zA-Z0-9_-] 字元替換為 `-`）",
        Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 18：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "快取、目錄結構、CLI", "Plugin 怎麼被隔離與管理")

    # ============================================================
    # Slide 19：Plugin 快取
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin 快取與檔案解析",
        "為什麼 plugin 不會就地使用？",
        slide_num=19, total=TOTAL, source="03 § Plugin 快取與檔案解析"
    )

    h.add_bullet_list(
        slide, [
            "兩種指定方式：`claude --plugin-dir` / `--plugin-url`（session 期間）vs marketplace（未來 session）",
            "**安全 + 驗證**：marketplace plugins **複製到 `~/.claude/plugins/cache`**，不是就地使用",
            "每個已安裝版本是快取中的單獨目錄，按 marketplace + plugin + 版本分組",
            "更新/卸載時，先前版本目錄**標記為孤立**，**14 天後**背景掃描移除",
            "卸載**最後一個** plugin 後，孤立即保留到下次安裝 plugin",
            "Claude 的 Glob/Grep 工具**跳過孤立目錄**（檔案結果不含過時程式碼）",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(2.8),
        font_size=13
    )

    h.add_text_block(
        slide, "⚠️ 路徑遍歷限制",
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "**已安裝的 plugins 無法參考其目錄外的檔案**",
            "`../shared-utils` 安裝後將無法運作（外部檔案不會被複製到快取）",
        ],
        Inches(0.7), Inches(5.0), Inches(12), Inches(0.8),
        font_size=12
    )

    h.add_text_block(
        slide, "🔗 在 Marketplace 內共享檔案（用 symlinks）",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**plugin 自身目錄內**：symlink 保留為**相對** symlink  /  **同 marketplace 內其他位置**：取消參考，內容複製到快取",
            "**marketplace 外**：因**安全考量被跳過**  /  用 `ln -s ../../shared-plugin/skills/foo ./skills/foo`",
        ],
        Inches(0.7), Inches(6.3), Inches(12), Inches(0.7),
        font_size=11
    )

    # ============================================================
    # Slide 20：完整目錄結構
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin 完整目錄結構",
        "標準 plugin 應該長什麼樣",
        slide_num=20, total=TOTAL, source="03 § Plugin 目錄結構"
    )

    h.add_code_block(
        slide, """enterprise-plugin/
├── .claude-plugin/           # Metadata directory (optional)
│   └── plugin.json             # plugin manifest
├── skills/                   # Skills
│   ├── code-reviewer/
│   │   └── SKILL.md
│   └── pdf-processor/
│       ├── SKILL.md
│       └── scripts/
├── commands/                 # Skills as flat .md files
│   ├── status.md
│   └── logs.md
├── agents/                   # Subagent definitions
│   ├── security-reviewer.md
│   ├── performance-tester.md
│   └── compliance-checker.md
├── output-styles/            # Output style definitions
│   └── terse.md
├── themes/                   # Color theme definitions
│   └── dracula.json
├── monitors/                 # Background monitor configurations
│   └── monitors.json
├── hooks/                    # Hook configurations
│   ├── hooks.json
│   └── security-hooks.json
├── bin/                      # Executables added to PATH
│   └── my-tool
├── settings.json             # Default settings
├── .mcp.json                 # MCP server definitions
├── .lsp.json                 # LSP server configurations
├── scripts/                  # Hook and utility scripts
│   ├── security-scan.sh
│   ├── format-code.py
│   └── deploy.js
├── LICENSE
└── CHANGELOG.md""",
        Inches(0.5), Inches(1.7), Inches(7.5), Inches(5.0),
        font_size=7
    )

    h.add_text_block(
        slide, "📌 重要結構規則",
        Inches(8.2), Inches(1.7), Inches(4.8), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`.claude-plugin/` **只放** `plugin.json`",
            "其他目錄（commands/、agents/、skills/…）**必須在 plugin 根目錄**",
            "Plugin 根目錄的 `CLAUDE.md` **不會**作為專案內容載入",
            "想提供指示 → **放在 skill 內**",
        ],
        Inches(8.2), Inches(2.1), Inches(4.8), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "📂 檔案位置速查",
        Inches(8.2), Inches(4.3), Inches(4.8), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "Manifest: `.claude-plugin/plugin.json`",
            "Skills: `skills/`  /  Commands: `commands/`",
            "Agents: `agents/`  /  Output styles: `output-styles/`",
            "Themes: `themes/`  /  Hooks: `hooks/hooks.json`",
            "MCP: `.mcp.json`  /  LSP: `.lsp.json`",
            "Monitors: `monitors/monitors.json`",
            "Executables: `bin/`  /  Settings: `settings.json`",
        ],
        Inches(8.2), Inches(4.7), Inches(4.8), Inches(2.2),
        font_size=10
    )

    # ============================================================
    # Slide 21：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "CLI 完整指令參考", "10 個 plugin 管理指令")

    # ============================================================
    # Slide 22：plugin init
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "plugin init — 快速搭建",
        "在 ~/.claude/skills/<name>/ 建立新 plugin",
        slide_num=22, total=TOTAL, source="03 § plugin init"
    )

    h.add_code_block(
        slide, """claude plugin init <name> [options]

# 選項
--description <desc>          Manifest 描述
--author <name>               作者名稱（預設 git config user.name）
--author-email <email>        作者電子郵件（預設 git config user.email）
--with <components...>        同時搭建元件資料夾
                              （skills/agents/hooks/mcp/lsp/output-style/channel）
-f, --force                   覆寫現有 .claude-plugin/
-h, --help                    顯示說明

# 別名：new

# 範例
claude plugin init my-helper                          # 最小
claude plugin init my-helper --with skills hooks       # 帶元件
claude plugin init my-helper --force                  # 覆寫""",
        Inches(0.5), Inches(1.7), Inches(8.0), Inches(4.5),
        font_size=11
    )

    h.add_text_block(
        slide, "🔧 --with 元件說明",
        Inches(8.7), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`skills`：額外的命名空間 `:example` skill",
            "`agents`：`agents/` subagent 定義",
            "`hooks`：`hooks/hooks.json` 範例事件",
            "`mcp`：`.mcp.json` HTTP/stdio 範例",
            "`lsp`：`.lsp.json` 語言伺服器範例",
            "`output-style`：`output-styles/.md`",
            "`channel`：基於 MCP 的 channel",
        ],
        Inches(8.7), Inches(2.1), Inches(4.5), Inches(3.5),
        font_size=11
    )

    # ============================================================
    # Slide 23：plugin install / uninstall
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "plugin install / uninstall",
        "從 marketplace 安裝與移除",
        slide_num=23, total=TOTAL, source="03 § install / uninstall"
    )

    h.add_text_block(
        slide, "📥 install",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_GREEN
    )

    h.add_code_block(
        slide, """claude plugin install <plugin> [options]
# 別名：add

# 選項
-s, --scope <scope>         user/project/local
--config <key=value>         設定 userConfig 選項
-y, --yes                   跳過 command source 確認

# 範例
claude plugin install formatter@my-marketplace
claude plugin install formatter@my-marketplace --scope project
claude plugin install formatter@my-marketplace --scope local""",
        Inches(0.5), Inches(2.1), Inches(6.0), Inches(2.7),
        font_size=11
    )

    h.add_text_block(
        slide, "📤 uninstall",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """claude plugin uninstall <plugin> [options]
# 別名：remove / rm

# 選項
-s, --scope <scope>         從範圍卸載
--keep-data                 保留 CLAUDE_PLUGIN_DATA
--prune                     移除自動安裝的依賴
-y, --yes                   跳過 --prune 確認

# 預設行為：從最後範圍卸載時
# 也刪除 ${CLAUDE_PLUGIN_DATA} 目錄
# 用 --keep-data 保留它""",
        Inches(6.733), Inches(2.1), Inches(6.0), Inches(2.7),
        font_size=11
    )

    h.add_callout(
        slide, "💡 卸載最後範圍時，`/plugin` 介面會顯示資料大小並在刪除前提示；CLI 預設刪除",
        Inches(0.5), Inches(5.1), Inches(12.333), Inches(0.5),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 24：plugin prune / enable / disable
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "prune / enable / disable",
        "依賤清理與啟停控制",
        slide_num=24, total=TOTAL, source="03 § prune / enable / disable"
    )

    h.add_code_block(
        slide, """# prune：移除自動安裝的孤立相依性
claude plugin prune [options]    # 別名：autoremove
  -s, --scope <scope>           在範圍進行修剪
  --dry-run                      列出將被移除的內容而不實際移除
  -y, --yes                      跳過確認

# enable：啟用已停用的 plugin
claude plugin enable <plugin> [options]
  -s, --scope <scope>           要啟用的範圍（自動偵測）
  # 若有 dependencies，會在相同範圍內傳遞啟用

# disable：停用 plugin 而不卸載
claude plugin disable [plugin] [options]
  -a, --all                     停用所有已啟用的 plugins
  -s, --scope <scope>           要停用的範圍
  # 當另一個 plugin depends on 目標時失敗""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.0),
        font_size=11
    )

    h.add_callout(
        slide, "💡 想在一個步驟中移除 plugin 並清理其依賴：`claude plugin uninstall --prune`",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    h.add_bullet_list(
        slide, [
            "`prune` 列出孤立依賴並在移除前要求確認  /  `--dry-run` 先看會移除什麼",
            "`enable` / `disable` 改變 plugin 狀態，**不**改變已安裝資料",
        ],
        Inches(0.7), Inches(6.4), Inches(12), Inches(0.7),
        font_size=12
    )

    # ============================================================
    # Slide 25：plugin update / list / details / tag
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "update / list / details / tag",
        "更新、檢視、發行標籤",
        slide_num=25, total=TOTAL, source="03 § update / list / details / tag"
    )

    h.add_comparison_table(
        slide,
        ["指令", "用途", "重要選項"],
        [
            ["`plugin update <plugin>`", "更新到最新版本", "`-s scope`（預設 user）、`-y` 跳過確認"],
            ["`plugin list`", "列出已安裝 plugin 與版本", "`--json` 輸出 JSON；`--available` 包含 marketplace"],
            ["`plugin details <name>`", "顯示元件清單與 token 成本", "包含 always-on 與 on-invoke 估算"],
            ["`plugin tag [path]`", "建立發行版 git 標籤", "`--push`、`--dry-run`、`-f`、`-m message`"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 plugin details 範例輸出",
        Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """dependency-guard 1.2.0
  Dependency analysis for Claude Code sessions

Component inventory
  Skills (2)  scan-dependencies, review-changes
  Agents (0)
  Hooks (1)  (harness-only — no model context cost)
  MCP servers (0)
  LSP servers (0)

Projected token cost
  Always-on:   ~180 tok   added to every session

Per-component (rounded)
  component            always-on  on-invoke
  scan-dependencies        ~100      ~2400
  review-changes            ~80      ~1800""",
        Inches(0.5), Inches(4.25), Inches(12.333), Inches(2.8),
        font_size=9
    )

    # ============================================================
    # Slide 26：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "偵錯與開發工具", "claude --debug 與常見問題")

    # ============================================================
    # Slide 27：偵錯命令
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "偵錯命令：claude --debug",
        "看到 plugin 載入的詳細資訊",
        slide_num=27, total=TOTAL, source="03 § 偵錯"
    )

    h.add_code_block(
        slide, """claude --debug""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.6),
        font_size=20
    )

    h.add_text_block(
        slide, "📋 會顯示",
        Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "正在載入哪些 plugins",
            "Plugin manifests 中的任何錯誤",
            "Skill、agent 和 hook 註冊",
            "MCP server 初始化",
            "LSP server 啟動狀態",
            "為什麼某個 LSP server 被跳過（如 `Executable not found`）",
        ],
        Inches(0.7), Inches(2.9), Inches(12), Inches(2.5),
        font_size=14
    )

    h.add_text_block(
        slide, "💡 想看警告是否實際影響：用 `claude plugin validate ./my-plugin --strict`",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        font_size=13, color=h.COLOR_GRAY_TXT, italic=True
    )

    h.add_text_block(
        slide, "詳細 hook 與 MCP 除錯見 06-hooks.md",
        Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.4),
        font_size=13, color=h.COLOR_GRAY_TXT, italic=True
    )

    # ============================================================
    # Slide 28：常見問題速查
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "常見問題與解決方案",
        "7 個最常踩的坑",
        slide_num=28, total=TOTAL, source="03 § 常見問題"
    )

    h.add_comparison_table(
        slide,
        ["問題", "原因", "解決方案"],
        [
            ["Plugin 未載入", "無效的 plugin.json", "`claude plugin validate` 或 `/plugin validate`"],
            ["Skills 未出現", "目錄結構錯誤", "skills/ 在 plugin 根目錄，不在 .claude-plugin/ 內"],
            ["Hooks 未觸發", "指令碼不可執行", "`chmod +x script.sh`"],
            ["MCP server 失敗", "缺少 ${CLAUDE_PLUGIN_ROOT}", "對所有 plugin 路徑使用變數"],
            ["路徑錯誤", "絕對路徑", "改為相對路徑，以 `./` 開頭"],
            ["LSP Executable not found", "語言伺服器未安裝", "安裝對應的二進位檔"],
            ["Commands 在 .claude-plugin/ 內", "目錄結構錯誤", "移到 plugin 根目錄"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.2),
        font_size=11
    )

    h.add_callout(
        slide, "💡 結構正確但仍失敗 → 看 06-hooks.md 的 hook 疑難排解章節（事件名稱、matcher、type 有效性）",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 29：範例錯誤訊息
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "範例錯誤訊息",
        "Manifest 驗證 vs Plugin 載入",
        slide_num=29, total=TOTAL, source="03 § 範例錯誤訊息"
    )

    h.add_text_block(
        slide, "📋 Manifest 驗證錯誤",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """Invalid JSON syntax: Unexpected token } in JSON at position 142
Plugin has an invalid manifest file at .claude-plugin/plugin.json.
  Validation errors: name: Required
Plugin has a corrupt manifest file at .claude-plugin/plugin.json.
  JSON parse error: ...""",
        Inches(0.5), Inches(2.1), Inches(12.333), Inches(1.6),
        font_size=11
    )

    h.add_text_block(
        slide, "🔌 Plugin 載入錯誤",
        Inches(0.5), Inches(3.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """Warning: No commands found in plugin my-plugin custom directory: ./cmds.
Plugin directory not found at path: ./plugins/my-plugin.
  Check that the marketplace entry has the correct path.
Plugin my-plugin has conflicting manifests:
  both plugin.json and marketplace entry specify components.""",
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(1.8),
        font_size=11
    )

    h.add_callout(
        slide, "💡 「conflicting manifests」→ marketplace 設了 `strict: false` 而 plugin 也有 plugin.json 宣告元件",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 30：Hook 疑難排解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 疑難排解速查",
        "指令碼未執行 vs 未觸發",
        slide_num=30, total=TOTAL, source="03 § Hook 疑難排解"
    )

    h.add_text_block(
        slide, "🔧 指令碼未執行",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "檢查指令碼是否可執行：`chmod +x ./scripts/your-script.sh`",
            "驗證 shebang：`#!/bin/bash` 或 `#!/usr/bin/env bash`",
            "路徑使用 `${CLAUDE_PLUGIN_ROOT}`",
            "手動測試：`./scripts/your-script.sh`",
        ],
        Inches(0.7), Inches(2.1), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "🔌 Hook 未觸發",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "驗證事件名稱（**區分大小寫**）：`PostToolUse`，不是 `postToolUse`",
            "檢查 matcher 模式：`\"matcher\": \"Write|Edit\"` 用於檔案操作",
            "確認 hook 類型有效：`command` / `http` / `mcp_tool` / `prompt` / `agent`",
        ],
        Inches(6.933), Inches(2.1), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "🔍 MCP Server 疑難排解",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**Server 未啟動**：檢查命令存在且可執行  /  驗證所有路徑使用 `${CLAUDE_PLUGIN_ROOT}` 變數",
            "看 `claude --debug` 輸出  /  **手動測試 server**：在 Claude Code 外執行 `command` 與 `args` 看是否能啟動",
        ],
        Inches(0.7), Inches(4.9), Inches(12), Inches(1.5),
        font_size=12
    )

    h.add_callout(
        slide, "💡 完整 hook 生命週期 + 自動回應見 06-hooks.md",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 31：目錄結構錯誤
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "目錄結構錯誤：最常見陷阱",
        "Plugin 載入但元件遺失？檢查這裡",
        slide_num=31, total=TOTAL, source="03 § 目錄結構錯誤"
    )

    h.add_text_block(
        slide, "✅ 正確結構",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_GREEN
    )

    h.add_code_block(
        slide, """my-plugin/
├── .claude-plugin/
│   └── plugin.json   ← Only manifest here
├── commands/         ← At root level
├── agents/           ← At root level
└── hooks/            ← At root level""",
        Inches(0.5), Inches(2.1), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "❌ 錯誤結構（症狀：元件遺失）",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """my-plugin/
└── .claude-plugin/   ← 元件不能放在這裡
    ├── plugin.json
    ├── commands/     ← 錯誤：應在根目錄
    ├── agents/       ← 錯誤：應在根目錄
    └── hooks/        ← 錯誤：應在根目錄""",
        Inches(6.733), Inches(2.1), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_bullet_list(
        slide, [
            "**症狀**：Plugin 載入但元件（skills、agents、hooks）**遺失**",
            "**唯一例外**：`.claude-plugin/` 內可以放 `plugin.json`（manifest）",
            "**所有其他目錄**必須在 plugin 根目錄  /  把元件從 `.claude-plugin/` 內移到根目錄",
        ],
        Inches(0.7), Inches(4.5), Inches(12), Inches(1.8),
        font_size=13
    )

    h.add_callout(
        slide, "💡 Plugin 根目錄的 `CLAUDE.md` 不會作為專案內容載入；想提供指示 → 放在 skill 內",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 32：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 6", "發佈與版本控制", "semver 與發行策略")

    # ============================================================
    # Slide 33：版本管理 4 種來源
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "版本管理：4 種解析來源",
        "Claude Code 用版本作為快取金鑰",
        slide_num=33, total=TOTAL, source="03 § 發佈與版本控制"
    )

    h.add_text_block(
        slide, "🔍 版本解析優先順序（第一個設定的項目勝出）",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "1️⃣ Plugin 的 `plugin.json` 中的 `version` 欄位",
            "2️⃣ Plugin 在 `marketplace.json` 中的 marketplace 項目的 `version` 欄位",
            "3️⃣ Plugin 來源的 git commit SHA（git 託管 marketplace）",
            "4️⃣ `unknown`（npm 來源或不在 git repo 內的本機目錄）",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(2.0),
        font_size=14
    )

    h.add_text_block(
        slide, "📊 兩種版本策略對比",
        Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["方法", "如何操作", "更新行為", "最適合"],
        [
            ["**明確版本**", "plugin.json 設 `\"version\": \"2.1.0\"`", "使用者只在 version 變更時收到更新", "穩定發行週期的已發佈 plugin"],
            ["**Commit-SHA 版本**", "plugin.json 與 marketplace 項目**省略** `version`", "每次新 commit 都是新版本", "積極開發中的內部/團隊 plugin"],
        ],
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(1.7),
        font_size=11
    )

    h.add_callout(
        slide, "📌 設 version 後每次想讓使用者收到變更都必須提升它；僅推 commit 不夠（會看到相同版本字串）",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 34：semver 與 CHANGELOG
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Semantic Versioning 與 CHANGELOG",
        "明確版本的最佳實踐",
        slide_num=34, total=TOTAL, source="03 § semver"
    )

    h.add_text_block(
        slide, "📐 Semver 三段式：MAJOR.MINOR.PATCH",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["段", "何時提升", "範例"],
        [
            ["**MAJOR**", "破壞性變更（不相容 API 變更）", "1.5.2 → 2.0.0"],
            ["**MINOR**", "新功能（向後相容）", "1.5.2 → 1.6.0"],
            ["**PATCH**", "錯誤修正（向後相容）", "1.5.2 → 1.5.3"],
        ],
        Inches(0.5), Inches(2.1), Inches(12.333), Inches(1.4),
        font_size=13
    )

    h.add_bullet_list(
        slide, [
            "遵循 [semver.org](https://semver.org) 規範  /  在 `CHANGELOG.md` 記錄所有變更",
            "**範例 CHANGELOG.md 結構**：",
        ],
        Inches(0.7), Inches(3.6), Inches(12), Inches(0.7),
        font_size=13
    )

    h.add_code_block(
        slide, """# Changelog

## [2.1.0] - 2026-01-15
### Added
- 新增 code-review skill
### Changed
- 升級 default 模型為 Sonnet
### Fixed
- 修復 race condition in hook execution

## [2.0.0] - 2025-12-01
### BREAKING
- 重新命名 skill: `review` → `code-review`
  （使用 `renames` 自動遷移舊名稱）""",
        Inches(0.5), Inches(4.35), Inches(12.333), Inches(2.7),
        font_size=9
    )

    # ============================================================
    # Slide 35：plugin tag 詳解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "plugin tag — 發行版標籤",
        "自動化 git tag 流程",
        slide_num=35, total=TOTAL, source="03 § plugin tag"
    )

    h.add_code_block(
        slide, """claude plugin tag [path] [options]

# 選項
--push                      建立標籤後推送到遠端
--dry-run                   列印將被標籤的內容而不建立
-f, --force                 即使工作樹髒污或標籤已存在也建立
-m, --message <msg>         標籤註解訊息
                            使用 %s 作為版本的佔位符
--remote <name>             使用 --push 時推送到的遠端（預設 origin）

# 範例
claude plugin tag
claude plugin tag --push
claude plugin tag -m "Release %s" --push
claude plugin tag --dry-run""",
        Inches(0.7), Inches(1.7), Inches(7.5), Inches(3.5),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 標籤命名慣例",
        Inches(8.5), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "依賴慣例：`{plugin-name}--v{version}`",
            "範例：`formatter--v2.1.0`",
            "用 `ref` 引用此標籤鎖定版本",
        ],
        Inches(8.5), Inches(2.1), Inches(4.5), Inches(1.5),
        font_size=12
    )

    h.add_text_block(
        slide, "🔄 完整發佈流程",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "1️⃣ 修改 plugin.json 的 `version` + 更新 `CHANGELOG.md`  /  2️⃣ commit 變更",
            "3️⃣ 跑 `claude plugin tag --push`  /  4️⃣ marketplace 偵測到新版本，使用者下次更新時收到",
        ],
        Inches(0.7), Inches(5.9), Inches(12), Inches(1.0),
        font_size=12
    )

    # ============================================================
    # Slide 36：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 7", "參考速查與重點回顧", "整合所有關鍵資訊")

    # ============================================================
    # Slide 37：7 種元件總結
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "7 種元件總結速查",
        "位置 + 用途 + 重要備註",
        slide_num=37, total=TOTAL, source="03 § 7 種元件"
    )

    h.add_comparison_table(
        slide,
        ["元件", "預設位置", "用途", "重要備註"],
        [
            ["Skills", "skills/", "可呼叫的 skill", "frontmatter 布林值用 yes/no/on/off/1/0"],
            ["Agents", "agents/", "專門 subagent", "**不支援** hooks、mcpServers、permissionMode"],
            ["Hooks", "hooks/hooks.json", "事件自動回應", "5 種 type：command/http/mcp_tool/prompt/agent"],
            ["MCP", ".mcp.json", "外部工具整合", "Plugin MCP 用 `plugin::<server>` 範圍"],
            ["LSP", ".lsp.json", "語言伺服器", "必須先安裝二進位檔，Plugin 不含"],
            ["Monitors", "monitors/monitors.json", "背景通知", "**無法**用 ${user_config.*}"],
            ["Themes", "themes/", "顏色主題", "實驗性；按 Ctrl+E 複製到 `~/.claude/themes/`"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 完整 manifest 欄位對照",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**9 個元件路徑欄位**：skills / commands / agents / hooks / mcpServers / lspServers / outputStyles / experimental.themes / experimental.monitors",
            "**3 個核心環境變數**：`${CLAUDE_PLUGIN_ROOT}` / `${CLAUDE_PLUGIN_DATA}` / `${CLAUDE_PROJECT_DIR}`",
            "**使用者設定**：`userConfig`（讓使用者啟用時輸入）",
        ],
        Inches(0.7), Inches(5.7), Inches(12), Inches(1.4),
        font_size=12
    )

    # ============================================================
    # Slide 38：CLI 指令速查表
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "CLI 完整指令速查表",
        "10 個 plugin 管理指令",
        slide_num=38, total=TOTAL, source="03 § CLI 速查"
    )

    h.add_comparison_table(
        slide,
        ["指令", "別名", "用途"],
        [
            ["`claude plugin init <name>`", "`new`", "在 skills/ 搭建新 plugin"],
            ["`claude plugin install <plugin>`", "—", "從 marketplace 安裝"],
            ["`claude plugin uninstall <plugin>`", "`remove`、`rm`", "移除已安裝的 plugin"],
            ["`claude plugin prune`", "`autoremove`", "移除孤立的自動安裝依賴"],
            ["`claude plugin enable <plugin>`", "—", "啟用已停用的 plugin"],
            ["`claude plugin disable [plugin]`", "—", "停用 plugin（不卸載）"],
            ["`claude plugin update <plugin>`", "—", "更新到最新版本"],
            ["`claude plugin list`", "—", "列出已安裝 plugin（`--json` / `--available`）"],
            ["`claude plugin details <name>`", "—", "顯示元件清單與 token 成本"],
            ["`claude plugin tag [path]`", "—", "建立發行版 git 標籤"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.5),
        font_size=11
    )

    h.add_callout(
        slide, "💡 還有 `validate`（驗證 JSON）與 `marketplace` 子命令群（add/list/remove/update），見 01-plugin-marketplaces.md",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 39：環境變數總結
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "環境變數總結",
        "Plugin 路徑與持久資料",
        slide_num=39, total=TOTAL, source="03 § 環境變數"
    )

    h.add_comparison_table(
        slide,
        ["變數", "解析為", "用途", "注意事項"],
        [
            ["`${CLAUDE_PLUGIN_ROOT}`", "plugin 安裝目錄絕對路徑", "指令碼、二進位檔、設定檔", "**最常用**"],
            ["`${CLAUDE_PLUGIN_DATA}`", "持久目錄（plugin 更新後保留）", "已安裝依賴（node_modules、Python venv）", "卸載時**預設刪除**（`--keep-data` 保留）"],
            ["`${CLAUDE_PROJECT_DIR}`", "專案根目錄", "專案本地指令碼和設定檔", "與 plugin 安裝位置無關"],
            ["`${user_config.KEY}`", "使用者啟用時輸入的值", "MCP/LSP 設定、hook 命令", "shell 拒絕（防注入）"],
            ["`CLAUDE_PLUGIN_OPTION_<KEY>`", "userConfig 環境變數版本", "hook 程序讀取", "**不適用** monitor"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.0),
        font_size=11
    )

    h.add_callout(
        slide, "💡 `${CLAUDE_PLUGIN_DATA}` 解析為 `~/.claude/plugins/data/{id}/`（非標準字元替換為 `-`）",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    h.add_bullet_list(
        slide, [
            "所有三個變數都會**匯出為環境變數**到 hook 程序與 MCP / LSP server 子程序",
            "Skill 與 agent 內容中的佔位符**也會被解析**（任何出現的地方）",
        ],
        Inches(0.7), Inches(6.4), Inches(12), Inches(0.7),
        font_size=12
    )

    # ============================================================
    # Slide 40：plugin.json 完整欄位表
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "plugin.json 欄位總表",
        "20+ 個欄位的權威參考",
        slide_num=40, total=TOTAL, source="03 § plugin.json"
    )

    h.add_comparison_table(
        slide,
        ["類別", "欄位", "說明"],
        [
            ["**核心**", "`name` ✅ 必需", "唯一識別碼（kebab-case）"],
            ["**中繼資料**", "`displayName` / `version` / `description`", "顯示名稱、版本、簡短說明"],
            ["**中繼資料**", "`author` / `homepage` / `repository` / `license`", "作者與授權"],
            ["**中繼資料**", "`keywords` / `defaultEnabled`", "探索標籤與預設啟用"],
            ["**元件路徑**", "`skills` / `commands` / `agents`", "skill / command / agent 路徑"],
            ["**元件路徑**", "`hooks` / `mcpServers` / `lspServers`", "hook / MCP / LSP 設定"],
            ["**元件路徑**", "`outputStyles`", "輸出樣式路徑"],
            ["**實驗性**", "`experimental.themes` / `experimental.monitors`", "主題與背景 monitor"],
            ["**進階**", "`userConfig`", "使用者啟用時輸入"],
            ["**進階**", "`dependencies`", "依賴其他 plugin（可選 semver）"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.5),
        font_size=11
    )

    h.add_callout(
        slide, "💡 Claude Code 忽略無法識別的頂層欄位（可保留 VS Code / Cursor / npm 中繼資料）",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 41：manifest 路徑行為規則
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "路徑行為規則速查",
        "取代 vs 新增 vs 自有合併",
        slide_num=41, total=TOTAL, source="03 § 路徑行為規則"
    )

    h.add_two_column_compare(
        slide,
        "🔄 取代預設值（6 個）",
        [
            "`commands`：預設 `commands/`",
            "`agents`：預設 `agents/`",
            "`outputStyles`：預設 `output-styles/`",
            "`experimental.themes`：預設 `themes/`",
            "`experimental.monitors`：預設 monitors 路徑",
            "若指定路徑 → **完全取代**預設目錄",
        ],
        "➕ 新增到預設值（1 個）",
        [
            "`skills` 欄位的路徑**新增**到預設 `skills/` 掃描",
            "預設 `skills/` 目錄**始終被掃描**",
            "多個 plugin 共享根目錄 skills → 列特定子目錄",
            "範例：`skills: [\"./skills/code-review\", \"./skills/docs\"]`",
        ],
        top=Inches(1.7), height=Inches(3.0)
    )

    h.add_text_block(
        slide, "🔧 自有合併規則（3 個）",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**hooks** / **MCP servers** / **LSP servers**：有自己的合併邏輯（marketplace 補充 vs plugin 內部）",
            "見 03 § Hooks / MCP / LSP 章節的詳細合併規則",
        ],
        Inches(0.7), Inches(5.4), Inches(12), Inches(0.8),
        font_size=12
    )

    h.add_callout(
        slide, "💡 所有路徑必須相對於 plugin 根目錄，並以 `./` 開頭；多個路徑可指定為陣列",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 42：plugin validate 完整指南
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "plugin validate 完整指南",
        "驗證是 CI 整合的關鍵",
        slide_num=42, total=TOTAL, source="03 § validate"
    )

    h.add_code_block(
        slide, """# 基本驗證
claude plugin validate .

# 從 Claude Code 內
/plugin validate .

# 驗證個別 plugin
claude plugin validate ./plugins/my-plugin

# CI 模式：將警告視為錯誤
claude plugin validate ./my-plugin --strict""",
        Inches(0.5), Inches(1.7), Inches(8.0), Inches(2.4),
        font_size=11
    )

    h.add_text_block(
        slide, "🔍 會檢查的項目",
        Inches(8.7), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "marketplace.json 架構錯誤",
            "重複的 plugin 名稱",
            "來源路徑遍歷",
            "每個 plugin 的 plugin.json",
            "項目 version 與 plugin.json 版本一致性",
        ],
        Inches(8.7), Inches(2.1), Inches(4.5), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "⚠️ 常見錯誤對照",
        Inches(0.5), Inches(4.4), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_comparison_table(
        slide,
        ["錯誤", "原因"],
        [
            ["`File not found: .claude-plugin/marketplace.json`", "缺少 manifest"],
            ["`Invalid JSON syntax: Unexpected token...`", "JSON 語法錯誤"],
            ["`Duplicate plugin name \"x\"`", "plugin 名稱重複"],
            ["`plugins[0].source: Path contains \"..\"`", "來源路徑含 `..`"],
            ["`YAML frontmatter failed to parse`", "skill/agent/command frontmatter 無效"],
            ["`Invalid JSON syntax` (hooks.json)", "hooks/hooks.json 格式錯誤，**阻止整個 plugin 載入**"],
        ],
        Inches(0.5), Inches(4.8), Inches(12.333), Inches(2.0),
        font_size=11
    )

    h.add_callout(
        slide, "💡 警告（非阻擋）：Marketplace has no plugins / No marketplace description / Plugin name not kebab-case",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 43：發佈完整流程
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "發佈完整流程（明確版本）",
        "從開發到使用者收到更新",
        slide_num=43, total=TOTAL, source="03 § 發佈流程"
    )

    flow_steps = [
        ("1", "修改程式碼", "新功能、修正"),
        ("2", "更新 plugin.json", "提升 version + 改 manifest"),
        ("3", "更新 CHANGELOG.md", "記錄變更"),
        ("4", "commit + push", "git 推送"),
        ("5", "claude plugin tag", "建立 + 推送 git 標籤"),
        ("6", "使用者收到更新", "下次 /reload-plugins 觸發"),
    ]

    step_w = Inches(1.95)
    step_h = Inches(2.5)
    h_gap = Inches(0.2)
    total_w = step_w * 6 + h_gap * 5
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.5)

    for i, (num, title, desc) in enumerate(flow_steps):
        x = start_x + i * (step_w + h_gap)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y, step_w, step_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + step_w / 2 - Inches(0.3), start_y + Inches(0.2),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = h.COLOR_PRIMARY
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            x + step_w / 2 - Inches(0.3), start_y + Inches(0.22),
            Inches(0.6), Inches(0.55),
            font_size=24, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, title,
            x + Inches(0.1), start_y + Inches(1.0), step_w - Inches(0.2), Inches(0.5),
            font_size=13, bold=True, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.1), start_y + Inches(1.55), step_w - Inches(0.2), Inches(0.8),
            font_size=10, align=PP_ALIGN.CENTER, color=h.COLOR_GRAY_TXT
        )

        # 箭頭
        if i < 5:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + step_w + Inches(0.02), start_y + Inches(1.0),
                Inches(0.16), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    h.add_callout(
        slide, "📌 若用 commit-SHA 版本（省略 version）→ 只需 commit + push，使用者**每次 commit** 收到更新",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    h.add_bullet_list(
        slide, [
            "**預先發佈標籤**：用 `1.0.0-rc.1`、`1.0.0-beta.1` 等 semver 預發版格式",
            "**依賴標籤**：依賴其他 plugin 時用 `{plugin-name}--v{version}` 慣例",
        ],
        Inches(0.7), Inches(6.0), Inches(12), Inches(0.9),
        font_size=12
    )

    # ============================================================
    # Slide 44：相關文件與資源
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "相關文件與資源",
        "延伸閱讀與官方連結",
        slide_num=44, total=TOTAL, source="03 § 參考資源"
    )

    h.add_bullet_list(
        slide, [
            "**系列內部**",
            "  - 02-plugins.md：教學和實際使用",
            "  - 01-plugin-marketplaces.md：建立和管理 marketplaces",
            "  - 04-skills.md：Skill 開發詳細資訊",
            "  - 05-subagents.md：Agent 設定和功能",
            "  - 06-hooks.md：事件處理和自動化",
            "**官方文件**",
            "  - MCP：外部工具整合",
            "  - Settings：Plugins 的設定選項",
            "  - Plugin 設定：settings.json 的 pluginConfigs 金鑰",
            "**重要版本**",
            "  - 當前適配 Claude Code v2.1.x",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(4.5),
        font_size=13
    )

    h.add_callout(
        slide, "💡 本章是技術參考權威；遇到具體問題先查 03，回到對應概念章節看教學",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 45：重點回顧
    # ============================================================
    h.add_summary_slide(
        slide=h.add_blank_slide(prs),
        title="重點回顧",
        key_points=[
            "**7 種元件**：Skills / Agents / Hooks / MCP / LSP / Monitors / Themes（位置 + 用途）",
            "**plugin.json 20+ 欄位**：`name` 必需；中繼資料、元件路徑、userConfig、dependencies",
            "**3 個關鍵環境變數**：`${CLAUDE_PLUGIN_ROOT}` / `${CLAUDE_PLUGIN_DATA}` / `${CLAUDE_PROJECT_DIR}`",
            "**快取機制**：marketplace plugins 複製到 `~/.claude/plugins/cache`，14 天後孤立版本清理",
            "**10 個 CLI 指令**：init / install / uninstall / prune / enable / disable / update / list / details / tag",
        ],
        next_steps=[
            "用 `claude plugin init my-helper --with skills hooks` 搭建你的 plugin",
            "用 `claude --debug` 觀察載入過程",
            "進階：試著把元件放在自訂路徑（如 `outputStyles`、`userConfig`）",
            "接著閱讀 `04-skills.md` 深入 skill 開發，或 `06-hooks.md` 學事件驅動！",
        ],
        source="03-plugins-reference.md"
    )

    # ============================================================
    # 儲存
    # ============================================================
    output = "/home/elan/pi-proj/03-plugins-reference.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
