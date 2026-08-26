"""
learn2deck 基礎形狀函式

每個函式都接受可選的 theme 參數，向後相容於原有 _pptx_helpers.py。
"""
from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt

from ..core import Theme
from .layout import get_font_size, get_slide_height, get_slide_width


# === 全域 helper：取得 theme 值或 fallback ===

def _color(theme: Theme | None, name: str, fallback_hex: str = "#000000") -> RGBColor:
    """取得 theme 顏色，沒設或沒 theme 就用 fallback"""
    if theme is None:
        return _hex_to_rgb(fallback_hex)
    return theme.get_color_or_default(name, fallback_hex)


def _font(theme: Theme | None, name: str, fallback: str = "Calibri") -> str:
    """取得 theme 字體，沒設或沒 theme 就用 fallback"""
    if theme is None:
        return fallback
    return theme.get_font(name)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# === 簡報與投影片 ===

def new_presentation(theme: Theme | None = None) -> Presentation:
    """建立 16:9 簡報"""
    prs = Presentation()
    prs.slide_width = get_slide_width(theme)
    prs.slide_height = get_slide_height(theme)
    return prs


def add_blank_slide(prs: Presentation) -> Slide:
    """新增空白頁（最常用，自己畫所有元素）"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_slide_bg(slide: Slide, color: RGBColor | None = None, theme: Theme | None = None) -> None:
    """設定投影片背景色"""
    if color is None:
        color = _color(theme, "bg_cream", "#FAF8F3")
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# === 標題列 ===

def add_title_bar(
    slide: Slide,
    title_text: str,
    subtitle_text: str | None = None,
    slide_num: int | None = None,
    total: int | None = None,
    source: str | None = None,
    theme: Theme | None = None,
) -> None:
    """統一的標題列設計
    - 橘色 accent line 在頂部
    - 主標題（大字）
    - 副標題（小字，灰色）
    - 右上角頁碼
    - 右下角來源（對應系列檔案）
    """
    primary = _color(theme, "primary", "#C75A1A")
    dark = _color(theme, "dark", "#2C2C2C")
    gray_text = _color(theme, "gray_text", "#6B6B6B")
    title_font = _font(theme, "title", "Calibri")
    body_font = _font(theme, "body", "Calibri")
    title_size = get_font_size(theme, "slide_title", 32)
    subtitle_size = get_font_size(theme, "slide_subtitle", 16)

    # 頂部 accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        get_slide_width(theme), Inches(0.15)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = primary
    line.line.fill.background()

    # 主標題
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(11), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = title_font
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = dark

    # 副標題
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.0),
            Inches(11), Inches(0.4)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = subtitle_text
        run.font.name = body_font
        run.font.size = Pt(subtitle_size)
        run.font.color.rgb = gray_text
        run.font.italic = True

    # 右上角頁碼
    if slide_num is not None and total is not None:
        page_box = slide.shapes.add_textbox(
            Inches(12), Inches(0.3),
            Inches(1.2), Inches(0.3)
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_num} / {total}"
        run.font.name = body_font
        run.font.size = Pt(10)
        run.font.color.rgb = gray_text

    # 右下角來源
    if source:
        src_box = slide.shapes.add_textbox(
            Inches(8), Inches(7.1),
            Inches(5), Inches(0.3)
        )
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"📖 來源：{source}"
        run.font.name = body_font
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = gray_text

    # 左下角品牌（從 theme 讀，預設為系列名稱）
    brand_text = "Claude Code Plugin 完整學習系列"  # TODO: theme.brand_text
    brand_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.1),
        Inches(6), Inches(0.3)
    )
    tf = brand_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = brand_text
    run.font.name = body_font
    run.font.size = Pt(9)
    run.font.color.rgb = primary
    run.font.bold = True


# === 文字區塊 ===

def add_text_block(
    slide: Slide,
    text: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor | None = None,
    align: Any = PP_ALIGN.LEFT,
    bg_color: RGBColor | None = None,
    font: str | None = None,
    theme: Theme | None = None,
):
    """新增純文字區塊（可選背景色）"""
    if color is None:
        color = _color(theme, "dark", "#2C2C2C")
    if font is None:
        font = _font(theme, "body", "Calibri")

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

    if bg_color:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
    return box


def add_bullet_list(
    slide: Slide,
    items: list,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    font_size: int = 14,
    color: RGBColor | None = None,
    font: str | None = None,
    bullet_char: str = "•",
    indent_levels: list | None = None,
    theme: Theme | None = None,
):
    """新增階層式項目清單
    items: list of (indent_level, text) tuples 或純字串清單
    """
    if color is None:
        color = _color(theme, "dark", "#2C2C2C")
    if font is None:
        font = _font(theme, "body", "Calibri")

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)

    # 正規化 items
    if isinstance(items[0], str):
        items = [(0, t) for t in items]
    elif isinstance(items[0], tuple) and len(items[0]) == 2:
        pass
    else:
        raise ValueError("items 應為 str list 或 (indent, text) tuple list")

    for i, (indent, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent

        # 縮排
        prefix = "    " * indent
        if indent == 0:
            prefix = f"{bullet_char} "

        run = p.add_run()
        run.text = f"{prefix}{text}"
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

    return box


# === 程式碼區塊 ===

def add_code_block(
    slide: Slide,
    code: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    language: str = "bash",
    font_size: int = 12,
    theme: Theme | None = None,
):
    """新增深色背景的程式碼區塊"""
    code_bg = _color(theme, "code_bg", "#1E1E1E")
    code_fg = _color(theme, "code_fg", "#E6E6E6")
    mono_font = _font(theme, "mono", "Consolas")

    # 背景矩形
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = code_bg
    bg.line.fill.background()

    # 程式碼文字
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.15)
    tf.margin_bottom = Inches(0.1)

    lines = code.strip("\n").split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = mono_font
        run.font.size = Pt(font_size)
        run.font.color.rgb = code_fg


# === 提示框 ===

def add_callout(
    slide: Slide,
    text: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    bg_color: RGBColor | None = None,
    border_color: RGBColor | None = None,
    icon: str = "💡",
    font_size: int = 13,
    theme: Theme | None = None,
):
    """新增提示框（圓角矩形）"""
    if bg_color is None:
        bg_color = _color(theme, "bg_gray", "#F3F0E9")
    if border_color is None:
        border_color = _color(theme, "primary", "#C75A1A")
    text_color = _color(theme, "dark", "#2C2C2C")
    body_font = _font(theme, "body", "Calibri")

    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = bg_color
    callout.line.color.rgb = border_color
    callout.line.width = Pt(1.5)

    tf = callout.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"{icon}  {text}" if icon else text
    run.font.name = body_font
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color


# === 表格 ===

def add_comparison_table(
    slide: Slide,
    headers: list,
    rows: list,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    header_bg: RGBColor | None = None,
    header_fg: RGBColor | None = None,
    alt_row_bg: RGBColor | None = None,
    font_size: int = 11,
    theme: Theme | None = None,
):
    """新增比較表格"""
    if header_bg is None:
        header_bg = _color(theme, "primary", "#C75A1A")
    if header_fg is None:
        header_fg = _color(theme, "white", "#FFFFFF")
    if alt_row_bg is None:
        alt_row_bg = _color(theme, "bg_gray", "#F3F0E9")
    white = _color(theme, "white", "#FFFFFF")
    dark = _color(theme, "dark", "#2C2C2C")
    title_font = _font(theme, "title", "Calibri")
    body_font = _font(theme, "body", "Calibri")

    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 標題列
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        tf = cell.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.name = title_font
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = header_fg

    # 資料列
    for row_idx, row in enumerate(rows, start=1):
        bg = alt_row_bg if row_idx % 2 == 0 else white
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = body_font
            run.font.size = Pt(font_size)
            run.font.color.rgb = dark


# === 流程圖 ===

def add_flow_box(
    slide: Slide,
    text: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    bg_color: RGBColor | None = None,
    border_color: RGBColor | None = None,
    font_size: int = 13,
    font_color: RGBColor | None = None,
    bold: bool = True,
    theme: Theme | None = None,
):
    """流程圖方塊（圓角矩形 + 文字）"""
    if bg_color is None:
        bg_color = _color(theme, "bg_gray", "#F3F0E9")
    if border_color is None:
        border_color = _color(theme, "primary", "#C75A1A")
    if font_color is None:
        font_color = _color(theme, "dark", "#2C2C2C")
    title_font = _font(theme, "title", "Calibri")

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = title_font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    return box


def add_arrow(
    slide: Slide,
    start_x: Any,
    start_y: Any,
    end_x: Any,
    end_y: Any,
    color: RGBColor | None = None,
    width: int = 2,
    theme: Theme | None = None,
):
    """新增箭頭（直線）"""
    if color is None:
        color = _color(theme, "gray_text", "#6B6B6B")
    connector = slide.shapes.add_connector(
        1,  # Straight arrow
        start_x, start_y,
        end_x - start_x, end_y - start_y
    )
    line = connector.line
    line.color.rgb = color
    line.width = Pt(width)
    return connector
