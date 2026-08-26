"""
learn2deck 特殊頁面建構函式

包含：cover, section_divider, summary, two_column
這些是較高階的組合，使用 shapes.py 的基礎函式。
"""
from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..core import Theme
from .shapes import (
    _color, _font,
    add_blank_slide, set_slide_bg, add_title_bar,
    add_text_block, add_bullet_list, add_code_block,
    get_font_size,
)


# === 封面頁 ===

def add_cover_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    tag: str = "Claude Code Plugin 完整學習系列",
    theme: Theme | None = None,
):
    """封面頁"""
    primary = _color(theme, "primary", "#C75A1A")
    dark = _color(theme, "dark", "#2C2C2C")
    bg_cream = _color(theme, "bg_cream", "#FAF8F3")
    gray_text = _color(theme, "gray_text", "#6B6B6B")
    title_font = _font(theme, "title", "Calibri")
    body_font = _font(theme, "body", "Calibri")
    cover_title_size = get_font_size(theme, "cover_title", 54)
    cover_subtitle_size = get_font_size(theme, "cover_subtitle", 22)

    slide = add_blank_slide(prs)
    set_slide_bg(slide, bg_cream, theme)

    # 大標題
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = title_font
    run.font.size = Pt(cover_title_size)
    run.font.bold = True
    run.font.color.rgb = primary

    # 副標題
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(11.333), Inches(1)
    )
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitle
    run.font.name = body_font
    run.font.size = Pt(cover_subtitle_size)
    run.font.color.rgb = dark

    # 標籤
    tag_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.0),
        Inches(11.333), Inches(0.5)
    )
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tag
    run.font.name = body_font
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = gray_text

    # 橘色裝飾條
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.666), Inches(5.6),
        Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    return slide


# === 章節分隔頁 ===

def add_section_divider(
    prs: Presentation,
    section_num: str,
    section_title: str,
    section_subtitle: str,
    theme: Theme | None = None,
):
    """章節分隔頁"""
    primary = _color(theme, "primary", "#C75A1A")
    dark = _color(theme, "dark", "#2C2C2C")
    bg_cream = _color(theme, "bg_cream", "#FAF8F3")
    gray_text = _color(theme, "gray_text", "#6B6B6B")
    title_font = _font(theme, "title", "Calibri")
    body_font = _font(theme, "body", "Calibri")
    section_num_size = get_font_size(theme, "section_num", 96)
    section_title_size = get_font_size(theme, "section_title", 40)
    section_subtitle_size = get_font_size(theme, "section_subtitle", 18)

    slide = add_blank_slide(prs)
    set_slide_bg(slide, bg_cream, theme)

    # 巨大章節編號
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(12.333), Inches(2)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_num
    run.font.name = title_font
    run.font.size = Pt(section_num_size)
    run.font.bold = True
    run.font.color.rgb = primary

    # 章節標題
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.0),
        Inches(12.333), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_title
    run.font.name = title_font
    run.font.size = Pt(section_title_size)
    run.font.bold = True
    run.font.color.rgb = dark

    # 副標題
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.2),
        Inches(12.333), Inches(0.8)
    )
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_subtitle
    run.font.name = body_font
    run.font.size = Pt(section_subtitle_size)
    run.font.italic = True
    run.font.color.rgb = gray_text

    return slide


# === 重點回顧頁 ===

def add_summary_slide(
    slide,
    title: str = "重點回顧",
    key_points: list | None = None,
    next_steps: list | None = None,
    source: str | None = None,
    theme: Theme | None = None,
):
    """重點回顧頁（直接畫在傳入的 slide 上）"""
    add_title_bar(slide, title, "重點整理與下一步行動", source=source, theme=theme)
    primary = _color(theme, "primary", "#C75A1A")

    if key_points:
        add_text_block(
            slide, "📌 關鍵要點",
            Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
            font_size=20, bold=True, color=primary, theme=theme
        )
        add_bullet_list(
            slide, key_points,
            Inches(0.7), Inches(2.2), Inches(12), Inches(2.5),
            font_size=15, theme=theme
        )

    if next_steps:
        add_text_block(
            slide, "🚀 下一步行動",
            Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
            font_size=20, bold=True, color=primary, theme=theme
        )
        add_bullet_list(
            slide, next_steps,
            Inches(0.7), Inches(5.3), Inches(12), Inches(1.5),
            font_size=15, theme=theme
        )


# === 雙欄對比 ===

def add_two_column_compare(
    slide,
    left_title: str,
    left_content: list,
    right_title: str,
    right_content: list,
    top: Any = Inches(1.7),
    height: Any = Inches(5.0),
    left_color: RGBColor | None = None,
    right_color: RGBColor | None = None,
    theme: Theme | None = None,
):
    """左右兩欄對比佈局"""
    if left_color is None:
        left_color = _color(theme, "blue", "#3B82F6")
    if right_color is None:
        right_color = _color(theme, "primary", "#C75A1A")
    bg_gray = _color(theme, "bg_gray", "#F3F0E9")
    white = _color(theme, "white", "#FFFFFF")
    body_font = _font(theme, "body", "Calibri")

    # 左欄
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top, Inches(6.1), height
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = bg_gray
    left_box.line.color.rgb = left_color
    left_box.line.width = Pt(2)

    # 左欄標題列
    left_title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), top, Inches(6.1), Inches(0.5)
    )
    left_title_bg.fill.solid()
    left_title_bg.fill.fore_color.rgb = left_color
    left_title_bg.line.fill.background()

    add_text_block(
        slide, left_title,
        Inches(0.5), top, Inches(6.1), Inches(0.5),
        font_size=16, bold=True, color=white,
        align=PP_ALIGN.CENTER, font=body_font
    )

    add_bullet_list(
        slide, left_content,
        Inches(0.7), top + Inches(0.7), Inches(5.7), height - Inches(0.8),
        font_size=13, theme=theme
    )

    # 右欄
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.733), top, Inches(6.1), height
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = bg_gray
    right_box.line.color.rgb = right_color
    right_box.line.width = Pt(2)

    # 右欄標題列
    right_title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.733), top, Inches(6.1), Inches(0.5)
    )
    right_title_bg.fill.solid()
    right_title_bg.fill.fore_color.rgb = right_color
    right_title_bg.line.fill.background()

    add_text_block(
        slide, right_title,
        Inches(6.733), top, Inches(6.1), Inches(0.5),
        font_size=16, bold=True, color=white,
        align=PP_ALIGN.CENTER, font=body_font
    )

    add_bullet_list(
        slide, right_content,
        Inches(6.933), top + Inches(0.7), Inches(5.7), height - Inches(0.8),
        font_size=13, theme=theme
    )
