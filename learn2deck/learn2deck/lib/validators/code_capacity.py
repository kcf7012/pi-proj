"""
R1: code 框容量驗證

規則：N 行 × 行高 + 0.2" margin ≤ 框高
"""
from __future__ import annotations

from typing import TYPE_CHECKING

from ..pptx_helpers.layout import LINE_HEIGHTS
from .base import BaseValidator, Issue, Severity

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide


# 黑色背景矩形（深色）的 RGB 閾值（小於此視為深色 = code 框）
DARK_BG_THRESHOLD = 0x40  # 64


class CodeCapacityValidator(BaseValidator):
    """code 框容量驗證

    計算邏輯：
    1. 找出每張投影片的「黑色矩形」（code 框背景）
    2. 找配對的 textbox（位置大小相同）
    3. 計算實際行數與字級 → 預估所需高度
    4. 若所需 > 實際 → 報告 ERROR
    """

    rule_id = "R1"

    def validate(self, prs: "Presentation") -> list[Issue]:
        issues: list[Issue] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            issues.extend(self._check_slide(slide, slide_idx))
        return issues

    def _check_slide(self, slide: "Slide", slide_num: int) -> list[Issue]:
        """檢查單張投影片的所有 code 框"""
        issues: list[Issue] = []

        # 找所有黑色矩形（code 框背景）
        for shape in slide.shapes:
            shape_type = str(shape.shape_type)
            if "AUTO_SHAPE" not in shape_type:
                continue

            try:
                rgb = shape.fill.fore_color.rgb
            except Exception:
                continue

            if not rgb or not self._is_dark(rgb):
                continue

            # 找配對的 textbox（含 TEXT_BOX 類型）
            top_in = shape.top / 914400 if shape.top else 0
            height_in = shape.height / 914400 if shape.height else 0

            textbox = self._find_paired_textbox(slide, shape)
            if textbox is None:
                continue

            # 計算行數與字級
            n_lines, font_size = self._analyze_textbox(textbox)
            if n_lines == 0:
                continue

            # 預估所需高度
            line_h = LINE_HEIGHTS.get(font_size, font_size / 72 * 1.2)
            needed = n_lines * line_h + 0.2  # 0.2" margin
            actual = height_in

            if actual < needed:
                margin = actual - needed
                issues.append(self.make_error(
                    f"Code 框裝不下：{n_lines} 行 @ {font_size}pt 需要 {needed:.2f}\"，實際 {actual:.2f}\"（不足 {-margin:.2f}\"）",
                    slide_num=slide_num,
                    lines=n_lines,
                    font_size=font_size,
                    needed_height=needed,
                    actual_height=actual,
                    margin=margin,
                ))

        return issues

    def _is_dark(self, rgb) -> bool:
        """判斷 RGB 是否為深色（code 框背景）"""
        try:
            r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
            return max(r, g, b) < DARK_BG_THRESHOLD
        except (TypeError, ValueError, IndexError):
            return False

    def _find_paired_textbox(self, slide: "Slide", bg_shape) -> "object | None":
        """找與背景矩形配對的 textbox（位置大小完全相同）

        注意：python-pptx 的 slide.shapes.add_shape() 會重換 XML，
        所以 bg_shape 參考可能 stale。這裡用型別判斷排除背景本身。
        """
        for shape in slide.shapes:
            # 跳過 AUTO_SHAPE 本身（背景矩形是 AUTO_SHAPE）
            shape_type = str(shape.shape_type)
            if "AUTO_SHAPE" in shape_type:
                continue
            if not shape.has_text_frame:
                continue
            if (shape.top == bg_shape.top
                and shape.left == bg_shape.left
                and shape.width == bg_shape.width
                and shape.height == bg_shape.height):
                return shape
        return None

    def _analyze_textbox(self, textbox) -> tuple[int, int]:
        """從 textbox 抽取行數與字級

        Returns:
            (n_lines, font_size) — 若無內容則 (0, 0)
        """
        text = textbox.text_frame.text
        if not text or not text.strip():
            return 0, 0

        # 行數 = 換行符 + 1
        n_lines = text.count("\n") + 1

        # 字級：取第一個 run 的字級
        font_size = 12  # 預設
        for para in textbox.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_size = int(run.font.size.pt)
                    break
            if font_size != 12:
                break

        return n_lines, font_size
