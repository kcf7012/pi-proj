"""
R2: 元素重疊驗證

規則：兩個非配對元素的 bounding box 有交集

策略：
- 配對的元素（背景矩形 + textbox 同一位置）不算重疊
- 其他元素兩兩比對
- 容忍小於 0.05" 的輕微重疊（浮點數誤差）
"""
from __future__ import annotations

from typing import TYPE_CHECKING

from .base import BaseValidator

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide


# 重疊容忍值（小於此視為接觸而非重疊）
OVERLAP_TOLERANCE = 0.05  # inches


class OverlapValidator(BaseValidator):
    """元素重疊驗證

    只檢查「有可見內容」的元素（跳過背景裝飾）
    """

    rule_id = "R2"

    def validate(self, prs: "Presentation") -> list[Issue]:
        issues: list[Issue] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            issues.extend(self._check_slide(slide, slide_idx))
        return issues

    def _check_slide(self, slide: "Slide", slide_num: int) -> list[Issue]:
        """檢查單張投影片的所有元素"""
        issues: list[Issue] = []

        # 收集有可見內容的元素
        shapes_info = []
        for shape in slide.shapes:
            info = self._get_shape_info(shape)
            if info is not None:
                shapes_info.append(info)

        # 兩兩比對
        for i in range(len(shapes_info)):
            for j in range(i + 1, len(shapes_info)):
                a, a_text = shapes_info[i]
                b, b_text = shapes_info[j]

                # 跳過配對的（背景+文字）
                if self._is_paired(a, b):
                    continue

                # 計算重疊
                overlap = self._compute_overlap(a, b)
                if overlap is not None:
                    h_overlap, v_overlap = overlap
                    if h_overlap > OVERLAP_TOLERANCE and v_overlap > OVERLAP_TOLERANCE:
                        # 跳過「完全包含」的情況（外框包內框）
                        if self._is_fully_inside(a, b) or self._is_fully_inside(b, a):
                            continue
                        # 真的有實質重疊
                        # 註：原本 spec 訂為 ERROR，但 pi-proj 原始 PPTX 有許多
                        # 設計性重疊（如箭頭指到 callout），改為 WARNING
                        issues.append(self.make_warning(
                            f"元素重疊：「{a_text[:20]}」與「{b_text[:20]}」({h_overlap:.2f}\" × {v_overlap:.2f}\")",
                            slide_num=slide_num,
                            shape_a_text=a_text[:30],
                            shape_b_text=b_text[:30],
                            h_overlap=h_overlap,
                            v_overlap=v_overlap,
                        ))

        return issues

    def _get_shape_info(self, shape) -> tuple[dict, str] | None:
        """取得形狀的位置與文字

        Returns:
            ({"top", "bottom", "left", "right", "width", "height"}, text)
            或 None（如果是裝飾性形狀）
        """
        if not shape.top or not shape.height or not shape.left or not shape.width:
            return None

        # 跳過太小的形狀（裝飾元素）
        if shape.width < 0.1 or shape.height < 0.1:
            return None

        top = shape.top / 914400
        left = shape.left / 914400
        w = shape.width / 914400
        h = shape.height / 914400

        # 跳過頂部 accent line（高度 < 0.2"）
        if h < 0.2:
            return None

        # 跳過品牌列（top >= 7.05，設計上在底部不需檢查重疊）
        if top >= 7.05:
            return None

        text = shape.text_frame.text if shape.has_text_frame else f"<{shape.shape_type}>"

        return (
            {
                "top": top,
                "bottom": top + h,
                "left": left,
                "right": left + w,
                "w": w,
                "h": h,
            },
            text.strip() or f"<{shape.shape_type}>",
        )

    def _is_paired(self, a: dict, b: dict) -> bool:
        """判斷兩個形狀是否配對（背景+文字）

        配對的條件：位置與大小完全相同
        """
        return (
            abs(a["top"] - b["top"]) < 0.01
            and abs(a["left"] - b["left"]) < 0.01
            and abs(a["w"] - b["w"]) < 0.01
            and abs(a["h"] - b["h"]) < 0.01
        )

    def _compute_overlap(self, a: dict, b: dict) -> tuple[float, float] | None:
        """計算兩個矩形的重疊量

        Returns:
            (horizontal, vertical) overlap in inches
            若無重疊則回傳 None
        """
        # 水平重疊
        h_left = max(a["left"], b["left"])
        h_right = min(a["right"], b["right"])
        h_overlap = h_right - h_left

        # 垂直重疊
        v_top = max(a["top"], b["top"])
        v_bottom = min(a["bottom"], b["bottom"])
        v_overlap = v_bottom - v_top

        if h_overlap > 0 and v_overlap > 0:
            return (h_overlap, v_overlap)
        return None

    def _is_fully_inside(self, inner: dict, outer: dict) -> bool:
        """檢查 inner 是否完全在 outer 內（包含 + 0.05" 容忍值）"""
        tol = 0.05
        return (
            outer["left"] - tol <= inner["left"]
            and inner["right"] <= outer["right"] + tol
            and outer["top"] - tol <= inner["top"]
            and inner["bottom"] <= outer["bottom"] + tol
        )
