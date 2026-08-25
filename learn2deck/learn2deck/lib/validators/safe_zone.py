"""
R3: 品牌列安全驗證

規則：任何內容元素 top + height > 7.0" → WARNING

品牌列在 y=7.1"，所有內容應在 7.0" 內。
"""
from __future__ import annotations

from typing import TYPE_CHECKING

from .base import BaseValidator

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide


# 品牌列安全邊界（內容必須在此 Y 之下）
SAFE_BOTTOM = 7.0  # inches

# 品牌列本身的高度（會出現在底部）
BRAND_BAR_HEIGHT = 0.4  # inches (7.1-7.5，給予充分容忍)


class SafeZoneValidator(BaseValidator):
    """品牌列安全驗證"""

    rule_id = "R3"

    def validate(self, prs: "Presentation") -> list[Issue]:
        issues: list[Issue] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            issues.extend(self._check_slide(slide, slide_idx))
        return issues

    def _check_slide(self, slide: "Slide", slide_num: int) -> list[Issue]:
        issues: list[Issue] = []

        for shape in slide.shapes:
            if not shape.top or not shape.height:
                continue

            bottom = (shape.top + shape.height) / 914400

            # 跳過品牌列本身（允許在 7.0-7.4 範圍）
            if bottom > SAFE_BOTTOM and bottom <= SAFE_BOTTOM + BRAND_BAR_HEIGHT + 0.05:
                continue

            if bottom > SAFE_BOTTOM:
                overflow = bottom - SAFE_BOTTOM
                text = ""
                if shape.has_text_frame:
                    text = shape.text_frame.text[:30].strip()
                issues.append(self.make_warning(
                    f"內容超出安全區：{text!r} 底部 y={bottom:.2f}\"（超出 {overflow:.2f}\"）",
                    slide_num=slide_num,
                    bottom=bottom,
                    overflow=overflow,
                ))

        return issues
