"""
learn2deck.validators - 品質驗證

公開 API：
- validate_deck(): 主入口
- ValidationReport: 報告物件
- Issue, Severity: 問題與嚴重度
- 個別 validator classes（給進階使用者）
"""
from __future__ import annotations

import subprocess
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation

from ..core import ValidationError
from .base import BaseValidator, Issue, Severity, ValidationReport
from .code_capacity import CodeCapacityValidator
from .file_format import FileFormatValidator
from .overlap import OverlapValidator
from .safe_zone import SafeZoneValidator

if TYPE_CHECKING:
    pass


# 所有內建 validator
BUILTIN_VALIDATORS: dict[str, type[BaseValidator]] = {
    "R1": CodeCapacityValidator,
    "R2": OverlapValidator,
    "R3": SafeZoneValidator,
    "R5": FileFormatValidator,
}


def validate_deck(
    pptx_path: str | Path,
    rules: list[str] | None = None,
    strict: bool = False,
) -> ValidationReport:
    """驗證已產出的 PPTX

    Args:
        pptx_path: PPTX 檔案路徑
        rules: 要跑的規則列表（預設全部）
        strict: 嚴格模式（WARNING 也算失敗）

    Returns:
        ValidationReport
    """
    path = Path(pptx_path)
    rules = rules or list(BUILTIN_VALIDATORS.keys())

    # 載入簡報
    try:
        prs = Presentation(str(path))
    except Exception as e:
        return ValidationReport(
            pptx_path=str(path),
            passed=False,
            issues=[
                Issue(
                    rule="R5",
                    severity=Severity.ERROR,
                    message=f"無法載入 PPTX：{e}",
                )
            ],
        )

    # 跑每個規則
    all_issues: list[Issue] = []

    # R5 特殊：接受檔案路徑
    if "R5" in rules:
        validator = FileFormatValidator()
        all_issues.extend(validator.validate(path))

    # R1, R2, R3 用 Presentation 物件
    for rule_id in rules:
        if rule_id == "R5":
            continue  # 已處理
        if rule_id not in BUILTIN_VALIDATORS:
            all_issues.append(Issue(
                rule=rule_id,
                severity=Severity.WARNING,
                message=f"未知的規則：{rule_id}",
            ))
            continue

        validator_cls = BUILTIN_VALIDATORS[rule_id]
        validator = validator_cls()
        all_issues.extend(validator.validate(prs))

    # 計算 passed
    has_error = any(i.severity == Severity.ERROR for i in all_issues)
    if strict:
        has_error = has_error or any(
            i.severity == Severity.WARNING for i in all_issues
        )

    # 統計
    stats = _compute_stats(prs, all_issues)

    return ValidationReport(
        pptx_path=str(path),
        passed=not has_error,
        issues=all_issues,
        stats=stats,
    )


def _compute_stats(prs: Presentation, issues: list[Issue]) -> dict:
    """統計資訊"""
    n_slides = len(prs.slides)
    n_code_blocks = 0
    n_tables = 0
    max_bottom = 0.0

    for slide in prs.slides:
        for shape in slide.shapes:
            # 統計 code 框（深色矩形）
            try:
                rgb = shape.fill.fore_color.rgb
                if rgb and max(int(rgb[0]), int(rgb[1]), int(rgb[2])) < 0x40:
                    n_code_blocks += 1
            except Exception:
                pass

            # 統計表格
            if shape.has_table:
                n_tables += 1

            # 最大底部
            if shape.top and shape.height:
                bottom = (shape.top + shape.height) / 914400
                max_bottom = max(max_bottom, bottom)

    return {
        "total_slides": n_slides,
        "code_blocks": n_code_blocks,
        "tables": n_tables,
        "max_content_bottom": round(max_bottom, 2),
        "n_errors": sum(1 for i in issues if i.severity == Severity.ERROR),
        "n_warnings": sum(1 for i in issues if i.severity == Severity.WARNING),
        "n_suggestions": sum(1 for i in issues if i.severity == Severity.SUGGESTION),
    }


# === CLI 整合輔助 ===

def print_report(report: ValidationReport) -> None:
    """用人類可讀格式印出報告"""
    from rich.console import Console
    from rich.table import Table

    console = Console()

    # 標題
    status = "✅ PASSED" if report.passed else "❌ FAILED"
    console.print(f"\n[bold]📋 Validation Report: {report.pptx_path}[/bold]")
    console.print(f"   Status: [bold]{status}[/bold]")
    console.print(f"   Total slides: {report.stats.get('total_slides', '?')}")
    console.print(f"   Code blocks: {report.stats.get('code_blocks', '?')}")
    console.print(f"   Tables: {report.stats.get('tables', '?')}")
    console.print(f"   Max content bottom: {report.stats.get('max_content_bottom', '?')}\"")

    if not report.issues:
        console.print("\n[green]✨ No issues found![/green]")
        return

    console.print(f"\n   Issues: {len(report.issues)}")

    # 分組
    errors = report.errors
    warnings = report.warnings
    suggestions = report.suggestions

    if errors:
        console.print(f"\n[red bold]Errors ({len(errors)}):[/red bold]")
        for issue in errors:
            loc = f"slide {issue.slide_num}" if issue.slide_num else "global"
            console.print(f"  [red]✗ [{issue.rule}] {loc}: {issue.message}[/red]")

    if warnings:
        console.print(f"\n[yellow bold]Warnings ({len(warnings)}):[/yellow bold]")
        for issue in warnings:
            loc = f"slide {issue.slide_num}" if issue.slide_num else "global"
            console.print(f"  [yellow]⚠ [{issue.rule}] {loc}: {issue.message}[/yellow]")

    if suggestions:
        console.print(f"\n[blue bold]Suggestions ({len(suggestions)}):[/blue bold]")
        for issue in suggestions:
            loc = f"slide {issue.slide_num}" if issue.slide_num else "global"
            console.print(f"  [blue]ℹ [{issue.rule}] {loc}: {issue.message}[/blue]")


__all__ = [
    "BaseValidator",
    "Issue",
    "Severity",
    "ValidationReport",
    "validate_deck",
    "print_report",
    "BUILTIN_VALIDATORS",
    "CodeCapacityValidator",
    "OverlapValidator",
    "SafeZoneValidator",
    "FileFormatValidator",
]
