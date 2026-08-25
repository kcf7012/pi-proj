"""
GridCardsBuilder - 網格卡片頁

對應 SlideType.GRID_CARDS
body schema: {
    "items": [
        {"icon": "🎯", "title": "...", "desc": "..."},
        ...
    ],
    "cols": 3  (optional, default 3)
}
"""
from __future__ import annotations

from pptx.util import Inches

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from ..core import MissingFieldError, SlideContent, SlideType
from ..pptx_helpers.shapes import _color, _font
from ..pptx_helpers import get_font_size
from .base import BaseBuilder


class GridCardsBuilder(BaseBuilder):
    """網格卡片 builder"""

    slide_type = SlideType.GRID_CARDS

    def build(
        self,
        slide: Slide,
        content: SlideContent,
        slide_num: int | None = None,
        total: int | None = None,
    ) -> None:
        # 1. 標題列
        self.add_title(slide, content, slide_num=slide_num, total=total)

        # 2. 卡片
        body = self.require_body(content, "items")
        items = body["items"]
        cols = body.get("cols", 3)

        if not isinstance(items, list) or len(items) == 0:
            raise MissingFieldError(
                f"GridCards '{content.title}' 的 items 應為非空 list"
            )

        # 計算版面
        n_items = len(items)
        rows = (n_items + cols - 1) // cols

        # 卡片尺寸（依欄數調整）
        if cols <= 3:
            card_w = 3.8
            card_h = 2.2
            v_gap = 0.3
            h_gap = 0.4
        else:
            card_w = 2.8
            card_h = 2.0
            v_gap = 0.2
            h_gap = 0.2

        grid_w = card_w * cols + h_gap * (cols - 1)
        start_x = (13.333 - grid_w) / 2
        start_y = 1.8

        for i, item in enumerate(items):
            if not isinstance(item, dict):
                continue
            row = i // cols
            col = i % cols
            x = start_x + col * (card_w + h_gap)
            y = start_y + row * (card_h + v_gap)
            self._draw_card(slide, item, x, y, card_w, card_h)

    def _draw_card(
        self,
        slide: Slide,
        item: dict,
        x: float,
        y: float,
        w: float,
        h: float,
    ) -> None:
        """畫單張卡片"""
        icon = item.get("icon", "")
        title = item.get("title", "")
        desc = item.get("desc", "")

        # 卡片背景
        if self.theme is not None:
            bg_color = _color(self.theme, "bg_gray", "#F3F0E9")
            border_color = _color(self.theme, "primary", "#C75A1A")
            title_color = _color(self.theme, "dark", "#2C2C2C")
            desc_color = _color(self.theme, "gray_text", "#6B6B6B")
        else:
            bg_color = RGBColor(0xF3, 0xF0, 0xE9)
            border_color = RGBColor(0xC7, 0x5A, 0x1A)
            title_color = RGBColor(0x2C, 0x2C, 0x2C)
            desc_color = RGBColor(0x6B, 0x6B, 0x6B)

        if self.theme is not None:
            title_font = _font(self.theme, "title", "Calibri")
            body_font = _font(self.theme, "body", "Calibri")
            title_size = get_font_size(self.theme, "body", 14)
            desc_size = 11
        else:
            title_font = "Calibri"
            body_font = "Calibri"
            title_size = 14
            desc_size = 11

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Inches(0.02)  # 1.5pt

        # 圖示
        if icon:
            from pptx.util import Pt
            icon_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 0.2), Inches(w), Inches(0.7)
            )
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = icon
            run.font.size = Pt(36)

        # 標題
        from pptx.util import Pt
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 1.0), Inches(w - 0.4), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = title_font
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color

        # 描述
        if desc:
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 1.5), Inches(w - 0.4), Inches(0.6)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = desc
            run.font.name = body_font
            run.font.size = Pt(desc_size)
            run.font.color.rgb = desc_color
