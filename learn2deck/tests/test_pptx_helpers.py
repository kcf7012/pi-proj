"""
learn2deck.pptx_helpers 單元測試

測試：
- layout 計算函式
- shapes 基礎形狀
- pages 特殊頁面
- 向後相容性（無 theme 也能用）
"""
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from learn2deck.lib.core import Theme
from learn2deck.lib.pptx_helpers import (
    # layout
    get_layout_value, get_slide_width, get_slide_height,
    safe_top, safe_bottom, safe_height, content_left, content_width,
    get_font_size, estimate_line_height,
    DEFAULT_LAYOUT, LINE_HEIGHTS,
    # shapes
    new_presentation, add_blank_slide, set_slide_bg, add_title_bar,
    add_text_block, add_bullet_list, add_code_block, add_callout,
    add_comparison_table, add_flow_box, add_arrow,
    # pages
    add_cover_slide, add_section_divider, add_summary_slide,
    add_two_column_compare,
)


# === Layout ===
class TestLayout:
    def test_default_slide_dimensions(self):
        """預設 16:9 簡報尺寸"""
        assert get_slide_width() == Inches(13.333)
        assert get_slide_height() == Inches(7.5)

    def test_safe_zones(self):
        """安全區域"""
        assert safe_top() == 1.3
        assert safe_bottom() == 7.0
        assert safe_height() == 5.7  # 7.0 - 1.3

    def test_content_dimensions(self):
        """內容區域"""
        assert content_left() == 0.5
        assert content_width() == 12.333  # 13.333 - 0.5 - 0.5

    def test_with_theme_override(self, tmp_path):
        """theme 覆寫版面參數"""
        yaml_content = """
name: custom
layout:
  slide_width: 10.0
  slide_height: 7.5
  content_top: 1.0
  content_bottom: 6.5
"""
        yaml_path = tmp_path / "custom.yaml"
        yaml_path.write_text(yaml_content)
        theme = Theme.from_yaml(yaml_path)

        assert get_slide_width(theme) == Inches(10.0)
        assert get_slide_height(theme) == Inches(7.5)
        assert safe_top(theme) == 1.0
        assert safe_bottom(theme) == 6.5
        assert safe_height(theme) == 5.5

    def test_get_font_size_with_theme(self, tmp_path):
        """theme 字級覆寫"""
        yaml_content = """
name: custom
font_sizes:
  cover_title: 60
  body: 16
"""
        yaml_path = tmp_path / "custom.yaml"
        yaml_path.write_text(yaml_content)
        theme = Theme.from_yaml(yaml_path)

        assert get_font_size(theme, "cover_title", 50) == 60
        assert get_font_size(theme, "body", 14) == 16
        assert get_font_size(theme, "nonexistent", 14) == 14  # default

    def test_estimate_line_height(self):
        """行高估算"""
        assert estimate_line_height(12) == 0.20
        assert estimate_line_height(10) == 0.17
        assert estimate_line_height(8) == 0.14
        assert estimate_line_height(7) == 0.13

    def test_estimate_line_height_fallback(self):
        """未知字級用公式 fallback"""
        h = estimate_line_height(15)  # 沒有 15 在 LINE_HEIGHTS
        # 15 / 72 * 1.2 = 0.25
        assert abs(h - 0.25) < 0.01


# === Shapes: 簡報與投影片 ===
class TestShapesBasics:
    def test_new_presentation_16_9(self):
        """預設 16:9 比例"""
        prs = new_presentation()
        assert prs.slide_width / prs.slide_height == pytest.approx(16/9, rel=0.01)

    def test_add_blank_slide(self):
        """新增空白頁"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        assert slide in prs.slides
        assert len(prs.slides) == 1

    def test_set_slide_bg(self):
        """設定背景色"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        set_slide_bg(slide)  # 用預設
        # 沒有直接 API 檢查，但執行不報錯就算通過


# === Shapes: 標題列 ===
class TestTitleBar:
    def test_basic_title(self):
        """基本標題列"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_title_bar(slide, "Test Title", slide_num=1, total=5)
        # 至少要有 3 個 shape：line + title + page num
        assert len(slide.shapes) >= 3

    def test_title_with_subtitle_and_source(self):
        """完整標題列"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_title_bar(
            slide, "Title", "Subtitle",
            slide_num=1, total=5, source="test.md"
        )
        # line + title + subtitle + page + source + brand = 6
        assert len(slide.shapes) >= 6

    def test_title_with_theme(self, tmp_path):
        """theme 套用到標題列"""
        yaml_content = """
name: custom
colors:
  primary: "#FF00FF"
fonts:
  title: "Arial"
  body: "Arial"
"""
        yaml_path = tmp_path / "custom.yaml"
        yaml_path.write_text(yaml_content)
        theme = Theme.from_yaml(yaml_path)

        prs = new_presentation(theme)
        slide = add_blank_slide(prs)
        add_title_bar(slide, "Title", theme=theme)
        # 應該用 theme 的顏色（#FF00FF）
        # 視覺驗證比較難，但執行不報錯即可
        assert len(slide.shapes) >= 3


# === Shapes: 文字 ===
class TestTextBlocks:
    def test_text_block(self):
        """純文字區塊"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_text_block(slide, "Hello", Inches(1), Inches(2), Inches(5), Inches(1))
        assert len(slide.shapes) >= 1

    def test_bullet_list_strings(self):
        """字串清單"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_bullet_list(
            slide, ["Item 1", "Item 2", "Item 3"],
            Inches(1), Inches(2), Inches(5), Inches(2)
        )
        # 1 個 textbox
        assert len(slide.shapes) >= 1

    def test_bullet_list_tuples(self):
        """階層式 tuples"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_bullet_list(
            slide, [(0, "Top"), (1, "  Sub 1"), (1, "  Sub 2")],
            Inches(1), Inches(2), Inches(5), Inches(2)
        )
        assert len(slide.shapes) >= 1

    def test_bullet_list_invalid_raises(self):
        """無效的 items 格式應拋出 ValueError"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        with pytest.raises(ValueError):
            add_bullet_list(slide, [123, 456], Inches(1), Inches(2), Inches(5), Inches(2))


# === Shapes: 程式碼 ===
class TestCodeBlock:
    def test_code_block_creates_pair(self):
        """code 框會建立 1 個矩形（背景）+ 1 個 textbox"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_code_block(
            slide, "print('hello')",
            Inches(1), Inches(2), Inches(5), Inches(1)
        )
        # 至少 2 個 shape
        assert len(slide.shapes) >= 2

    def test_code_block_multiline(self):
        """多行程式碼"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        code = "line 1\nline 2\nline 3"
        add_code_block(slide, code, Inches(1), Inches(2), Inches(5), Inches(1))
        assert len(slide.shapes) >= 2


# === Shapes: 提示框 ===
class TestCallout:
    def test_callout_creates(self):
        """提示框建立"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_callout(slide, "Important!", Inches(1), Inches(2), Inches(5), Inches(1))
        # 至少有 1 個 shape
        assert len(slide.shapes) >= 1

    def test_callout_no_icon(self):
        """無 icon 的 callout"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_callout(
            slide, "Plain text", Inches(1), Inches(2), Inches(5), Inches(1),
            icon=""
        )
        assert len(slide.shapes) >= 1


# === Shapes: 表格 ===
class TestTable:
    def test_table_basic(self):
        """基本表格"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_comparison_table(
            slide,
            headers=["Col1", "Col2"],
            rows=[["a", "b"], ["c", "d"]],
            left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(2)
        )
        # 表格也算 1 個 shape
        assert len(slide.shapes) >= 1


# === Shapes: 流程圖 ===
class TestFlowBox:
    def test_flow_box(self):
        """流程圖方塊"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        box = add_flow_box(
            slide, "Step 1",
            Inches(1), Inches(2), Inches(2), Inches(1)
        )
        assert box is not None
        assert len(slide.shapes) >= 1

    def test_arrow(self):
        """箭頭"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_arrow(
            slide, Inches(1), Inches(2),
            Inches(3), Inches(3)
        )
        assert len(slide.shapes) >= 1


# === Pages: 特殊頁面 ===
class TestPages:
    def test_cover_slide(self):
        """封面頁"""
        prs = new_presentation()
        slide = add_cover_slide(prs, "Title", "Subtitle", tag="# Tag")
        assert slide in prs.slides

    def test_section_divider(self):
        """章節分隔頁"""
        prs = new_presentation()
        slide = add_section_divider(prs, "Part 1", "Title", "Subtitle")
        assert slide in prs.slides

    def test_summary_slide(self):
        """重點回顧頁"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_summary_slide(
            slide,
            title="Summary",
            key_points=["Point 1", "Point 2"],
            next_steps=["Step 1", "Step 2"],
            source="test.md"
        )
        assert len(slide.shapes) >= 1

    def test_summary_no_next_steps(self):
        """沒有 next_steps 的 summary"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_summary_slide(
            slide,
            key_points=["Point 1"]
        )
        assert len(slide.shapes) >= 1

    def test_two_column_compare(self):
        """雙欄對比"""
        prs = new_presentation()
        slide = add_blank_slide(prs)
        add_two_column_compare(
            slide,
            left_title="Pros",
            left_content=["Good 1", "Good 2"],
            right_title="Cons",
            right_content=["Bad 1", "Bad 2"],
        )
        # 至少要有 4 個 shape（2 個方塊 + 2 個標題列）
        assert len(slide.shapes) >= 4


# === 向後相容性 ===
class TestBackwardCompatibility:
    def test_no_theme_works(self):
        """不傳 theme 也能用（向後相容於 _pptx_helpers.py）"""
        prs = new_presentation()  # 沒 theme
        slide = add_blank_slide(prs)
        add_title_bar(slide, "Title", slide_num=1, total=5)
        add_text_block(slide, "Text", Inches(1), Inches(2), Inches(5), Inches(1))
        add_code_block(slide, "code", Inches(1), Inches(3), Inches(5), Inches(1))
        add_callout(slide, "Tip", Inches(1), Inches(4), Inches(5), Inches(1))
        add_cover_slide(prs, "Cover", "Sub")
        add_section_divider(prs, "Part 1", "Title", "Sub")
        add_two_column_compare(
            slide, "L", ["a"], "R", ["b"]
        )
        # 能執行到這裡就算成功
        assert len(prs.slides) >= 3

    def test_save_and_reload(self, tmp_path):
        """產出 + 重新讀取"""
        prs = new_presentation()
        add_cover_slide(prs, "My Deck", "Subtitle")
        add_section_divider(prs, "Part 1", "Title", "Subtitle")

        out_path = tmp_path / "test.pptx"
        prs.save(str(out_path))

        # 重新讀取
        reloaded = Presentation(str(out_path))
        assert len(reloaded.slides) == 2
        assert reloaded.slide_width / reloaded.slide_height == pytest.approx(16/9, rel=0.01)
