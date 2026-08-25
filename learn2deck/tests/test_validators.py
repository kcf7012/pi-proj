"""
learn2deck.validators 單元測試

測試 4 條規則與 validate_deck() 主入口。
"""
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation

from learn2deck.lib.parsers import parse_content
from learn2deck.lib.builders import build_full_deck
from learn2deck.lib.validators import (
    validate_deck, print_report,
    Issue, Severity, ValidationReport,
    CodeCapacityValidator, OverlapValidator,
    SafeZoneValidator, FileFormatValidator,
    BUILTIN_VALIDATORS,
)


# === Base ===
class TestBase:
    def test_severity_enum(self):
        """嚴重度 enum 包含 3 個值"""
        assert Severity.ERROR.value == "error"
        assert Severity.WARNING.value == "warning"
        assert Severity.SUGGESTION.value == "suggestion"

    def test_issue_str(self):
        """Issue 可字串化"""
        i = Issue(rule="R1", severity=Severity.ERROR, slide_num=5, message="Test")
        s = str(i)
        assert "R1" in s
        assert "error" in s
        assert "slide 5" in s
        assert "Test" in s

    def test_issue_str_global(self):
        """Issue 沒指定 slide_num 時顯示 global"""
        i = Issue(rule="R5", severity=Severity.ERROR, message="Bad file")
        assert "global" in str(i)

    def test_validation_report_properties(self):
        """Report 自動分類 issues"""
        report = ValidationReport(
            pptx_path="test.pptx",
            passed=True,
            issues=[
                Issue(rule="R1", severity=Severity.ERROR, message="e1"),
                Issue(rule="R2", severity=Severity.WARNING, message="w1"),
                Issue(rule="R3", severity=Severity.SUGGESTION, message="s1"),
            ],
        )
        assert len(report.errors) == 1
        assert len(report.warnings) == 1
        assert len(report.suggestions) == 1
        assert report.errors[0].message == "e1"

    def test_validation_report_to_dict(self):
        """Report 可序列化為 dict"""
        report = ValidationReport(
            pptx_path="test.pptx",
            passed=True,
            issues=[Issue(rule="R1", severity=Severity.ERROR, message="e1")],
        )
        d = report.to_dict()
        assert d["pptx_path"] == "test.pptx"
        assert d["passed"] is True
        assert d["issues"][0]["rule"] == "R1"
        assert d["issues"][0]["severity"] == "error"

    def test_validation_report_print_human(self):
        """Report 有人類可讀格式"""
        report = ValidationReport(
            pptx_path="test.pptx",
            passed=True,
            issues=[],
        )
        text = report.print_human()
        assert "test.pptx" in text
        assert "PASSED" in text


# === R1: Code Capacity ===
class TestR1CodeCapacity:
    def test_built_in_validator_registry(self):
        """R1 在內建清單中"""
        assert "R1" in BUILTIN_VALIDATORS
        assert BUILTIN_VALIDATORS["R1"] is CodeCapacityValidator

    def test_validator_id(self):
        """Validator 設定正確的 rule_id"""
        v = CodeCapacityValidator()
        assert v.rule_id == "R1"

    def test_no_code_blocks_no_issues(self):
        """沒有 code 框的簡報不應有 R1 問題"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        v = CodeCapacityValidator()
        issues = v.validate(prs)
        assert issues == []

    def test_oversized_code_block_detected(self, tmp_path):
        """code 框裝不下會被偵測"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 取得 slide 的 shapes list（不要保存 bg 參考，因為 fill.solid() 會重換 XML）
        shapes = slide.shapes

        # 建立一個很小的 code 框
        bg = shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.5), Inches(2.0), Inches(8), Inches(0.5)  # 太矮
        )
        bg.fill.solid()
        # 設定深色
        bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)

        # 加配對的 textbox
        tb = shapes.add_textbox(
            Inches(0.5), Inches(2.0), Inches(8), Inches(0.5)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        # 加很多行
        for i in range(10):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"line {i+1}"
            run.font.size = Pt(12)

        # 重取 shapes（因為上面 fill 動作可能重換了）
        # 重新驗證取最新的 shapes
        v = CodeCapacityValidator()
        issues = v.validate(prs)
        assert len(issues) >= 1, f"預期有 R1 問題，但沒有找到"
        assert issues[0].rule == "R1"
        assert "裝不下" in issues[0].message

    def test_fitting_code_block_no_issue(self):
        """code 框裝得下不應有問題"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.shapes.add_shape(
            1, Inches(0.5), Inches(2.0), Inches(8), Inches(2.0)
        )
        from pptx.dml.color import RGBColor
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)

        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.0), Inches(8), Inches(2.0)
        )
        for i in range(5):
            p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"line {i+1}"
            run.font.size = Pt(12)

        v = CodeCapacityValidator()
        issues = v.validate(prs)
        assert issues == []


# === R2: Overlap ===
class TestR2Overlap:
    def test_built_in(self):
        assert "R2" in BUILTIN_VALIDATORS
        assert BUILTIN_VALIDATORS["R2"] is OverlapValidator

    def test_validator_id(self):
        v = OverlapValidator()
        assert v.rule_id == "R2"

    def test_no_overlap(self):
        """無重疊的投影片不應有問題"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 兩個分離的方塊
        b1 = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(3), Inches(1))
        b2 = slide.shapes.add_shape(1, Inches(4), Inches(1), Inches(3), Inches(1))

        v = OverlapValidator()
        issues = v.validate(prs)
        assert issues == []

    def test_overlap_detected(self):
        """明顯重疊應被偵測"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 兩個重疊的方塊（不是嵌套關係）
        b1 = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(3), Inches(1))
        b1.text_frame.text = "Block A"
        b2 = slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(3), Inches(1))
        b2.text_frame.text = "Block B"

        v = OverlapValidator()
        issues = v.validate(prs)
        assert len(issues) >= 1
        assert issues[0].rule == "R2"

    def test_nested_not_flagged(self):
        """嵌套關係（外框包含內框）不應被當作重疊"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 外框
        outer = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(6), Inches(3))
        outer.text_frame.text = "Outer"
        # 內框（完全在外框內）
        inner = slide.shapes.add_shape(1, Inches(1), Inches(1.5), Inches(4), Inches(1))
        inner.text_frame.text = "Inner"

        v = OverlapValidator()
        issues = v.validate(prs)
        # 嵌套不應被視為重疊
        assert issues == []

    def test_brand_bar_skipped(self):
        """品牌列（top >= 7.05）不應被檢查"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 品牌列
        b = slide.shapes.add_shape(1, Inches(0.5), Inches(7.1), Inches(5), Inches(0.3))
        b.text_frame.text = "Brand"
        # 故意重疊的方塊
        b2 = slide.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(5), Inches(0.3))
        b2.text_frame.text = "Above"

        v = OverlapValidator()
        issues = v.validate(prs)
        # 兩個都在品牌列區域，應被跳過
        assert issues == []


# === R3: Safe Zone ===
class TestR3SafeZone:
    def test_built_in(self):
        assert "R3" in BUILTIN_VALIDATORS
        assert BUILTIN_VALIDATORS["R3"] is SafeZoneValidator

    def test_validator_id(self):
        v = SafeZoneValidator()
        assert v.rule_id == "R3"

    def test_content_within_safe_zone(self):
        """內容在 7.0" 內不應有問題"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        b = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(5), Inches(2))
        b.text_frame.text = "OK"

        v = SafeZoneValidator()
        issues = v.validate(prs)
        assert issues == []

    def test_content_beyond_safe_zone(self):
        """內容超出 7.0" 應被警告"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        b = slide.shapes.add_shape(1, Inches(0.5), Inches(6.5), Inches(5), Inches(1))
        b.text_frame.text = "Overflow"

        v = SafeZoneValidator()
        issues = v.validate(prs)
        assert len(issues) >= 1
        assert issues[0].rule == "R3"
        assert "超出" in issues[0].message

    def test_brand_bar_allowed(self):
        """品牌列（top=7.1, h=0.3）應被允許"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        b = slide.shapes.add_shape(1, Inches(0.5), Inches(7.1), Inches(5), Inches(0.3))
        b.text_frame.text = "Brand bar"

        v = SafeZoneValidator()
        issues = v.validate(prs)
        assert issues == []


# === R5: File Format ===
class TestR5FileFormat:
    def test_built_in(self):
        assert "R5" in BUILTIN_VALIDATORS
        assert BUILTIN_VALIDATORS["R5"] is FileFormatValidator

    def test_validator_id(self):
        v = FileFormatValidator()
        assert v.rule_id == "R5"

    def test_valid_pptx_file(self, tmp_path):
        """有效的 PPTX 檔案不應有問題"""
        # 建立有效的 PPTX
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        out = tmp_path / "valid.pptx"
        prs.save(str(out))

        v = FileFormatValidator()
        issues = v.validate(out)
        assert issues == []

    def test_wrong_extension(self, tmp_path):
        """副檔名錯誤"""
        out = tmp_path / "wrong.txt"
        out.write_text("hello")

        v = FileFormatValidator()
        issues = v.validate(out)
        assert len(issues) >= 1
        assert "副檔名" in issues[0].message

    def test_nonexistent_file(self, tmp_path):
        """不存在的檔案"""
        out = tmp_path / "doesnt.pptx"

        v = FileFormatValidator()
        issues = v.validate(out)
        assert len(issues) >= 1

    def test_corrupted_zip(self, tmp_path):
        """損壞的 ZIP 檔"""
        out = tmp_path / "bad.pptx"
        out.write_text("This is not a valid PPTX")

        v = FileFormatValidator()
        issues = v.validate(out)
        assert len(issues) >= 1

    def test_validates_presentation_object(self):
        """也接受 Presentation 物件"""
        prs = Presentation()
        v = FileFormatValidator()
        issues = v.validate(prs)
        # 空 PPTX 會有 warning
        assert all(i.rule == "R5" for i in issues)


# === validate_deck 主入口 ===
class TestValidateDeck:
    @pytest.fixture
    def valid_pptx(self, tmp_path):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        out = tmp_path / "test.pptx"
        prs.save(str(out))
        return out

    def test_validate_valid_file_passes(self, valid_pptx):
        """有效檔案通過驗證"""
        report = validate_deck(valid_pptx)
        assert report.passed is True
        assert report.pptx_path == str(valid_pptx)
        assert "total_slides" in report.stats

    def test_validate_nonexistent_raises(self):
        """不存在的檔案不應崩潰"""
        report = validate_deck("/tmp/nonexistent.pptx")
        assert report.passed is False
        assert len(report.issues) >= 1

    def test_validate_specific_rules(self, valid_pptx):
        """指定只跑某些規則"""
        report = validate_deck(valid_pptx, rules=["R5"])
        # 只有 R5 應該有結果
        for issue in report.issues:
            assert issue.rule in ["R5", "unknown"]

    def test_validate_unknown_rule(self, valid_pptx):
        """未知規則會給 warning"""
        report = validate_deck(valid_pptx, rules=["R99"])
        assert any(i.rule == "R99" for i in report.issues)

    def test_strict_mode_treats_warnings_as_errors(self, tmp_path):
        """strict 模式把 warning 當 error"""
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        b = slide.shapes.add_shape(1, Inches(0.5), Inches(6.5), Inches(5), Inches(1))
        b.text_frame.text = "Warning"  # 觸發 R3 warning

        out = tmp_path / "test.pptx"
        prs.save(str(out))

        # 非 strict：passed=True
        report = validate_deck(out, strict=False)
        assert report.passed is True  # 只有 warning
        assert report.stats["n_warnings"] >= 1

        # strict：passed=False
        report = validate_deck(out, strict=True)
        assert report.passed is False


# === End-to-end 整合測試 ===
class TestEndToEnd:
    """用 parser + builder + validator 完整跑一次"""

    @pytest.fixture
    def md_dir(self):
        return Path("/home/elan/pi-proj")

    @pytest.mark.parametrize("filename", [
        "00-claude-code-plugins-series.md",
        "02-plugins.md",
        "04-skills.md",
        "07-discover-plugins.md",
    ])
    def test_parse_build_validate_cycle(self, md_dir, filename, tmp_path):
        """parse → build → validate 完整 cycle"""
        deck = parse_content(md_dir / filename)
        out = tmp_path / f"{filename.replace('.md', '.pptx')}"
        build_full_deck(deck, str(out))

        report = validate_deck(out)
        # 至少有 0 個 error（warning 可有可無）
        assert report.stats["n_errors"] == 0, \
            f"{filename} 有 {report.stats['n_errors']} 個錯誤: {report.issues}"
        # 至少有 1 張投影片
        assert report.stats["total_slides"] >= 1
