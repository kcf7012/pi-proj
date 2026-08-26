"""
learn2deck.builders 單元測試

測試 9 種 builder 與 build_full_deck() 端到端整合。
"""
import pytest

from learn2deck.lib.builders import (
    BaseBuilder, build_slide, build_full_deck,
    CoverBuilder, ObjectivesBuilder, SectionDividerBuilder,
    TitleContentBuilder, TitleTableBuilder, TitleCodeBuilder,
    TwoColumnBuilder, GridCardsBuilder, SummaryBuilder,
)
from learn2deck.lib.core import (
    DeckSpec, MissingFieldError, SlideContent, SlideType, BuildError,
    load_theme,
)
from learn2deck.lib.pptx_helpers import new_presentation, add_blank_slide


@pytest.fixture
def prs():
    """空的 16:9 簡報"""
    return new_presentation()


@pytest.fixture
def theme():
    return load_theme("claude-orange")


# === 9 種 Builder 基本存在性 ===
class TestBuildersExist:
    def test_all_9_builders_importable(self):
        """9 種 builder 都能 import"""
        builders = [
            CoverBuilder, ObjectivesBuilder, SectionDividerBuilder,
            TitleContentBuilder, TitleTableBuilder, TitleCodeBuilder,
            TwoColumnBuilder, GridCardsBuilder, SummaryBuilder,
        ]
        for b in builders:
            assert issubclass(b, BaseBuilder)

    def test_all_builders_have_slide_type(self):
        """每個 builder 都有對應的 slide_type"""
        mapping = {
            CoverBuilder: SlideType.COVER,
            ObjectivesBuilder: SlideType.OBJECTIVES,
            SectionDividerBuilder: SlideType.SECTION_DIVIDER,
            TitleContentBuilder: SlideType.TITLE_CONTENT,
            TitleTableBuilder: SlideType.TITLE_TABLE,
            TitleCodeBuilder: SlideType.TITLE_CODE,
            TwoColumnBuilder: SlideType.TWO_COLUMN,
            GridCardsBuilder: SlideType.GRID_CARDS,
            SummaryBuilder: SlideType.SUMMARY,
        }
        for builder_cls, expected_type in mapping.items():
            assert builder_cls.slide_type == expected_type


# === 個別 Builder 測試 ===
class TestSummaryBuilder:
    def test_summary_with_key_points(self, prs, theme):
        """summary 有 key_points"""
        slide = add_blank_slide(prs)
        builder = SummaryBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.SUMMARY,
            title="Summary",
            body={"key_points": ["p1", "p2"]},
        )
        builder.build(slide, content)
        assert len(slide.shapes) > 0

    def test_summary_missing_body_raises(self, prs, theme):
        """summary 沒有 body 應拋出 MissingFieldError"""
        slide = add_blank_slide(prs)
        builder = SummaryBuilder(theme=theme)
        content = SlideContent(type=SlideType.SUMMARY, title="Empty")
        with pytest.raises(MissingFieldError):
            builder.build(slide, content)


class TestTitleContentBuilder:
    def test_with_bullets(self, prs, theme):
        """bullets 模式"""
        slide = add_blank_slide(prs)
        builder = TitleContentBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Title",
            subtitle="Subtitle",
            body={"items": ["a", "b", "c"]},
        )
        builder.build(slide, content)
        # 標題列 + bullets
        assert len(slide.shapes) > 0

    def test_with_text(self, prs, theme):
        """純文字模式"""
        slide = add_blank_slide(prs)
        builder = TitleContentBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Title",
            body={"text": "Hello world"},
        )
        builder.build(slide, content)
        assert len(slide.shapes) > 0

    def test_missing_text_and_items_raises(self, prs, theme):
        """既沒 text 也沒 items 應拋出 MissingFieldError"""
        slide = add_blank_slide(prs)
        builder = TitleContentBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Empty",
            body={},
        )
        with pytest.raises(MissingFieldError):
            builder.build(slide, content)


class TestTitleTableBuilder:
    def test_basic_table(self, prs, theme):
        """基本表格"""
        slide = add_blank_slide(prs)
        builder = TitleTableBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_TABLE,
            title="Table",
            body={
                "headers": ["A", "B"],
                "rows": [["1", "2"], ["3", "4"]],
            },
        )
        builder.build(slide, content)
        # 標題列 + 表格
        assert len(slide.shapes) > 0

    def test_missing_headers_raises(self, prs, theme):
        """缺少 headers 應拋出"""
        slide = add_blank_slide(prs)
        builder = TitleTableBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_TABLE,
            title="Empty",
            body={"rows": []},
        )
        with pytest.raises(MissingFieldError):
            builder.build(slide, content)


class TestTitleCodeBuilder:
    def test_basic_code(self, prs, theme):
        """基本程式碼"""
        slide = add_blank_slide(prs)
        builder = TitleCodeBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_CODE,
            title="Code",
            body={"code": "print('hi')", "language": "python"},
        )
        builder.build(slide, content)
        # 標題列 + 背景 + textbox
        assert len(slide.shapes) >= 2

    def test_multiline_code(self, prs, theme):
        """多行程式碼"""
        slide = add_blank_slide(prs)
        builder = TitleCodeBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_CODE,
            title="Multi",
            body={"code": "line 1\nline 2\nline 3"},
        )
        builder.build(slide, content)
        assert len(slide.shapes) >= 2

    def test_missing_code_raises(self, prs, theme):
        """缺少 code 欄位應拋出"""
        slide = add_blank_slide(prs)
        builder = TitleCodeBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TITLE_CODE,
            title="Empty",
            body={"language": "python"},
        )
        with pytest.raises(MissingFieldError):
            builder.build(slide, content)


class TestTwoColumnBuilder:
    def test_basic_two_column(self, prs, theme):
        """基本雙欄"""
        slide = add_blank_slide(prs)
        builder = TwoColumnBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TWO_COLUMN,
            title="Compare",
            body={
                "left": {"title": "Pros", "items": ["a"]},
                "right": {"title": "Cons", "items": ["b"]},
            },
        )
        builder.build(slide, content)
        # 標題列 + 2 個方塊 + 2 個標題列 + 2 個 bullets
        assert len(slide.shapes) >= 4

    def test_missing_left_or_right_raises(self, prs, theme):
        """缺少 left 或 right 應拋出"""
        slide = add_blank_slide(prs)
        builder = TwoColumnBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.TWO_COLUMN,
            title="Empty",
            body={"left": {"title": "P", "items": []}},
        )
        with pytest.raises(MissingFieldError):
            builder.build(slide, content)


class TestGridCardsBuilder:
    def test_with_3_items(self, prs, theme):
        """3 個項目（預設 3 欄）"""
        slide = add_blank_slide(prs)
        builder = GridCardsBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.GRID_CARDS,
            title="Cards",
            body={
                "items": [
                    {"icon": "🎯", "title": "A", "desc": "a"},
                    {"icon": "📦", "title": "B", "desc": "b"},
                    {"icon": "🧪", "title": "C", "desc": "c"},
                ]
            },
        )
        builder.build(slide, content)
        # 標題列 + 3 張卡片（每張 = 1 shape）
        assert len(slide.shapes) >= 3

    def test_with_6_items_2_rows(self, prs, theme):
        """6 個項目（2 列 3 欄）"""
        slide = add_blank_slide(prs)
        builder = GridCardsBuilder(theme=theme)
        items = [
            {"icon": str(i), "title": f"T{i}", "desc": f"d{i}"}
            for i in range(1, 7)
        ]
        content = SlideContent(
            type=SlideType.GRID_CARDS,
            title="Many",
            body={"items": items},
        )
        builder.build(slide, content)
        # 標題列 + 6 張卡片
        assert len(slide.shapes) >= 6

    def test_missing_items_raises(self, prs, theme):
        """缺少 items 應拋出"""
        slide = add_blank_slide(prs)
        builder = GridCardsBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.GRID_CARDS,
            title="Empty",
            body={"cols": 3},
        )
        with pytest.raises(MissingFieldError):
            builder.build(slide, content)


class TestObjectivesBuilder:
    def test_inherits_from_grid_cards(self):
        """ObjectivesBuilder 繼承自 GridCardsBuilder"""
        assert issubclass(ObjectivesBuilder, GridCardsBuilder)

    def test_default_title(self, prs, theme):
        """沒給 title 時用預設「本章你會學到」"""
        slide = add_blank_slide(prs)
        builder = ObjectivesBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.OBJECTIVES,
            title="",  # 空標題
            body={
                "items": [
                    {"icon": "🎯", "title": "A", "desc": "a"},
                ]
            },
        )
        builder.build(slide, content)
        # 標題應該被自動填入
        assert content.title == "本章你會學到"


class TestCoverBuilder:
    def test_build_raises_build_error(self, prs, theme):
        """Cover 必須透過 build_full_deck() 處理"""
        slide = add_blank_slide(prs)
        builder = CoverBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.COVER,
            title="Test",
            body={"tag": "# T"},
        )
        with pytest.raises(BuildError):
            builder.build(slide, content)


class TestSectionDividerBuilder:
    def test_build_raises_build_error(self, prs, theme):
        """Section divider 必須透過 build_full_deck() 處理"""
        slide = add_blank_slide(prs)
        builder = SectionDividerBuilder(theme=theme)
        content = SlideContent(
            type=SlideType.SECTION_DIVIDER,
            title="Title",
            body={"section_num": "P1"},
        )
        with pytest.raises(BuildError):
            builder.build(slide, content)


# === build_full_deck 端到端 ===
class TestBuildFullDeck:
    def test_build_simple_deck(self, theme, tmp_path):
        """建立一個簡單的 3 張簡報"""
        deck = DeckSpec(title="Test", subtitle="Sub")
        deck.add_slide(SlideContent(
            type=SlideType.COVER,
            title="Test",
            subtitle="Sub",
            body={"tag": "# T"},
        ))
        deck.add_slide(SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Slide 2",
            body={"text": "Hello"},
        ))
        deck.add_slide(SlideContent(
            type=SlideType.SUMMARY,
            title="End",
            body={"key_points": ["p1"]},
        ))

        out = tmp_path / "simple.pptx"
        build_full_deck(deck, str(out))
        assert out.exists()
        assert out.stat().st_size > 0

    def test_build_all_9_types(self, theme, tmp_path):
        """所有 9 種版型都能正常建立"""
        deck = DeckSpec(title="All 9", subtitle="Test")

        # Cover
        deck.add_slide(SlideContent(
            type=SlideType.COVER,
            title="All 9",
            body={"tag": "# All"},
        ))
        # Section
        deck.add_slide(SlideContent(
            type=SlideType.SECTION_DIVIDER,
            title="Section",
            body={"section_num": "P1", "section_subtitle": "Sub"},
        ))
        # Objectives
        deck.add_slide(SlideContent(
            type=SlideType.OBJECTIVES,
            title="Goals",
            body={"items": [
                {"icon": "🎯", "title": "A", "desc": "a"},
            ]},
        ))
        # Title content
        deck.add_slide(SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Content",
            body={"items": ["a", "b"]},
        ))
        # Title table
        deck.add_slide(SlideContent(
            type=SlideType.TITLE_TABLE,
            title="Table",
            body={"headers": ["A"], "rows": [["1"]]},
        ))
        # Title code
        deck.add_slide(SlideContent(
            type=SlideType.TITLE_CODE,
            title="Code",
            body={"code": "x = 1"},
        ))
        # Two column
        deck.add_slide(SlideContent(
            type=SlideType.TWO_COLUMN,
            title="Compare",
            body={
                "left": {"title": "L", "items": ["a"]},
                "right": {"title": "R", "items": ["b"]},
            },
        ))
        # Grid cards
        deck.add_slide(SlideContent(
            type=SlideType.GRID_CARDS,
            title="Cards",
            body={"items": [
                {"icon": "1", "title": "A", "desc": "a"},
            ]},
        ))
        # Summary
        deck.add_slide(SlideContent(
            type=SlideType.SUMMARY,
            title="End",
            body={"key_points": ["p"]},
        ))

        out = tmp_path / "all9.pptx"
        build_full_deck(deck, str(out))

        # 重新讀取，應該有 9 張
        from pptx import Presentation
        reloaded = Presentation(str(out))
        assert len(reloaded.slides) == 9

    def test_build_with_different_theme(self, tmp_path):
        """用 minimal-bw 主題也能正常建立"""
        deck = DeckSpec(title="Test", theme="minimal-bw")
        deck.add_slide(SlideContent(
            type=SlideType.COVER,
            title="BW Test",
            body={"tag": "# BW"},
        ))
        deck.add_slide(SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Slide",
            body={"text": "Hello"},
        ))

        out = tmp_path / "bw.pptx"
        build_full_deck(deck, str(out))
        assert out.exists()

    def test_build_slide_helper(self, prs, theme):
        """build_slide() 便利函式"""
        slide = add_blank_slide(prs)
        content = SlideContent(
            type=SlideType.TITLE_CONTENT,
            title="Helper Test",
            body={"items": ["a", "b"]},
        )
        build_slide(slide, content, theme=theme, slide_num=1, total=1)
        assert len(slide.shapes) > 0

    def test_build_slide_with_invalid_type_raises(self, prs):
        """未知的 slide_type 應拋出 BuildError"""
        # 強制建立一個無效的 SlideContent
        content = SlideContent(
            type=SlideType.COVER,
            title="X",
        )
        # 改 type 為無效值
        content.type = "invalid_type"
        slide = add_blank_slide(prs)
        with pytest.raises((BuildError, ValueError)):
            build_slide(slide, content)
