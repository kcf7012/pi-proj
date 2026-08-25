"""
詳細檢視單份 PPTX。

比 diff_pptx.py 更詳細，列出每張 slide：
- 所有 shape 的類型、座標、尺寸、文字內容、字級、顏色
- 用「人類可讀」格式，方便比對
"""
from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor


def emu_to_inch(emu) -> float:
    if emu is None:
        return 0.0
    return round(int(emu) / 914400.0, 2)


def color_to_hex(color) -> str | None:
    try:
        if color and color.type is not None:
            if hasattr(color, 'rgb') and color.rgb is not None:
                return str(color.rgb)
    except Exception:
        pass
    return None


def shape_detail(shape, idx: int) -> str:
    """產出 shape 的詳細描述。"""
    t = str(shape.shape_type) if shape.shape_type else "?"
    left = emu_to_inch(shape.left)
    top = emu_to_inch(shape.top)
    w = emu_to_inch(shape.width)
    h = emu_to_inch(shape.height)

    fill_color = None
    try:
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            fill_color = color_to_hex(shape.fill.fore_color)
    except Exception:
        pass

    lines = [
        f"  [{idx:02d}] {t:25s} ({left},{top}) {w}x{h} fill={fill_color}"
    ]

    if shape.has_text_frame:
        for pi, para in enumerate(shape.text_frame.paragraphs):
            for ri, run in enumerate(para.runs):
                if not run.text.strip():
                    continue
                size = run.font.size
                size_pt = round(size.pt, 1) if size else "?"
                bold = "B" if run.font.bold else " "
                font_name = run.font.name or "?"
                color = color_to_hex(run.font.color) if run.font.color and run.font.color.type else "?"
                text = run.text[:50] + ("…" if len(run.text) > 50 else "")
                lines.append(
                    f"        p{pi}r{ri} {size_pt}pt {bold} {font_name} {color}: {text!r}"
                )
    return "\n".join(lines)


def inspect(path: Path, slide_num: int | None = None) -> None:
    p = Presentation(str(path))
    print(f"# {path}")
    print(f"# slides: {len(p.slides)}")
    print()

    slides = p.slides
    if slide_num is not None:
        slides = [p.slides[slide_num - 1]]
        print(f"# showing slide {slide_num} only\n")

    for i, slide in enumerate(slides):
        if slide_num is None:
            print(f"=== Slide {i+1} ===")
        shapes = list(slide.shapes)
        for j, sh in enumerate(shapes):
            print(shape_detail(sh, j))
        print()


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    path = Path(sys.argv[1])
    slide_num = int(sys.argv[2]) if len(sys.argv) >= 3 else None
    if not path.exists():
        print(f"Not found: {path}")
        return 2
    inspect(path, slide_num)
    return 0


if __name__ == "__main__":
    sys.exit(main())
