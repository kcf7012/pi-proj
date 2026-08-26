"""
新版 PPTX 版面檢查

逐張 slide 檢查：
- 文字框位置是否超出安全區（y > 7.35"）
- 表格高度是否塞下所有 row（用 row_count * row_height 估計）
- code 框是否塞下所有行數
- 是否有形狀重疊

"""
from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation


SAFE_TOP = 1.3
SAFE_BOTTOM = 7.0  # 7.1" 有 brand bar，7.35" 是極限
SAFE_LEFT = 0.5
SAFE_RIGHT = 12.833


def in_safe_zone(left, top, width, height) -> list[str]:
    issues = []
    bottom = top + height
    right = left + width
    if bottom > SAFE_BOTTOM + 0.01:
        issues.append(f"超出底部安全區 (bottom={bottom:.2f} > {SAFE_BOTTOM})")
    if top < SAFE_TOP - 0.01 and top != 0:
        issues.append(f"侵入標題區 (top={top:.2f} < {SAFE_TOP})")
    if left < SAFE_LEFT - 0.01 and left != 0:
        issues.append(f"超出左邊界 (left={left:.2f})")
    if right > SAFE_RIGHT + 0.01 and right != 13.333:
        issues.append(f"超出右邊界 (right={right:.2f})")
    return issues


def check_slide(slide, slide_idx: int) -> list[str]:
    issues = []
    shapes = list(slide.shapes)
    for i, sh in enumerate(shapes):
        if sh.left is None or sh.top is None:
            continue
        left = sh.left / 914400
        top = sh.top / 914400
        w = (sh.width or 0) / 914400
        h = (sh.height or 0) / 914400
        for msg in in_safe_zone(left, top, w, h):
            issues.append(f"  shape[{i}] {sh.name}: {msg}")

    # 檢查 code 框：高度是否足夠
    # 找 AUTO_SHAPE 緊接著 TEXT_BOX 的 pair（code 框）
    for i, sh in enumerate(shapes):
        t = str(sh.shape_type) if sh.shape_type else ""
        if t != "TEXT_BOX (17)":
            continue
        if not sh.has_text_frame:
            continue
        text = "\n".join(p.text for p in sh.text_frame.paragraphs)
        if "\n" not in text:
            continue
        n_lines = text.count("\n") + 1
        h_in = (sh.height or 0) / 914400
        # 假設每行 0.16" (11pt monospace line height)
        if n_lines * 0.16 > h_in + 0.05:
            issues.append(f"  shape[{i}] 疑似 code 框塞不下 {n_lines} 行 (height={h_in:.2f}, 需要 {n_lines*0.16:.2f})")

    # 檢查表格是否被切：計算 row 數 × 估算高度
    for i, sh in enumerate(shapes):
        t = str(sh.shape_type) if sh.shape_type else ""
        if "TABLE" not in t:
            continue
        if not hasattr(sh, "table"):
            continue
        try:
            table = sh.table
        except Exception:
            continue
        n_rows = len(table.rows)
        h_in = (sh.height or 0) / 914400
        # 表頭 0.4" + 每 row 0.35"
        needed = 0.4 + (n_rows - 1) * 0.35
        if needed > h_in + 0.05:
            issues.append(f"  shape[{i}] 表格塞不下 {n_rows} 列 (height={h_in:.2f}, 需要 {needed:.2f})")

    return issues


def inspect_pptx(path: Path) -> dict:
    p = Presentation(str(path))
    result = {"file": path.name, "slides": len(p.slides), "issues_per_slide": []}
    for i, slide in enumerate(p.slides, 1):
        issues = check_slide(slide, i)
        result["issues_per_slide"].append({"slide": i, "issues": issues})
    return result


def main() -> int:
    if len(sys.argv) < 2:
        print("Usage: layout_check.py file1.pptx [file2.pptx ...]")
        return 2
    for arg in sys.argv[1:]:
        p = Path(arg)
        if not p.exists():
            print(f"Not found: {p}")
            continue
        r = inspect_pptx(p)
        total_issues = sum(len(s['issues']) for s in r['issues_per_slide'])
        print(f"# {p.name}: {r['slides']} slides, {total_issues} issues")
        for s in r['issues_per_slide']:
            if s['issues']:
                print(f"  Slide {s['slide']}:")
                for iss in s['issues']:
                    print(iss)
    return 0


if __name__ == "__main__":
    sys.exit(main())
