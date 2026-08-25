"""
比對兩份 PPTX 的結構特徵。

用途：Phase 9 視覺驗證的程式化輔助。
- 列出每張 slide 的 shape 數、文字內容、座標範圍、填色
- 兩份並排比對，輸出差異

使用：
    python diff_pptx.py OLD.pptx NEW.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


def emu_to_inch(emu: int) -> float:
    return round(emu / 914400.0, 3)


def shape_signature(shape) -> dict:
    """取出 shape 的結構特徵（不含樣式細節）。"""
    sig = {
        "type": str(shape.shape_type) if shape.shape_type else "UNKNOWN",
        "name": shape.name,
        "left_in": emu_to_inch(shape.left) if shape.left is not None else None,
        "top_in": emu_to_inch(shape.top) if shape.top is not None else None,
        "width_in": emu_to_inch(shape.width) if shape.width is not None else None,
        "height_in": emu_to_inch(shape.height) if shape.height is not None else None,
        "text": "",
    }
    if shape.has_text_frame:
        texts = []
        for para in shape.text_frame.paragraphs:
            t = "".join(r.text for r in para.runs)
            if t.strip():
                texts.append(t.strip())
        sig["text"] = " | ".join(texts)
    return sig


def slide_signature(slide) -> dict:
    """取出 slide 的結構特徵。"""
    shapes = list(slide.shapes)
    sig = {
        "shape_count": len(shapes),
        "shapes": [shape_signature(s) for s in shapes],
    }
    # 抓 layout name
    try:
        sig["layout"] = slide.slide_layout.name
    except Exception:
        sig["layout"] = "?"
    return sig


def signature_text(sig: dict, indent: int = 0) -> str:
    """把 signature 轉成人類可讀文字。"""
    pad = "  " * indent
    lines = [f"{pad}layout={sig['layout']} shapes={sig['shape_count']}"]
    for i, sh in enumerate(sig["shapes"]):
        text_preview = sh["text"][:60] + ("…" if len(sh["text"]) > 60 else "")
        lines.append(
            f"{pad}  [{i:02d}] {sh['type']:25s} "
            f"pos=({sh['left_in']},{sh['top_in']}) "
            f"size=({sh['width_in']}x{sh['height_in']}) "
            f"text={text_preview!r}"
        )
    return "\n".join(lines)


def diff_two(old_path: Path, new_path: Path) -> int:
    old_p = Presentation(str(old_path))
    new_p = Presentation(str(new_path))

    print(f"# OLD: {old_path}")
    print(f"# NEW: {new_path}")
    print(f"# OLD slides: {len(old_p.slides)}  |  NEW slides: {len(new_p.slides)}")
    print()

    if len(old_p.slides) != len(new_p.slides):
        print(f"⚠️  slide count differs: {len(old_p.slides)} vs {len(new_p.slides)}")
        print()

    n = min(len(old_p.slides), len(new_p.slides))
    diff_count = 0
    for i in range(n):
        old_sig = slide_signature(old_p.slides[i])
        new_sig = slide_signature(new_p.slides[i])

        # 簡單比對：shape 數、text 內容總長度、layout
        old_text_total = sum(len(s["text"]) for s in old_sig["shapes"])
        new_text_total = sum(len(s["text"]) for s in new_sig["shapes"])

        differs = (
            old_sig["shape_count"] != new_sig["shape_count"]
            or old_sig["layout"] != new_sig["layout"]
            or old_text_total != new_text_total
        )

        marker = "❌ DIFF" if differs else "✅ match"
        if differs:
            diff_count += 1

        print(f"--- Slide {i+1:02d} {marker} ---")
        print(f"  OLD: shapes={old_sig['shape_count']} layout={old_sig['layout']} text_chars={old_text_total}")
        print(f"  NEW: shapes={new_sig['shape_count']} layout={new_sig['layout']} text_chars={new_text_total}")

        if differs:
            # 詳細列出 NEW 的 shapes
            print(f"  --- NEW shapes ---")
            print(signature_text(new_sig, indent=2))
            print()

    print(f"# Summary: {n - diff_count}/{n} slides match structurally, {diff_count} differ")
    return diff_count


def main() -> int:
    if len(sys.argv) != 3:
        print(__doc__)
        return 2
    old = Path(sys.argv[1])
    new = Path(sys.argv[2])
    if not old.exists():
        print(f"OLD not found: {old}")
        return 2
    if not new.exists():
        print(f"NEW not found: {new}")
        return 2
    return diff_two(old, new)


if __name__ == "__main__":
    sys.exit(main())
