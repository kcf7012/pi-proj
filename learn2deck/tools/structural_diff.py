"""
Phase 9 結構性視覺驗證報告

對 8 份新舊 .pptx 做結構性比對，產出 markdown 報告。
- 比較 slide 數、shape 數、版型分布、文字內容
- 不假裝是真正的視覺比對
- 輸出供人工決策的差異清單
"""
from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation


def slide_summary(slide) -> dict:
    """slide 結構摘要"""
    shapes = list(slide.shapes)
    text_chars = 0
    type_counts: dict[str, int] = {}
    for s in shapes:
        t = str(s.shape_type) if s.shape_type else "?"
        type_counts[t] = type_counts.get(t, 0) + 1
        if s.has_text_frame:
            text_chars += sum(len(p.text) for p in s.text_frame.paragraphs)

    # 第一個非空文字
    title = ""
    for s in shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                if p.text.strip():
                    title = p.text.strip()
                    break
            if title:
                break

    return {
        "shape_count": len(shapes),
        "type_counts": type_counts,
        "text_chars": text_chars,
        "title": title[:50],
    }


def diff_one(old: Path, new: Path) -> dict:
    old_p = Presentation(str(old))
    new_p = Presentation(str(new))

    old_summary = [slide_summary(s) for s in old_p.slides]
    new_summary = [slide_summary(s) for s in new_p.slides]

    return {
        "old_slides": len(old_summary),
        "new_slides": len(new_summary),
        "old_summary": old_summary,
        "new_summary": new_summary,
    }


def main() -> int:
    if len(sys.argv) != 3:
        print("Usage: structural_diff.py OLD_DIR NEW_DIR")
        return 2
    old_dir = Path(sys.argv[1])
    new_dir = Path(sys.argv[2])

    pairs = [
        ("00-overview.pptx", "new_00-claude-code-plugins-series.pptx"),
        ("01-plugin-marketplaces.pptx", "new_01-plugin-marketplaces.pptx"),
        ("02-plugins.pptx", "new_02-plugins.pptx"),
        ("03-plugins-reference.pptx", "new_03-plugins-reference.pptx"),
        ("04-skills.pptx", "new_04-skills.pptx"),
        ("05-subagents.pptx", "new_05-subagents.pptx"),
        ("06-hooks.pptx", "new_06-hooks.pptx"),
        ("07-discover-plugins.pptx", "new_07-discover-plugins.pptx"),
    ]

    print("# Phase 9 結構性視覺驗證報告\n")
    print(f"OLD: {old_dir}  (現有 pi-proj .pptx)")
    print(f"NEW: {new_dir}  (learn2deck build 產出)\n")

    for old_name, new_name in pairs:
        old_p = old_dir / old_name
        new_p = new_dir / new_name
        if not old_p.exists():
            print(f"⚠️  Missing old: {old_p}")
            continue
        if not new_p.exists():
            print(f"⚠️  Missing new: {new_p}")
            continue

        d = diff_one(old_p, new_p)
        print(f"## {new_name}")
        print(f"- OLD slides: **{d['old_slides']}**")
        print(f"- NEW slides: **{d['new_slides']}**")
        diff = d['new_slides'] - d['old_slides']
        sign = "+" if diff > 0 else ""
        print(f"- Δ: {sign}{diff}")

        old_shapes = sum(s['shape_count'] for s in d['old_summary'])
        new_shapes = sum(s['shape_count'] for s in d['new_summary'])
        old_chars = sum(s['text_chars'] for s in d['old_summary'])
        new_chars = sum(s['text_chars'] for s in d['new_summary'])

        print(f"- Total shapes: OLD={old_shapes} / NEW={new_shapes} (Δ={new_shapes-old_shapes:+d})")
        print(f"- Total text chars: OLD={old_chars} / NEW={new_chars} (Δ={new_chars-old_chars:+d})")

        # 列出前 3 張 NEW 的 title
        print(f"- First 3 NEW titles:")
        for s in d['new_summary'][:3]:
            print(f"  - `{s['title']}` ({s['shape_count']} shapes, {s['text_chars']} chars)")

        print()

    return 0


if __name__ == "__main__":
    sys.exit(main())
