"""
Claude Code Plugin 系列簡報產生器 - 共用 Helper
設計系統：Claude 品牌色 + 自學友善的標準型 + 詳細型密度
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 設計系統：顏色與字體
# ============================================================
COLOR_PRIMARY   = RGBColor(0xC7, 0x5A, 0x1A)   # Claude 橘
COLOR_DARK      = RGBColor(0x2C, 0x2C, 0x2C)   # 深灰（主要文字）
COLOR_BG_CREAM   = RGBColor(0xFA, 0xF8, 0xF3)   # 米白（背景）
COLOR_BG_GRAY    = RGBColor(0xF3, 0xF0, 0xE9)   # 淺米（卡片背景）
COLOR_BLUE      = RGBColor(0x3B, 0x82, 0xF6)   # 輔助藍
COLOR_GREEN     = RGBColor(0x16, 0xA3, 0x4A)   # 強調綠
COLOR_RED       = RGBColor(0xDC, 0x26, 0x26)   # 警告紅
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY_TXT  = RGBColor(0x6B, 0x6B, 0x6B)   # 次要文字
COLOR_CODE_BG   = RGBColor(0x1E, 0x1E, 0x1E)   # 程式碼區塊背景
COLOR_CODE_FG   = RGBColor(0xE6, 0xE6, 0xE6)   # 程式碼前景

FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri"
FONT_MONO    = "Consolas"

# ============================================================
# 簡報尺寸：16:9
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    """建立 16:9 簡報"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    """新增空白頁（最常用，自己畫所有元素）"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_slide_bg(slide, color=COLOR_BG_CREAM):
    """設定投影片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================================
# 標題列（每張投影片頂部）
# ============================================================
def add_title_bar(slide, title_text, subtitle_text=None, slide_num=None, total=None, source=None):
    """
    統一的標題列設計
    - 橘色 accent line 在頂部
    - 主標題（大字）
    - 副標題（小字，灰色）
    - 右上角頁碼
    - 右下角來源（對應系列檔案）
    """
    # 橘色頂部 accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.15)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()

    # 主標題
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(11), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_TITLE
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLOR_DARK

    # 副標題
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.0),
            Inches(11), Inches(0.4)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = subtitle_text
        run.font.name = FONT_BODY
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_GRAY_TXT
        run.font.italic = True

    # 右上角頁碼
    if slide_num is not None and total is not None:
        page_box = slide.shapes.add_textbox(
            Inches(12), Inches(0.3),
            Inches(1.2), Inches(0.3)
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_num} / {total}"
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = COLOR_GRAY_TXT

    # 右下角來源
    if source:
        src_box = slide.shapes.add_textbox(
            Inches(8), Inches(7.1),
            Inches(5), Inches(0.3)
        )
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"📖 來源：{source}"
        run.font.name = FONT_BODY
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = COLOR_GRAY_TXT

    # 左下角品牌
    brand_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.1),
        Inches(6), Inches(0.3)
    )
    tf = brand_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Claude Code Plugin 完整學習系列"
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.color.rgb = COLOR_PRIMARY
    run.font.bold = True


# ============================================================
# 內容區塊
# ============================================================
def add_text_block(slide, text, left, top, width, height,
                   font_size=14, bold=False, italic=False, color=COLOR_DARK,
                   align=PP_ALIGN.LEFT, bg_color=None, font=FONT_BODY):
    """新增純文字區塊（可選背景色）"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

    if bg_color:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
    return box


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=14, color=COLOR_DARK, font=FONT_BODY,
                    bullet_char="•", indent_levels=None):
    """
    新增階層式項目清單
    items: list of (indent_level, text) tuples 或純字串清單
    indent_levels: None 表示全部用 bullet_char
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)

    # 正規化 items
    if isinstance(items[0], str):
        items = [(0, t) for t in items]
    elif isinstance(items[0], tuple) and len(items[0]) == 2:
        pass
    else:
        raise ValueError("items 應為 str list 或 (indent, text) tuple list")

    for i, (indent, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent

        # 縮排
        prefix = "    " * indent
        if indent == 0:
            prefix = f"{bullet_char} "

        run = p.add_run()
        run.text = f"{prefix}{text}"
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

    return box


def add_code_block(slide, code, left, top, width, height,
                   language="bash", font_size=12):
    """新增深色背景的程式碼區塊"""
    # 背景矩形
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_CODE_BG
    bg.line.fill.background()

    # 程式碼文字
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.15)
    tf.margin_bottom = Inches(0.1)

    lines = code.strip("\n").split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(font_size)
        run.font.color.rgb = COLOR_CODE_FG


def add_callout(slide, text, left, top, width, height,
                bg_color=COLOR_BG_GRAY, border_color=None,
                icon="💡", font_size=13):
    """新增提示框（圓角矩形）"""
    if border_color is None:
        border_color = COLOR_PRIMARY

    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = bg_color
    callout.line.color.rgb = border_color
    callout.line.width = Pt(1.5)

    tf = callout.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"{icon}  {text}" if icon else text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.color.rgb = COLOR_DARK


# ============================================================
# 表格
# ============================================================
def add_comparison_table(slide, headers, rows, left, top, width, height,
                         header_bg=COLOR_PRIMARY, header_fg=COLOR_WHITE,
                         alt_row_bg=COLOR_BG_GRAY, font_size=11):
    """
    新增比較表格
    headers: list of column titles
    rows: list of list of cell values (string)
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 標題列
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        tf = cell.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.name = FONT_TITLE
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = header_fg

    # 資料列
    for row_idx, row in enumerate(rows, start=1):
        bg = alt_row_bg if row_idx % 2 == 0 else COLOR_WHITE
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = COLOR_DARK


# ============================================================
# 流程圖
# ============================================================
def add_flow_box(slide, text, left, top, width, height,
                 bg_color=COLOR_BG_GRAY, border_color=None,
                 font_size=13, font_color=COLOR_DARK, bold=True):
    """流程圖方塊（圓角矩形 + 文字）"""
    if border_color is None:
        border_color = COLOR_PRIMARY

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TITLE
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    return box


def add_arrow(slide, start_x, start_y, end_x, end_y,
              color=None, width=2):
    """新增箭頭（直線）"""
    if color is None:
        color = COLOR_GRAY_TXT
    connector = slide.shapes.add_connector(
        1,  # Straight arrow
        start_x, start_y,
        end_x - start_x, end_y - start_y
    )
    line = connector.line
    line.color.rgb = color
    line.width = Pt(width)
    return connector


# ============================================================
# 特殊頁面
# ============================================================
def add_cover_slide(prs, title, subtitle, tag="Claude Code Plugin 學習系列"):
    """封面頁"""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, COLOR_BG_CREAM)

    # 大標題
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = COLOR_PRIMARY

    # 副標題
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(11.333), Inches(1)
    )
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitle
    run.font.name = FONT_BODY
    run.font.size = Pt(22)
    run.font.color.rgb = COLOR_DARK

    # 標籤
    tag_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.0),
        Inches(11.333), Inches(0.5)
    )
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tag
    run.font.name = FONT_BODY
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = COLOR_GRAY_TXT

    # 橘色裝飾條
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.666), Inches(5.6),
        Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    return slide


def add_section_divider(prs, section_num, section_title, section_subtitle):
    """章節分隔頁"""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, COLOR_BG_CREAM)

    # 巨大章節編號
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(12.333), Inches(2)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_num
    run.font.name = FONT_TITLE
    run.font.size = Pt(96)
    run.font.bold = True
    run.font.color.rgb = COLOR_PRIMARY

    # 章節標題
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.0),
        Inches(12.333), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_title
    run.font.name = FONT_TITLE
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_DARK

    # 副標題
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.2),
        Inches(12.333), Inches(0.8)
    )
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_subtitle
    run.font.name = FONT_BODY
    run.font.size = Pt(18)
    run.font.italic = True
    run.font.color.rgb = COLOR_GRAY_TXT

    return slide


def add_summary_slide(slide, title="重點回顧", key_points=None,
                       next_steps=None, source=None):
    """重點回顧頁"""
    add_title_bar(slide, title, "重點整理與下一步行動", source=source)

    if key_points:
        add_text_block(
            slide, "📌 關鍵要點",
            Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
            font_size=20, bold=True, color=COLOR_PRIMARY
        )
        add_bullet_list(
            slide, key_points,
            Inches(0.7), Inches(2.2), Inches(12), Inches(2.5),
            font_size=15
        )

    if next_steps:
        add_text_block(
            slide, "🚀 下一步行動",
            Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
            font_size=20, bold=True, color=COLOR_PRIMARY
        )
        add_bullet_list(
            slide, next_steps,
            Inches(0.7), Inches(5.3), Inches(12), Inches(1.5),
            font_size=15
        )


# ============================================================
# 特殊內容：兩欄對比
# ============================================================
def add_two_column_compare(slide, left_title, left_content, right_title, right_content,
                            top=Inches(1.7), height=Inches(5.0),
                            left_color=COLOR_BLUE, right_color=COLOR_PRIMARY):
    """左右兩欄對比佈局"""
    # 左欄
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top, Inches(6.1), height
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_BG_GRAY
    left_box.line.color.rgb = left_color
    left_box.line.width = Pt(2)

    # 左欄標題列
    left_title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), top, Inches(6.1), Inches(0.5)
    )
    left_title_bg.fill.solid()
    left_title_bg.fill.fore_color.rgb = left_color
    left_title_bg.line.fill.background()

    add_text_block(
        slide, left_title,
        Inches(0.5), top, Inches(6.1), Inches(0.5),
        font_size=16, bold=True, color=COLOR_WHITE,
        align=PP_ALIGN.CENTER
    )

    add_bullet_list(
        slide, left_content,
        Inches(0.7), top + Inches(0.7), Inches(5.7), height - Inches(0.8),
        font_size=13
    )

    # 右欄
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.733), top, Inches(6.1), height
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_BG_GRAY
    right_box.line.color.rgb = right_color
    right_box.line.width = Pt(2)

    # 右欄標題列
    right_title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.733), top, Inches(6.1), Inches(0.5)
    )
    right_title_bg.fill.solid()
    right_title_bg.fill.fore_color.rgb = right_color
    right_title_bg.line.fill.background()

    add_text_block(
        slide, right_title,
        Inches(6.733), top, Inches(6.1), Inches(0.5),
        font_size=16, bold=True, color=COLOR_WHITE,
        align=PP_ALIGN.CENTER
    )

    add_bullet_list(
        slide, right_content,
        Inches(6.933), top + Inches(0.7), Inches(5.7), height - Inches(0.8),
        font_size=13
    )
