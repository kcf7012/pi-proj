"""
簡報 1/4：建立並分發 Plugin Marketplace (01-plugin-marketplaces.pptx)
約 35 張
對應：01-plugin-marketplaces.md
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 35

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "建立並分發 Plugin Marketplace",
        "把你的 plugin 變成可分享的應用程式商店",
        tag="#01 · Marketplace"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "從建立第一個 marketplace 到企業級部署",
        slide_num=2, total=TOTAL, source="01-plugin-marketplaces.md"
    )

    objectives = [
        ("🏪", "什麼是 Marketplace", "理解「plugin 目錄」的核心概念"),
        ("🚀", "建立本機 Marketplace", "6 步驟完成第一個 marketplace"),
        ("📋", "完整 JSON Schema", "marketplace.json 與 plugin 項目的所有欄位"),
        ("🌐", "5 種 Plugin 來源", "本機路徑、GitHub、git、git-subdir、npm"),
        ("🛡️", "Strict Mode 與託管", "企業環境的 marketplace 限制"),
        ("🚀", "發佈與分發", "GitHub 託管、團隊預配置、容器預填"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(2.2)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_cols = 3
    grid_rows = 2
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.8)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_PRIMARY
        card.line.width = Pt(2)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.2), box_w, Inches(0.7),
            font_size=40, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.2), y + Inches(1.0), box_w - Inches(0.4), Inches(0.5),
            font_size=15, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.5), box_w - Inches(0.4), Inches(0.6),
            font_size=11, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Marketplace 基礎", "概念、流程、必備檔案")

    # ============================================================
    # Slide 4：什麼是 Marketplace
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "什麼是 Plugin Marketplace？",
        "把 plugin 分發給其他人的目錄",
        slide_num=4, total=TOTAL, source="01 § 什麼是 Plugin Marketplace"
    )

    h.add_text_block(
        slide, "✨ Marketplace 提供",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "✅ **集中式探索**：讓人容易找到",
            "✅ **版本追蹤**：清楚知道用的是哪一版",
            "✅ **自動更新**：推送後使用者能拉新版本",
            "✅ **多種來源支援**：git 儲存庫、本機路徑、npm…",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(2.0),
        font_size=15
    )

    h.add_callout(
        slide, "一句話：Marketplace = 別人加入來源、安裝你寫的 plugin 的入口",
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.5),
        icon="💡", font_size=15
    )

    h.add_text_block(
        slide, "🔄 整體流程（4 步）",
        Inches(0.5), Inches(5.1), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    flow_steps = ["建立 plugin", "建立 marketplace.json", "託管（GitHub / GitLab）", "與使用者分享"]
    box_w = Inches(2.7)
    box_h = Inches(1.0)
    h_gap = Inches(0.3)
    total_w = box_w * 4 + h_gap * 3
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(5.6)

    for i, text in enumerate(flow_steps):
        x = start_x + i * (box_w + h_gap)
        h.add_flow_box(
            slide, f"{i+1}. {text}",
            x, start_y, box_w, box_h,
            bg_color=h.COLOR_BG_GRAY, border_color=h.COLOR_PRIMARY,
            font_size=14
        )
        # 箭頭
        if i < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.02), start_y + Inches(0.3),
                Inches(0.26), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    # ============================================================
    # Slide 5：使用者端指令
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "使用者端 3 個關鍵指令",
        "使用者怎麼使用你的 marketplace？",
        slide_num=5, total=TOTAL, source="01 § 整體流程"
    )

    h.add_comparison_table(
        slide,
        ["指令", "作用"],
        [
            ["/plugin marketplace add <source>", "新增 marketplace 來源"],
            ["/plugin install <plugin>@<marketplace>", "安裝特定 plugin"],
            ["/plugin marketplace update", "重新整理本機副本"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.5),
        font_size=15
    )

    h.add_callout(
        slide, "假設你已經有要分發的 plugin。還沒？見 02-plugins.md",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5),
        icon="💡", font_size=13
    )

    h.add_text_block(
        slide, "📌 為什麼需要 Marketplace？",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "不用 marketplace 也能用 `--plugin-dir` 直接載入（個人開發用）",
            "但要**分享給團隊或公開** → 必須透過 marketplace",
            "marketplace 提供**版本管理、自動更新**、**安全管控**",
        ],
        Inches(0.7), Inches(5.7), Inches(12), Inches(1.5),
        font_size=14
    )

    # ============================================================
    # Slide 6：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "快速入門：建立本機 Marketplace", "6 步驟完成第一個 marketplace")

    # ============================================================
    # Slide 7：6 步驟總覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "建立 Marketplace 的 6 個步驟",
        "從目錄結構到使用 skill",
        slide_num=7, total=TOTAL, source="01 § 快速入門"
    )

    steps = [
        ("1", "建立目錄結構", "mkdir my-marketplace/...", "建立 plugin 根目錄"),
        ("2", "建立 SKILL.md", "skills/quality-review/", "定義 skill 行為"),
        ("3", "建立 plugin manifest", "plugin.json", "plugin 中繼資料"),
        ("4", "建立 marketplace 目錄", "marketplace.json", "marketplace 中繼資料"),
        ("5", "新增並安裝", "/plugin marketplace add...", "本地測試"),
        ("6", "使用 skill", "/plugin:skill", "驗證功能"),
    ]

    step_w = Inches(1.95)
    step_h = Inches(3.0)
    h_gap = Inches(0.2)
    total_w = step_w * 6 + h_gap * 5
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.3)

    for i, (num, title, code, desc) in enumerate(steps):
        x = start_x + i * (step_w + h_gap)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y, step_w, step_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        # 編號圓
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + step_w / 2 - Inches(0.3), start_y + Inches(0.2),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = h.COLOR_PRIMARY
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            x + step_w / 2 - Inches(0.3), start_y + Inches(0.22),
            Inches(0.6), Inches(0.55),
            font_size=24, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

        # 標題
        h.add_text_block(
            slide, title,
            x + Inches(0.1), start_y + Inches(0.95), step_w - Inches(0.2), Inches(0.6),
            font_size=12, bold=True, align=PP_ALIGN.CENTER
        )
        # code
        h.add_text_block(
            slide, code,
            x + Inches(0.1), start_y + Inches(1.65), step_w - Inches(0.2), Inches(0.8),
            font_size=9, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.1), start_y + Inches(2.45), step_w - Inches(0.2), Inches(0.5),
            font_size=10, align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 8：範例完整程式碼
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "範例：完整程式碼",
        "建立 quality-review marketplace",
        slide_num=8, total=TOTAL, source="01 § 快速入門"
    )

    h.add_text_block(
        slide, "📁 Step 1: 建立目錄",
        Inches(0.5), Inches(1.65), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """mkdir -p my-marketplace/.claude-plugin
mkdir -p my-marketplace/plugins/quality-review-plugin/.claude-plugin
mkdir -p my-marketplace/plugins/quality-review-plugin/skills/quality-review""",
        Inches(0.7), Inches(2.05), Inches(12), Inches(0.75),
        font_size=10
    )

    h.add_text_block(
        slide, "📝 Step 2: SKILL.md（plugin 的 skills/quality-review/SKILL.md）",
        Inches(0.5), Inches(2.9), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """---
description: 檢查程式碼中的錯誤、安全性和效能問題
---

檢查我選擇的程式碼或最近的變更，查找：
- 潛在的錯誤或邊界情況
- 安全性問題
- 效能問題
- 可讀性改進

簡潔且可行動。""",
        Inches(0.7), Inches(3.3), Inches(12), Inches(2.15),
        font_size=10
    )

    h.add_text_block(
        slide, "📦 Step 3: plugin manifest（.claude-plugin/plugin.json）",
        Inches(0.5), Inches(5.45), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """{
  "name": "quality-review-plugin",
  "description": "新增 quality-review skill 以進行快速程式碼審查",
  "version": "1.0.0"
}""",
        Inches(0.7), Inches(5.85), Inches(12), Inches(1.15),
        font_size=10
    )

    # ============================================================
    # Slide 9：Step 4-6
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "範例：完整程式碼（續）",
        "marketplace.json 與測試",
        slide_num=9, total=TOTAL, source="01 § 快速入門"
    )

    h.add_text_block(
        slide, "🏪 Step 4: marketplace 目錄（.claude-plugin/marketplace.json）",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """{
  "name": "my-plugins",
  "owner": {
    "name": "Your Name"
  },
  "plugins": [
    {
      "name": "quality-review-plugin",
      "source": "./plugins/quality-review-plugin",
      "description": "新增 quality-review skill 以進行快速程式碼審查"
    }
  ]
}""",
        Inches(0.7), Inches(2.1), Inches(12), Inches(2.7),
        font_size=10
    )

    h.add_text_block(
        slide, "🚀 Step 5-6: 新增、安裝並使用",
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """# Step 5: 新增並安裝
/plugin marketplace add ./my-marketplace
/plugin install quality-review-plugin@my-plugins

# Step 6: 選取一些程式碼後執行
/quality-review-plugin:quality-review""",
        Inches(0.7), Inches(5.3), Inches(12), Inches(1.5),
        font_size=11
    )

    h.add_callout(
        slide, "📌 Plugin skills 使用 **plugin 名稱做命名空間**",
        Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.35),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 10：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "Marketplace 與 Plugin Schema", "完整欄位速查")

    # ============================================================
    # Slide 11：Marketplace 完整範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Marketplace 完整範例",
        "多個 plugin 與完整 metadata",
        slide_num=11, total=TOTAL, source="01 § Marketplace 檔案結構"
    )

    h.add_code_block(
        slide, """{
  "name": "company-tools",
  "owner": {
    "name": "DevTools Team",
    "email": "devtools@example.com"
  },
  "plugins": [
    {
      "name": "code-formatter",
      "source": "./plugins/formatter",
      "description": "在保存時自動格式化程式碼",
      "version": "2.1.0",
      "author": { "name": "DevTools Team" }
    },
    {
      "name": "deployment-tools",
      "source": {
        "source": "github",
        "repo": "company/deploy-plugin"
      },
      "description": "部署自動化工具"
    }
  ]
}""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(4.6),
        font_size=11
    )

    h.add_callout(
        slide, "💡 儲存庫根目錄建立 `.claude-plugin/marketplace.json`，定義名稱、擁有者、plugin 清單",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 12：Marketplace 必需與選用欄位
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Marketplace Schema",
        "3 個必需欄位 + 7 個選用欄位",
        slide_num=12, total=TOTAL, source="01 § Marketplace 架構"
    )

    h.add_text_block(
        slide, "⭐ 必需欄位",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_comparison_table(
        slide,
        ["欄位", "類型", "描述"],
        [
            ["`name`", "string", "識別碼（kebab-case、無空格）"],
            ["`owner`", "object", "維護者資訊"],
            ["`plugins`", "array", "可用 plugin 清單"],
        ],
        Inches(0.5), Inches(2.1), Inches(6.0), Inches(2.2),
        font_size=11
    )

    h.add_text_block(
        slide, "📋 選用欄位",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["欄位", "描述"],
        [
            ["`$schema`", "JSON Schema URL（編輯器自動完成）"],
            ["`description`", "簡短描述"],
            ["`version`", "Marketplace 版本"],
            ["`metadata.pluginRoot`", "相對路徑的基底（例：`./plugins`）"],
            ["`allowCrossMarketplaceDependenciesOn`", "允許依賴的其他 marketplace 名單"],
            ["`renames`", "舊名稱 → 新名稱對應（自動遷移）"],
        ],
        Inches(6.733), Inches(2.1), Inches(6.0), Inches(2.5),
        font_size=11
    )

    h.add_callout(
        slide, "⚠️ 保留名稱（不可使用）：claude-code-marketplace、claude-plugins-official、anthropic-plugins 等",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5),
        icon="", font_size=12
    )

    h.add_text_block(
        slide, "👤 owner 欄位",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`name` ✅ 必需：維護者或團隊名稱  /  `email` ❌ 選用：聯絡電子郵件",
        ],
        Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
        font_size=12
    )

    # ============================================================
    # Slide 13：Plugin 項目欄位（完整）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin 項目欄位",
        "2 個必需 + 多個選用",
        slide_num=13, total=TOTAL, source="01 § Plugin 項目欄位"
    )

    h.add_text_block(
        slide, "⭐ 必需",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )
    h.add_bullet_list(
        slide, [
            "`name`（kebab-case，無空格）  /  `source`（string 或 object）",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
        font_size=13
    )

    h.add_text_block(
        slide, "📋 標準中繼資料（選用）",
        Inches(0.5), Inches(2.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["欄位", "說明"],
        [
            ["`displayName`", "UI 顯示名稱（可含空格，v2.1.143+）"],
            ["`description`", "簡短描述"],
            ["`version`", "Plugin 版本（設了之後 plugin 會固定）"],
            ["`author` / `homepage` / `repository` / `license`", "作者與授權資訊"],
            ["`keywords` / `category` / `tags`", "分類與可搜尋標籤"],
            ["`strict`", "控制 plugin.json 是否為元件定義權威（預設 true）"],
            ["`defaultEnabled`", "安裝後是否啟用（預設 true）"],
        ],
        Inches(0.5), Inches(3.1), Inches(6.0), Inches(2.7),
        font_size=10
    )

    h.add_text_block(
        slide, "🔌 元件配置（選用）",
        Inches(6.733), Inches(2.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["欄位", "說明"],
        [
            ["`skills`", "skill 目錄的自訂路徑"],
            ["`commands`", "命令檔案或目錄的自訂路徑"],
            ["`agents`", "agent 檔案路徑"],
            ["`hooks`", "hooks 配置或檔案路徑"],
            ["`mcpServers`", "MCP server 配置或路徑"],
            ["`lspServers`", "LSP server 配置或路徑"],
        ],
        Inches(6.733), Inches(3.1), Inches(6.0), Inches(2.7),
        font_size=10
    )

    h.add_callout(
        slide, "💡 `${CLAUDE_PLUGIN_ROOT}` 引用 plugin 安裝目錄；想保留更新狀態用 `${CLAUDE_PLUGIN_DATA}`",
        Inches(0.5), Inches(5.95), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 14：進階 plugin 範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "進階 Plugin 項目範例",
        "完整 hooks + MCP + 多元件",
        slide_num=14, total=TOTAL, source="01 § Plugin 項目欄位"
    )

    h.add_code_block(
        slide, """{
  "name": "enterprise-tools",
  "source": {
    "source": "github",
    "repo": "company/enterprise-plugin"
  },
  "description": "企業工作流程自動化工具",
  "version": "2.1.0",
  "author": {
    "name": "Enterprise Team",
    "email": "enterprise@example.com"
  },
  "homepage": "https://docs.example.com/enterprise",
  "license": "MIT",
  "commands": ["./commands/core/", "./commands/enterprise/"],
  "agents": ["./agents/security-reviewer.md", "./agents/compliance-checker.md"],
  "hooks": {
    "PostToolUse": [{
      "matcher": "Write|Edit",
      "hooks": [{
        "type": "command",
        "command": "${CLAUDE_PLUGIN_ROOT}/scripts/validate.sh"
      }]
    }]
  },
  "mcpServers": {
    "enterprise-db": {
      "command": "${CLAUDE_PLUGIN_ROOT}/servers/db-server",
      "args": ["--config", "${CLAUDE_PLUGIN_ROOT}/config.json"]
    }
  },
  "strict": false
}""",
        Inches(0.5), Inches(1.7), Inches(8.0), Inches(5.0),
        font_size=8
    )

    h.add_text_block(
        slide, "📌 重點提醒",
        Inches(8.7), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`commands` / `agents`：可指定多個目錄或檔案",
            "路徑**相對於 plugin 根目錄**",
            "`${CLAUDE_PLUGIN_ROOT}` 引用 plugin 安裝目錄",
            "`${CLAUDE_PLUGIN_DATA}` 保留 plugin 更新後狀態",
            "`strict: false` → marketplace 項目是完整定義",
        ],
        Inches(8.7), Inches(2.1), Inches(4.5), Inches(3.0),
        font_size=12
    )

    # ============================================================
    # Slide 15：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "Plugin 來源（5 種）", "從本機路徑到 npm 套件")

    # ============================================================
    # Slide 16：5 種來源速覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "5 種 Plugin 來源",
        "決定 plugin 從哪裡取得",
        slide_num=16, total=TOTAL, source="01 § Plugin 來源"
    )

    sources = [
        ("📁", "相對路徑", '"./my-plugin"', "marketplace 儲存庫內的本機目錄"),
        ("🐙", "GitHub", '{"source":"github","repo":"owner/repo"}', "最常用，自動 clone"),
        ("🔗", "Git URL", '{"source":"url","url":"https://..."}', "任意 git 主機（GitLab、Bitbucket）"),
        ("📦", "Git 子目錄", '{"source":"git-subdir","path":"..."}', "monorepo 用，稀疏複製"),
        ("📚", "npm", '{"source":"npm","package":"@acme/..."}', "透過 `npm install` 安裝"),
    ]

    box_w = Inches(6.1)
    box_h = Inches(0.92)
    box_gap_x = Inches(0.13)
    box_gap_y = Inches(0.08)
    grid_start_x = Inches(0.5)
    grid_start_y = Inches(1.7)

    for i, (icon, name, example, desc) in enumerate(sources):
        x = grid_start_x
        y = grid_start_y + i * (box_h + box_gap_y)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(1.5)

        h.add_text_block(
            slide, f"{icon}  {name}",
            x + Inches(0.2), y + Inches(0.12), Inches(2.0), Inches(0.4),
            font_size=14, bold=True
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(0.5), Inches(2.0), Inches(0.4),
            font_size=10, color=h.COLOR_GRAY_TXT, italic=True
        )
        h.add_code_block(
            slide, example,
            x + Inches(2.3), y + Inches(0.18), Inches(3.7), Inches(0.6),
            font_size=9
        )

    h.add_callout(
        slide, "📌 同時設 `ref` + `sha` 時，`sha` 是有效的固定（即使分支被刪，只要 commit 還能取得就成功）",
        Inches(0.5), Inches(6.75), Inches(12.333), Inches(0.3),
        icon="", font_size=10
    )

    # ============================================================
    # Slide 17：相對路徑陷阱
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "相對路徑的陷阱",
        "URL-based marketplace 會失效",
        slide_num=17, total=TOTAL, source="01 § 相對路徑"
    )

    h.add_text_block(
        slide, "⚠️ 失敗情境",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "**症狀**：用 URL 加 marketplace，plugin 報「path not found」",
            "**原因**：URL-based marketplace 只下載 `marketplace.json`，不下載其他檔案",
            "**限制**：必須以 `./` 開頭，不能用 `../` 跳出根目錄",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(1.8),
        font_size=14
    )

    h.add_text_block(
        slide, "✅ 解決方案",
        Inches(0.5), Inches(4.1), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_GREEN
    )

    h.add_bullet_list(
        slide, [
            "改用外部來源（GitHub / npm / git URL）",
            "或用 **git-based marketplace**（會 clone 整個 repo）",
            "對於企業內部使用：考慮把整個 marketplace 託管在 GitHub",
        ],
        Inches(0.7), Inches(4.5), Inches(12), Inches(1.5),
        font_size=14
    )

    h.add_callout(
        slide, "💡 本機路徑相對於 **marketplace 根目錄**（包含 `.claude-plugin/` 的那層）解析",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 18：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "Strict Mode 與託管", "企業環境的安全管控")

    # ============================================================
    # Slide 19：Strict Mode
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Strict Mode 嚴格模式",
        "控制 plugin.json 是否為元件定義的權威",
        slide_num=19, total=TOTAL, source="01 § Strict Mode"
    )

    h.add_two_column_compare(
        slide,
        "✅ strict: true（預設）",
        [
            "`plugin.json` 為**權威**",
            "marketplace 項目可補充元件",
            "兩者**合併**生效",
            "大多數情境適用",
            "plugin 有自己的 plugin.json 並管理元件",
        ],
        "⚠️ strict: false",
        [
            "marketplace 項目是**完整定義**",
            "若 plugin 也有 plugin.json 宣告元件 → **衝突**",
            "衝突 → **無法載入**",
            "marketplace 運營商要完全控制時用",
            "重組或策劃 plugin 的元件時適用",
        ]
    )

    h.add_callout(
        slide, "💡 經驗法則：用 strict: true 讓 plugin 自己管理，市場只補充",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        icon="", font_size=13
    )

    # ============================================================
    # Slide 20：託管選項
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "託管並分發 Marketplace",
        "4 種託管選項",
        slide_num=20, total=TOTAL, source="01 § 託管"
    )

    h.add_two_column_compare(
        slide,
        "🐙 1. GitHub（推薦）",
        [
            "建立新 repo",
            "加 `.claude-plugin/marketplace.json`",
            "分享指令：",
            "`/plugin marketplace add owner/repo`",
            "免費、CI/CD 整合完善",
        ],
        "🌐 2. 其他 Git 主機",
        [
            "GitLab、Bitbucket、Gitea…",
            "/plugin marketplace add https://gitlab.com/.../plugins.git",
            "支援 SSH 與 HTTPS",
            "私人 repo 用既有 git 認證助手",
        ]
    )

    h.add_text_block(
        slide, "🔐 3. 私人 repo：背景自動更新的挑戰",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_text_block(
        slide, "預設 background refresh 對 `git pull` 停用認證助手 → HTTPS 認證失敗。解法：CLAUDE_CODE_PLUGIN_KEEP_MARKETPLACE_ON_FAILURE=1、gh auth setup-git、URL token 重寫",
        Inches(0.7), Inches(7.25), Inches(12), Inches(0.5),
        font_size=11
    )

    # ============================================================
    # Slide 21：個人 repo 解決方案
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "私人 Repo 認證：3 種解法",
        "處理 background refresh 問題",
        slide_num=21, total=TOTAL, source="01 § 私人 repo"
    )

    h.add_text_block(
        slide, "解法 1：保留最後副本（環境變數）",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """export CLAUDE_CODE_PLUGIN_KEEP_MARKETPLACE_ON_FAILURE=1""",
        Inches(0.7), Inches(2.1), Inches(12), Inches(0.6),
        font_size=13
    )

    h.add_text_block(
        slide, "解法 2：全域 URL 重寫（嵌入 token）",
        Inches(0.5), Inches(2.9), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """# GitHub 範例
git config --global url."https://x-access-token:YOUR_TOKEN@github.com/acme-corp/plugins".insteadOf "https://github.com/acme-corp/plugins"

# GitLab
git config --global url."https://oauth2:YOUR_TOKEN@gitlab.com/acme-corp/plugins".insteadOf "https://gitlab.com/acme-corp/plugins"

# Bitbucket
git config --global url."https://x-token-auth:YOUR_TOKEN@bitbucket.org/acme-corp/plugins".insteadOf "https://bitbucket.org/acme-corp/plugins" """,
        Inches(0.7), Inches(3.3), Inches(12), Inches(2.3),
        font_size=11
    )

    h.add_callout(
        slide, "⚠️ 重寫以**純文字**儲存在 gitconfig，請用**唯讀** token",
        Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    h.add_text_block(
        slide, "💡 CI/CD：GitHub Actions 匯出 GH_TOKEN → `gh auth setup-git`",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True
    )

    # ============================================================
    # Slide 22：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 6", "為團隊與容器預先配置", "pre-config 與 seed dir")

    # ============================================================
    # Slide 23：為團隊預先配置
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "為團隊預先配置 Marketplace",
        ".claude/settings.json 自動提示安裝",
        slide_num=23, total=TOTAL, source="01 § 為團隊預先配置"
    )

    h.add_code_block(
        slide, """// .claude/settings.json
{
  "extraKnownMarketplaces": {
    "company-tools": {
      "source": {
        "source": "github",
        "repo": "your-org/claude-plugins"
      }
    }
  },
  "enabledPlugins": {
    "code-formatter@company-tools": true,
    "deployment-tools@company-tools": true
  }
}""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(3.0),
        font_size=11
    )

    h.add_bullet_list(
        slide, [
            "團隊成員信任專案資料夾時**自動被提示安裝**",
            "`enabledPlugins` 指定**預設啟用**的 plugin",
            "Marketplace 狀態存於 `~/.claude/plugins/known_marketplaces.json`（每個使用者一份）",
        ],
        Inches(0.7), Inches(4.85), Inches(12), Inches(1.3),
        font_size=13
    )

    h.add_callout(
        slide, "📌 v2.1.195+：來自外部來源的 plugin 在團隊成員安裝前**不會載入**",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 24：為容器預先填充
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "為容器預先填充 Plugin",
        "CI/CD 環境的 seed dir 機制",
        slide_num=24, total=TOTAL, source="01 § 為容器預先填充"
    )

    h.add_code_block(
        slide, """$CLAUDE_CODE_PLUGIN_SEED_DIR/
├── known_marketplaces.json
├── marketplaces/<name>/...
└── cache/<marketplace>/<plugin>/<version>/...""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(1.4),
        font_size=13
    )

    h.add_text_block(
        slide, "🛠️ 兩種建置方式",
        Inches(0.5), Inches(3.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "方式 A：建置時跑安裝再複製",
        [
            "在建置中執行：",
            "`claude plugin marketplace add ...`",
            "`claude plugin install my-tool@...`",
            "然後把 `~/.claude/plugins` 複製到映像",
            "設 SEED_DIR 指向它",
        ],
        "方式 B：直接安裝到目標路徑",
        [
            "用 `CLAUDE_CODE_PLUGIN_CACHE_DIR` 環境變數：",
            "`CLAUDE_CODE_PLUGIN_CACHE_DIR=/opt/claude-seed \\`",
            "  `claude plugin install my-tool@...`",
            "Runtime 設 `CLAUDE_CODE_PLUGIN_SEED_DIR`",
            "跳過複製步驟",
        ],
        top=Inches(3.7), height=Inches(2.3),
        left_color=h.COLOR_BLUE, right_color=h.COLOR_GREEN
    )

    h.add_bullet_list(
        slide, [
            "**唯讀**：種子目錄永遠不會被寫入；自動更新被停用  /  **種子優先**：種子 marketplace 會覆寫使用者配置",
            "**變更被擋**：`/plugin marketplace remove` 或 `update` 會失敗",
        ],
        Inches(0.7), Inches(6.2), Inches(12), Inches(0.8),
        font_size=12
    )

    # ============================================================
    # Slide 25：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 7", "受管 Marketplace 限制", "企業級安全管控")

    # ============================================================
    # Slide 26：strictKnownMarketplaces
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "受管 Marketplace 限制",
        "用 strictKnownMarketplaces 控管來源",
        slide_num=26, total=TOTAL, source="01 § 受管 Marketplace 限制"
    )

    h.add_text_block(
        slide, "🛡️ 限制如何運作",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "在 marketplace **新增**以及 plugin 安裝、更新、重新整理、**自動更新**時檢查",
            "在任何網路或檔案系統操作**之前**檢查",
            "大多數來源類型用**精確匹配**（URL 結尾的 `/`、`.git` 後綴、`ssh://` vs `https://` 視為不同）",
            "受管設定中，**使用者與專案配置無法覆蓋**",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(2.5),
        font_size=13
    )

    h.add_text_block(
        slide, "📋 4 種配置模式",
        Inches(0.5), Inches(4.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["配置", "行為"],
        [
            ["未定義（預設）", "無限制"],
            ["空陣列 `[]`", "**完全鎖定**"],
            ["來源清單", "**只允許白名單**內的 marketplace"],
            ["`hostPattern` / `pathPattern`", "內部 git 伺服器 / 檔案系統路徑"],
        ],
        Inches(0.5), Inches(5.1), Inches(12.333), Inches(1.8),
        font_size=12
    )

    h.add_callout(
        slide, "💡 `strictKnownMarketplaces` 只限制能新增的內容，不會自動註冊；搭配 `extraKnownMarketplaces` 才能自動供應",
        Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3),
        icon="", font_size=10
    )

    # ============================================================
    # Slide 27：實戰配置範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "strictKnownMarketplaces 實戰配置",
        "4 種典型場景",
        slide_num=27, total=TOTAL, source="01 § 限制"
    )

    h.add_text_block(
        slide, "🔒 完全鎖定",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )
    h.add_code_block(
        slide, """{
  "strictKnownMarketplaces": []
}""",
        Inches(0.5), Inches(2.1), Inches(6.0), Inches(0.8),
        font_size=12
    )

    h.add_text_block(
        slide, "📋 只允許白名單",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """{
  "strictKnownMarketplaces": [
    {
      "source": "github",
      "repo": "acme-corp/approved-plugins"
    },
    {
      "source": "github",
      "repo": "acme-corp/security-tools",
      "ref": "v2.0"
    }
  ]
}""",
        Inches(6.733), Inches(2.1), Inches(6.0), Inches(2.3),
        font_size=9
    )

    h.add_text_block(
        slide, "🏢 內部 git 伺服器（hostPattern）",
        Inches(0.5), Inches(4.5), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """{
  "strictKnownMarketplaces": [
    {
      "source": "hostPattern",
      "hostPattern": "^github\\\\.example\\\\.com$"
    }
  ]
}""",
        Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.4),
        font_size=10
    )

    h.add_text_block(
        slide, "📂 檔案系統路徑（pathPattern）",
        Inches(6.733), Inches(4.5), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """{
  "strictKnownMarketplaces": [
    {
      "source": "pathPattern",
      "pathPattern": "^/opt/approved/"
    }
  ]
}""",
        Inches(6.733), Inches(4.9), Inches(6.0), Inches(1.4),
        font_size=10
    )

    h.add_callout(
        slide, "💡 允許任何檔案路徑但仍控管網路來源：pathPattern 用 `.*`  /  搭配 `disableSideloadFlags` (拒絕 CLI 旗標側載) 與 `pluginSuggestionMarketplaces` (目錄相關建議)",
        Inches(0.5), Inches(6.45), Inches(12.333), Inches(0.6),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 28：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 8", "版本、發佈、驗證", "semver、發行通道、CLI 指令")

    # ============================================================
    # Slide 29：版本解析
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "版本解析與發行通道",
        "穩定 / 最新 分流",
        slide_num=29, total=TOTAL, source="01 § 版本解析"
    )

    h.add_text_block(
        slide, "🔍 解析優先順序",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "1️⃣ `plugin.json` 的 version",
            "2️⃣ `marketplace.json` 項目的 version",
            "3️⃣ Plugin 來源的 git commit SHA",
            "⚠️ 同時設 plugin.json 和 marketplace.json 的 version → **前者優先**（且會靜默使用）",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(2.0),
        font_size=14
    )

    h.add_text_block(
        slide, "📊 固定版本的兩種方式",
        Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["方法", "更新行為", "最適合"],
        [
            ["**明確版本**（plugin.json 設 version）", "使用者只在 version 變更時收到更新", "穩定發行週期的已發佈 plugin"],
            ["**Commit-SHA 版本**（省略 version）", "每次新 commit 都被視為新版本", "積極開發中的內部/團隊 plugin"],
        ],
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(1.5),
        font_size=12
    )

    h.add_callout(
        slide, "📌 若用明確版本，遵循 semver（MAJOR.MINOR.PATCH）並在 CHANGELOG.md 記錄",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 30：穩定 vs 最新 範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "穩定 vs 最新 發行通道",
        "用 ref 把使用者分流",
        slide_num=30, total=TOTAL, source="01 § 發行通道"
    )

    h.add_text_block(
        slide, "🟢 stable-tools marketplace",
        Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_GREEN
    )

    h.add_code_block(
        slide, """{
  "name": "stable-tools",
  "plugins": [
    {
      "name": "code-formatter",
      "source": {
        "source": "github",
        "repo": "acme-corp/code-formatter",
        "ref": "stable"
      }
    }
  ]
}""",
        Inches(0.5), Inches(2.1), Inches(6.0), Inches(2.7),
        font_size=11
    )

    h.add_text_block(
        slide, "🔥 latest-tools marketplace",
        Inches(6.733), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """{
  "name": "latest-tools",
  "plugins": [
    {
      "name": "code-formatter",
      "source": {
        "source": "github",
        "repo": "acme-corp/code-formatter",
        "ref": "latest"
      }
    }
  ]
}""",
        Inches(6.733), Inches(2.1), Inches(6.0), Inches(2.7),
        font_size=11
    )

    h.add_bullet_list(
        slide, [
            "每個通道必須解析為**不同版本**（用 `version` → plugin.json 在每個 ref 都要不同；省略 → 不同 SHA）",
            "指派給使用者群組：穩定群組 vs 早期存取群組（用 `extraKnownMarketplaces`）",
            "Plugin 可用 semver 範圍限制依賴（慣例為 `{plugin-name}--v{version}` git 標籤）",
        ],
        Inches(0.7), Inches(4.8), Inches(12), Inches(2.0),
        font_size=12
    )

    # ============================================================
    # Slide 31：重新命名與移除
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "重新命名或移除 Plugin",
        "用 renames 自動遷移",
        slide_num=31, total=TOTAL, source="01 § 重新命名或移除"
    )

    h.add_bullet_list(
        slide, [
            "Plugin 的 `name` 是**穩定識別碼**，更動會破壞現有安裝",
            "想改 UI 顯示但保留識別 → 設 `displayName`、**保持 `name` 不變**",
            "真的必須改 `name` 或從 `plugins` 移除 → 用頂層 `renames` 自動遷移",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(1.5),
        font_size=13
    )

    h.add_code_block(
        slide, """{
  "name": "acme-tools",
  "owner": { "name": "Acme" },
  "plugins": [
    { "name": "code-formatter", "source": "./plugins/code-formatter" }
  ],
  "renames": {
    "formatter": "code-formatter",
    "legacy-linter": null
  }
}""",
        Inches(0.7), Inches(3.4), Inches(12), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "📌 renames 行為",
        Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "指向新名稱 → 載入新名稱並重寫 `enabledPlugins` / `pluginConfigs` 內的舊金鑰  /  指向 `null` → 刪除舊金鑰，通知已從 marketplace 移除",
            "把 `renames` 視為**僅附加歷史**  /  跑 `claude plugin validate .` 避免鏈形成循環",
        ],
        Inches(0.7), Inches(6.4), Inches(12), Inches(0.7),
        font_size=12
    )

    # ============================================================
    # Slide 32：驗證與測試
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "驗證與測試",
        "claude plugin validate 與常見錯誤",
        slide_num=32, total=TOTAL, source="01 § 驗證與測試"
    )

    h.add_code_block(
        slide, """# 驗證 marketplace
claude plugin validate .

# 從 Claude Code 內
/plugin validate .

# 驗證個別 plugin
claude plugin validate ./plugins/my-plugin""",
        Inches(0.7), Inches(1.7), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "🔍 validate 會檢查",
        Inches(6.933), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "marketplace.json 架構錯誤",
            "重複的 plugin 名稱",
            "來源路徑遍歷",
            "每個 plugin 的 plugin.json",
            "項目 version 與 plugin.json 版本一致性",
        ],
        Inches(6.933), Inches(2.1), Inches(6.0), Inches(2.5),
        font_size=12
    )

    h.add_text_block(
        slide, "⚠️ 常見錯誤",
        Inches(0.5), Inches(4.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_comparison_table(
        slide,
        ["錯誤", "原因"],
        [
            ["`File not found: .claude-plugin/marketplace.json`", "缺少 manifest"],
            ["`Invalid JSON syntax: Unexpected token...`", "JSON 語法錯誤"],
            ["`Duplicate plugin name \"x\"`", "plugin 名稱重複"],
            ["`plugins[0].source: Path contains \"..\"`", "來源路徑含 `..`"],
            ["`YAML frontmatter failed to parse`", "skill/agent/command frontmatter 無效"],
        ],
        Inches(0.5), Inches(5.1), Inches(12.333), Inches(1.8),
        font_size=11
    )

    # ============================================================
    # Slide 33：CLI 指令速查
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "CLI 指令速查",
        "非互動式指令（用於腳本/CI）",
        slide_num=33, total=TOTAL, source="01 § CLI 指令參考"
    )

    h.add_code_block(
        slide, """# 新增 marketplace
claude plugin marketplace add <source> [options]
  -s, --scope <user|project|local>   位置（預設 user）
  --sparse <paths...>                 限制 sparse-checkout

# 範例
claude plugin marketplace add acme-corp/claude-plugins         # GitHub
claude plugin marketplace add acme-corp/claude-plugins@v2.0    # 固定 tag
claude plugin marketplace add https://gitlab.example.com/.../plugins.git
claude plugin marketplace add ./my-marketplace                 # 本機
claude plugin marketplace add acme-corp/monorepo --sparse .claude-plugin plugins

# 其他指令
claude plugin marketplace list [--json]                         # 列出
claude plugin marketplace remove <name>                        # 移除（別名 rm）
claude plugin marketplace update [name]                        # 更新""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(4.8),
        font_size=11
    )

    h.add_callout(
        slide, "⚠️ v2.1.196+ 起，未帶協議的主機（如 gitlab.example.com/...）會被拒絕為無效的 owner/repo",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 34：疑難排解速查
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "疑難排解速查",
        "8 個常見問題與解法",
        slide_num=34, total=TOTAL, source="01 § 疑難排解"
    )

    h.add_comparison_table(
        slide,
        ["問題", "解法"],
        [
            ["1. Marketplace 未載入", "驗證 URL 可存取 + 檢查 `.claude-plugin/marketplace.json` 存在"],
            ["2. 驗證錯誤", "跑 `claude plugin validate .` 取得具體錯誤"],
            ["3. Plugin 安裝失敗", "驗證 plugin 來源 URL + repo 公開或有權限"],
            ["4. 私人 repo 驗證失敗", "`gh auth status` 確認登入 + `git config credential.helper`"],
            ["5. 離線環境更新失敗", "`export CLAUDE_CODE_PLUGIN_KEEP_MARKETPLACE_ON_FAILURE=1`"],
            ["6. Git 操作逾時", "`export CLAUDE_CODE_PLUGIN_GIT_TIMEOUT_MS=300000`"],
            ["7. URL-based 相對路徑失敗", "改用 GitHub / npm / git URL，或 git-based marketplace"],
            ["8. 安裝後找不到檔案", "plugin 被複製到快取，參考 plugin 目錄外檔案失效"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.8),
        font_size=11
    )

    h.add_callout(
        slide, "💡 環境變數速查：CLAUDE_CODE_PLUGIN_PREFER_HTTPS、KEEP_MARKETPLACE_ON_FAILURE、GIT_TIMEOUT_MS、SEED_DIR、CACHE_DIR",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="", font_size=11
    )

    # ============================================================
    # Slide 35：重點回顧
    # ============================================================
    h.add_summary_slide(
        slide=h.add_blank_slide(prs),
        title="重點回顧",
        key_points=[
            "**Marketplace** = 別人加入來源、安裝你寫 plugin 的入口",
            "**6 步驟建立本機 marketplace**：目錄 → SKILL.md → plugin.json → marketplace.json → 安裝 → 使用",
            "**5 種 plugin 來源**：相對路徑、GitHub、git URL、git-subdir、npm",
            "**Strict Mode** 控制 plugin.json 是否為元件定義的權威（預設 true）",
            "**strictKnownMarketplaces** 是企業級安全管控的關鍵設定",
        ],
        next_steps=[
            "建立你的第一個本機 marketplace（從 6 步驟範例開始）",
            "把它推上 GitHub，用 `/plugin marketplace add owner/repo` 測試",
            "進階：試著建立 `strict` 模式的多 plugin marketplace",
            "接著閱讀 `02-plugins.md` 深入學習 plugin 開發本身！",
        ],
        source="01-plugin-marketplaces.md"
    )

    # ============================================================
    # 儲存
    # ============================================================
    output = "/home/elan/pi-proj/01-plugin-marketplaces.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
