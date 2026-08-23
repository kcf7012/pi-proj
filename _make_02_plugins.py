"""
簡報 2/4：Plugin 開發入門 (02-plugins.pptx)
約 25 張
對應：02-plugins.md
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 25

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "Plugin 開發入門",
        "從零建立你的第一個 Claude Code Plugin",
        tag="#02 · Plugin 開發"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "建立、測試、除錯你的第一個 plugin",
        slide_num=2, total=TOTAL, source="02-plugins.md"
    )

    objectives = [
        ("🎯", "理解 plugin 的核心概念", "為什麼要用 plugin？什麼時候該用？"),
        ("📦", "建立第一個 plugin", "從零開始的完整步驟"),
        ("🧪", "本地測試 plugin", "用 --plugin-dir 快速迭代"),
        ("🔧", "加入更多元件", "Skills、Agents、Hooks、MCP"),
        ("🚀", "共享給團隊", "轉換獨立配置 + Marketplace"),
        ("🐛", "除錯常見問題", "結構錯誤、權限、路徑問題"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(2.2)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_cols = 3
    grid_rows = 2
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.8)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_PRIMARY
        card.line.width = Pt(2)

        # Icon
        h.add_text_block(
            slide, icon,
            x, y + Inches(0.2), box_w, Inches(0.7),
            font_size=40, align=PP_ALIGN.CENTER
        )
        # 標題
        h.add_text_block(
            slide, title,
            x + Inches(0.2), y + Inches(1.0), box_w - Inches(0.4), Inches(0.5),
            font_size=15, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.5), box_w - Inches(0.4), Inches(0.6),
            font_size=11, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Plugin vs 獨立配置", "什麼時候該升級成 plugin？")

    # ============================================================
    # Slide 4：Plugin vs 獨立配置
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin vs 獨立配置",
        "兩種擴展 Claude Code 的方式",
        slide_num=4, total=TOTAL, source="02 § Plugin vs 獨立配置"
    )

    h.add_comparison_table(
        slide,
        ["面向", "獨立配置（.claude/）", "Plugin"],
        [
            ["Skill 名稱", "/hello（簡短）", "/my-plugin:hello（命名空間）"],
            ["可用範圍", "僅當前專案", "跨專案、跨團隊"],
            ["分享方式", "手動複製", "透過 marketplace 一鍵安裝"],
            ["版本管理", "無", "semver 或 git SHA"],
            ["適合情境", "個人、實驗、單一專案", "團隊、正式、跨專案"],
            ["建議起點", "✓ 先用這個", "確認要共享再升級"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.2),
        font_size=12
    )

    h.add_callout(
        slide, "建議路徑：在 .claude/ 中從獨立配置開始快速迭代，準備好共享時再轉成 plugin",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.5),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 5：什麼時候該用 Plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "什麼時候該升級成 Plugin？",
        "明確的決策依據",
        slide_num=5, total=TOTAL, source="02 § 獨立 vs Plugin"
    )

    h.add_two_column_compare(
        slide,
        "✅ 用 Plugin 的時機",
        [
            "想跟團隊成員共享",
            "需要在多個專案重用",
            "需要版本控制和更新機制",
            "透過 marketplace 發布",
            "接受命名空間化（/plugin-name:hello）",
            "預期要長期維護",
        ],
        "⏸️ 維持獨立的時機",
        [
            "只是個人實驗",
            "只在單一專案使用",
            "不需要分享",
            "希望有簡短 skill 名稱",
            "還在快速迭代階段",
            "短期一次性使用",
        ]
    )

    # ============================================================
    # Slide 6：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "5 步建立第一個 Plugin", "從零到完成的完整流程")

    # ============================================================
    # Slide 7：5 步驟總覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "建立 Plugin 的 5 個步驟",
        "從目錄結構到本地測試",
        slide_num=7, total=TOTAL, source="02 § 快速開始"
    )

    steps = [
        ("1", "建立目錄", "mkdir my-first-plugin", "建立 plugin 根目錄"),
        ("2", "建立 manifest", ".claude-plugin/plugin.json", "定義身份中繼資料"),
        ("3", "加入 skill", "skills/hello/SKILL.md", "第一個可呼叫的指令"),
        ("4", "本地測試", "claude --plugin-dir ./...", "用 CLI 旗標載入"),
        ("5", "嘗試使用", "/my-first-plugin:hello", "驗證功能正常"),
    ]

    step_w = Inches(2.3)
    step_h = Inches(2.5)
    h_gap = Inches(0.2)
    total_w = step_w * 5 + h_gap * 4
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.5)

    for i, (num, title, code, desc) in enumerate(steps):
        x = start_x + i * (step_w + h_gap)

        # 步驟方塊
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y, step_w, step_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        # 步驟編號
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + step_w/2 - Inches(0.3), start_y - Inches(0.25),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = h.COLOR_PRIMARY
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            x + step_w/2 - Inches(0.3), start_y - Inches(0.25),
            Inches(0.6), Inches(0.6),
            font_size=22, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 標題
        h.add_text_block(
            slide, title,
            x, start_y + Inches(0.5), step_w, Inches(0.4),
            font_size=15, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        # 程式碼
        h.add_text_block(
            slide, code,
            x + Inches(0.1), start_y + Inches(1.0),
            step_w - Inches(0.2), Inches(0.6),
            font_size=10, color=h.COLOR_BLUE, font="Consolas",
            align=PP_ALIGN.CENTER
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.1), start_y + Inches(1.7),
            step_w - Inches(0.2), Inches(0.7),
            font_size=10, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

        # 連接箭頭
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + step_w + Inches(0.02), start_y + step_h/2 - Inches(0.1),
                h_gap - Inches(0.04), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    h.add_text_block(
        slide, "⏱ 整個流程 5 分鐘以內",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    # ============================================================
    # Slide 8：Step 1 - 建立目錄
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Step 1：建立 Plugin 目錄",
        "Plugin 從一個資料夾開始",
        slide_num=8, total=TOTAL, source="02 § Step 1"
    )

    h.add_code_block(
        slide, """$ mkdir my-first-plugin
$ cd my-first-plugin

# 預先建立所有需要的子目錄
$ mkdir -p my-first-plugin/.claude-plugin
$ mkdir -p my-first-plugin/skills/hello""",
        Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.0),
        font_size=14
    )

    h.add_text_block(
        slide, "📁 最終目錄結構",
        Inches(1.0), Inches(4.0), Inches(11.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """my-first-plugin/
└── .claude-plugin/        ← 稍後放 plugin.json
└── skills/
    └── hello/             ← 稍後放 SKILL.md""",
        Inches(1.0), Inches(4.5), Inches(11.333), Inches(1.8),
        font_size=14
    )

    h.add_callout(
        slide, "位置不重要，本章用 ./my-first-plugin/，production 時用任何位置都行",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="📍", font_size=13
    )

    # ============================================================
    # Slide 9：Step 2 - 建立 manifest
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Step 2：建立 Plugin Manifest",
        "plugin.json 定義你的身份",
        slide_num=9, total=TOTAL, source="02 § Step 2"
    )

    h.add_code_block(
        slide, """# my-first-plugin/.claude-plugin/plugin.json
{
  "name": "my-first-plugin",
  "description": "A greeting plugin to learn the basics",
  "version": "1.0.0",
  "author": {
    "name": "Your Name"
  }
}""",
        Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5),
        font_size=14
    )

    h.add_text_block(
        slide, "📋 關鍵欄位說明",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "name：唯一識別碼 + skill 命名空間（如 /my-first-plugin:hello）",
            "description：顯示在 /plugin 管理器",
            "version：語義版本，1.0.0 開始；省略則用 git commit SHA",
            "author：選用，標示作者歸屬",
        ],
        Inches(0.7), Inches(4.9), Inches(12), Inches(2.0),
        font_size=13
    )

    # ============================================================
    # Slide 10：Step 3 - 加入 Skill
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Step 3：加入第一個 Skill",
        "SKILL.md 決定 skill 的行為",
        slide_num=10, total=TOTAL, source="02 § Step 3"
    )

    h.add_code_block(
        slide, """# my-first-plugin/skills/hello/SKILL.md
---
description: Greet the user with a friendly message
disable-model-invocation: true
---

Greet the user warmly and ask how you can help them today.""",
        Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5),
        font_size=13
    )

    h.add_text_block(
        slide, "💡 重點：資料夾名稱 = skill 名稱",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "hello/ 資料夾 → /my-first-plugin:hello",
            "前綴：plugin name（my-first-plugin）",
            "後綴：skill 目錄名（hello）",
            "為什麼要命名空間？避免多個 plugin skill 衝突",
            "完整名稱：plugin-name:skill-name",
        ],
        Inches(0.7), Inches(4.9), Inches(12), Inches(2.0),
        font_size=13
    )

    # ============================================================
    # Slide 11：Step 4 - 本地測試
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Step 4：用 --plugin-dir 本地測試",
        "免安裝，直接載入開發中的 plugin",
        slide_num=11, total=TOTAL, source="02 § 本地測試"
    )

    h.add_code_block(
        slide, """# 用 --plugin-dir 旗標啟動
$ claude --plugin-dir ./my-first-plugin

# 一次載入多個 plugins
$ claude --plugin-dir ./plugin-one --plugin-dir ./plugin-two

# 載入 .zip 套件（v2.1.128+）
$ claude --plugin-dir ./my-plugin.zip

# 從 URL 載入
$ claude --plugin-url https://example.com/my-plugin.zip""",
        Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.0),
        font_size=13
    )

    h.add_callout(
        slide, "注意：當 --plugin-dir plugin 與已安裝的市場 plugin 同名時，本地副本優先",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    h.add_text_block(
        slide, "🔄 變更後立即生效",
        Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """/reload-plugins    ← 重新載入，無需重啟""",
        Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.5),
        font_size=14
    )

    # ============================================================
    # Slide 12：Step 5 - 嘗試使用
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Step 5：嘗試使用你的 Skill",
        "看到 /my-first-plugin:hello 成功就算完成！",
        slide_num=12, total=TOTAL, source="02 § Step 5"
    )

    h.add_two_column_compare(
        slide,
        "❶ 試試 skill",
        [
            "在 Claude Code 中執行：",
            "",
            "/my-first-plugin:hello",
            "",
            "或讓 Claude 自動觸發：",
            "",
            "「跟我打個招呼」",
        ],
        "❷ 加上引數",
        [
            "更新 SKILL.md，加入 $ARGUMENTS：",
            "",
            "Greet the user named \"$ARGUMENTS\"",
            "warmly and ask how you can help.",
            "",
            "然後執行：",
            "",
            "/my-first-plugin:hello Alex",
        ]
    )

    h.add_callout(
        slide, "🎉 完成！你已經建立並執行了你的第一個 plugin",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="🎉", font_size=14
    )

    # ============================================================
    # Slide 13：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "Plugin 完整結構", "理解每個目錄的用途")

    # ============================================================
    # Slide 14：完整目錄結構
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "完整 Plugin 目錄結構",
        "標準 enterprise-plugin 的樣貌",
        slide_num=14, total=TOTAL, source="02 § Plugin 結構"
    )

    h.add_code_block(
        slide, """enterprise-plugin/
├── .claude-plugin/           ← 只有 plugin.json
│   └── plugin.json
├── skills/                   ← /name 觸發
│   ├── code-reviewer/
│   └── pdf-processor/
├── commands/                 ← 平面 .md skills
├── agents/                   ← 專門代理
├── hooks/                    ← 事件處理
├── .mcp.json                 ← MCP servers
├── .lsp.json                 ← LSP servers
├── monitors/                 ← 背景監視器
├── bin/                      ← 可執行檔
└── settings.json             ← 預設設定""",
        Inches(0.7), Inches(1.7), Inches(6.5), Inches(5.0),
        font_size=12
    )

    # 右側說明
    h.add_text_block(
        slide, "⚠️ 常見錯誤",
        Inches(7.5), Inches(1.7), Inches(5.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "不要把元件放在 .claude-plugin/ 內",
            "components 必須在 plugin 根目錄",
            "plugin 根目錄 ≠ ~/.claude/",
            "  （例：~/.claude/.mcp.json 沒用）",
            "Plugin 根目錄 = 含 plugin.json 的那層",
        ],
        Inches(7.5), Inches(2.1), Inches(5.5), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 簡化版",
        Inches(7.5), Inches(4.3), Inches(5.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_GREEN
    )

    h.add_bullet_list(
        slide, [
            "可以只放 SKILL.md 在根目錄",
            "不需建立 skills/ 資料夾",
            "Claude Code 自動當作單一 skill 載入",
            "建議：可能擴充時用 skills/ 結構",
        ],
        Inches(7.5), Inches(4.7), Inches(5.5), Inches(2.0),
        font_size=12
    )

    # ============================================================
    # Slide 15：Plugin Manifest 完整欄位
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin Manifest 完整欄位",
        "plugin.json 的所有可用欄位",
        slide_num=15, total=TOTAL, source="03 § Manifest"
    )

    h.add_code_block(
        slide, """{
  "name": "plugin-name",              ← 必需：唯一識別碼
  "version": "1.2.0",                 ← 語義版本
  "description": "...",               ← 顯示用
  "author": { "name": "...", "email": "..." },
  "homepage": "https://...",
  "repository": "https://github.com/...",
  "license": "MIT",
  "keywords": ["deployment", "ci-cd"],
  "defaultEnabled": true,             ← 預設啟用

  // 元件路徑
  "skills": "./custom/skills/",
  "agents": ["./agents/reviewer.md"],
  "hooks": "./config/hooks.json",
  "mcpServers": "./mcp-config.json",
  "lspServers": "./.lsp.json"
}""",
        Inches(0.7), Inches(1.7), Inches(7), Inches(5.2),
        font_size=11
    )

    # 右側速查
    h.add_text_block(
        slide, "📋 欄位分類",
        Inches(8.0), Inches(1.7), Inches(5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "識別：name（必需）",
            "中繼資料：version、author、license",
            "元件路徑：skills/agents/hooks/mcp/lsp",
            "路徑規則：必須以 ./ 開頭",
            "skills 為附加，其他為取代預設",
        ],
        Inches(8.0), Inches(2.1), Inches(5), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 完整 Manifest 架構",
        Inches(8.0), Inches(4.3), Inches(5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_text_block(
        slide, "見 03-plugins-reference.md § Plugin manifest 架構",
        Inches(8.0), Inches(4.7), Inches(5), Inches(0.4),
        font_size=11, color=h.COLOR_GRAY_TXT
    )

    h.add_text_block(
        slide, "完整範例含 userConfig、channels、dependencies 等進階欄位",
        Inches(8.0), Inches(5.0), Inches(5), Inches(0.4),
        font_size=10, color=h.COLOR_GRAY_TXT, italic=True
    )

    h.add_callout(
        slide, "name 必須是 kebab-case（無空格、無大寫）",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="⚠️", font_size=12
    )

    # ============================================================
    # Slide 16：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "加入更多元件", "Skills、Agents、Hooks、MCP")

    # ============================================================
    # Slide 17：加入 Skills
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "加入更多 Skills",
        "Plugin 內可以有多個 skill",
        slide_num=17, total=TOTAL, source="02 § 加入 Skills"
    )

    h.add_code_block(
        slide, """my-plugin/
├── .claude-plugin/
│   └── plugin.json
└── skills/
    ├── code-review/        ← /my-plugin:code-review
    │   ├── SKILL.md
    │   └── examples/
    ├── test-runner/        ← /my-plugin:test-runner
    │   ├── SKILL.md
    │   └── scripts/
    └── deploy-helper/      ← /my-plugin:deploy-helper
        └── SKILL.md""",
        Inches(0.7), Inches(1.7), Inches(7), Inches(3.5),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 設計建議",
        Inches(8.0), Inches(1.7), Inches(5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "每個 skill 是一個獨立資料夾",
            "資料夾名 → skill 命名空間",
            "description 決定 Claude 何時自動觸發",
            "用 disable-model-invocation 控制",
            "保持 SKILL.md 簡潔（< 500 行）",
            "支援檔案放 reference.md、examples/",
        ],
        Inches(8.0), Inches(2.1), Inches(5), Inches(3.5),
        font_size=12
    )

    h.add_text_block(
        slide, "完整 SKILL.md 撰寫指南見 04-skills.md",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True,
        align=PP_ALIGN.CENTER
    )

    h.add_callout(
        slide, "每個 skill 必須有 SKILL.md，支援檔案是選用的",
        Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
        icon="📁", font_size=12
    )

    # ============================================================
    # Slide 18：加入 Agents 與 Hooks
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "加入 Agents 和 Hooks",
        "完整的 plugin 元件組合",
        slide_num=18, total=TOTAL, source="02 § 加入 Agents/Hooks"
    )

    h.add_two_column_compare(
        slide,
        "🤖 Agents",
        [
            "位置：agents/ 目錄",
            "格式：Markdown + frontmatter",
            "範例：",
            "  agents/code-reviewer.md",
            "  agents/security-auditor.md",
            "特性：",
            "  ✓ 隔離 context",
            "  ✓ 專門工具",
            "  ✓ 獨立權限",
        ],
        "🎣 Hooks",
        [
            "位置：hooks/hooks.json",
            "或內聯在 plugin.json",
            "範例：",
            "  hooks/hooks.json",
            "  → 編輯後跑 lint",
            "事件類型：",
            "  PreToolUse / PostToolUse",
            "  SessionStart / Stop 等",
        ]
    )

    h.add_callout(
        slide, "完整 Agents 指南見 05-subagents.md，完整 Hooks 指南見 06-hooks.md",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="📚", font_size=12
    )

    # ============================================================
    # Slide 19：加入 MCP 與 LSP
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "加入 MCP Servers 和 LSP",
        "連接外部服務與程式碼智慧",
        slide_num=19, total=TOTAL, source="02 § 加入 MCP/LSP"
    )

    h.add_code_block(
        slide, """# .mcp.json
{
  "mcpServers": {
    "plugin-database": {
      "command": "${CLAUDE_PLUGIN_ROOT}/servers/db-server",
      "args": ["--config", "${CLAUDE_PLUGIN_ROOT}/config.json"],
      "env": {
        "DB_PATH": "${CLAUDE_PLUGIN_ROOT}/data"
      }
    }
  }
}""",
        Inches(0.7), Inches(1.7), Inches(6), Inches(2.5),
        font_size=11
    )

    h.add_code_block(
        slide, """# .lsp.json
{
  "go": {
    "command": "gopls",
    "args": ["serve"],
    "extensionToLanguage": {
      ".go": "go"
    }
  }
}""",
        Inches(0.7), Inches(4.3), Inches(6), Inches(2.0),
        font_size=12
    )

    # 右側說明
    h.add_text_block(
        slide, "🔌 MCP 整合",
        Inches(7.0), Inches(1.7), Inches(6), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "Plugin 啟用時自動啟動 server",
            "MCP 工具以 mcp__server__tool 命名",
            "Plugin 提供者用 plugin:: 範圍",
            "範例：github、gitlab、slack",
        ],
        Inches(7.0), Inches(2.1), Inches(6), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "🔍 LSP 程式碼智慧",
        Inches(7.0), Inches(4.3), Inches(6), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "即時診斷（編輯後）",
            "跳轉到定義、尋找參考",
            "使用前須安裝語言伺服器",
            "  npm install -g typescript-language-server",
            "  pip install pyright",
        ],
        Inches(7.0), Inches(4.7), Inches(6), Inches(2.0),
        font_size=11
    )

    # ============================================================
    # Slide 20：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "轉換獨立配置 + 除錯", "從 .claude/ 升級到 plugin")

    # ============================================================
    # Slide 21：轉換獨立配置
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "將獨立配置轉換為 Plugin",
        "從 .claude/ 升級到 plugin",
        slide_num=21, total=TOTAL, source="02 § 轉換"
    )

    h.add_code_block(
        slide, """# 1. 建立 plugin 目錄
$ mkdir -p my-plugin/.claude-plugin
$ mkdir my-plugin/hooks

# 2. 建立 manifest
$ cat > my-plugin/.claude-plugin/plugin.json <<EOF
{
  "name": "my-plugin",
  "description": "Migrated from standalone",
  "version": "1.0.0"
}
EOF

# 3. 複製現有配置
$ cp -r .claude/commands my-plugin/
$ cp -r .claude/agents my-plugin/
$ cp -r .claude/skills my-plugin/

# 4. 移動 hooks 到 hooks/hooks.json
# 從 .claude/settings.json 複製 hooks 物件

# 5. 驗證
$ claude --plugin-dir ./my-plugin""",
        Inches(0.7), Inches(1.7), Inches(7.5), Inches(5.0),
        font_size=11
    )

    h.add_text_block(
        slide, "⚠️ 遷移後清理",
        Inches(8.5), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "從 .claude/ 移除原始檔案",
            "避免重複（獨立 vs plugin）",
            "Plugin agents 由 .claude/agents 覆蓋",
            "Plugin skills 命名空間化",
            "  /plugin-name:skill",
            "  與原本的 /skill 各自可用",
        ],
        Inches(8.5), Inches(2.1), Inches(4.5), Inches(2.5),
        font_size=11
    )

    h.add_callout(
        slide, "建議：準備好要分享時才升級，否則維持 .claude/ 較簡單",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    # ============================================================
    # Slide 22：常見問題除錯
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "常見問題與除錯",
        "5 個最常見的失敗原因",
        slide_num=22, total=TOTAL, source="02 § 除錯"
    )

    issues = [
        ("❌ Plugin 未載入",
         "Invalid plugin.json",
         "執行 claude plugin validate ./my-plugin"),
        ("❌ Skills 未出現",
         "目錄結構錯誤",
         "確認 skills/ 在 plugin 根目錄，非 .claude-plugin/ 內"),
        ("❌ Hooks 未觸發",
         "腳本無執行權限",
         "chmod +x ./scripts/your-script.sh"),
        ("❌ MCP 連線失敗",
         "絕對路徑問題",
         "改用 ${CLAUDE_PLUGIN_ROOT} 變數"),
        ("❌ 啟動後找不到 plugin",
         "路徑問題",
         "使用相對路徑以 ./ 開頭"),
    ]

    for i, (title, cause, solution) in enumerate(issues):
        y = Inches(1.7) + i * Inches(0.95)

        # 標題列
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), y, Inches(4.5), Inches(0.8)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = h.COLOR_RED
        title_box.line.fill.background()

        h.add_text_block(
            slide, title,
            Inches(0.7), y, Inches(4.3), Inches(0.8),
            font_size=15, bold=True, color=h.COLOR_WHITE
        )

        # 原因 + 解法
        info_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.0), y, Inches(7.833), Inches(0.8)
        )
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        info_box.line.color.rgb = h.COLOR_GRAY_TXT
        info_box.line.width = Pt(0.5)

        h.add_text_block(
            slide, f"原因：{cause}",
            Inches(5.2), y + Inches(0.05), Inches(7.5), Inches(0.3),
            font_size=12, color=h.COLOR_DARK, bold=True
        )
        h.add_text_block(
            slide, f"解法：{solution}",
            Inches(5.2), y + Inches(0.4), Inches(7.5), Inches(0.4),
            font_size=11, color=h.COLOR_GREEN
        )

    h.add_callout(
        slide, "用 claude --debug 啟動可看到 plugin 載入詳情",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="🔍", font_size=13
    )

    # ============================================================
    # Slide 23：提交到 Marketplace
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "提交到 Marketplace",
        "從個人專案到社群分享",
        slide_num=23, total=TOTAL, source="02 § 提交 Marketplace"
    )

    h.add_text_block(
        slide, "📋 提交前檢查清單",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    checklist = [
        "✓ claude plugin validate 通過（無錯誤）",
        "✓ README.md 含安裝與使用說明",
        "✓ 完整的 plugin.json 欄位",
        "✓ 所有元件測試正常",
        "✓ .claude-plugin/marketplace.json 正確",
        "✓ 在團隊內部試用過",
    ]

    h.add_bullet_list(
        slide, checklist,
        Inches(0.7), Inches(2.2), Inches(12), Inches(2.0),
        font_size=14
    )

    h.add_text_block(
        slide, "🌐 兩個公開 Marketplace",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["Marketplace", "策展者", "提交方式"],
        [
            ["claude-plugins-official", "Anthropic 策展", "不接受提交，Anthropic 自行決定"],
            ["claude-plugins-community", "社群 + Anthropic 審查", "用應用內表單提交"],
        ],
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.5),
        font_size=12
    )

    h.add_text_block(
        slide, "提交表單：claude.ai/admin-settings/directory/submissions/plugins/new",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER, italic=True
    )

    # ============================================================
    # Slide 24：重點回顧
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_summary_slide(
        slide,
        title="重點回顧",
        key_points=[
            "Plugin = 可分享、可版本化、可團隊協作的擴展包裝",
            "5 步建立 plugin：目錄 → manifest → skill → 測試 → 使用",
            "--plugin-dir 是開發測試的關鍵旗標",
            "Plugin 元件：skills、agents、hooks、MCP、LSP、monitors、themes",
            "結構陷阱：元件必須在根目錄，不在 .claude-plugin/ 內",
        ],
        next_steps=[
            "🎯 立即：用 claude plugin init my-helper 建立第一個 plugin",
            "📚 進階：閱讀 03-plugins-reference.md 掌握完整 Manifest 架構",
            "🛠 擴充：嘗試加入 agents 或 hooks",
            "🚀 分享：建立團隊 marketplace.json",
        ],
        source="02-plugins.md"
    )

    # ============================================================
    # Slide 25：結束頁
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide, h.COLOR_BG_CREAM)

    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "建立你的第一個 Plugin！📦"
    run.font.name = h.FONT_TITLE
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = h.COLOR_PRIMARY

    h.add_text_block(
        slide, "從 5 步驟開始，動手做最重要",
        Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
        font_size=20, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.666), Inches(5.0),
        Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = h.COLOR_PRIMARY
    bar.line.fill.background()

    h.add_text_block(
        slide, "下一份簡報：04-skills.pptx（Skills 完整指南）",
        Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
        font_size=14, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "Claude Code Plugin 完整學習系列 · #02",
        Inches(1), Inches(6.1), Inches(11.333), Inches(0.5),
        font_size=12, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    # 儲存
    output = "/home/elan/pi-proj/02-plugins.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
