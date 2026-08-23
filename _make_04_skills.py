"""
簡報 3/4：Skills 完整指南 (04-skills.pptx)
約 40 張
對應：04-skills.md
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 40

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "Skills 完整指南",
        "可重複使用的指令、知識、工作流",
        tag="#04 · Skills"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "從基礎到進階的完整 Skills 知識",
        slide_num=2, total=TOTAL, source="04-skills.md"
    )

    objectives = [
        ("📚", "Skill 是什麼", "理解兩種角色：參考型 vs 任務型"),
        ("🚀", "建立第一個 Skill", "5 分鐘可完成的完整流程"),
        ("📍", "Skill 範圍管理", "企業/個人/專案/Plugin 的優先級"),
        ("🎨", "Frontmatter 完整語法", "所有可用欄位與用途"),
        ("🔧", "字串替換與引數", "$ARGUMENTS、$name、動態變數"),
        ("🔄", "動態上下文注入", "!`command` 預處理技巧"),
        ("🎯", "控制叫用權限", "user-invocable / disable-model-invocation"),
        ("⚙️", "Subagent 整合", "context: fork、agent 欄位"),
        ("📊", "評估與改進", "skill-creator 與基準比較"),
    ]

    box_w = Inches(3.0)
    box_h = Inches(1.5)
    h_gap = Inches(0.3)
    v_gap = Inches(0.25)
    grid_cols = 3
    grid_rows = 3
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap * (grid_rows - 1)
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.7)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_PRIMARY
        card.line.width = Pt(1.5)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.1), Inches(0.6), Inches(0.5),
            font_size=24, align=PP_ALIGN.LEFT
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.7), y + Inches(0.1),
            box_w - Inches(0.8), Inches(0.5),
            font_size=13, bold=True, color=h.COLOR_DARK
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(0.7),
            box_w - Inches(0.4), Inches(0.7),
            font_size=10, color=h.COLOR_GRAY_TXT
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Skill 基礎概念", "理解 skill 是什麼")

    # ============================================================
    # Slide 4：什麼是 Skill
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "什麼是 Skill？",
        "一個 SKILL.md 檔案改變 Claude 的行為",
        slide_num=4, total=TOTAL, source="04 § 什麼是 Skill"
    )

    h.add_text_block(
        slide, "📚 Skill 是一個資料夾，包含 SKILL.md 和支援檔案",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.5),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_code_block(
        slide, """skills/
└── summarize-changes/
    ├── SKILL.md          ← 必需：YAML frontmatter + 說明
    ├── reference.md       ← 選用：詳細參考
    └── scripts/           ← 選用：可執行腳本
        └── helper.py""",
        Inches(2.0), Inches(2.5), Inches(9.333), Inches(2.0),
        font_size=12
    )

    h.add_two_column_compare(
        slide,
        "✅ Skill 的好處",
        [
            "可重用：寫一次，到處用",
            "知識庫：Claude 隨時可參考",
            "SOP 標準化：團隊流程統一",
            "可分享：透過 plugin 發布",
            "低 context 成本：按需載入",
        ],
        "🆚 與 CLAUDE.md 差異",
        [
            "CLAUDE.md：每次 session 自動載入",
            "Skill：手動觸發或按需載入",
            "CLAUDE.md：放「always do X」",
            "Skill：放「有時需要的知識」",
            "建議：CLAUDE.md < 200 行",
        ],
        top=Inches(4.7), height=Inches(2.3)
    )

    # ============================================================
    # Slide 5：兩種角色
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Skill 的兩種角色",
        "參考型 vs 任務型",
        slide_num=5, total=TOTAL, source="04 § Skill 類型"
    )

    # 兩個大卡片
    cards = [
        ("📖", "參考型（Reference）", h.COLOR_BLUE,
         "新增 Claude 應用於工作的知識",
         [
             "API 慣例、命名規範",
             "風格指南、領域知識",
             "Claude 在整個 session 內聯使用",
             "通常不設 disable-model-invocation",
             "description 描述「何時使用」",
         ],
         """---
name: api-conventions
description: API design patterns
---

When writing API endpoints:
- Use RESTful naming
- Return consistent error formats
- Include request validation"""),
        ("⚡", "任務型（Action）", h.COLOR_PRIMARY,
         "為 Claude 提供特定動作的逐步說明",
         [
             "部署、提交、程式碼生成",
             "用 /skill-name 直接叫用",
             "通常設 disable-model-invocation: true",
             "防止 Claude 自動觸發",
             "只有你想執行時才跑",
         ],
         """---
name: deploy
description: Deploy to production
disable-model-invocation: true
context: fork
---

Deploy $ARGUMENTS to production:
1. Run the test suite
2. Build the application
3. Push to deployment target""")
    ]

    for i, (icon, title, color, role, points, code) in enumerate(cards):
        x = Inches(0.5) + i * Inches(6.4)
        y = Inches(1.7)
        w = Inches(6.0)
        h_card = Inches(5.0)

        # 主卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h_card)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # 標題列
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.6))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = color
        title_bar.line.fill.background()
        h.add_text_block(
            slide, f"{icon}  {title}",
            x, y, w, Inches(0.6),
            font_size=18, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 角色描述
        h.add_text_block(
            slide, role,
            x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(0.4),
            font_size=13, color=h.COLOR_DARK, bold=True
        )

        # 特點清單
        h.add_bullet_list(
            slide, points,
            x + Inches(0.2), y + Inches(1.1), w - Inches(0.4), Inches(1.5),
            font_size=11
        )

        # 範例
        h.add_text_block(
            slide, "範例：",
            x + Inches(0.2), y + Inches(2.7), w - Inches(0.4), Inches(0.3),
            font_size=11, bold=True, color=h.COLOR_PRIMARY
        )
        h.add_code_block(
            slide, code,
            x + Inches(0.2), y + Inches(3.0), w - Inches(0.4), Inches(1.8),
            font_size=9
        )

    # ============================================================
    # Slide 6：快速開始
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "快速開始：建立你的第一個 Skill",
        "總共 3 步，5 分鐘完成",
        slide_num=6, total=TOTAL, source="04 § 快速開始"
    )

    steps = [
        ("1", "建立目錄",
         "mkdir -p ~/.claude/skills/summarize-changes",
         "個人 skill 在 ~/.claude/skills/ 下，所有專案可用"),
        ("2", "建立 SKILL.md",
         "見右側程式碼",
         "YAML frontmatter + Markdown 內容"),
        ("3", "測試",
         "/summarize-changes 或讓 Claude 自動觸發",
         "用 git diff HEAD 預處理即時資料"),
    ]

    step_w = Inches(3.8)
    step_h = Inches(3.0)
    h_gap = Inches(0.3)
    total_w = step_w * 3 + h_gap * 2
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(1.7)

    for i, (num, title, code, desc) in enumerate(steps):
        x = start_x + i * (step_w + h_gap)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        # 編號圓
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + step_w/2 - Inches(0.35), start_y - Inches(0.2), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = h.COLOR_PRIMARY
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            x + step_w/2 - Inches(0.35), start_y - Inches(0.2), Inches(0.7), Inches(0.7),
            font_size=24, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

        # 標題
        h.add_text_block(
            slide, title,
            x, start_y + Inches(0.6), step_w, Inches(0.5),
            font_size=16, bold=True, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )

        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), start_y + Inches(1.1), step_w - Inches(0.4), Inches(0.5),
            font_size=11, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )

        # 程式碼（若提供）
        if code:
            h.add_code_block(
                slide, code,
                x + Inches(0.2), start_y + Inches(1.7), step_w - Inches(0.4), Inches(1.2),
                font_size=10
            )

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_w + Inches(0.05), start_y + step_h/2 - Inches(0.15), h_gap - Inches(0.1), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    # 下方：完整 SKILL.md 範例
    h.add_text_block(
        slide, "完整 SKILL.md 範例",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """---
description: Summarizes uncommitted changes and flags anything risky. Use when the user asks what changed.
---

## Current changes
!`git diff HEAD`           ← 預處理：執行命令並替換

## Instructions
Summarize the changes above in 2-3 bullets, then list risks
such as missing error handling or hardcoded values.""",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(1.3),
        font_size=11
    )

    h.add_callout(
        slide, "完成！/summarize-changes 即可用",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="✅", font_size=13
    )

    # ============================================================
    # Slide 7：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "Skill 位置與範圍", "4 個層級與優先級規則")

    # ============================================================
    # Slide 8：Skill 範圍
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Skill 的 4 個位置",
        "決定 skill 誰能用、優先順序如何",
        slide_num=8, total=TOTAL, source="04 § Skill 位置"
    )

    scopes = [
        ("🏢", "企業", "受管設定", "組織內所有使用者", h.COLOR_RED, 1),
        ("👤", "個人", "~/.claude/skills/", "你的所有專案", h.COLOR_BLUE, 2),
        ("📁", "專案", ".claude/skills/", "僅此專案", h.COLOR_GREEN, 3),
        ("📦", "外掛", "plugin/skills/", "啟用 plugin 的位置", h.COLOR_PRIMARY, 4),
    ]

    box_w = Inches(2.9)
    box_h = Inches(3.0)
    h_gap = Inches(0.3)
    total_w = box_w * 4 + h_gap * 3
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(1.7)

    for i, (icon, name, path, scope, color, priority) in enumerate(scopes):
        x = start_x + i * (box_w + h_gap)
        y = start_y

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # 標題
        h.add_text_block(
            slide, f"{icon} {name}",
            x, y + Inches(0.2), box_w, Inches(0.5),
            font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER
        )
        # 優先級
        h.add_text_block(
            slide, f"優先級 #{priority}",
            x, y + Inches(0.7), box_w, Inches(0.3),
            font_size=11, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
        )
        # 路徑
        h.add_code_block(
            slide, path,
            x + Inches(0.2), y + Inches(1.1), box_w - Inches(0.4), Inches(0.6),
            font_size=10
        )
        # 範圍
        h.add_text_block(
            slide, scope,
            x + Inches(0.2), y + Inches(1.9), box_w - Inches(0.4), Inches(0.9),
            font_size=12, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "優先級規則：企業 > 個人 > 專案。同名時，較高優先級的 skill 勝出",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        icon="📌", font_size=13
    )

    h.add_text_block(
        slide, "Plugin skills 用 plugin-name:skill-name 命名空間，永不與其他層級衝突",
        Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "巢狀目錄：子目錄中的 .claude/skills/ 也會載入（用相對路徑限定）",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        font_size=11, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
    )

    # ============================================================
    # Slide 9：即時變更偵測
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "即時變更偵測",
        "改完即生效，無需重啟",
        slide_num=9, total=TOTAL, source="04 § 即時變更"
    )

    h.add_text_block(
        slide, "⚡ Claude Code 會監視 skill 目錄",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.5),
        font_size=20, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_two_column_compare(
        slide,
        "✅ 自動生效",
        [
            "~/.claude/skills/ 內的變更",
            "專案 .claude/skills/ 內的變更",
            "--add-dir 目錄的 .claude/skills/",
            "新增、編輯、移除 skill",
        ],
        "⚠️ 需重啟",
        [
            "建立全新的頂級 skills 目錄",
            "Plugin 元件：hooks、.mcp.json",
            "agents、output-styles 變更",
            "（需 /reload-plugins）",
        ]
    )

    h.add_text_block(
        slide, "即時變更偵測僅涵蓋 SKILL.md 文字",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True,
        align=PP_ALIGN.CENTER
    )

    h.add_callout(
        slide, "想即時生效 plugin 元件：/reload-plugins",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="🔄", font_size=13
    )

    # ============================================================
    # Slide 10：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "Frontmatter 完整語法", "所有可用欄位")

    # ============================================================
    # Slide 11：Frontmatter 完整參考（上）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Frontmatter 完整欄位（上）",
        "基礎設定與可見性",
        slide_num=11, total=TOTAL, source="04 § Frontmatter"
    )

    h.add_comparison_table(
        slide,
        ["欄位", "必需", "描述"],
        [
            ["name", "❌", "顯示名稱（預設為目錄名）"],
            ["description", "✅ 建議", "Skill 功能與何時使用。Claude 用此決定自動觸發"],
            ["when_to_use", "❌", "額外觸發上下文，附加到 description"],
            ["argument-hint", "❌", "自動完成提示，例如 [issue-number]"],
            ["disable-model-invocation", "❌", "true = 防止 Claude 自動載入（任務型用）"],
            ["user-invocable", "❌", "false = 從 / 功能表隱藏（背景知識用）"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.5),
        font_size=12
    )

    h.add_callout(
        slide, "description 是最重要的欄位！決定 Claude 是否自動觸發。1,536 字元上限",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        icon="📌", font_size=12
    )

    # ============================================================
    # Slide 12：Frontmatter 完整參考（下）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Frontmatter 完整欄位（下）",
        "進階控制選項",
        slide_num=12, total=TOTAL, source="04 § Frontmatter"
    )

    h.add_comparison_table(
        slide,
        ["欄位", "描述"],
        [
            ["allowed-tools", "skill 啟用時免批准的工具（Bash(git *)）"],
            ["disallowed-tools", "skill 啟用時移除的工具"],
            ["model", "覆蓋模型（sonnet/opus/haiku/inherit）"],
            ["effort", "覆蓋努力級別（low/medium/high/xhigh/max）"],
            ["context", "fork = 在 subagent 中執行"],
            ["agent", "context: fork 時使用的 subagent 類型"],
            ["hooks", "限定 skill 生命週期的 hooks"],
            ["paths", "Glob 模式限制自動載入時機"],
            ["shell", "bash 或 powershell（用於 !`command`）"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.8),
        font_size=11
    )

    h.add_callout(
        slide, "建議：先用最少欄位（description），需要時再加 disable-model-invocation",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    # ============================================================
    # Slide 13：Frontmatter 範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Frontmatter 實戰範例",
        "不同用途的 frontmatter 寫法",
        slide_num=13, total=TOTAL, source="04 § Frontmatter"
    )

    # 4 個範例
    examples = [
        ("📖 標準 reference skill", h.COLOR_BLUE,
         """---
description: 撰寫 API 端點的慣例
---

# API Conventions
- Use RESTful naming
- Return consistent errors"""),
        ("⚡ 任務型 skill", h.COLOR_PRIMARY,
         """---
description: 部署到 production
disable-model-invocation: true
---

Deploy $ARGUMENTS to production:
1. Run tests
2. Build
3. Push"""),
        ("🔧 帶引數", h.COLOR_GREEN,
         """---
description: 修復 GitHub issue
argument-hint: [issue-number]
arguments: [issue]
---

Fix issue $issue following
our coding standards."""),
        ("🤖 在 subagent 中", RGBColor(0x7C, 0x3A, 0xED),
         """---
description: 深度研究
context: fork
agent: Explore
---

Research $ARGUMENTS thoroughly
and summarize findings.""")
    ]

    box_w = Inches(5.8)
    box_h = Inches(2.4)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)

    for i, (label, color, code) in enumerate(examples):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * (box_w + h_gap)
        y = Inches(1.7) + row * (box_h + v_gap)

        # 標題
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.4))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = color
        title_bar.line.fill.background()
        h.add_text_block(
            slide, label,
            x, y, box_w, Inches(0.4),
            font_size=14, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 程式碼
        h.add_code_block(
            slide, code,
            x, y + Inches(0.5), box_w, box_h - Inches(0.5),
            font_size=10
        )

    # ============================================================
    # Slide 14：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "字串替換與動態內容", "$ARGUMENTS 與 !`command`")

    # ============================================================
    # Slide 15：字串替換
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "可用的字串替換",
        "動態內容的 8 個變數",
        slide_num=15, total=TOTAL, source="04 § 字串替換"
    )

    h.add_comparison_table(
        slide,
        ["變數", "描述"],
        [
            ["$ARGUMENTS", "skill 叫用時的所有引數"],
            ["$ARGUMENTS[N] / $N", "0-based 索引的第 N 個引數"],
            ["$name", "frontmatter arguments 清單宣告的具名引數"],
            ["${CLAUDE_SESSION_ID}", "目前 session ID（用於記錄）"],
            ["${CLAUDE_EFFORT}", "目前努力級別（low/medium/high/xhigh/max）"],
            ["${CLAUDE_SKILL_DIR}", "skill 的 SKILL.md 所在目錄"],
            ["${CLAUDE_PROJECT_DIR}", "專案根目錄（與 hooks 相同）"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.5),
        font_size=12
    )

    h.add_code_block(
        slide, """---
name: session-logger
description: Log activity for this session
---

Log to logs/${CLAUDE_SESSION_ID}.log:

$ARGUMENTS""",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.7),
        font_size=12
    )

    h.add_callout(
        slide, "轉義：$1.00 用 \\$1.00。雙引號包多字引數：/my-skill \"hello world\" second",
        Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.3),
        icon="💡", font_size=11
    )

    # ============================================================
    # Slide 16：引數與具名引數
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "傳遞引數給 Skill",
        "基本、索引、具名三種方式",
        slide_num=16, total=TOTAL, source="04 § 引數"
    )

    h.add_text_block(
        slide, "❶ 基本：$ARGUMENTS",
        Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """Fix GitHub issue $ARGUMENTS
following our standards.

→ /fix-issue 123
  "Fix GitHub issue 123..." """,
        Inches(0.5), Inches(2.1), Inches(6), Inches(1.6),
        font_size=11
    )

    h.add_text_block(
        slide, "❷ 索引：$ARGUMENTS[N] / $N",
        Inches(0.5), Inches(3.8), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """Migrate $0 from $1 to $2.

→ /migrate SearchBar React Vue
  $0=SearchBar, $1=React, $2=Vue""",
        Inches(0.5), Inches(4.2), Inches(6), Inches(1.6),
        font_size=11
    )

    h.add_text_block(
        slide, "❸ 具名：$name",
        Inches(6.833), Inches(1.7), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """---
arguments: [issue, branch]
---

Review PR for issue $issue
on branch $branch.

→ /review-pr 421 main
  $issue=421, $branch=main""",
        Inches(6.833), Inches(2.1), Inches(6), Inches(1.6),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 堆疊多個 skills",
        Inches(6.833), Inches(3.8), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """/code-review /fix-issue 123
  載入兩個 skills，
  123 作為 $ARGUMENTS
  傳遞給每個""",
        Inches(6.833), Inches(4.2), Inches(6), Inches(1.6),
        font_size=11
    )

    # ============================================================
    # Slide 17：動態上下文注入
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "動態上下文注入：!`command`",
        "執行 shell 命令，結果替換佔位符",
        slide_num=17, total=TOTAL, source="04 § 動態上下文"
    )

    h.add_text_block(
        slide, "⚡ 預處理機制",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_text_block(
        slide, "Claude Code 在看到 skill 內容之前執行 !`command`，輸出替換佔位符。Claude 永遠只看到最終結果。",
        Inches(0.5), Inches(2.2), Inches(12.333), Inches(0.6),
        font_size=14, color=h.COLOR_DARK
    )

    # 範例 1
    h.add_text_block(
        slide, "❶ 內聯單行",
        Inches(0.5), Inches(3.0), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )
    h.add_code_block(
        slide, """## Pull request context
- PR diff: !`gh pr diff`
- Comments: !`gh pr view --comments`
- Changed files: !`gh pr diff --name-only`

→ gh CLI 執行，結果插入""",
        Inches(0.5), Inches(3.4), Inches(6), Inches(2.2),
        font_size=10
    )

    # 範例 2
    h.add_text_block(
        slide, "❷ 多行圍欄區塊",
        Inches(6.833), Inches(3.0), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )
    h.add_code_block(
        slide, """## Environment
```!
node --version
npm --version
git status --short
```

→ 多行命令用 ```! 開啟""",
        Inches(6.833), Inches(3.4), Inches(6), Inches(2.2),
        font_size=10
    )

    h.add_text_block(
        slide, "⚠️ 限制：替換對原始檔案執行一次，命令輸出不會被重新掃描",
        Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_RED
    )

    h.add_text_block(
        slide, "⚠️ ! 必須在行首或緊接在空白後（KEY=!`cmd` 不會執行）",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        font_size=11, color=h.COLOR_GRAY_TXT
    )

    h.add_callout(
        slide, "用 disableSkillShellExecution: true 停用 shell 執行（受管設定用）",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="🔒", font_size=12
    )

    # ============================================================
    # Slide 18：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "支援檔案與目錄結構", "建立可維護的 skill")

    # ============================================================
    # Slide 19：支援檔案
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "新增支援檔案",
        "用輔助檔案保持 SKILL.md 簡潔",
        slide_num=19, total=TOTAL, source="04 § 支援檔案"
    )

    h.add_code_block(
        slide, """my-skill/
├── SKILL.md           ← 必需 - 概覽與導航
├── reference.md       ← 詳細 API 文件（必要時載入）
├── examples.md        ← 使用範例（必要時載入）
└── scripts/
    └── validate.sh    ← 工具腳本（執行而非載入）""",
        Inches(0.7), Inches(1.7), Inches(7), Inches(3.0),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 設計原則",
        Inches(8.0), Inches(1.7), Inches(5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "SKILL.md 保持在 500 行以下",
            "詳細資料移至 reference.md",
            "範例移至 examples.md",
            "腳本放 scripts/ 供 Claude 執行",
            "從 SKILL.md 明確參考支援檔案",
        ],
        Inches(8.0), Inches(2.1), Inches(5), Inches(2.5),
        font_size=12
    )

    h.add_code_block(
        slide, """## Additional resources

- For complete API details, see [reference.md](reference.md)
- For usage examples, see [examples.md](examples.md)""",
        Inches(0.7), Inches(5.0), Inches(12.333), Inches(1.0),
        font_size=12
    )

    h.add_callout(
        slide, "每多一行都會增加 token 成本，支援檔案只在需要時載入",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        icon="💰", font_size=13
    )

    # ============================================================
    # Slide 20：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 6", "控制誰能叫用 Skill", "user-invocable 與 disable-model-invocation")

    # ============================================================
    # Slide 21：控制叫用權限
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "控制誰能叫用 Skill",
        "user-invocable 與 disable-model-invocation",
        slide_num=21, total=TOTAL, source="04 § 控制叫用"
    )

    h.add_comparison_table(
        slide,
        ["Frontmatter", "你可叫用", "Claude 可叫用", "何時載入"],
        [
            ["（預設）", "✅", "✅", "描述始終在 context；叫用時載入完整內容"],
            ["disable-model-invocation: true", "✅", "❌", "描述不在 context；你叫用時才載入完整內容"],
            ["user-invocable: false", "❌", "✅", "描述始終在 context；叫用時載入完整內容"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=12
    )

    h.add_text_block(
        slide, "📖 參考型 vs ⚡ 任務型",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "參考型（知識）",
        [
            "（無特殊設定）",
            "Claude 在需要時自動載入",
            "user-invocable: false",
            "  → /legacy-context 對使用者無意義",
        ],
        "任務型（動作）",
        [
            "disable-model-invocation: true",
            "Claude 永遠不會自動跑",
            "只有你 /deploy 才會執行",
            "避免「程式碼看起來 OK 就部署」",
        ],
        top=Inches(5.8), height=Inches(1.4)
    )

    # ============================================================
    # Slide 22：隱藏與覆寫
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "隱藏與覆寫 Skill 可見性",
        "skillOverrides 與 /permissions 控制",
        slide_num=22, total=TOTAL, source="04 § skillOverrides"
    )

    h.add_text_block(
        slide, "從設定覆寫可見性",
        Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """{
  "skillOverrides": {
    "legacy-context": "name-only",
    "deploy": "off"
  }
}""",
        Inches(0.5), Inches(2.1), Inches(6), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "四種狀態",
        Inches(0.5), Inches(4.2), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            '"on" — 名稱 + 描述都給 Claude',
            '"name-only" — 只給名稱',
            '"user-invocable-only" — 給你但藏 Claude',
            '"off" — 完全隱藏',
        ],
        Inches(0.5), Inches(4.6), Inches(6), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "限制 Claude 叫用哪些 skills",
        Inches(6.833), Inches(1.7), Inches(6), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """# 在 /permissions deny 規則：
Skill(deploy *)        # 拒絕特定 skill
Skill                 # 拒絕所有 skills

# 允許特定 skills：
Skill(commit)
Skill(review-pr *)""",
        Inches(6.833), Inches(2.1), Inches(6), Inches(2.5),
        font_size=11
    )

    h.add_callout(
        slide, "Plugin skills 不受 skillOverrides 影響，請用 /plugin 管理",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="⚠️", font_size=12
    )

    # ============================================================
    # Slide 23：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 7", "進階功能", "工具權限、Subagent、視覺輸出")

    # ============================================================
    # Slide 24：預先批准工具
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "為 Skill 預先批准工具",
        "allowed-tools 的使用",
        slide_num=24, total=TOTAL, source="04 § 預先批准工具"
    )

    h.add_code_block(
        slide, """---
name: commit
description: Stage and commit current changes
disable-model-invocation: true
allowed-tools: Bash(git add *) Bash(git commit *) Bash(git status *)
---

# 自動批准 git 命令，無需每次都按確認""",
        Inches(0.7), Inches(1.7), Inches(12.333), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "🔧 allowed-tools vs disallowed-tools",
        Inches(0.5), Inches(3.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "allowed-tools（白名單）",
        [
            "列出 skill 可用的工具",
            "Bash(git add *)",
            "Bash(gh pr create *)",
            "Read、Grep、Glob",
            "會覆蓋 session 的工具設定",
        ],
        "disallowed-tools（黑名單）",
        [
            "從可用池中移除工具",
            "AskUserQuestion",
            "WebSearch",
            "Write、Edit（限制寫入）",
            "限制在下次訊息後清除",
        ],
        top=Inches(4.4), height=Inches(2.0)
    )

    h.add_callout(
        slide, "專案簽入的 skills 中，allowed-tools 需在信任工作區後才生效",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="🔒", font_size=12
    )

    # ============================================================
    # Slide 25：Skill 在 Subagent 中執行
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "在 Subagent 中執行 Skill",
        "context: fork 隔離執行",
        slide_num=25, total=TOTAL, source="04 § Subagent 整合"
    )

    h.add_code_block(
        slide, """---
name: deep-research
description: Research a topic thoroughly
context: fork          ← 在 subagent 中執行
agent: Explore        ← 使用 Explore subagent
---

Research $ARGUMENTS thoroughly:

1. Find relevant files using Glob and Grep
2. Read and analyze the code
3. Summarize findings with specific file references""",
        Inches(0.7), Inches(1.7), Inches(12.333), Inches(2.5),
        font_size=12
    )

    h.add_text_block(
        slide, "🔄 Skills 與 Subagents 雙向整合",
        Inches(0.5), Inches(4.4), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["方向", "系統提示", "任務", "也載入"],
        [
            ["Skill 帶 context: fork", "來自 agent 類型", "SKILL.md 內容", "CLAUDE.md（除非是 Explore/Plan）"],
            ["Subagent 帶 skills 欄位", "Subagent 的 markdown", "Claude 的委派訊息", "預載入的 skills + CLAUDE.md"],
        ],
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(1.8),
        font_size=11
    )

    h.add_callout(
        slide, "context: fork 僅對「有明確任務」的 skill 有意義；純知識的 skill 會浪費 subagent 呼叫",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    # ============================================================
    # Slide 26：視覺輸出範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "視覺輸出：生成互動式 HTML",
        "skills 可以執行腳本並輸出視覺化結果",
        slide_num=26, total=TOTAL, source="04 § 視覺輸出"
    )

    h.add_code_block(
        slide, """---
name: codebase-visualizer
description: Generate interactive tree visualization
allowed-tools: Bash(python3 *)
---

# Codebase Visualizer

Generate an interactive HTML tree view of your project structure.

## Usage

```bash
python3 ${CLAUDE_SKILL_DIR}/scripts/visualize.py .
```

Creates codebase-map.html and opens in browser.""",
        Inches(0.7), Inches(1.7), Inches(7), Inches(3.0),
        font_size=11
    )

    h.add_text_block(
        slide, "📊 視覺化的應用場景",
        Inches(8.0), Inches(1.7), Inches(5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "📂 程式碼結構視覺化",
            "📈 相依性圖表",
            "🧪 測試涵蓋率報告",
            "📋 API 文件",
            "🗄️ 資料庫架構圖",
        ],
        Inches(8.0), Inches(2.1), Inches(5), Inches(2.5),
        font_size=11
    )

    h.add_callout(
        slide, "腳本用 ${CLAUDE_SKILL_DIR} 引用支援檔案，無論安裝在哪都正確解析",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        icon="📁", font_size=13
    )

    h.add_text_block(
        slide, "完整範例：04-skills.md § 視覺輸出範例（含 250 行 Python 程式碼）",
        Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
    )

    h.add_callout(
        slide, "skill 可以生成任何格式：HTML、SVG、PDF、Markdown 報告",
        Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
        icon="🎨", font_size=13
    )

    # ============================================================
    # Slide 27：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 8", "技能生命週期與評估", "如何知道 skill 有效？")

    # ============================================================
    # Slide 28：Skill 內容生命週期
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Skill 內容生命週期",
        "載入後會一直在 context 中",
        slide_num=28, total=TOTAL, source="04 § 內容生命週期"
    )

    h.add_text_block(
        slide, "📌 一旦載入，內容保持整個 session",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.5),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_bullet_list(
        slide, [
            "叫用 skill → 內容作為單一訊息進入對話",
            "整個 session 保持有效（不重新讀取檔案）",
            "重複叫用相同內容 → 簡短註記「已載入」",
            "重複叫用不同內容（不同引數）→ 附加新副本",
            "Auto-compact 保留每個 skill 前 5,000 tokens",
            "所有 skill 共享 25,000 tokens 預算",
            "compaction 後重新叫用可恢復完整內容",
        ],
        Inches(0.7), Inches(2.5), Inches(12), Inches(3.5),
        font_size=13
    )

    h.add_callout(
        slide, "建議：保持 skill 簡潔，每行都是重複的 token 成本",
        Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
        icon="💰", font_size=13
    )

    # ============================================================
    # Slide 29：評估與改進
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "評估與改進 Skill",
        "用 skill-creator 量化效果",
        slide_num=29, total=TOTAL, source="04 § 評估"
    )

    h.add_text_block(
        slide, "🎯 觸發 ≠ 有效",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "看到 skill 觸發告訴你 Claude 找到了它，不是它做了你想要的",
        Inches(0.5), Inches(2.2), Inches(12.333), Inches(0.4),
        font_size=13, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "用 skill-creator 量化評估",
        Inches(0.5), Inches(2.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """/plugin install skill-creator@claude-plugins-official
/reload-plugins
evaluate my summarize-changes skill with skill-creator""",
        Inches(0.5), Inches(3.3), Inches(12.333), Inches(1.3),
        font_size=12
    )

    h.add_text_block(
        slide, "skill-creator 提供的功能",
        Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "📊 量化指標",
        [
            "測試案例（evals/evals.json）",
            "隔離執行（每個 subagent）",
            "通過率 / 時間 / tokens",
            "版本 A/B 盲測",
            "HTML 報告",
        ],
        "🔍 描述調整",
        [
            "生成應觸發/不應觸發的提示",
            "測試命中率",
            "建議描述編輯",
            "改進 description 提升匹配度",
        ],
        top=Inches(5.3), height=Inches(1.6)
    )

    # ============================================================
    # Slide 30：疑難排解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "疑難排解",
        "4 個最常見的 skill 問題",
        slide_num=30, total=TOTAL, source="04 § 疑難排解"
    )

    issues = [
        ("🔇 Skill 未觸發",
         "原因：描述模糊、缺少關鍵字、YAML 格式錯",
         "解法：加關鍵字 / 修 YAML / /debug 看錯誤"),
        ("🔊 Skill 觸發過於頻繁",
         "原因：描述太寬泛",
         "解法：更具體 / 加 disable-model-invocation"),
        ("✂️ 描述被截斷",
         "原因：太多 skills 擠壓清單",
         "解法：提高 skillListingBudgetFraction"),
        ("🛑 Skill 載入後無效",
         "原因：compaction 後內容被丟",
         "解法：compaction 後重新叫用恢復"),
    ]

    for i, (title, cause, solution) in enumerate(issues):
        y = Inches(1.7) + i * Inches(1.2)

        # 標題列
        title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(4.0), Inches(1.0))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = h.COLOR_RED
        title_box.line.fill.background()
        h.add_text_block(
            slide, title,
            Inches(0.5), y, Inches(4.0), Inches(1.0),
            font_size=15, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 原因 + 解法
        info_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), y, Inches(8.333), Inches(1.0))
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        info_box.line.color.rgb = h.COLOR_GRAY_TXT
        info_box.line.width = Pt(0.5)

        h.add_text_block(
            slide, cause,
            Inches(4.7), y + Inches(0.1), Inches(8.0), Inches(0.4),
            font_size=12, color=h.COLOR_DARK, bold=True
        )
        h.add_text_block(
            slide, solution,
            Inches(4.7), y + Inches(0.55), Inches(8.0), Inches(0.4),
            font_size=11, color=h.COLOR_GREEN
        )

    h.add_callout(
        slide, "/debug 啟用 debug 日誌，可看到 skill 載入的所有錯誤",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="🔍", font_size=13
    )

    # ============================================================
    # Slide 31：分享 Skills
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "分享 Skills 給其他人",
        "3 種分享範圍",
        slide_num=31, total=TOTAL, source="04 § 分享"
    )

    h.add_text_block(
        slide, "📢 分享的 3 種方式",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    methods = [
        ("📁", "專案", "將 .claude/skills/ 提交到版本控制", "團隊協作", h.COLOR_BLUE),
        ("📦", "Plugin", "在 plugin 中建立 skills/ 目錄", "跨專案重用", h.COLOR_GREEN),
        ("🏢", "受管", "透過受管設定部署組織範圍", "企業級", h.COLOR_PRIMARY),
    ]

    box_w = Inches(3.8)
    box_h = Inches(3.0)
    h_gap = Inches(0.4)
    total_w = box_w * 3 + h_gap * 2
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.4)

    for i, (icon, name, method, use, color) in enumerate(methods):
        x = start_x + i * (box_w + h_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # Icon + 標題
        h.add_text_block(
            slide, f"{icon} {name}",
            x, start_y + Inches(0.3), box_w, Inches(0.5),
            font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER
        )

        # 方式
        h.add_text_block(
            slide, method,
            x + Inches(0.3), start_y + Inches(1.0), box_w - Inches(0.6), Inches(0.9),
            font_size=12, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )

        # 用途
        label_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), start_y + Inches(2.1), box_w - Inches(0.6), Inches(0.5))
        label_box.fill.solid()
        label_box.fill.fore_color.rgb = color
        label_box.line.fill.background()
        h.add_text_block(
            slide, f"適用：{use}",
            x + Inches(0.3), start_y + Inches(2.1), box_w - Inches(0.6), Inches(0.5),
            font_size=11, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "Plugin 是最靈活的方式：可包含 skills + agents + hooks + MCP",
        Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.4),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 32-39：所有 Frontmatter 速查表
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "常用 Frontmatter 模式速查",
        "複製貼上就能用",
        slide_num=32, total=TOTAL, source="04 § 速查表"
    )

    patterns = [
        ("標準 reference skill", "description", h.COLOR_BLUE),
        ("任務型 skill（防自動觸發）", "description + disable-model-invocation: true", h.COLOR_GREEN),
        ("使用者專用任務 skill", "description + disable-model-invocation: true + argument-hint", h.COLOR_PRIMARY),
        ("預先批准工具", "allowed-tools", h.COLOR_RED),
        ("背景知識（Claude 專用）", "user-invocable: false", RGBColor(0x7C, 0x3A, 0xED)),
        ("帶引數的 skill", "arguments: [name] + $name", h.COLOR_BLUE),
        ("在 subagent 中執行", "context: fork + agent: Explore", h.COLOR_GREEN),
        ("限制自動觸發的時機", "paths: [\"*.ts\"]", h.COLOR_RED),
    ]

    for i, (label, fields, color) in enumerate(patterns):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(1.7) + row * Inches(1.25)

        # 標題列
        title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(6.0), Inches(0.35))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = color
        title_box.line.fill.background()
        h.add_text_block(
            slide, label,
            x, y, Inches(6.0), Inches(0.35),
            font_size=12, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 內容
        content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.35), Inches(6.0), Inches(0.8))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        content_box.line.color.rgb = h.COLOR_GRAY_TXT
        content_box.line.width = Pt(0.5)

        h.add_text_block(
            slide, fields,
            x + Inches(0.2), y + Inches(0.4), Inches(5.6), Inches(0.7),
            font_size=11, color=h.COLOR_DARK, font="Consolas"
        )

    # ============================================================
    # Slide 33：與其他元件的比較
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Skills vs 其他元件",
        "選擇正確的工具",
        slide_num=33, total=TOTAL, source="00 § 比較"
    )

    h.add_comparison_table(
        slide,
        ["元件", "載入時機", "上下文影響", "最佳用途"],
        [
            ["CLAUDE.md", "Session 開始", "每個請求都消耗", "「永遠要做的規則」"],
            ["Skill（參考型）", "需要時自動載入", "描述始終在 context", "API 慣例、模式"],
            ["Skill（任務型）", "/name 觸發", "叫用時才載入", "deploy、commit、review"],
            ["Subagent", "委派時", "獨立 context", "需要隔離的研究任務"],
            ["Hook", "事件觸發", "零消耗（除非回傳）", "lint 自動跑、阻擋命令"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=12
    )

    h.add_text_block(
        slide, "Skills + Subagents 組合",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "Subagent 預載入 skills：把專業知識塞進 subagent 啟動時",
            "Skill 啟動 subagent：context: fork 讓 skill 跑在 subagent 中",
            "兩種方向互補，視場景選擇",
        ],
        Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2),
        font_size=12
    )

    # ============================================================
    # Slide 34：實戰：建立第一個 Skill
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "實戰：30 分鐘建立你的第一個實用 Skill",
        "git-commit-helper 完整教學",
        slide_num=34, total=TOTAL, source="綜合實戰"
    )

    h.add_text_block(
        slide, "🎯 目標：自動產生符合團隊風格的 commit 訊息",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    # Step 1
    h.add_text_block(
        slide, "Step 1：建立目錄與 SKILL.md",
        Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )
    h.add_code_block(
        slide, """$ mkdir -p ~/.claude/skills/git-commit-helper""",
        Inches(0.5), Inches(2.7), Inches(6), Inches(0.5),
        font_size=11
    )

    # Step 2
    h.add_text_block(
        slide, "Step 2：撰寫 SKILL.md",
        Inches(0.5), Inches(3.3), Inches(12), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )
    h.add_code_block(
        slide, """---
name: git-commit-helper
description: 產生 conventional commit 風格的訊息
disable-model-invocation: true
allowed-tools: Bash(git *)
---

## Current status
!`git status --short`
!`git diff --staged --stat`

## Recent commits
!`git log --oneline -5`

## Instructions
Generate a commit message following Conventional Commits:
<type>(<scope>): <subject>

Types: feat, fix, docs, refactor, test, chore
Subject: imperative mood, lowercase, no period
Body: explain what and why, not how""",
        Inches(0.5), Inches(3.7), Inches(12.333), Inches(2.6),
        font_size=10
    )

    # Step 3
    h.add_text_block(
        slide, "Step 3：使用",
        Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )
    h.add_code_block(
        slide, """/git-commit-helper   ← 自動產生訊息，按確認即可""",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        font_size=11
    )

    # ============================================================
    # Slide 35：實戰：建立多 skill Plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "實戰：建立團隊多 Skill Plugin",
        "把常用 skills 打包成 plugin",
        slide_num=35, total=TOTAL, source="綜合實戰"
    )

    h.add_code_block(
        slide, """team-workflow/
├── .claude-plugin/
│   └── plugin.json
├── skills/
│   ├── code-review/
│   │   ├── SKILL.md
│   │   └── checklist.md        ← 團隊審查清單
│   ├── run-tests/
│   │   └── SKILL.md
│   ├── deploy-staging/
│   │   └── SKILL.md
│   └── generate-changelog/
│       ├── SKILL.md
│       └── scripts/
│           └── changelog.py
├── hooks/
│   └── hooks.json              ← 編輯後自動跑 lint
├── agents/
│   └── code-reviewer.md
└── README.md                   ← 團隊使用說明""",
        Inches(0.5), Inches(1.7), Inches(7), Inches(5.0),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 設計重點",
        Inches(8.0), Inches(1.7), Inches(5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "一個 plugin 集中多個相關 skill",
            "hooks 補強 skills（自動化）",
            "agents 處理複雜任務",
            "文件清楚說明使用方式",
            "透過 marketplace 給團隊",
        ],
        Inches(8.0), Inches(2.1), Inches(5), Inches(2.5),
        font_size=12
    )

    h.add_callout(
        slide, "Plugin 是 skill + agents + hooks + MCP 的組合體",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="📦", font_size=13
    )

    # ============================================================
    # Slide 36：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 9", "下一步行動", "把所學變成實戰")

    # ============================================================
    # Slide 37：練習題
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "練習題：5 個實作挑戰",
        "由淺入深，鞏固所學",
        slide_num=37, total=TOTAL, source="練習"
    )

    challenges = [
        ("🟢 入門", "建立 summarize skill", "用 !`git diff HEAD` 摘要變更", "10 分鐘"),
        ("🟢 入門", "建立 deploy 任務型 skill", "用 disable-model-invocation + context: fork", "15 分鐘"),
        ("🟡 中等", "Skill 帶具名引數", "用 arguments 欄位與 $name 替換", "20 分鐘"),
        ("🟡 中等", "Skill 加上 allowed-tools", "讓 git 命令免批准", "20 分鐘"),
        ("🔴 進階", "建立完整 plugin", "包含 3 個 skills + 1 個 hook + 1 個 agent", "1 小時"),
    ]

    for i, (level, title, desc, time) in enumerate(challenges):
        y = Inches(1.8) + i * Inches(1.0)

        # 等級
        level_color = h.COLOR_GREEN if "🟢" in level else (h.COLOR_PRIMARY if "🟡" in level else h.COLOR_RED)
        level_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(1.2), Inches(0.85))
        level_box.fill.solid()
        level_box.fill.fore_color.rgb = level_color
        level_box.line.fill.background()
        h.add_text_block(
            slide, level,
            Inches(0.5), y, Inches(1.2), Inches(0.85),
            font_size=14, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 標題
        h.add_text_block(
            slide, title,
            Inches(1.8), y + Inches(0.05), Inches(6.5), Inches(0.4),
            font_size=15, bold=True, color=h.COLOR_DARK
        )
        # 描述
        h.add_text_block(
            slide, desc,
            Inches(1.8), y + Inches(0.45), Inches(6.5), Inches(0.4),
            font_size=11, color=h.COLOR_GRAY_TXT
        )

        # 時間
        time_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y, Inches(4.333), Inches(0.85))
        time_box.fill.solid()
        time_box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        time_box.line.color.rgb = h.COLOR_PRIMARY
        time_box.line.width = Pt(1.5)
        h.add_text_block(
            slide, f"⏱ {time}",
            Inches(8.5), y, Inches(4.333), Inches(0.85),
            font_size=15, bold=True, color=h.COLOR_PRIMARY,
            align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "完成 5 題後，你已經具備建立任何複雜度的 skill 的能力",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="🎓", font_size=14
    )

    # ============================================================
    # Slide 38：重點回顧
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_summary_slide(
        slide,
        title="重點回顧",
        key_points=[
            "Skill = SKILL.md 檔案 + 可選支援檔案",
            "兩種角色：參考型（知識）vs 任務型（動作）",
            "description 是最重要欄位，決定自動觸發",
            "disable-model-invocation + user-invocable 控制誰能叫用",
            "動態注入：!`command` 預處理、$ARGUMENTS 替換",
            "支援 files：reference.md、examples.md、scripts/",
            "評估用 skill-creator 量化效果",
        ],
        next_steps=[
            "🎯 立即：建立你的第一個 skill（10 分鐘）",
            "📚 進階：嘗試所有 Frontmatter 欄位",
            "🛠 整合：把 skills 包成 plugin 分享給團隊",
            "🚀 評估：用 skill-creator 測試你的 skill",
        ],
        source="04-skills.md"
    )

    # ============================================================
    # Slide 39：本章總結
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Skills 是 Plugin 的核心",
        "把它包進 plugin，分享給世界",
        slide_num=39, total=TOTAL, source="04 § 總結"
    )

    h.add_text_block(
        slide, "💎 你學到了 Claude Code 最強的擴展機制",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.5),
        font_size=20, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "Skills 涵蓋所有擴展類型：知識、動作、自動化、視覺化",
        Inches(0.5), Inches(2.3), Inches(12.333), Inches(0.4),
        font_size=14, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "下一步：把這些 skills 包進 plugin",
        Inches(0.5), Inches(3.0), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_GREEN,
        align=PP_ALIGN.CENTER
    )

    # 視覺化：Skills → Plugin
    # 左：多個 skills
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(5.0), Inches(2.5))
    slide.shapes[-1].fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = h.COLOR_BG_GRAY
    slide.shapes[-1].line.color.rgb = h.COLOR_BLUE
    slide.shapes[-1].line.width = Pt(2)

    h.add_text_block(
        slide, "多個獨立 Skills",
        Inches(0.5), Inches(3.9), Inches(5.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE, align=PP_ALIGN.CENTER
    )

    skills_list = ["summarize-changes", "code-review", "deploy", "run-tests", "changelog"]
    for i, s in enumerate(skills_list):
        y = Inches(4.4) + i * 0.32
        h.add_text_block(
            slide, f"📚 {s}",
            Inches(0.7), y, Inches(4.6), Inches(0.3),
            font_size=12, color=h.COLOR_DARK, font="Consolas"
        )

    # 箭頭
    h.add_text_block(
        slide, "→",
        Inches(5.7), Inches(5.0), Inches(1.0), Inches(0.6),
        font_size=48, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    # 右：Plugin
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(3.8), Inches(6.0), Inches(2.5))
    slide.shapes[-1].fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = h.COLOR_BG_GRAY
    slide.shapes[-1].line.color.rgb = h.COLOR_PRIMARY
    slide.shapes[-1].line.width = Pt(2.5)

    h.add_text_block(
        slide, "📦 Plugin（打包）",
        Inches(6.833), Inches(3.9), Inches(6.0), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY, align=PP_ALIGN.CENTER
    )

    plugin_items = [
        "✓ 所有 skills 命名空間化",
        "✓ + agents、hooks",
        "✓ + MCP servers",
        "✓ + Marketplace 分發",
        "✓ 版本管理 + 自動更新",
    ]
    for i, item in enumerate(plugin_items):
        y = Inches(4.4) + i * 0.32
        h.add_text_block(
            slide, item,
            Inches(7.033), y, Inches(5.7), Inches(0.3),
            font_size=12, color=h.COLOR_DARK
        )

    h.add_callout(
        slide, "Plugin 是 skills 的載體；先寫好 skills，再考慮 plugin 化",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="💡", font_size=14
    )

    # ============================================================
    # Slide 40：結束頁
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide, h.COLOR_BG_CREAM)

    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Skills 大師之路 ✨"
    run.font.name = h.FONT_TITLE
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = h.COLOR_PRIMARY

    h.add_text_block(
        slide, "從第一個 skill 開始，逐步建立你的技能庫",
        Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
        font_size=20, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.666), Inches(5.0),
        Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = h.COLOR_PRIMARY
    bar.line.fill.background()

    h.add_text_block(
        slide, "下一份簡報：06-hooks.pptx（Hooks 自動化指南）",
        Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
        font_size=14, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "Claude Code Plugin 完整學習系列 · #04",
        Inches(1), Inches(6.1), Inches(11.333), Inches(0.5),
        font_size=12, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    # 儲存
    output = "/home/elan/pi-proj/04-skills.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
