"""
簡報 4/4：探索並安裝 Plugin (07-discover-plugins.pptx)
約 22 張
對應：07-discover-plugins.md
一般使用者視角：如何使用 /plugin 介面找到、安裝、管理 plugin
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 22

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "探索並安裝 Plugin",
        "一般使用者指南：如何找到、安裝、管理好用工具",
        tag="#07 · Plugin 探索"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "從找 plugin 到用 plugin 的完整路徑",
        slide_num=2, total=TOTAL, source="07-discover-plugins.md"
    )

    objectives = [
        ("🛒", "Plugin 與 Marketplace 概念", "理解「商店 vs 商品」的關係"),
        ("🏛️", "官方 Anthropic Marketplace", "最權威、安全的 plugin 來源"),
        ("🧩", "六類實用 Plugin", "程式碼智慧、安全審查、開發流程…"),
        ("🎛️", "/plugin 互動介面", "4 個標籤頁的完整操作"),
        ("⚙️", "CLI 安裝與管理", "指令列路線 + 4 種安裝範圍"),
        ("🔄", "更新與重載機制", "/reload-plugins 與自動更新"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(2.2)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_cols = 3
    grid_rows = 2
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.8)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_PRIMARY
        card.line.width = Pt(2)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.2), box_w, Inches(0.7),
            font_size=40, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.2), y + Inches(1.0), box_w - Inches(0.4), Inches(0.5),
            font_size=15, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.5), box_w - Inches(0.4), Inches(0.6),
            font_size=11, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Plugin 與 Marketplace", "先搞懂「商店 vs 商品」")

    # ============================================================
    # Slide 4：Plugin vs Marketplace
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Plugin 與 Marketplace 的關係",
        "一句話：Marketplace 是商店，Plugin 是商品",
        slide_num=4, total=TOTAL, source="07 § 基本概念"
    )

    h.add_two_column_compare(
        slide,
        "🛒 Marketplace（市集）",
        [
            "別人建立和分享的 plugin 目錄",
            "像「應用程式商店」",
            "用 `/plugin marketplace add` 新增來源",
            "可加多個（官方、社群、本機…）",
            "通常託管在 GitHub / GitLab",
            "更新一次，所有使用者都收到",
        ],
        "🧩 Plugin（套件）",
        [
            "實際擴充 Claude 功能的元件",
            "用 skills / agents / hooks / MCP 組成",
            "用 `/plugin install` 安裝",
            "從 marketplace 中挑選",
            "安裝後 Claude Code 自動載入",
            "可用 `/plugin disable` 暫時關閉",
        ]
    )

    h.add_callout(
        slide, "工作流程：新增 marketplace（加商店）→ 從中安裝 plugin（買商品）",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 5：官方 Anthropic Marketplace
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "官方 Anthropic Marketplace",
        "claude-plugins-official：啟動時自動可用",
        slide_num=5, total=TOTAL, source="07 § 官方"
    )

    h.add_text_block(
        slide, "✅ 由 Anthropic 策劃，分類清楚",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "啟動 Claude Code 時自動可用，無需手動新增",
            "執行 `/plugin` → 切到 **Discover** 標籤瀏覽",
            "命令列安裝：`/plugin install github@claude-plugins-official`",
            "若沒看到，執行 `/plugin marketplace update claude-plugins-official`",
        ],
        Inches(0.7), Inches(2.2), Inches(12), Inches(1.6),
        font_size=14
    )

    h.add_text_block(
        slide, "📂 主要分類",
        Inches(0.7), Inches(4.0), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    categories = [
        ("🧠", "程式碼智慧", "LSP plugin（pyright、ts、rust…）"),
        ("🔌", "外部整合", "MCP 預設配置（github、jira、figma…）"),
        ("🛡️", "安全審查", "security-guidance：每變更自動掃漏洞"),
        ("🛠️", "開發流程", "commit-commands、pr-review-toolkit…"),
        ("🎨", "輸出樣式", "explanatory / learning style"),
    ]

    cat_w = Inches(2.3)
    cat_h = Inches(1.8)
    cat_gap = Inches(0.2)
    cat_total = cat_w * 5 + cat_gap * 4
    cat_start_x = (h.SLIDE_W - cat_total) / 2
    cat_y = Inches(4.5)

    for i, (icon, title, desc) in enumerate(categories):
        x = cat_start_x + i * (cat_w + cat_gap)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, cat_y, cat_w, cat_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(1.5)

        h.add_text_block(
            slide, icon,
            x, cat_y + Inches(0.15), cat_w, Inches(0.5),
            font_size=28, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.1), cat_y + Inches(0.75), cat_w - Inches(0.2), Inches(0.4),
            font_size=13, bold=True, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.1), cat_y + Inches(1.15), cat_w - Inches(0.2), Inches(0.6),
            font_size=10, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 6：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "六類實用 Plugin 速覽", "從語言伺服器到安全審查")

    # ============================================================
    # Slide 7：程式碼智慧 Plugin（LSP）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "程式碼智慧 Plugin（LSP）",
        "讓 Claude 即時看到型別錯誤、跳轉定義",
        slide_num=7, total=TOTAL, source="07 § 程式碼智慧"
    )

    h.add_text_block(
        slide, "🔍 兩大能力",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**自動診斷**：每次編輯後，語言伺服器回報錯誤和警告，Claude 同一輪修正",
            "**程式碼導航**：跳轉定義、找參考、懸停型別、列符號、追蹤呼叫層次",
        ],
        Inches(0.7), Inches(2.2), Inches(12), Inches(1.0),
        font_size=13
    )

    h.add_text_block(
        slide, "🌐 各語言 Plugin 對照",
        Inches(0.7), Inches(3.4), Inches(12), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["語言", "Plugin 名稱", "需要的二進位檔"],
        [
            ["Python", "pyright-lsp", "pyright-langserver"],
            ["TypeScript", "typescript-lsp", "typescript-language-server"],
            ["Rust", "rust-analyzer-lsp", "rust-analyzer"],
            ["Go", "gopls-lsp", "gopls"],
            ["Java", "jdtls-lsp", "jdtls"],
            ["C/C++", "clangd-lsp", "clangd"],
            ["C#", "csharp-lsp", "csharp-ls"],
        ],
        Inches(0.5), Inches(3.9), Inches(12.333), Inches(2.8),
        font_size=12
    )

    h.add_callout(
        slide, "⚠️ 必須先安裝語言伺服器二進位檔！Plugin 只是配置連接，不包含伺服器本體",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 8：外部整合 Plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "外部整合 Plugin（MCP 預設配置）",
        "連接外部服務，無需手動設定",
        slide_num=8, total=TOTAL, source="07 § 外部整合"
    )

    h.add_bullet_list(
        slide, [
            "捆綁預先配置的 MCP servers，**一鍵安裝就能用**",
            "無需手動編輯 `.mcp.json` 設定",
            "Plugin 啟用時 MCP server 自動啟動",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(1.0),
        font_size=14
    )

    h.add_text_block(
        slide, "🔌 整合對照表",
        Inches(0.7), Inches(2.8), Inches(12), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["類別", "Plugin"],
        [
            ["**原始碼控制**", "`github`、`gitlab`"],
            ["**專案管理**", "`atlassian`（Jira/Confluence）、`asana`、`linear`、`notion`"],
            ["**設計**", "`figma`"],
            ["**基礎設施**", "`vercel`、`firebase`、`supabase`"],
            ["**通訊**", "`slack`"],
            ["**監控**", "`sentry`"],
        ],
        Inches(0.5), Inches(3.3), Inches(12.333), Inches(3.3),
        font_size=13
    )

    # ============================================================
    # Slide 9：其他四類 Plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "其他四類實用 Plugin",
        "安全、開發流程、輸出樣式",
        slide_num=9, total=TOTAL, source="07 § 安全/開發/輸出"
    )

    # 4 個分類卡片
    cards = [
        ("🛡️", "自動安全審查", "security-guidance",
         ["審查 Claude 每項變更", "檢查常見漏洞", "發現問題**同工作階段**修復"]),
        ("🛠️", "開發工作流程", "commit-commands / pr-review-toolkit / plugin-dev",
         ["Git 提交、PR 建立", "PR 審查專用 agents", "Plugin 開發工具組"]),
        ("🎨", "輸出樣式", "explanatory / learning style",
         ["對實作選擇提供教育性見解", "用於 skill 建立的互動學習模式"]),
        ("🌐", "社群 Marketplace", "anthropics/claude-plugins-community",
         ["通過 Anthropic 自動驗證", "安全篩選的第三方 plugin", "固定到 commit SHA"]),
    ]

    card_w = Inches(6.1)
    card_h = Inches(2.4)
    card_gap_x = Inches(0.13)
    card_gap_y = Inches(0.2)
    grid_start_x = Inches(0.5)
    grid_start_y = Inches(1.7)

    for i, (icon, title, plugin, bullets) in enumerate(cards):
        row = i // 2
        col = i % 2
        x = grid_start_x + col * (card_w + card_gap_x)
        y = grid_start_y + row * (card_h + card_gap_y)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, card_w, card_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        # icon
        h.add_text_block(
            slide, icon,
            x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.6),
            font_size=32
        )
        # 標題
        h.add_text_block(
            slide, title,
            x + Inches(1.0), y + Inches(0.25), card_w - Inches(1.2), Inches(0.5),
            font_size=18, bold=True, color=h.COLOR_DARK
        )
        # plugin 名稱
        h.add_text_block(
            slide, plugin,
            x + Inches(0.3), y + Inches(0.95), card_w - Inches(0.5), Inches(0.4),
            font_size=11, color=h.COLOR_GRAY_TXT, italic=True
        )
        # bullets
        h.add_bullet_list(
            slide, bullets,
            x + Inches(0.3), y + Inches(1.35), card_w - Inches(0.5), card_h - Inches(1.5),
            font_size=12
        )

    # ============================================================
    # Slide 10：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "使用 /plugin 介面", "互動式管理 plugin")

    # ============================================================
    # Slide 11：/plugin 介面 4 個標籤
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "/plugin 介面：4 個標籤",
        "用 Tab（或 Shift+Tab）切換",
        slide_num=11, total=TOTAL, source="07 § 使用 /plugin 介面"
    )

    tabs = [
        ("🛒", "Discover", "瀏覽", "從所有 marketplaces 瀏覽可用 plugin"),
        ("📦", "Installed", "管理", "檢視與管理已安裝的 plugin"),
        ("🏪", "Marketplaces", "來源", "新增、移除或更新已新增的 marketplace"),
        ("❌", "Errors", "排錯", "檢視 plugin 載入錯誤"),
    ]

    tab_w = Inches(2.9)
    tab_h = Inches(3.5)
    tab_gap = Inches(0.3)
    total_w = tab_w * 4 + tab_gap * 3
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.0)

    for i, (icon, name, action, desc) in enumerate(tabs):
        x = start_x + i * (tab_w + tab_gap)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y, tab_w, tab_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        # icon
        h.add_text_block(
            slide, icon,
            x, start_y + Inches(0.3), tab_w, Inches(0.9),
            font_size=48, align=PP_ALIGN.CENTER
        )
        # 名稱
        h.add_text_block(
            slide, name,
            x + Inches(0.1), start_y + Inches(1.3), tab_w - Inches(0.2), Inches(0.5),
            font_size=22, bold=True, align=PP_ALIGN.CENTER, color=h.COLOR_DARK
        )
        # 動作
        h.add_text_block(
            slide, action,
            x + Inches(0.1), start_y + Inches(1.8), tab_w - Inches(0.2), Inches(0.4),
            font_size=12, color=h.COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), start_y + Inches(2.3), tab_w - Inches(0.4), Inches(1.0),
            font_size=11, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "選擇 plugin 後的詳細資訊窗格：Context cost（v2.1.143+）、Last updated（v2.1.144+）、Will install（v2.1.145+）",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.5),
        icon="💡", font_size=12
    )

    # ============================================================
    # Slide 12：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "新增 Marketplace 來源", "從 6 種來源中挑選")

    # ============================================================
    # Slide 13：6 種 marketplace 來源
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "新增 Marketplace：6 種來源",
        "使用 `/plugin marketplace add <source>`",
        slide_num=13, total=TOTAL, source="07 § 新增 Marketplace 來源"
    )

    h.add_comparison_table(
        slide,
        ["來源類型", "範例指令"],
        [
            ["**GitHub 簡寫**", "`/plugin marketplace add anthropics/claude-code`"],
            ["**Git URL（HTTPS）**", "`/plugin marketplace add https://gitlab.com/company/plugins.git`"],
            ["**Git URL（SSH）**", "`/plugin marketplace add git@gitlab.com:company/plugins.git`"],
            ["**本機路徑**", "`/plugin marketplace add ./my-marketplace`"],
            ["**本機 JSON 檔案**", "`/plugin marketplace add ./path/to/marketplace.json`"],
            ["**遠端 URL**", "`/plugin marketplace add https://example.com/marketplace.json`"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=13
    )

    h.add_text_block(
        slide, "⚠️ 注意事項",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "GitHub 簡寫用 `owner/repo` 格式（v2.1.196+ 起拒絕無 `https://` 前綴的 URL）",
            "加 `.git` 後綴讓 Claude Code 複製整個 repo（而非當作 marketplace.json 直接連結）",
            "加 `#ref` 固定分支或標籤：`add https://.../repo.git#v1.0.0`",
        ],
        Inches(0.7), Inches(5.7), Inches(12), Inches(1.3),
        font_size=12
    )

    # ============================================================
    # Slide 14：URL-based vs Git-based
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "URL-based vs Git-based",
        "直接指向 marketplace.json 有個限制",
        slide_num=14, total=TOTAL, source="07 § URL vs Git"
    )

    h.add_two_column_compare(
        slide,
        "🌐 URL-based（直接指向 JSON）",
        [
            "只下載 `marketplace.json` 本身",
            "Plugin 項目用相對路徑（`./plugins/x`）會**找不到檔案**",
            "⚠️ 不適合內含本機 plugin 的 marketplace",
            "優點：簡單、快",
        ],
        "📦 Git-based（推薦）",
        [
            "Clone 整個 repo 到本機",
            "Plugin 項目相對路徑可正常解析",
            "可加 `#ref` 固定版本",
            "缺點：repo 要先存在",
        ]
    )

    h.add_callout(
        slide, "解決方案：改用 GitHub / npm / git URL 來源，或用 git-based marketplace",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 15：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "安裝與管理", "4 種範圍 + 6 個指令")

    # ============================================================
    # Slide 16：安裝 plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "安裝 Plugin",
        "從 Discover 或命令列",
        slide_num=16, total=TOTAL, source="07 § 安裝 Plugin"
    )

    h.add_text_block(
        slide, "🎯 語法",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """/plugin install <plugin-name>@<marketplace-name>""",
        Inches(0.7), Inches(2.2), Inches(12), Inches(0.6),
        font_size=18
    )

    h.add_text_block(
        slide, "🖱️ 從 Discover 標籤（互動）",
        Inches(0.7), Inches(3.1), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "執行 `/plugin` → 切到 **Discover** 標籤",
            "在 plugin 上按 **Enter** → 選擇範圍 → 安裝",
        ],
        Inches(0.7), Inches(3.5), Inches(12), Inches(0.7),
        font_size=13
    )

    h.add_text_block(
        slide, "💻 從命令列（腳本友好）",
        Inches(0.7), Inches(4.4), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """/plugin install commit-commands@claude-code-plugins
# 或 CLI（預設 user 範圍）
claude plugin install formatter@your-org""",
        Inches(0.7), Inches(4.85), Inches(12), Inches(1.1),
        font_size=13
    )

    h.add_callout(
        slide, "安裝時會打開 plugin 詳細資訊窗格，讓你預覽 components 後選擇安裝範圍",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    # ============================================================
    # Slide 17：4 種安裝範圍
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "4 種安裝範圍",
        "決定 plugin 給誰用、存到哪",
        slide_num=17, total=TOTAL, source="07 § 安裝範圍"
    )

    h.add_comparison_table(
        slide,
        ["範圍", "設定檔", "使用案例"],
        [
            ["**user**", "`~/.claude/settings.json`", "個人 plugin，所有專案可用（**預設**）"],
            ["**project**", ".claude/settings.json", "團隊共享，透過版本控制"],
            ["**local**", ".claude/settings.local.json", "專案特定，不共享（gitignored）"],
            ["**managed**", "受管設定", "管理員部署，**唯讀**"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.0),
        font_size=13
    )

    h.add_text_block(
        slide, "💡 解除安裝時的範圍處理（v2.1.203+）",
        Inches(0.7), Inches(4.8), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "解除 `.claude/settings.json` 啟用的 plugin 時會問：",
            "  - 只為自己停用（寫到 `.claude/settings.local.json`）",
            "  - 為所有人解除安裝（從 `.claude/settings.json` 移除）",
            "**managed** 範圍的 plugin 你無法修改（管理員控制）",
        ],
        Inches(0.7), Inches(5.2), Inches(12), Inches(1.7),
        font_size=13
    )

    # ============================================================
    # Slide 18：管理已安裝的 plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "管理已安裝的 Plugin",
        "互動介面 + CLI 指令",
        slide_num=18, total=TOTAL, source="07 § 管理已安裝"
    )

    h.add_text_block(
        slide, "🎛️ 互動式介面（Installed 標籤）",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "按 `f` 加入/取消最愛",
            "輸入文字篩選（按名稱或描述）",
            "排序：載入錯誤 → 未解決依賴 → 最愛 → 停用",
            "v2.1.187+：**Not used recently** 自動列出 2 週未用 plugin",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(1.4),
        font_size=13
    )

    h.add_text_block(
        slide, "💻 CLI 指令",
        Inches(0.7), Inches(3.6), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """# 列出已安裝
/plugin list                          # 全部
/plugin list --enabled                # 只列啟用
/plugin list --disabled               # 只列停用

# 停用（不移除）
/plugin disable <plugin>@<marketplace>

# 重新啟用
/plugin enable <plugin>@<marketplace>

# 完全移除
/plugin uninstall <plugin>@<marketplace>""",
        Inches(0.7), Inches(4.0), Inches(12), Inches(2.9),
        font_size=11
    )

    # ============================================================
    # Slide 19：/reload-plugins
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "不重新啟動就套用變更",
        "/reload-plugins：工作階段中重載所有 plugin",
        slide_num=19, total=TOTAL, source="07 § /reload-plugins"
    )

    h.add_code_block(
        slide, """/reload-plugins""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.6),
        font_size=20
    )

    h.add_text_block(
        slide, "📋 會重新載入",
        Inches(0.7), Inches(2.6), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_GREEN
    )

    h.add_bullet_list(
        slide, [
            "所有 plugin、skills、agents、hooks",
            "Plugin MCP servers 與 LSP servers",
            "顯示各類元件的計數",
        ],
        Inches(0.7), Inches(3.0), Inches(12), Inches(1.2),
        font_size=13
    )

    h.add_text_block(
        slide, "⚠️ Token 成本與快取失效",
        Inches(0.7), Inches(4.4), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "下個請求會產生 token 成本（新元件會宣告自己）",
            "提供 MCP servers 的 plugin 若其工具未被 tool search 延遲 → 快取失效",
            "快取失效時 `/reload-plugins` 顯示警告且**不套用**",
            "傳 `--force` 強制套用",
        ],
        Inches(0.7), Inches(4.8), Inches(12), Inches(1.8),
        font_size=13
    )

    # ============================================================
    # Slide 20：自動更新
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "自動更新",
        "工作階段後背景更新，不打斷進行中的工作",
        slide_num=20, total=TOTAL, source="07 § 自動更新"
    )

    h.add_text_block(
        slide, "⚙️ 啟用後的行為",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "工作階段開始後檢查更新，隨機延遲最多 **10 分鐘**",
            "執行中的工作階段仍使用啟動時載入的版本",
            "若任何 plugin 已更新 → 提示執行 `/reload-plugins` 或下次啟動時載入",
        ],
        Inches(0.7), Inches(2.1), Inches(12), Inches(1.5),
        font_size=13
    )

    h.add_text_block(
        slide, "🎛️ 預設行為",
        Inches(0.7), Inches(3.8), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "官方 Anthropic marketplace：**預設啟用**",
            "第三方與本機開發 marketplace：**預設停用**",
        ],
        Inches(0.7), Inches(4.2), Inches(12), Inches(0.8),
        font_size=13
    )

    h.add_text_block(
        slide, "🛠️ 環境變數控制",
        Inches(0.7), Inches(5.1), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """# 完全停用所有自動更新
export DISABLE_AUTOUPDATER=1

# Claude Code 不自動更新，但 plugin 仍自動更新
export DISABLE_AUTOUPDATER=1
export FORCE_AUTOUPDATE_PLUGINS=1""",
        Inches(0.7), Inches(5.5), Inches(12), Inches(1.4),
        font_size=11
    )

    # ============================================================
    # Slide 21：團隊 marketplace 與安全性
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "團隊配置與安全性",
        "pre-config + 安全注意",
        slide_num=21, total=TOTAL, source="07 § 團隊/安全"
    )

    h.add_text_block(
        slide, "👥 為團隊預先配置 marketplace",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """// .claude/settings.json
{
  "extraKnownMarketplaces": {
    "my-team-tools": {
      "source": {
        "source": "github",
        "repo": "your-org/claude-plugins"
      }
    }
  },
  "enabledPlugins": {
    "code-formatter@my-team-tools": true
  }
}""",
        Inches(0.7), Inches(2.1), Inches(12), Inches(2.8),
        font_size=11
    )

    h.add_text_block(
        slide, "🛡️ 安全性提醒",
        Inches(0.7), Inches(5.05), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "**Plugin 和 marketplace 是高度受信任元件**（用你的權限執行任意程式碼）",
            "僅從你信任的來源安裝；檢查每個 plugin 的首頁",
            "組織可用 `strictKnownMarketplaces` 限制使用者能新增的 marketplace",
            "注意 plugin 內含的 MCP servers / 外部軟體（Anthropic 無法控制）",
        ],
        Inches(0.7), Inches(5.45), Inches(12), Inches(1.5),
        font_size=13
    )

    # ============================================================
    # Slide 22：重點回顧
    # ============================================================
    h.add_summary_slide(
        slide=h.add_blank_slide(prs),
        title="重點回顧",
        key_points=[
            "**Marketplace = 商店**，**Plugin = 商品**：先加商店再買商品",
            "官方 Anthropic marketplace 啟動時自動可用，分 5 大類：智慧、整合、安全、開發、輸出",
            "`/plugin` 介面有 4 個標籤：Discover、Installed、Marketplaces、Errors",
            "4 種安裝範圍：user / project / local / managed",
            "執行 `/reload-plugins` 不重啟就套用新安裝的 plugin",
        ],
        next_steps=[
            "執行 `/plugin` → 切到 Discover 標籤，先裝 `commit-commands`（實用、低風險）",
            "重啟 Claude Code（或執行 `/reload-plugins`）",
            "試用 skill，例如 `/commit-commands:commit`",
            "接著閱讀 `02-plugins.md` 學習如何從零建立自己的 plugin！",
        ],
        source="07-discover-plugins.md"
    )

    # ============================================================
    # 儲存
    # ============================================================
    output = "/home/elan/pi-proj/07-discover-plugins.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
