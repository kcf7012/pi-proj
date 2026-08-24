"""
簡報 3/4：Subagents 自訂指南 (05-subagents.pptx)
約 30 張
對應：05-subagents.md
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 30

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "Subagents 自訂指南",
        "打造專門化的 AI 助手，保護主對話 context",
        tag="#05 · Subagents"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "從內建 subagent 到自訂專門 agent",
        slide_num=2, total=TOTAL, source="05-subagents.md"
    )

    objectives = [
        ("🤖", "什麼是 Subagent", "獨立 context、專門化的 AI 助手"),
        ("⚙️", "內建 4 種 Subagent", "Explore、Plan、general-purpose…"),
        ("🚀", "建立第一個 Subagent", "3 步驟：描述 → 檔案 → 叫用"),
        ("🎛️", "完整 frontmatter 配置", "tools、model、memory、hooks…"),
        ("🧠", "選擇模型與能力控制", "白名單、黑名單、persistence"),
        ("💡", "常見模式與陷阱", "隔離高容量、平行、鏈接"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(2.2)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_cols = 3
    grid_rows = 2
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.8)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_PRIMARY
        card.line.width = Pt(2)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.2), box_w, Inches(0.7),
            font_size=40, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.2), y + Inches(1.0), box_w - Inches(0.4), Inches(0.5),
            font_size=15, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.5), box_w - Inches(0.4), Inches(0.6),
            font_size=11, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "什麼是 Subagent", "獨立 context 的專門助手")

    # ============================================================
    # Slide 4：Subagent 核心概念
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Subagent 核心概念",
        "在獨立 context 視窗中執行的專門助手",
        slide_num=4, total=TOTAL, source="05 § 什麼是 Subagent"
    )

    h.add_bullet_list(
        slide, [
            "**獨立 context**：在自己的 context 視窗執行，主對話不會被搜尋結果/日誌淹沒",
            "**自訂系統提示**：針對特定領域的專注行為",
            "**特定工具存取**：可限制只能用白名單或排除黑名單工具",
            "**獨立權限**：每個 subagent 都有自己的權限判斷",
            "**自動委派**：Claude 根據 subagent 的 `description` 欄位決定何時使用",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(2.5),
        font_size=14
    )

    h.add_text_block(
        slide, "✨ Subagents 幫你做到",
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "✅ **保留 context**：將探索和實作保留在主對話之外",
            "✅ **強制約束**：限制 subagent 可以使用哪些工具",
            "✅ **跨專案重用**：使用者層級 subagents 全專案可用",
            "✅ **專門化行為**：特定領域的專注系統提示",
            "✅ **控制成本**：可路由到更快更便宜的模型（如 Haiku）",
        ],
        Inches(0.7), Inches(4.7), Inches(12), Inches(2.0),
        font_size=13
    )

    # ============================================================
    # Slide 5：使用時機
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "什麼時候該用 Subagent？",
        "主對話 vs Subagent 的決策依據",
        slide_num=5, total=TOTAL, source="05 § 什麼是 Subagent"
    )

    h.add_two_column_compare(
        slide,
        "🗨️ 使用主對話的情境",
        [
            "任務需要**頻繁的來回**或迭代改進",
            "多個階段**共享重要 context**（如計劃、實作、測試）",
            "快速、有針對性的更改",
            "**延遲很重要**（非 fork subagent 從頭開始）",
            "需要查閱先前對話內容",
        ],
        "🤖 使用 Subagent 的情境",
        [
            "任務**產生你不需要的詳細輸出**（測試、日誌）",
            "想強制執行**特定工具限制或權限**",
            "工作**自包含**且可返回摘要",
            "要保護主對話 context 不被淹沒",
            "想用更便宜模型處理（Haiku）",
        ]
    )

    h.add_callout(
        slide, "💡 黃金法則：旁支任務用 subagent，核心對話留給主對話",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        icon="", font_size=13
    )

    # ============================================================
    # Slide 6：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "內建 Subagents", "Claude Code 內建的 4 種助手")

    # ============================================================
    # Slide 7：4 種內建 subagent
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "4 種內建 Subagents",
        "Claude 自動使用，你也可以明確叫用",
        slide_num=7, total=TOTAL, source="05 § 內建 Subagents"
    )

    builtins = [
        ("🔍", "Explore", "唯讀搜尋", "Model: 從主對話繼承（上限 Opus）", "Tools: 唯讀（拒絕 Write/Edit）", "檔案探索、程式碼搜尋、程式碼庫探索"),
        ("📋", "Plan", "plan mode 研究", "Model: 從主對話繼承", "Tools: 唯讀", "在 plan mode 中理解程式碼庫"),
        ("🎯", "general-purpose", "複雜多步任務", "Model: 從主對話繼承", "Tools: 全部可用", "複雜研究、多步驟操作、程式碼修改"),
        ("🤖", "claude（預設）", "通用助手", "Model: 從主對話繼承", "Tools: 全部", "任務不適合更專門的代理時"),
    ]

    box_w = Inches(6.1)
    box_h = Inches(2.4)
    box_gap_x = Inches(0.13)
    box_gap_y = Inches(0.2)
    grid_start_x = Inches(0.5)
    grid_start_y = Inches(1.7)

    for i, (icon, name, role, model, tools, purpose) in enumerate(builtins):
        row = i // 2
        col = i % 2
        x = grid_start_x + col * (box_w + box_gap_x)
        y = grid_start_y + row * (box_h + box_gap_y)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        h.add_text_block(
            slide, f"{icon}  {name}",
            x + Inches(0.2), y + Inches(0.15), box_w - Inches(0.4), Inches(0.5),
            font_size=18, bold=True
        )
        h.add_text_block(
            slide, role,
            x + Inches(0.2), y + Inches(0.65), box_w - Inches(0.4), Inches(0.3),
            font_size=11, italic=True, color=h.COLOR_GRAY_TXT
        )
        h.add_bullet_list(
            slide, [model, tools, f"**用途**：{purpose}"],
            x + Inches(0.3), y + Inches(1.05), box_w - Inches(0.5), box_h - Inches(1.2),
            font_size=11
        )

    # ============================================================
    # Slide 8：限制與停用內建 subagent
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "限制與停用內建 Subagent",
        "想關掉某些 subagent 怎麼做？",
        slide_num=8, total=TOTAL, source="05 § 內建 Subagents"
    )

    h.add_bullet_list(
        slide, [
            "**阻擋特定內建類型**：加到 `permissions.deny`（如 `deny: [\"Agent(Explore)\"]`）",
            "**防止任何委派**：用 `permissions.deny` 拒絕 `Agent` 工具本身",
            "**只停用 Explore/Plan**：設 `CLAUDE_CODE_DISABLE_EXPLORE_PLAN_AGENTS=1`（v2.1.198+）",
            "**完全移除內建類型**（非互動模式 / Agent SDK）：設 `CLAUDE_AGENT_SDK_DISABLE_BUILTIN_AGENTS=1`",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(2.0),
        font_size=13
    )

    h.add_text_block(
        slide, "💡 其他內建助手",
        Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["Agent", "Model", "使用時機"],
        [
            ["`statusline-setup`", "Sonnet", "執行 `/statusline` 配置狀態列時"],
            ["`claude-code-guide`", "Haiku", "詢問 Claude Code 功能問題時"],
        ],
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(1.5),
        font_size=13
    )

    h.add_callout(
        slide, "⚠️ 重要例外：Explore 和 Plan 跳過 CLAUDE.md 與 git status（保持研究快速便宜）",
        Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 9：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "建立第一個 Subagent", "3 步驟：描述 → 檔案 → 叫用")

    # ============================================================
    # Slide 10：3 步驟總覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "建立 Subagent 的 3 個步驟",
        "從描述到使用",
        slide_num=10, total=TOTAL, source="05 § 快速開始"
    )

    steps = [
        ("1", "向 Claude 描述", "在 Claude Code 描述想要的 subagent", "「Create a personal code-improver subagent...」"),
        ("2", "檢查檔案", "Claude 寫入 markdown 檔案", "`~/.claude/agents/code-improver.md`"),
        ("3", "叫用", "重新啟動（或直接使用）", "「Use the code-improver agent...」"),
    ]

    step_w = Inches(3.8)
    step_h = Inches(3.2)
    h_gap = Inches(0.4)
    total_w = step_w * 3 + h_gap * 2
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.0)

    for i, (num, title, action, example) in enumerate(steps):
        x = start_x + i * (step_w + h_gap)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y, step_w, step_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        # 編號圓圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + step_w / 2 - Inches(0.4), start_y + Inches(0.3),
            Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = h.COLOR_PRIMARY
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            x + step_w / 2 - Inches(0.4), start_y + Inches(0.35),
            Inches(0.8), Inches(0.7),
            font_size=32, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

        # 標題
        h.add_text_block(
            slide, title,
            x + Inches(0.2), start_y + Inches(1.3), step_w - Inches(0.4), Inches(0.5),
            font_size=18, bold=True, align=PP_ALIGN.CENTER
        )
        # 動作
        h.add_text_block(
            slide, action,
            x + Inches(0.2), start_y + Inches(1.85), step_w - Inches(0.4), Inches(0.5),
            font_size=13, align=PP_ALIGN.CENTER, color=h.COLOR_DARK
        )
        # 範例
        h.add_text_block(
            slide, example,
            x + Inches(0.2), start_y + Inches(2.45), step_w - Inches(0.4), Inches(0.7),
            font_size=11, align=PP_ALIGN.CENTER, italic=True, color=h.COLOR_GRAY_TXT
        )

    # ============================================================
    # Slide 11：範例 subagent 檔案
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "範例：code-improver subagent",
        "使用者層級，所有專案可用",
        slide_num=11, total=TOTAL, source="05 § 快速開始"
    )

    h.add_code_block(
        slide, """# ~/.claude/agents/code-improver.md
---
name: code-improver
description: Scans files and suggests improvements
  for readability, performance, and best practices.
  Use after writing or modifying code.
tools: Read, Grep, Glob
model: sonnet
---

You are a code improvement specialist. For each
issue you find, explain the problem, show the
current code, and provide an improved version.""",
        Inches(0.7), Inches(1.7), Inches(7.5), Inches(4.5),
        font_size=11
    )

    h.add_text_block(
        slide, "📂 檔案位置決定範圍",
        Inches(8.5), Inches(1.7), Inches(4.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "`~/.claude/agents/` → 所有專案可用（**個人**）",
            "`<project>/.claude/agents/` → 該專案（**團隊**）",
            "**專案 subagents 簽入版本控制**讓團隊協作",
            "Claude Code 監視資料夾，**幾秒內偵測變更**（無需重啟）",
        ],
        Inches(8.5), Inches(2.1), Inches(4.5), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 重要提示",
        Inches(8.5), Inches(4.3), Inches(4.5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "如果 Claude 找不到新 subagent，**重啟 Claude Code** 後再試",
            "（只發生在 `~/.claude/agents/` 在 session 開始前不存在的情況）",
        ],
        Inches(8.5), Inches(4.7), Inches(4.5), Inches(1.5),
        font_size=11
    )

    # ============================================================
    # Slide 12：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "配置 Subagent", "frontmatter 16 個欄位全解析")

    # ============================================================
    # Slide 13：5 種範圍優先級
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "5 種 Subagent 範圍",
        "優先級 1（最高）到 5（最低）",
        slide_num=13, total=TOTAL, source="05 § 配置 Subagents"
    )

    h.add_comparison_table(
        slide,
        ["優先級", "位置", "範圍", "如何建立"],
        [
            ["1（最高）", "Managed settings", "組織範圍", "透過 managed settings 部署"],
            ["2", "`--agents` CLI flag", "目前 session", "啟動時傳遞 JSON"],
            ["3", "`.claude/agents/`", "目前專案", "向 Claude 詢問或手動建立"],
            ["4", "`~/.claude/agents/`", "所有專案", "向 Claude 詢問或手動建立"],
            ["5（最低）", "Plugin 的 `agents/`", "Plugin 啟用處", "隨 plugins 安裝"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.0),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 重點補充",
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "Claude Code 會**向上走掃描** `.claude/agents/`，所以可組織成子資料夾（`agents/review/`）",
            "v2.1.178+：巢狀目錄若定義相同 `name` → 使用**最接近工作目錄**的定義",
            "子目錄路徑不影響識別（身份僅來自 `name` frontmatter）",
            "`--agents` flag 接受 JSON 與 `prompt` 欄位（**不是** frontmatter 格式）",
        ],
        Inches(0.7), Inches(5.3), Inches(12), Inches(1.6),
        font_size=12
    )

    # ============================================================
    # Slide 14：16 個 frontmatter 欄位（速查表）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Frontmatter 欄位速查表",
        "只有 `name` 和 `description` 必需",
        slide_num=14, total=TOTAL, source="05 § frontmatter"
    )

    h.add_comparison_table(
        slide,
        ["欄位", "必需", "說明"],
        [
            ["`name`", "✅", "唯一識別碼（kebab-case），hooks 收到此值"],
            ["`description`", "✅", "Claude 何時委派給此 subagent"],
            ["`tools`", "❌", "可用工具清單（白名單）"],
            ["`disallowedTools`", "❌", "拒絕的工具（黑名單）"],
            ["`model`", "❌", "sonnet / opus / haiku / inherit，預設 inherit"],
            ["`permissionMode`", "❌", "default / acceptEdits / auto / plan…"],
            ["`maxTurns`", "❌", "停止前的最大代理輪次"],
            ["`skills`", "❌", "啟動時預載入的 skills"],
            ["`mcpServers`", "❌", "此 subagent 可用的 MCP servers"],
            ["`hooks`", "❌", "限定此 subagent 的生命週期 hooks"],
        ],
        Inches(0.5), Inches(1.7), Inches(6.1), Inches(5.0),
        font_size=11
    )

    h.add_comparison_table(
        slide,
        ["欄位", "必需", "說明"],
        [
            ["`memory`", "❌", "持久記憶範圍：user / project / local"],
            ["`background`", "❌", "true = 強制背景執行"],
            ["`effort`", "❌", "努力級別（low / medium / high）"],
            ["`isolation`", "❌", "worktree = 在臨時 worktree 執行"],
            ["`color`", "❌", "任務清單中的顯示顏色"],
            ["`initialPrompt`", "❌", "作為主 session 代理時的自動首訊"],
        ],
        Inches(6.733), Inches(1.7), Inches(6.1), Inches(3.5),
        font_size=11
    )

    h.add_callout(
        slide, "⚠️ Plugin 限制：plugin subagents **不支援** hooks、mcpServers、permissionMode",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 15：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "選擇模型與控制能力", "模型路由 + 工具限制")

    # ============================================================
    # Slide 16：選擇模型
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "選擇模型：4 種方式",
        "`model` 欄位控制 subagent 用哪個 AI",
        slide_num=16, total=TOTAL, source="05 § 選擇模型"
    )

    h.add_comparison_table(
        slide,
        ["寫法", "意義", "用途"],
        [
            ["**Model 別名**", "`sonnet` / `opus` / `haiku` / `fable`", "快速指定模型家族"],
            ["**完整模型 ID**", "`claude-opus-5` / `claude-sonnet-5`", "指定特定版本"],
            ["**inherit**", "與主對話相同模型", "預設（v2.1.198+ 也繼承擴展思考）"],
            ["**Omitted**", "預設為 inherit", "與 inherit 相同"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.5),
        font_size=13
    )

    h.add_text_block(
        slide, "🔍 模型解析優先順序（每次叫用時）",
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "1️⃣ 環境變數 `CLAUDE_CODE_SUBAGENT_MODEL`",
            "2️⃣ 每次叫用的 `model` 參數",
            "3️⃣ Subagent 定義的 `model` frontmatter",
            "4️⃣ 主對話的模型",
        ],
        Inches(0.7), Inches(4.7), Inches(12), Inches(2.0),
        font_size=14
    )

    # ============================================================
    # Slide 17：工具限制：白名單 / 黑名單
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "工具限制：白名單與黑名單",
        "`tools` 與 `disallowedTools` 互補",
        slide_num=17, total=TOTAL, source="05 § 控制能力"
    )

    h.add_text_block(
        slide, "✅ 白名單（tools）",
        Inches(0.7), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_GREEN
    )

    h.add_code_block(
        slide, """---
name: safe-researcher
description: Research agent with
  restricted capabilities
tools: Read, Grep, Glob, Bash
---""",
        Inches(0.7), Inches(2.1), Inches(6.0), Inches(1.8),
        font_size=12
    )

    h.add_text_block(
        slide, "🚫 黑名單（disallowedTools）",
        Inches(6.933), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """---
name: no-writes
description: Inherits available tools
  except file writes
disallowedTools: Write, Edit
---""",
        Inches(6.933), Inches(2.1), Inches(6.0), Inches(1.8),
        font_size=12
    )

    h.add_text_block(
        slide, "💡 重要規則",
        Inches(0.5), Inches(4.1), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "兩者都設 → `disallowedTools` **先套用**，再 `tools` 解析剩餘池",
            "`tools` 清單無項目解析為工具 → subagent **通常拒絕啟動**",
            "支援 MCP 模式：`mcp__<server>` 或 `mcp__<server>__*`",
            "`disallowedTools` 還支援 `mcp__*`（從任何 server 移除）",
            "範例：`disallowedTools: mcp__github` → 排除整個 GitHub MCP server",
        ],
        Inches(0.7), Inches(4.5), Inches(12), Inches(2.2),
        font_size=13
    )

    # ============================================================
    # Slide 18：限制可生成的 subagent
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "限制可生成的 Subagent 類型",
        "Agent(agent_type) 白名單語法",
        slide_num=18, total=TOTAL, source="05 § 控制能力"
    )

    h.add_code_block(
        slide, """---
name: coordinator
description: Coordinates work across
  specialized agents
tools: Agent(worker, researcher), Read, Bash
---""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(1.55),
        font_size=12
    )

    h.add_bullet_list(
        slide, [
            "**白名單**：只能生成 `worker` 和 `researcher` 兩種 subagent",
            "**無限制**：`tools: Agent`（不加括號）允許生成任何 subagent",
            "**僅限** `claude --agent` 作為主執行緒的代理",
        ],
        Inches(0.7), Inches(3.4), Inches(12), Inches(1.0),
        font_size=13
    )

    h.add_text_block(
        slide, "🔌 範圍 MCP Servers 到 Subagent",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """---
name: browser-tester
mcpServers:
  - playwright:        # inline 定義，僅限此 subagent
      type: stdio
      command: npx
      args: ["-y", "@playwright/mcp@latest"]
  - github             # 按名稱引用，重用已配置的 server
---""",
        Inches(0.7), Inches(4.95), Inches(12), Inches(2.05),
        font_size=10
    )

    # ============================================================
    # Slide 19：預載入 skills 與持久記憶
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "預載入 Skills 與持久記憶",
        "讓 subagent 啟動時就有完整 context",
        slide_num=19, total=TOTAL, source="05 § 預載入/memory"
    )

    h.add_text_block(
        slide, "📚 預載入 Skills",
        Inches(0.7), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """---
name: api-developer
description: Implement API endpoints
  following team conventions
skills:
  - api-conventions
  - error-handling-patterns
---""",
        Inches(0.7), Inches(2.1), Inches(6.0), Inches(2.2),
        font_size=12
    )

    h.add_text_block(
        slide, "🧠 持久記憶（memory）",
        Inches(6.933), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """---
name: code-reviewer
description: Reviews code for quality
memory: user   # 跨所有專案
# project / local / 不寫
---""",
        Inches(6.933), Inches(2.1), Inches(6.0), Inches(2.2),
        font_size=12
    )

    h.add_text_block(
        slide, "記憶範圍選擇",
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**`user`**（`~/.claude/agent-memory/<name>/`）：跨所有專案記住學習內容",
            "**`project`**（`.claude/agent-memory/<name>/`）：專案特定，簽入版本控制",
            "**`local`**（`.claude/agent-memory-local/<name>/`）：專案特定，不簽入",
            "啟用時：系統提示包含讀寫指令 + `MEMORY.md` 前 200 行（25KB）",
            "Read/Write/Edit 工具**自動啟用**讓 subagent 管理記憶",
        ],
        Inches(0.7), Inches(5.0), Inches(12), Inches(1.8),
        font_size=12
    )

    # ============================================================
    # Slide 20：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 6", "條件式規則與 Hooks", "用 PreToolUse 動態控制")

    # ============================================================
    # Slide 21：唯讀 DB 範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "條件式規則：唯讀 DB 查詢",
        "用 PreToolUse hook 阻擋 SQL 寫入",
        slide_num=21, total=TOTAL, source="05 § 條件式規則"
    )

    h.add_code_block(
        slide, """---
name: db-reader
description: Execute read-only database queries
tools: Bash
hooks:
  PreToolUse:
    - matcher: "Bash"
      hooks:
        - type: command
          command: "./scripts/validate-readonly-query.sh"
---""",
        Inches(0.7), Inches(1.7), Inches(12), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "📜 驗證腳本",
        Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """#!/bin/bash
# ./scripts/validate-readonly-query.sh
INPUT=$(cat)
COMMAND=$(echo "$INPUT" | jq -r '.tool_input.command // empty')

# 阻擋 SQL 寫入操作（不區分大小寫）
if echo "$COMMAND" | grep -iE '\\b(INSERT|UPDATE|DELETE|DROP|CREATE|ALTER|TRUNCATE)\\b' > /dev/null; then
  echo "Blocked: Only SELECT queries are allowed" >&2
  exit 2
fi
exit 0""",
        Inches(0.7), Inches(4.6), Inches(12), Inches(2.2),
        font_size=10
    )

    h.add_callout(
        slide, "⚠️ macOS/Linux 必須 `chmod +x script.sh`，否則 hook 會失敗（不是阻擋）",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.3),
        icon="", font_size=10
    )

    # ============================================================
    # Slide 22：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 7", "使用 Subagents", "叫用、平行、深度、context")

    # ============================================================
    # Slide 23：3 種叫用方式
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "3 種叫用 Subagent 的方式",
        "從一次性建議到 session 範圍預設",
        slide_num=23, total=TOTAL, source="05 § 使用 Subagents"
    )

    h.add_text_block(
        slide, "1️⃣ 自然語言（建議性）",
        Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """Use the test-runner subagent to fix failing tests
Have the code-reviewer subagent look at my recent changes""",
        Inches(0.7), Inches(2.1), Inches(12), Inches(0.7),
        font_size=12
    )

    h.add_text_block(
        slide, "2️⃣ @-mention（保證叫用）",
        Inches(0.7), Inches(2.95), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, '''@"code-reviewer (agent)" look at the auth changes''',
        Inches(0.7), Inches(3.35), Inches(12), Inches(0.6),
        font_size=12
    )

    h.add_text_block(
        slide, "3️⃣ Session 範圍（整個 session）",
        Inches(0.7), Inches(4.1), Inches(12), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """# 命令列啟動
claude --agent code-reviewer

# 或在 .claude/settings.json 設定
{
  "agent": "code-reviewer"
}""",
        Inches(0.7), Inches(4.5), Inches(12), Inches(1.6),
        font_size=12
    )

    h.add_callout(
        slide, "💡 @-mention 控制 Claude 叫用哪個 subagent，不是 subagent 接收什麼 prompt",
        Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 24：平行、深度、context
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "平行、深度與 Context 管理",
        "資源限制與隔離策略",
        slide_num=24, total=TOTAL, source="05 § 平行/深度/context"
    )

    h.add_text_block(
        slide, "🔁 並行限制",
        Inches(0.7), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "預設 20 個 subagents 同時運行",
            "超過 → 報錯 `Concurrent subagent limit reached`",
            "用 `CLAUDE_CODE_MAX_CONCURRENT_SUBAGENTS` 調整",
            "ultracode session **豁免**此限制",
        ],
        Inches(0.7), Inches(2.1), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "📐 深度限制",
        Inches(6.933), Inches(1.7), Inches(6.0), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "預設 subagent 可生成 subagent，**最多 3 層深**",
            "用 `CLAUDE_CODE_MAX_SUBAGENT_SPAWN_DEPTH` 調整",
            "v2.1.219+ 預設 3 層（v2.1.217 為 1）",
            "深度限制時扣留 `Agent` 工具，subagent 自己委派",
        ],
        Inches(6.933), Inches(2.1), Inches(6.0), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "🪟 Context 管理",
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "**非 fork subagent** 從**全新隔離 context** 開始，看不到對話歷史",
            "**例外**：`fork` subagent 繼承父對話",
            "載入：system prompt + CLAUDE.md + git status snapshot + preload skills + sibling roster",
            "**永不載入**：output style、auto memory、context window size",
            "**恢復 subagent**：保留完整對話歷史，從停止處繼續（v2.1.191+ 自己停止的 subagent 不會自動恢復）",
        ],
        Inches(0.7), Inches(4.7), Inches(12), Inches(2.0),
        font_size=12
    )

    # ============================================================
    # Slide 25：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 8", "範例與常見模式", "4 個實戰範例 + 3 個核心模式")

    # ============================================================
    # Slide 26：4 個實戰範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "4 個實戰範例 Subagent",
        "從 code reviewer 到 db-reader",
        slide_num=26, total=TOTAL, source="05 § 範例 Subagents"
    )

    examples = [
        ("🔍", "Code Reviewer", "唯讀審查", "tools: Read, Grep, Glob, Bash", "model: inherit"),
        ("🐛", "Debugger", "分析與修復", "tools: Read, Edit, Bash, Grep, Glob", "model: inherit"),
        ("📊", "Data Scientist", "SQL/BigQuery 分析", "tools: Bash, Read, Write", "model: sonnet"),
        ("🗃️", "DB Reader", "唯讀 SQL 查詢", "tools: Bash + PreToolUse hook", "model: inherit"),
    ]

    box_w = Inches(6.1)
    box_h = Inches(2.4)
    box_gap_x = Inches(0.13)
    box_gap_y = Inches(0.2)
    grid_start_x = Inches(0.5)
    grid_start_y = Inches(1.7)

    for i, (icon, name, role, tools, model) in enumerate(examples):
        row = i // 2
        col = i % 2
        x = grid_start_x + col * (box_w + box_gap_x)
        y = grid_start_y + row * (box_h + box_gap_y)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = h.COLOR_PRIMARY
        box.line.width = Pt(2)

        h.add_text_block(
            slide, f"{icon}  {name}",
            x + Inches(0.2), y + Inches(0.15), box_w - Inches(0.4), Inches(0.5),
            font_size=18, bold=True
        )
        h.add_text_block(
            slide, role,
            x + Inches(0.2), y + Inches(0.7), box_w - Inches(0.4), Inches(0.4),
            font_size=12, italic=True, color=h.COLOR_GRAY_TXT
        )
        h.add_bullet_list(
            slide, [f"**tools**: {tools}", f"**model**: {model}"],
            x + Inches(0.3), y + Inches(1.2), box_w - Inches(0.5), box_h - Inches(1.4),
            font_size=12
        )

    # ============================================================
    # Slide 27：3 個常見模式
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "3 個常見模式",
        "隔離、平行、鏈接",
        slide_num=27, total=TOTAL, source="05 § 常見模式"
    )

    h.add_two_column_compare(
        slide,
        "🔇 隔離高容量操作",
        [
            "**最佳用途**：執行測試、抓文件、處理日誌",
            "詳細輸出保留在 subagent context",
            "只有**相關摘要**返回主對話",
            "範例：跑測試套件，只回報失敗的測試",
        ],
        "🔗 鏈接 Subagents",
        [
            "**多步驟工作流**",
            "每個 subagent 完成後，Claude 把結果給下一個",
            "範例：reviewer 找問題 → optimizer 修復",
            "適合需要**順序依賴**的工作",
        ]
    )

    h.add_text_block(
        slide, "⚡ 平行運行研究",
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_text_block(
        slide, "獨立調查 → 同時運行多個 subagent → Claude 綜合發現。研究路徑不互相依賴時效果最好",
        Inches(0.7), Inches(7.25), Inches(12), Inches(0.4),
        font_size=12
    )

    # ============================================================
    # Slide 28：子代理輸出掃描
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Subagent 輸出掃描",
        "Claude Code 自動防禦 prompt injection",
        slide_num=28, total=TOTAL, source="05 § 輸出掃描"
    )

    h.add_bullet_list(
        slide, [
            "**掃描時機**：每個 subagent 最終報告被 Claude 讀取**之前**",
            "**為什麼**：subagent 可能讀取你未審查的檔案/網頁/命令輸出，文字可能帶有對主對話的指令",
            "**兩種變更**：",
            "  - **反斜線插入**：模仿 Claude Code 自己的輸出（如 `<tool_use>`、`Human:`、`Assistant:`）",
            "  - **Marker line**：報告提到權限設定（`bypassPermissions` 等）時在前面加上一行",
            "**保證**：掃描**永不刪除或重新措辭**任何內容",
        ],
        Inches(0.7), Inches(1.7), Inches(12), Inches(3.5),
        font_size=14
    )

    h.add_callout(
        slide, "⚠️ 掃描需要 Claude Code v2.1.210+",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        icon="", font_size=13
    )

    h.add_text_block(
        slide, "💡 完整 subagent 生命週期看 06-hooks.md（事件觸發 + 自動回應）",
        Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.4),
        font_size=14, color=h.COLOR_GRAY_TXT, italic=True
    )

    # ============================================================
    # Slide 29：疑難排解
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "疑難排解速查",
        "常見問題與解決方案",
        slide_num=29, total=TOTAL, source="05 § 速查"
    )

    h.add_comparison_table(
        slide,
        ["問題", "可能原因", "解決方案"],
        [
            ["Subagent 沒被叫用", "description 不夠清楚", "加入「use proactively」等詞，並明確何時使用"],
            ["新 subagent 找不到", "session 在建立資料夾前已啟動", "重啟 Claude Code"],
            ["工具限制沒生效", "tools/disallowedTools 衝突", "先套用 disallowedTools，再讓 tools 解析剩餘池"],
            ["MCP server 沒連接", "沒在 subagent 內 inline 定義", "在 `mcpServers` 欄位 inline 定義，僅此 subagent 可用"],
            ["Preload skills 沒生效", "skill 設了 disable-model-invocation", "移除該標記（預載入從同組 skills 取）"],
            ["記憶目錄沒內容", "auto memory 全域關閉", "啟用 auto memory 或不使用 memory 欄位"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.0),
        font_size=12
    )

    h.add_callout(
        slide, "💡 詳細除錯：`claude --debug` 顯示 plugin 載入與 subagent 註冊詳細資訊",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        icon="", font_size=12
    )

    # ============================================================
    # Slide 30：重點回顧
    # ============================================================
    h.add_summary_slide(
        slide=h.add_blank_slide(prs),
        title="重點回顧",
        key_points=[
            "**Subagent** = 獨立 context 的專門助手，保護主對話 context",
            "**5 種範圍**：managed → CLI flag → 專案 → 使用者 → plugin",
            "**frontmatter** 16 個欄位，最常用：name、description、tools、model",
            "**白名單 vs 黑名單**：`disallowedTools` 先套用，再讓 `tools` 解析",
            "**常見模式**：隔離高容量、平行研究、鏈接 subagents",
        ],
        next_steps=[
            "從內建 `Explore` 開始體驗 subagent 的威力",
            "建立你的第一個自訂 subagent（個人層級 `~/.claude/agents/`）",
            "用 `memory: user` 讓它跨專案累積知識",
            "接著閱讀 `06-hooks.md` 學習事件驅動自動化！",
        ],
        source="05-subagents.md"
    )

    # ============================================================
    # 儲存
    # ============================================================
    output = "/home/elan/pi-proj/05-subagents.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
