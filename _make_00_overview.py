"""
簡報 1/4：系列總覽 (00-overview.pptx)
約 30 張
對應：00-claude-code-plugins-series.md
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 30  # 預估總張數（會動態更新）

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "Claude Code Plugin 學習系列",
        "完整學習路徑 · 8 份文件 · 6,000+ 行",
        tag="系列總覽 · #00"
    )

    # ============================================================
    # Slide 2：為什麼需要這系列
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "為什麼需要了解 Claude Code Plugin？",
        "擴展 Claude 的能力 · 打造個人化 AI 工作流",
        slide_num=2, total=TOTAL, source="00-claude-code-plugins-series.md"
    )

    h.add_two_column_compare(
        slide,
        "😐 預設 Claude Code",
        [
            "通用對話能力",
            "內建工具：檔案、搜尋、Bash",
            "每次 session 都要重新解釋",
            "無法自動化重複任務",
            "無法團隊共享設定",
            "個人風格無法持久化"
        ],
        "🚀 使用 Plugin 後",
        [
            "可重複使用的 Skills（/command）",
            "隔離上下文的 Subagents",
            "事件驅動的 Hooks（自動化）",
            "團隊共享的 Marketplace",
            "自動化的格式化與檢查",
            "個人化 AI 助手工作流"
        ]
    )

    # ============================================================
    # Slide 3：系列規模
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "這套系列的規模",
        "完整覆蓋 Claude Code Plugin 系統的各個面向",
        slide_num=3, total=TOTAL, source="00-claude-code-plugins-series.md"
    )

    # 三大數據卡
    cards = [
        ("8", "份", "完整文件", h.COLOR_PRIMARY),
        ("6,682", "行", "繁體中文內容", h.COLOR_BLUE),
        ("272", "KB", "總檔案大小", h.COLOR_GREEN),
    ]
    card_w = Inches(3.8)
    card_h = Inches(2.5)
    spacing = Inches(0.3)
    total_w = card_w * 3 + spacing * 2
    start_x = (h.SLIDE_W - total_w) / 2

    for i, (num, unit, label, color) in enumerate(cards):
        x = start_x + i * (card_w + spacing)
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(2.2), card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # 數字
        h.add_text_block(
            slide, num,
            x, Inches(2.4), card_w, Inches(1.2),
            font_size=72, bold=True, color=color,
            align=PP_ALIGN.CENTER
        )
        # 單位
        h.add_text_block(
            slide, unit,
            x, Inches(3.7), card_w, Inches(0.4),
            font_size=18, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        # 標籤
        h.add_text_block(
            slide, label,
            x, Inches(4.1), card_w, Inches(0.4),
            font_size=14, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide,
        "涵蓋 Claude Code v2.1.x（到 v2.1.236），所有內容整理自官方文件",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.6),
        icon="📊", font_size=14
    )

    h.add_text_block(
        slide,
        "建立日期：2026/01　|　授權：整理自官方文件，繁體中文教學用途",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER, bold=False
    )

    # ============================================================
    # Slide 4：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Claude Code 擴展生態系全景", "了解整體架構與各元件的角色")

    # ============================================================
    # Slide 5：擴展元件總覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Claude Code 擴展生態系：7 個核心元件",
        "從基礎規則到進階自動化，層層堆疊的擴展能力",
        slide_num=5, total=TOTAL, source="00 § 擴展全景"
    )

    # 7 個元件以圓角矩形呈現，2x4 網格
    components = [
        ("CLAUDE.md", "專案說明書", "每次 session 自動載入的規則", h.COLOR_PRIMARY),
        ("Skills", "可重用知識庫", "/name 觸發，隨叫隨到", h.COLOR_BLUE),
        ("Subagents", "隔離代理人", "獨立 context 的子任務", h.COLOR_GREEN),
        ("Hooks", "事件自動化", "確定性觸發的腳本", h.COLOR_RED),
        ("MCP", "外部服務", "連接資料庫、API、瀏覽器", RGBColor(0x7C, 0x3A, 0xED)),
        ("Plugins", "包裝箱", "把上述元件打包成可發布單元", RGBColor(0xEA, 0x58, 0x0C)),
        ("Marketplaces", "商店目錄", "多個 plugin 的集合 + 版本管理", RGBColor(0x0E, 0x7C, 0x66)),
    ]

    grid_cols = 3
    grid_rows = 3  # 多一個空位
    box_w = Inches(3.8)
    box_h = Inches(1.5)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap * (grid_rows - 1)
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.7)

    for i, (name, role, desc, color) in enumerate(components):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        # 元件方塊
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_w, box_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(2.5)

        # 標題條
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, box_w, Inches(0.4)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = color
        title_bar.line.fill.background()

        # 名稱
        h.add_text_block(
            slide, name,
            x, y, box_w, Inches(0.4),
            font_size=14, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 角色
        h.add_text_block(
            slide, role,
            x, y + Inches(0.45), box_w, Inches(0.35),
            font_size=12, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )

        # 描述
        h.add_text_block(
            slide, desc,
            x, y + Inches(0.85), box_w, Inches(0.6),
            font_size=10, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide,
        "學習順序建議：CLAUDE.md → Skills → Plugins（其他元件視需求）",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 6：每個元件的使用時機
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "何時該用哪個元件？",
        "依需求選擇最合適的擴展元件",
        slide_num=6, total=TOTAL, source="00 § 元件使用時機"
    )

    h.add_comparison_table(
        slide,
        ["你的需求", "推薦元件", "原因"],
        [
            ["每次 session 都要遵守的規則", "CLAUDE.md", "自動載入，零成本"],
            ["重複使用的 SOP 或專業知識", "Skill", "隨時叫用，可重用"],
            ["需要上下文隔離的任務", "Subagent", "不污染主對話"],
            ["每次檔案編輯後自動跑 X", "Hook", "確定性觸發"],
            ["連接外部資料庫 / API", "MCP", "標準化介面"],
            ["想把上述組合成可發布的套件", "Plugin", "統一打包"],
            ["想分享給團隊 / 社群", "Marketplace", "版本管理 + 發布"],
        ],
        Inches(0.5), Inches(1.8), Inches(12.333), Inches(4.8),
        font_size=13
    )

    h.add_callout(
        slide,
        "元件之間可以組合：例如 Plugin 可以包含 Skills + Subagents + Hooks + MCP servers",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="🧩", font_size=13
    )

    # ============================================================
    # Slide 7：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "8 份系列文件導覽", "完整檔案地圖與學習路徑")

    # ============================================================
    # Slide 8：8 份文件總覽（表格）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "8 份系列文件總覽",
        "每份對應一個主題，獨立可閱讀",
        slide_num=8, total=TOTAL, source="00 § 系列目錄"
    )

    h.add_comparison_table(
        slide,
        ["編號", "檔名", "主題", "大小"],
        [
            ["00", "00-claude-code-plugins-series.md", "系列總覽（本文件）", "8 KB"],
            ["01", "01-plugin-marketplaces.md", "Plugin 集合的建立與發布", "32 KB"],
            ["02", "02-plugins.md", "Plugin 開發入門（從零開始）", "20 KB"],
            ["03", "03-plugins-reference.md", "完整技術規格與 CLI 參考", "44 KB"],
            ["04", "04-skills.md", "Skills 設計、撰寫、評估", "44 KB"],
            ["05", "05-subagents.md", "隔離上下文與平行任務", "40 KB"],
            ["06", "06-hooks.md", "事件驅動自動化", "64 KB"],
            ["07", "07-discover-plugins.md", "從市場找到並使用 plugin", "20 KB"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(5.0),
        font_size=11
    )

    h.add_callout(
        slide,
        "所有檔案在 /home/elan/pi-proj/ 目錄下，副檔名都是 .md",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="📁", font_size=12
    )

    # ============================================================
    # Slide 9：學習路徑 - 新手
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "🟢 新手入門路徑（建議 1-2 天）",
        "從零開始，快速建立第一個 plugin",
        slide_num=9, total=TOTAL, source="00 § 學習路徑"
    )

    # 4 步驟橫向流程圖
    steps = [
        ("1", "先讀 #07", "探索並安裝", "實際裝幾個 plugin 玩玩", h.COLOR_BLUE),
        ("2", "讀 #04", "Skills", "學會寫簡單的 skill", h.COLOR_GREEN),
        ("3", "讀 #02", "Plugin 開發", "把 skill 包成 plugin", h.COLOR_PRIMARY),
        ("4", "動手做", "claude plugin init", "建立第一個 plugin！", h.COLOR_RED),
    ]

    step_w = Inches(2.8)
    step_h = Inches(3.5)
    gap = Inches(0.3)
    total_w = step_w * 4 + gap * 3
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.2)

    for i, (num, title, subtitle, desc, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # 步驟方塊
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y, step_w, step_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(2.5)

        # 步驟編號圓圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + step_w/2 - Inches(0.4), start_y - Inches(0.2),
            Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        h.add_text_block(
            slide, num,
            x + step_w/2 - Inches(0.4), start_y - Inches(0.2),
            Inches(0.8), Inches(0.8),
            font_size=28, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 步驟標題
        h.add_text_block(
            slide, title,
            x, start_y + Inches(0.8), step_w, Inches(0.6),
            font_size=18, bold=True, color=color,
            align=PP_ALIGN.CENTER
        )

        # 副標題
        h.add_text_block(
            slide, subtitle,
            x, start_y + Inches(1.5), step_w, Inches(0.5),
            font_size=14, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER, bold=True
        )

        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), start_y + Inches(2.2),
            step_w - Inches(0.4), Inches(1.0),
            font_size=11, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

        # 箭頭（除了最後一個）
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + step_w + Inches(0.05), start_y + step_h/2 - Inches(0.15),
                gap - Inches(0.1), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    h.add_callout(
        slide,
        "重點：動手做比讀完更重要！每讀完一份，立刻實際操作一次。",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5),
        icon="🎯", font_size=14
    )

    # ============================================================
    # Slide 10：學習路徑 - 進階
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "🟡 進階使用者路徑（建議 3-5 天）",
        "掌握 Subagents、Hooks、Marketplace",
        slide_num=10, total=TOTAL, source="00 § 學習路徑"
    )

    steps = [
        ("1", "讀 #05", "Subagents", "隔離 context、平行任務", h.COLOR_BLUE),
        ("2", "讀 #06", "Hooks", "事件驅動自動化", h.COLOR_GREEN),
        ("3", "讀 #01", "Marketplaces", "建立團隊 marketplace", h.COLOR_PRIMARY),
        ("4", "實戰", "完整 plugin", "結合所有元件的專案", h.COLOR_RED),
    ]

    for i, (num, title, subtitle, desc, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(2.5)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + step_w/2 - Inches(0.4), start_y - Inches(0.2), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        h.add_text_block(slide, num, x + step_w/2 - Inches(0.4), start_y - Inches(0.2), Inches(0.8), Inches(0.8), font_size=28, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER)
        h.add_text_block(slide, title, x, start_y + Inches(0.8), step_w, Inches(0.6), font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        h.add_text_block(slide, subtitle, x, start_y + Inches(1.5), step_w, Inches(0.5), font_size=14, color=h.COLOR_DARK, align=PP_ALIGN.CENTER, bold=True)
        h.add_text_block(slide, desc, x + Inches(0.2), start_y + Inches(2.2), step_w - Inches(0.4), Inches(1.0), font_size=11, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_w + Inches(0.05), start_y + step_h/2 - Inches(0.15), gap - Inches(0.1), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    h.add_callout(
        slide, "挑戰：建立一個 Hook 自動跑 ESLint + Subagent 做程式碼審查 + 上傳到 Marketplace",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5),
        icon="🏆", font_size=13
    )

    # ============================================================
    # Slide 11：學習路徑 - 專家
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "🔴 專家 / 團隊負責人路徑（建議 1 週+）",
        "建立企業級 plugin 系統",
        slide_num=11, total=TOTAL, source="00 § 學習路徑"
    )

    steps = [
        ("1", "讀 #03", "技術參考", "掌握所有細節", h.COLOR_BLUE),
        ("2", "回頭讀 #01", "Marketplace 進階", "企業 marketplace", h.COLOR_GREEN),
        ("3", "整合", "所有元件", "複雜多功能 plugin", h.COLOR_PRIMARY),
        ("4", "部署", "企業級", "跨團隊 plugin 系統", h.COLOR_RED),
    ]

    for i, (num, title, subtitle, desc, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(2.5)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + step_w/2 - Inches(0.4), start_y - Inches(0.2), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        h.add_text_block(slide, num, x + step_w/2 - Inches(0.4), start_y - Inches(0.2), Inches(0.8), Inches(0.8), font_size=28, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER)
        h.add_text_block(slide, title, x, start_y + Inches(0.8), step_w, Inches(0.6), font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        h.add_text_block(slide, subtitle, x, start_y + Inches(1.5), step_w, Inches(0.5), font_size=14, color=h.COLOR_DARK, align=PP_ALIGN.CENTER, bold=True)
        h.add_text_block(slide, desc, x + Inches(0.2), start_y + Inches(2.2), step_w - Inches(0.4), Inches(1.0), font_size=11, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_w + Inches(0.05), start_y + step_h/2 - Inches(0.15), gap - Inches(0.1), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_PRIMARY
            arrow.line.fill.background()

    h.add_callout(
        slide, "進階：了解 strictKnownMarketplaces、managed settings、企業部署",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5),
        icon="👑", font_size=13
    )

    # ============================================================
    # Slide 12：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "任務導向索引", "想做 X？看哪一份")

    # ============================================================
    # Slide 13：任務索引
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "任務導向索引",
        "想做某個特定的事？從這裡開始找",
        slide_num=13, total=TOTAL, source="00 § 任務索引"
    )

    h.add_comparison_table(
        slide,
        ["你的目標", "該看哪份"],
        [
            ["找別人寫好的 plugin 來用", "→ #07-discover-plugins.md"],
            ["寫第一個 plugin", "→ #02-plugins.md"],
            ["寫可重用的 /command 指令", "→ #04-skills.md"],
            ["寫可重用的「代理人」", "→ #05-subagents.md"],
            ["編輯檔案後自動跑測試", "→ #06-hooks.md"],
            ["阻擋危險命令（rm -rf）", "→ #06-hooks.md § PreToolUse"],
            ["部署到團隊/公司", "→ #01-plugin-marketplaces.md"],
            ["提交到 Anthropic 官方", "→ #02-plugins.md § 提交"],
            ["理解完整技術規格", "→ #03-plugins-reference.md"],
            ["查特定 CLI 指令用法", "→ #03-plugins-reference.md § CLI"],
            ["解決疑難問題", "→ 各文件最後的「疑難排解」章節"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(5.2),
        font_size=13
    )

    # ============================================================
    # Slide 14：必裝官方 Plugin
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "推薦先裝的官方 Plugin",
        "從 Marketplace 開始體驗",
        slide_num=14, total=TOTAL, source="07-discover-plugins.md"
    )

    plugins = [
        ("commit-commands", "Git 提交工作流", "適合所有人的日常必備", h.COLOR_GREEN),
        ("security-guidance", "自動安全審查", "每次變更都檢查常見漏洞", h.COLOR_RED),
        ("pr-review-toolkit", "PR 審查工具", "團隊協作的利器", h.COLOR_BLUE),
        ("typescript-lsp", "TypeScript 程式碼智慧", "即時型別錯誤檢查", h.COLOR_PRIMARY),
        ("pyright-lsp", "Python 程式碼智慧", "即時型別錯誤檢查", h.COLOR_PRIMARY),
        ("explanatory-output-style", "教育性輸出", "學習 Claude 決策的推理", RGBColor(0x7C, 0x3A, 0xED)),
    ]

    box_w = Inches(5.8)
    box_h = Inches(1.2)
    h_gap = Inches(0.4)
    v_gap = Inches(0.25)
    grid_cols = 2
    grid_rows = 3
    grid_w = box_w * grid_cols + h_gap
    grid_h = box_h * grid_rows + v_gap * 2
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.8)

    for i, (name, role, desc, color) in enumerate(plugins):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        # 插件卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # 左側色條
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.15), box_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 名稱
        h.add_text_block(
            slide, name,
            x + Inches(0.3), y + Inches(0.1), box_w - Inches(0.4), Inches(0.4),
            font_size=15, bold=True, color=color
        )
        # 角色
        h.add_text_block(
            slide, role,
            x + Inches(0.3), y + Inches(0.5), box_w - Inches(0.4), Inches(0.3),
            font_size=11, bold=True, color=h.COLOR_DARK
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.3), y + Inches(0.8), box_w - Inches(0.4), Inches(0.4),
            font_size=10, color=h.COLOR_GRAY_TXT
        )

    h.add_callout(
        slide, "新手建議：先裝 commit-commands + security-guidance，體驗 10 分鐘就有感",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        icon="🚀", font_size=13
    )

    # ============================================================
    # Slide 15：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "7 大元件速覽", "每個元件的核心概念與範例")

    # ============================================================
    # Slide 16：CLAUDE.md
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 1/7：CLAUDE.md",
        "每次 session 都會看到的「專案說明書」",
        slide_num=16, total=TOTAL, source="memory + best-practices"
    )

    h.add_two_column_compare(
        slide,
        "📖 是什麼",
        [
            "Markdown 檔案，放在專案根目錄",
            "每次 Claude Code 啟動時自動載入",
            "作為「系統提示的一部分」",
            "可放在 .claude/CLAUDE.md 或子目錄",
            "可使用 @path 引用其他檔案"
        ],
        "✨ 適合放什麼",
        [
            "專案慣例與風格指南",
            "「Always do X」型規則",
            "建構/測試指令",
            "禁止事項的明確規定",
            "新成員 onboarding 須知"
        ]
    )

    h.add_callout(
        slide,
        "💡 規則：CLAUDE.md 保持在 200 行以下。複雜內容放 Skills 或 .claude/rules/",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="📏", font_size=12
    )

    # ============================================================
    # Slide 17：Skills
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 2/7：Skills",
        "可重複使用的指令知識庫（/command 觸發）",
        slide_num=17, total=TOTAL, source="04-skills.md"
    )

    h.add_text_block(
        slide, "基本結構",
        Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """~/.claude/skills/
└── summarize-changes/
    ├── SKILL.md          ← 必需
    ├── reference.md       ← 選用
    └── scripts/           ← 選用
        └── helper.py""",
        Inches(0.5), Inches(2.0), Inches(5.5), Inches(2.0),
        font_size=12
    )

    h.add_text_block(
        slide, "SKILL.md 範例",
        Inches(0.5), Inches(4.2), Inches(6), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """---
description: 總結未提交的變更
  並標記風險內容
---

## Current changes
!`git diff HEAD`

## Instructions
Summarize in 2-3 bullets,
list risks you notice...""",
        Inches(0.5), Inches(4.6), Inches(5.5), Inches(2.2),
        font_size=11
    )

    # 右側說明
    h.add_text_block(
        slide, "💡 兩種角色",
        Inches(6.5), Inches(1.6), Inches(6.5), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "參考型：API 慣例、命名規範",
            "任務型：/deploy、/commit 流程",
            "任務型加 disable-model-invocation: true",
            "  → 防止 Claude 自動觸發",
            "  → 只在你想執行時才會跑",
        ],
        Inches(6.5), Inches(2.0), Inches(6.5), Inches(2.0),
        font_size=13
    )

    h.add_text_block(
        slide, "🔧 動態上下文",
        Inches(6.5), Inches(4.2), Inches(6.5), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "!`git diff HEAD` 預處理",
            "呼叫 shell 取得即時資料",
            "結果直接替換佔位符",
            "Claude 看到的是最終結果",
        ],
        Inches(6.5), Inches(4.6), Inches(6.5), Inches(2.2),
        font_size=13
    )

    # ============================================================
    # Slide 18：Subagents
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 3/7：Subagents",
        "隔離上下文的子代理人",
        slide_num=18, total=TOTAL, source="05-subagents.md"
    )

    h.add_two_column_compare(
        slide,
        "🎯 核心特性",
        [
            "獨立 context 視窗",
            "專門的系統提示",
            "獨立的工具權限",
            "可指定模型（便宜模型省錢）",
            "完成後返回摘要給主對話"
        ],
        "🚀 使用時機",
        [
            "旁支任務會淹沒主對話",
            "需要強制工具限制",
            "可平行執行的獨立任務",
            "想重用某個專門設定",
            "想用 Haiku 省成本"
        ]
    )

    h.add_callout(
        slide, "3 個內建 Subagent：Explore（唯讀）、Plan（計畫模式）、general-purpose",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="🛠", font_size=12
    )

    # ============================================================
    # Slide 19：Hooks
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 4/7：Hooks",
        "事件驅動的確定性自動化",
        slide_num=19, total=TOTAL, source="06-hooks.md"
    )

    h.add_text_block(
        slide, "⚡ Hook 5 種類型",
        Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "command：執行 shell 命令",
            "http：POST 到 URL endpoint",
            "mcp_tool：呼叫 MCP 工具",
            "prompt：用 LLM 評估（Haiku）",
            "agent：subagent 驗證器（實驗性）",
        ],
        Inches(0.5), Inches(2.0), Inches(6), Inches(2.5),
        font_size=13
    )

    h.add_text_block(
        slide, "🎯 熱門使用場景",
        Inches(0.5), Inches(4.7), Inches(6), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "PostToolUse 自動跑 formatter",
            "PreToolUse 阻擋 rm -rf",
            "SessionStart 載入專案 context",
            "Notification 發送桌面通知",
            "Stop 自動跑測試再停止",
        ],
        Inches(0.5), Inches(5.1), Inches(6), Inches(1.8),
        font_size=13
    )

    # 右側：完整生命週期事件
    h.add_text_block(
        slide, "🔄 主要生命週期事件",
        Inches(6.5), Inches(1.6), Inches(6.5), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """SessionStart      ← session 開始
UserPromptSubmit   ← 提交 prompt
PreToolUse         ← 工具執行前（可阻擋）
PostToolUse        ← 工具成功後
PostToolUseFailure ← 工具失敗後
Notification       ← 系統通知
SubagentStart/Stop ← subagent 生命週期
PreCompact         ← context 壓縮前
Stop               ← Claude 完成回應
SessionEnd         ← session 結束

（總共 28 個事件）""",
        Inches(6.5), Inches(2.0), Inches(6.5), Inches(4.5),
        font_size=12
    )

    # ============================================================
    # Slide 20：MCP
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 5/7：MCP (Model Context Protocol)",
        "連接外部服務的標準化介面",
        slide_num=20, total=TOTAL, source="MCP 文件"
    )

    h.add_text_block(
        slide, "🔌 MCP = 連接外部世界的標準協議",
        Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.5),
        font_size=20, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_two_column_compare(
        slide,
        "✅ MCP 提供",
        [
            "外部工具和資料存取",
            "資料庫查詢",
            "Slack 訊息發送",
            "瀏覽器控制（Playwright）",
            "GitHub/GitLab 整合",
            "統一認證和連線管理"
        ],
        "🎁 範例 plugin",
        [
            "github：原始碼控制",
            "gitlab：原始碼控制",
            "atlassian：Jira/Confluence",
            "figma：設計整合",
            "vercel：基礎設施",
            "sentry：監控"
        ]
    )

    h.add_callout(
        slide, "Skills 教你怎麼用好 MCP！MCP 提供工具 + Skills 教策略 = 完美組合",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="💎", font_size=13
    )

    # ============================================================
    # Slide 21：Plugins
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 6/7：Plugins",
        "把所有元件打包成可發布的單元",
        slide_num=21, total=TOTAL, source="02-plugins.md"
    )

    h.add_text_block(
        slide, "📦 Plugin = 上述所有元件的包裝箱",
        Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.5),
        font_size=20, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "標準目錄結構",
        Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """my-plugin/
├── .claude-plugin/
│   └── plugin.json     ← manifest
├── skills/             ← /name 觸發
│   └── hello/
│       └── SKILL.md
├── agents/             ← 專門代理
├── hooks/              ← 事件處理
│   └── hooks.json
├── .mcp.json           ← MCP servers
├── .lsp.json           ← LSP servers
└── bin/                ← 可執行檔""",
        Inches(0.5), Inches(2.7), Inches(5.8), Inches(4.0),
        font_size=11
    )

    h.add_text_block(
        slide, "🎯 為什麼用 Plugin？",
        Inches(6.5), Inches(2.3), Inches(6.5), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "與團隊共享",
            "跨專案重用",
            "版本控制（明確版本 / git SHA）",
            "輕鬆更新",
            "社群分享",
            "⚠️ 命名空間：",
            "  /my-plugin:hello",
            "  避免 plugin 間衝突",
        ],
        Inches(6.5), Inches(2.7), Inches(6.5), Inches(4.0),
        font_size=13
    )

    # ============================================================
    # Slide 22：Marketplaces
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "元件 7/7：Marketplaces",
        "Plugin 集合的「商店目錄」",
        slide_num=22, total=TOTAL, source="01-plugin-marketplaces.md"
    )

    h.add_text_block(
        slide, "🏪 Marketplace = 多個 Plugin 的版本化目錄",
        Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.5),
        font_size=20, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "三種 Marketplace 來源",
        Inches(0.5), Inches(2.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_DARK
    )

    sources = [
        ("📦 官方", "claude-plugins-official", "Anthropic 策展", "自動註冊", h.COLOR_PRIMARY),
        ("🌐 社群", "claude-plugins-community", "第三方提交", "手動加入", h.COLOR_BLUE),
        ("🛠 自建", "你的 marketplace.json", "完全自訂", "手動加入", h.COLOR_GREEN),
    ]

    box_w = Inches(4.0)
    box_h = Inches(1.5)
    h_gap = Inches(0.3)
    total_w = box_w * 3 + h_gap * 2
    start_x = (h.SLIDE_W - total_w) / 2

    for i, (icon_label, name, desc, install, color) in enumerate(sources):
        x = start_x + i * (box_w + h_gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.8), box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        h.add_text_block(
            slide, f"{icon_label}",
            x, Inches(2.9), box_w, Inches(0.4),
            font_size=20, bold=True, color=color,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, name,
            x, Inches(3.3), box_w, Inches(0.4),
            font_size=13, bold=True, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER, font=h.COLOR_MONO if hasattr(h, 'COLOR_MONO') else "Consolas"
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), Inches(3.7), box_w - Inches(0.4), Inches(0.3),
            font_size=11, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, f"安裝方式：{install}",
            x + Inches(0.2), Inches(4.0), box_w - Inches(0.4), Inches(0.3),
            font_size=10, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER
        )

    h.add_text_block(
        slide, "🔐 受管設定",
        Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    h.add_code_block(
        slide, """strictKnownMarketplaces: 限制使用者能新增的 marketplace
disableSideloadFlags: 拒絕 CLI 側載 plugin
pluginSuggestionMarketplaces: 允許內容相關建議""",
        Inches(2.0), Inches(5.2), Inches(9.333), Inches(1.4),
        font_size=12
    )

    # ============================================================
    # Slide 23：元件組合示意圖
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "實戰：所有元件的組合應用",
        "Plugin 可以包含 Skills + Subagents + Hooks + MCP",
        slide_num=23, total=TOTAL, source="00 § 組合"
    )

    # 中央 Plugin 大方塊
    center_x = Inches(4.0)
    center_y = Inches(2.0)
    center_w = Inches(5.333)
    center_h = Inches(3.0)

    main = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_x, center_y, center_w, center_h)
    main.fill.solid()
    main.fill.fore_color.rgb = h.COLOR_BG_GRAY
    main.line.color.rgb = h.COLOR_PRIMARY
    main.line.width = Pt(3)

    h.add_text_block(
        slide, "📦 Plugin",
        center_x, center_y + Inches(0.1), center_w, Inches(0.4),
        font_size=22, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    # Plugin 內的 4 個子元件（小方塊）
    inner_w = Inches(2.3)
    inner_h = Inches(0.9)
    inner_gap = Inches(0.2)

    inner_components = [
        ("Skills", "/my-plugin:hello", h.COLOR_BLUE),
        ("Subagents", "code-reviewer", h.COLOR_GREEN),
        ("Hooks", "PostToolUse", h.COLOR_RED),
        ("MCP servers", "github, slack", RGBColor(0x7C, 0x3A, 0xED)),
    ]

    inner_x_start = center_x + (center_w - (inner_w * 2 + inner_gap)) / 2
    inner_y_start = center_y + Inches(0.7)

    for i, (name, detail, color) in enumerate(inner_components):
        row = i // 2
        col = i % 2
        x = inner_x_start + col * (inner_w + inner_gap)
        y = inner_y_start + row * (inner_h + inner_gap)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, inner_w, inner_h)
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        h.add_text_block(
            slide, name,
            x, y + Inches(0.05), inner_w, Inches(0.3),
            font_size=13, bold=True, color=color,
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, detail,
            x, y + Inches(0.4), inner_w, Inches(0.4),
            font_size=10, color=h.COLOR_GRAY_TXT,
            align=PP_ALIGN.CENTER, font="Consolas"
        )

    # 周圍的 Marketplace
    mkt_x = Inches(10.0)
    mkt_y = Inches(3.0)
    mkt = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mkt_x, mkt_y, Inches(2.5), Inches(1.0))
    mkt.fill.solid()
    mkt.fill.fore_color.rgb = h.COLOR_WHITE
    mkt.line.color.rgb = RGBColor(0x0E, 0x7C, 0x66)
    mkt.line.width = Pt(2.5)
    h.add_text_block(
        slide, "🏪 Marketplace",
        mkt_x, mkt_y + Inches(0.1), Inches(2.5), Inches(0.3),
        font_size=13, bold=True, color=RGBColor(0x0E, 0x7C, 0x66),
        align=PP_ALIGN.CENTER
    )
    h.add_text_block(
        slide, "多個 plugins\n+ 版本管理",
        mkt_x, mkt_y + Inches(0.45), Inches(2.5), Inches(0.5),
        font_size=10, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    # 連接箭頭
    h.add_arrow(slide, Inches(9.3), Inches(3.5), Inches(10.0), Inches(3.5), color=RGBColor(0x0E, 0x7C, 0x66), width=2)

    # CLAUDE.md 在左邊
    claude_x = Inches(0.5)
    claude_y = Inches(3.0)
    claude = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, claude_x, claude_y, Inches(3.0), Inches(1.0))
    claude.fill.solid()
    claude.fill.fore_color.rgb = h.COLOR_WHITE
    claude.line.color.rgb = h.COLOR_DARK
    claude.line.width = Pt(2.5)
    h.add_text_block(
        slide, "📋 CLAUDE.md",
        claude_x, claude_y + Inches(0.1), Inches(3.0), Inches(0.3),
        font_size=13, bold=True, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )
    h.add_text_block(
        slide, "每次 session 載入\n的專案說明書",
        claude_x, claude_y + Inches(0.45), Inches(3.0), Inches(0.5),
        font_size=10, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    # 連接箭頭（CLAUDE.md → Plugin）
    h.add_arrow(slide, Inches(3.5), Inches(3.5), Inches(4.0), Inches(3.5), color=h.COLOR_DARK, width=2)

    h.add_callout(
        slide, "完整 Plugin 範例：typeform-feedback 插件（整合 Form、Slack、GitHub Issue）",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        icon="🌟", font_size=13
    )

    h.add_text_block(
        slide, "真實世界：每個你看到的「整合 Claude + X」背後，都是這個組合",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER, bold=False, font="Consolas"
    )

    # ============================================================
    # Slide 24：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "實際案例與最佳實踐", "從理論到實戰")

    # ============================================================
    # Slide 25：真實使用案例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "4 個真實世界的 Plugin 使用案例",
        "看看別人怎麼用",
        slide_num=25, total=TOTAL, source="各 plugin 範例"
    )

    cases = [
        ("🛡️", "安全審查",
         "security-guidance plugin",
         "每次 Claude 寫程式時自動檢查 SQL injection、XSS、敏感資料洩漏。發現問題就指示 Claude 在同一輪修復。",
         h.COLOR_RED),
        ("🔍", "程式碼審查",
         "code-reviewer subagent",
         "使用唯讀工具掃描 PR 變更，輸出 Critical/Warning/Suggestion 三級分類。隔離 context 讓主對話保持乾淨。",
         h.COLOR_BLUE),
        ("🚀", "PR 工作流",
         "commit-commands plugin",
         "/commit-commands:commit 自動暫存 → 產生訊息 → 建立 commit。/pr 自動推 PR + 填寫模板。",
         h.COLOR_GREEN),
        ("📊", "資料視覺化",
         "codebase-visualizer skill",
         "!`find . -type f` 預處理 → Python 腳本生成互動式 HTML → 在瀏覽器打開可摺疊樹狀圖。",
         h.COLOR_PRIMARY),
    ]

    box_w = Inches(5.8)
    box_h = Inches(2.4)
    h_gap = Inches(0.4)
    v_gap = Inches(0.3)
    grid_cols = 2
    grid_rows = 2

    for i, (icon, title, subtitle, desc, color) in enumerate(cases):
        row = i // grid_cols
        col = i % grid_cols
        x = Inches(0.5) + col * (box_w + h_gap)
        y = Inches(1.7) + row * (box_h + v_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # 左側大 icon
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.7))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = color
        icon_bg.line.fill.background()
        h.add_text_block(
            slide, icon,
            x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.7),
            font_size=24, align=PP_ALIGN.CENTER
        )

        # 標題
        h.add_text_block(
            slide, title,
            x + Inches(1.1), y + Inches(0.2), box_w - Inches(1.2), Inches(0.4),
            font_size=18, bold=True, color=color
        )
        # 副標題
        h.add_text_block(
            slide, subtitle,
            x + Inches(1.1), y + Inches(0.6), box_w - Inches(1.2), Inches(0.3),
            font_size=11, color=h.COLOR_DARK, font="Consolas"
        )
        # 描述
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.0), box_w - Inches(0.4), Inches(1.3),
            font_size=12, color=h.COLOR_DARK
        )

    # ============================================================
    # Slide 26：最佳實踐
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "開發 Plugin 的最佳實踐",
        "從一開始就做對，減少重構",
        slide_num=26, total=TOTAL, source="02 § 開發指南 + 03 § 技術參考"
    )

    h.add_two_column_compare(
        slide,
        "✅ 應該做",
        [
            "✅ 從獨立配置開始（.claude/）",
            "✅ 確定能用後再包成 Plugin",
            "✅ Plugin 命名統一 kebab-case",
            "✅ 寫詳細的 description",
            "✅ 完整 README.md（安裝、使用）",
            "✅ 用 !`command` 預處理動態資料",
            "✅ 在 SKILL.md 中參考支援檔案",
            "✅ Skills 控制在 500 行以內",
        ],
        "❌ 不該做",
        [
            "❌ 元件放在 .claude-plugin/ 內",
            "❌ 在 .claude-plugin/ 放其他目錄",
            "❌ 跨 plugin 目錄共享（../）",
            "❌ 絕對路徑",
            "❌ 一次寫太大，慢慢組合",
            "❌ Hook 沒給可執行權限",
            "❌ 沒驗證就 commit",
            "❌ 忽略 validate 警告",
        ]
    )

    # ============================================================
    # Slide 27：常見陷阱
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "新手常見的 5 個陷阱",
        "別讓這些小錯誤浪費你的時間",
        slide_num=27, total=TOTAL, source="各章節疑難排解"
    )

    pitfalls = [
        ("🚫 元件放錯位置",
         "commands/agents/skills/hooks 必須在 plugin 根目錄，不是 .claude-plugin/ 內。只有 plugin.json 屬於 .claude-plugin/。",
         h.COLOR_RED),
        ("🚫 沒 chmod",
         "Hook 腳本沒設可執行權限 → 靜默失敗。記得 chmod +x script.sh。",
         h.COLOR_ORANGE if hasattr(h, 'COLOR_ORANGE') else h.COLOR_RED),
        ("🚫 絕對路徑",
         "Plugin 必須用相對路徑（./），絕對路徑在其他機器會壞掉。",
         h.COLOR_RED),
        ("🚫 ../ 路徑",
         "Plugin 被複製到快取，../shared 不會被複製。需要共享用 symlinks。",
         h.COLOR_RED),
        ("🚫 跳過驗證",
         "claude plugin validate 警告也要看。拼錯欄位名是常見錯誤。",
         h.COLOR_RED),
    ]

    for i, (title, desc, color) in enumerate(pitfalls):
        y = Inches(1.8) + i * Inches(0.95)

        # 標題
        h.add_text_block(
            slide, title,
            Inches(0.5), y, Inches(12.333), Inches(0.4),
            font_size=16, bold=True, color=color
        )
        # 描述
        h.add_text_block(
            slide, desc,
            Inches(0.7), y + Inches(0.4), Inches(12.133), Inches(0.5),
            font_size=12, color=h.COLOR_DARK
        )

    h.add_callout(
        slide, "養成習慣：寫完先 claude plugin validate，跑過再測試",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="🔍", font_size=13
    )

    # ============================================================
    # Slide 28：重點回顧
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_summary_slide(
        slide,
        title="重點回顧",
        key_points=[
            "7 個核心元件：CLAUDE.md、Skills、Subagents、Hooks、MCP、Plugins、Marketplaces",
            "8 份文件、6,682 行、272 KB — 涵蓋 Claude Code v2.1.x 完整功能",
            "學習路徑：新手（07→04→02）→ 進階（05→06→01）→ 專家（03 細讀）",
            "推薦起手：先裝官方 plugin（commit-commands）→ 寫第一個 skill → 包成 plugin",
            "設計原則：標準型 + 詳細型密度，每頁一個核心概念",
        ],
        next_steps=[
            "🎯 立刻：執行 /plugin，安裝 commit-commands 試試",
            "📚 30 分鐘：翻完 #04-skills.md 的目錄",
            "🛠 2 小時：跟著 #02-plugins.md 步驟建立第一個 plugin",
            "🚀 一週：完成學習路徑中的「實戰」項目",
        ],
        source="00-claude-code-plugins-series.md"
    )

    # ============================================================
    # Slide 29：4 份簡報總覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "📊 本簡報系列總覽",
        "4 份 PPTX 對應 4 個學習主題",
        slide_num=29, total=TOTAL, source="00 § 簡報系列"
    )

    ppts = [
        ("00-overview.pptx", "本檔案", "系列總覽、生態系、學習路徑", "30 張", h.COLOR_PRIMARY),
        ("02-plugins.pptx", "Plugin 開發", "從零建立第一個 plugin", "25 張", h.COLOR_BLUE),
        ("04-skills.pptx", "Skills", "Skills 設計、撰寫、評估", "40 張", h.COLOR_GREEN),
        ("06-hooks.pptx", "Hooks", "事件驅動自動化", "50 張", h.COLOR_RED),
    ]

    box_w = Inches(2.9)
    box_h = Inches(2.8)
    h_gap = Inches(0.25)
    total_w = box_w * 4 + h_gap * 3
    start_x = (h.SLIDE_W - total_w) / 2

    for i, (filename, title, desc, pages, color) in enumerate(ppts):
        x = start_x + i * (box_w + h_gap)
        y = Inches(2.2)

        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # 頂部色條
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.6))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()

        h.add_text_block(
            slide, f"📄 {title}",
            x, y + Inches(0.1), box_w, Inches(0.4),
            font_size=14, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, filename,
            x + Inches(0.2), y + Inches(0.8), box_w - Inches(0.4), Inches(0.5),
            font_size=11, color=h.COLOR_DARK, font="Consolas",
            align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.4), box_w - Inches(0.4), Inches(0.8),
            font_size=11, color=h.COLOR_DARK,
            align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, pages,
            x + Inches(0.2), y + Inches(2.3), box_w - Inches(0.4), Inches(0.4),
            font_size=14, bold=True, color=color,
            align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "剩下 4 份（marketplaces、reference、subagents、discover-plugins）依需求再做",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        icon="📅", font_size=13
    )

    h.add_text_block(
        slide, "所有 PPTX 都在 /home/elan/pi-proj/ 目錄下",
        Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER, bold=False
    )

    # ============================================================
    # Slide 30：結束頁
    # ============================================================
    slide = h.add_blank_slide(prs)
    set_slide_bg = lambda s: h.set_slide_bg(s, h.COLOR_BG_CREAM)
    set_slide_bg(slide)

    # 大標題
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "開始你的 Plugin 之旅 🚀"
    run.font.name = h.FONT_TITLE
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = h.COLOR_PRIMARY

    h.add_text_block(
        slide, "從 #02-plugins.md 開始，建立你的第一個 plugin",
        Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
        font_size=20, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    # 橘色裝飾條
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.666), Inches(5.0),
        Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = h.COLOR_PRIMARY
    bar.line.fill.background()

    h.add_text_block(
        slide, "Claude Code Plugin 完整學習系列 · 2026/01",
        Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
        font_size=14, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "整理自官方文件（code.claude.com/docs）· 繁體中文教學用途",
        Inches(1), Inches(6.1), Inches(11.333), Inches(0.5),
        font_size=12, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    # 儲存
    output = "/home/elan/pi-proj/00-overview.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
