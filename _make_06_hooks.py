"""
簡報 4/4：Hooks 自動化指南 (06-hooks.pptx)
約 50 張
對應：06-hooks.md
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import _pptx_helpers as h


def build():
    prs = h.new_presentation()
    TOTAL = 50

    # ============================================================
    # 封面
    # ============================================================
    h.add_cover_slide(
        prs,
        "Hooks 自動化指南",
        "事件驅動的確定性工作流",
        tag="#06 · Hooks"
    )

    # ============================================================
    # Slide 2：本章學習目標
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "本章你會學到",
        "從基礎到進階的 Hook 完整知識",
        slide_num=2, total=TOTAL, source="06-hooks.md"
    )

    objectives = [
        ("🎯", "Hook 是什麼", "確定性 vs LLM 判斷的差異"),
        ("⏱", "Hook 生命週期", "28 個事件完整解析"),
        ("📍", "Hook 設定位置", "使用者/專案/Plugin/Skill"),
        ("🎨", "Matcher 模式", "正則、字串、混用規則"),
        ("🔧", "5 種 Handler 類型", "command、http、mcp_tool、prompt、agent"),
        ("📥", "Hook 輸入/輸出", "JSON 格式與退出碼"),
        ("🎯", "決策控制", "allow/deny/ask/改寫輸入"),
        ("🛡", "安全最佳實踐", "避免踩坑"),
        ("🐛", "除錯技巧", "debug 模式與日誌"),
    ]

    box_w = Inches(3.0)
    box_h = Inches(1.5)
    h_gap = Inches(0.3)
    v_gap = Inches(0.25)
    grid_cols = 3
    grid_rows = 3
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap * (grid_rows - 1)
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.7)

    for i, (icon, title, desc) in enumerate(objectives):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_RED
        card.line.width = Pt(1.5)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.1), Inches(0.6), Inches(0.5),
            font_size=22, align=PP_ALIGN.LEFT
        )
        h.add_text_block(
            slide, title,
            x + Inches(0.7), y + Inches(0.1),
            box_w - Inches(0.8), Inches(0.5),
            font_size=13, bold=True, color=h.COLOR_DARK
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(0.7),
            box_w - Inches(0.4), Inches(0.7),
            font_size=10, color=h.COLOR_GRAY_TXT
        )

    # ============================================================
    # Slide 3：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 1", "Hook 基礎概念", "理解確定性 vs 判斷性")

    # ============================================================
    # Slide 4：什麼是 Hook
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "什麼是 Hook？",
        "使用者定義的自動化執行點",
        slide_num=4, total=TOTAL, source="06 § 什麼是 Hook"
    )

    h.add_text_block(
        slide, "🎣 Hook = 在生命週期事件上自動觸發的腳本/HTTP/LLM 評估",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.5),
        font_size=20, bold=True, color=h.COLOR_RED,
        align=PP_ALIGN.CENTER
    )

    h.add_two_column_compare(
        slide,
        "✅ Hook 提供（確定性）",
        [
            "使用者定義的 shell 命令",
            "HTTP endpoint",
            "LLM prompt（半確定）",
            "Subagent 驗證器",
            "特定點必觸發",
            "適合：lint、阻擋、通知",
        ],
        "🧠 由 Claude 判斷（不確定）",
        [
            "Claude 自行決定何時做",
            "依賴 LLM 理解",
            "可能漏掉或做錯",
            "適合：開放性任務",
            "不可強制",
        ]
    )

    h.add_callout(
        slide, "口訣：「必須做的事」用 Hook，「希望做的事」用 CLAUDE.md/Skill",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="💡", font_size=14
    )

    # ============================================================
    # Slide 5：Hook 生命週期
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 生命週期：3 個節奏",
        "28 個事件分屬三類",
        slide_num=5, total=TOTAL, source="06 § 生命週期"
    )

    categories = [
        ("⏱ 每個 Session 一次", h.COLOR_BLUE,
         "SessionStart", "SessionEnd",
         "Session 開始或結束時各執行一次"),
        ("🔄 每輪一次", h.COLOR_GREEN,
         "UserPromptSubmit", "Stop", "StopFailure",
         "每次 prompt 提交或 Claude 回應完成"),
        ("🛠 每次工具呼叫", h.COLOR_RED,
         "PreToolUse", "PostToolUse", "PostToolUseFailure",
         "在 agentic loop 中每次工具呼叫時觸發"),
    ]

    for i, (label, color, *events_desc) in enumerate(categories):
        x = Inches(0.5) + i * Inches(4.3)
        y = Inches(1.8)
        w = Inches(4.0)
        h_box = Inches(3.5)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h_box)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # 標題
        h.add_text_block(
            slide, label,
            x, y + Inches(0.2), w, Inches(0.5),
            font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER
        )

        # 事件
        events = events_desc[:2]
        for j, event in enumerate(events):
            h.add_text_block(
                slide, event,
                x + Inches(0.2), y + Inches(0.9 + j * 0.4), w - Inches(0.4), Inches(0.3),
                font_size=12, color=h.COLOR_DARK, font="Consolas"
            )

        # 描述
        h.add_text_block(
            slide, events_desc[2],
            x + Inches(0.2), y + Inches(2.0), w - Inches(0.4), Inches(1.3),
            font_size=11, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "完整 28 個事件見 06-hooks.md § 完整事件表",
        Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.4),
        icon="📋", font_size=13
    )

    # ============================================================
    # Slide 6：Hook 解析流程
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 解析流程",
        "從事件觸發到決策執行的完整路徑",
        slide_num=6, total=TOTAL, source="06 § 解析流程"
    )

    h.add_text_block(
        slide, "範例：阻擋 rm -rf 命令",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED,
        align=PP_ALIGN.CENTER
    )

    # 流程圖（5 步驟垂直）
    steps = [
        ("1", "PreToolUse 事件觸發", "Claude 送 JSON 到 stdin", h.COLOR_RED),
        ("2", "Matcher 匹配 Bash", "工具名稱過濾", h.COLOR_BLUE),
        ("3", "if 條件匹配 rm *", "細部命令過濾", h.COLOR_GREEN),
        ("4", "腳本執行 + 輸出 JSON", "permissionDecision: deny", h.COLOR_PRIMARY),
        ("5", "Claude Code 阻擋工具", "向 Claude 顯示原因", h.COLOR_DARK),
    ]

    step_h = Inches(0.8)
    step_gap = Inches(0.2)
    total_h = step_h * len(steps) + step_gap * (len(steps) - 1)
    start_y = Inches(2.3)

    for i, (num, title, desc, color) in enumerate(steps):
        y = start_y + i * (step_h + step_gap)

        # 編號圓
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.15), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            Inches(0.5), y + Inches(0.15), Inches(0.5), Inches(0.5),
            font_size=18, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

        # 方塊
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(7.0), step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        h.add_text_block(
            slide, title,
            Inches(1.4), y + Inches(0.1), Inches(6.6), Inches(0.4),
            font_size=14, bold=True, color=color
        )
        h.add_text_block(
            slide, desc,
            Inches(1.4), y + Inches(0.45), Inches(6.6), Inches(0.3),
            font_size=10, color=h.COLOR_GRAY_TXT
        )

        # 連接線（除了最後一個）
        if i < len(steps) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(0.6), y + Inches(0.7), Inches(0.3), step_gap + Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = h.COLOR_GRAY_TXT
            line.line.fill.background()

    # 右側：範例 JSON 輸出
    h.add_code_block(
        slide, """{
  "hookSpecificOutput": {
    "hookEventName": "PreToolUse",
    "permissionDecision": "deny",
    "permissionDecisionReason": "Destructive command blocked"
  }
}""",
        Inches(8.5), Inches(2.3), Inches(4.5), Inches(2.0),
        font_size=10
    )

    h.add_text_block(
        slide, "JSON 決策輸出",
        Inches(8.5), Inches(4.4), Inches(4.5), Inches(0.3),
        font_size=12, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_callout(
        slide, "完整 hook 腳本見 06-hooks.md § 解析流程",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="📖", font_size=12
    )

    # ============================================================
    # Slide 7：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 2", "Hook 設定與位置", "6 個位置決定範圍")

    # ============================================================
    # Slide 8：Hook 位置
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 設定的 6 個位置",
        "範圍決定誰能觸發",
        slide_num=8, total=TOTAL, source="06 § Hook 位置"
    )

    locations = [
        ("🏠", "~/.claude/settings.json", "你的所有專案", "個人", h.COLOR_BLUE),
        ("📁", ".claude/settings.json", "單一專案", "可 commit", h.COLOR_GREEN),
        ("🔒", ".claude/settings.local.json", "單一專案", "gitignored", h.COLOR_RED),
        ("🏢", "Managed policy settings", "組織範圍", "管理員控制", h.COLOR_PRIMARY),
        ("📦", "Plugin hooks/hooks.json", "啟用 plugin 時", "隨 plugin 打包", RGBColor(0x7C, 0x3A, 0xED)),
        ("📝", "Skill/Subagent frontmatter", "限定範圍", "跟著元件走", RGBColor(0x0E, 0x7C, 0x66)),
    ]

    box_w = Inches(3.8)
    box_h = Inches(2.2)
    h_gap = Inches(0.3)
    v_gap = Inches(0.3)
    grid_cols = 3
    grid_rows = 2
    grid_w = box_w * grid_cols + h_gap * (grid_cols - 1)
    grid_h = box_h * grid_rows + v_gap
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.7)

    for i, (icon, path, scope, share, color) in enumerate(locations):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2)

        h.add_text_block(
            slide, icon,
            x, y + Inches(0.15), box_w, Inches(0.4),
            font_size=24, align=PP_ALIGN.CENTER
        )
        h.add_code_block(
            slide, path,
            x + Inches(0.2), y + Inches(0.7), box_w - Inches(0.4), Inches(0.7),
            font_size=9
        )
        h.add_text_block(
            slide, f"範圍：{scope}",
            x + Inches(0.2), y + Inches(1.4), box_w - Inches(0.4), Inches(0.3),
            font_size=10, color=h.COLOR_DARK
        )
        h.add_text_block(
            slide, f"特性：{share}",
            x + Inches(0.2), y + Inches(1.7), box_w - Inches(0.4), Inches(0.3),
            font_size=10, color=h.COLOR_DARK
        )

    h.add_callout(
        slide, "Plugin hooks 自動跟著 plugin 啟用/停用，最容易管理",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 9：Matcher 模式
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Matcher 模式詳解",
        "控制 hook 何時觸發",
        slide_num=9, total=TOTAL, source="06 § Matcher"
    )

    h.add_text_block(
        slide, "🎯 三種匹配方式",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_comparison_table(
        slide,
        ["Matcher 值", "評估為", "範例"],
        [
            ["\"*\"、\"\" 或省略", "匹配所有", "任何工具"],
            ["Bash、Edit、Write", "精確字串", "只匹配該工具"],
            ["Edit|Write、Bash,Read", "清單（\\| 或 ,）", "匹配任一工具"],
            ["mcp__.*", "JavaScript regex", "所有 MCP 工具"],
            ["^Edit$", "錨定 regex", "只匹配 Edit（不匹配 NotebookEdit）"],
        ],
        Inches(0.5), Inches(2.3), Inches(12.333), Inches(3.0),
        font_size=12
    )

    h.add_two_column_compare(
        slide,
        "常見工具名 matchers",
        [
            "Bash — shell 命令",
            "Edit|Write — 檔案修改",
            "Read — 檔案讀取",
            "Agent — subagent 委派",
            "mcp__github__.* — GitHub 工具",
        ],
        "事件特定 matchers",
        [
            "SessionStart: startup|resume|clear",
            "Notification: permission_prompt|idle",
            "SubagentStart: Explore|Plan",
            "StopFailure: rate_limit|overloaded",
            "ConfigChange: user|project",
        ],
        top=Inches(5.5), height=Inches(1.5)
    )

    # ============================================================
    # Slide 10：常見 Matcher 範例
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "常見 Matcher 範例",
        "可直接複製貼上",
        slide_num=10, total=TOTAL, source="06 § Matcher 範例"
    )

    h.add_code_block(
        slide, """// Bash 工具的所有呼叫
"matcher": "Bash"

// 編輯或寫入檔案
"matcher": "Edit|Write"

// 特定的 MCP 工具
"matcher": "mcp__github__search_repositories"

// 特定 server 的所有工具
"matcher": "mcp__memory__.*"

// 寫入所有 mcp 工具（用 regex）
"matcher": "mcp__.*__write.*"

// 排除字首匹配（用 regex 錨定）
"matcher": "^Edit$"  // 只匹配 Edit，不匹配 NotebookEdit""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=12
    )

    h.add_callout(
        slide, "MCP server 名稱含連字號（如 brave-search）需用 mcp__brave-search__.* 形式",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4),
        icon="⚠️", font_size=12
    )

    h.add_text_block(
        slide, "Plugin 提供的 MCP server 用 plugin:: 範圍（如 plugin:my-plugin:db）",
        Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.4),
        font_size=11, color=h.COLOR_GRAY_TXT, italic=True
    )

    h.add_callout(
        slide, "Plugin 範圍名稱含 : → matcher 走 regex 路徑 → 用 ^ $ 錨定",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="🔧", font_size=12
    )

    # ============================================================
    # Slide 11：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 3", "5 種 Hook Handler 類型", "從簡單到強大")

    # ============================================================
    # Slide 12：5 種類型概覽
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook Handler 5 種類型",
        "選擇最適合的工具",
        slide_num=12, total=TOTAL, source="06 § Handler 類型"
    )

    types = [
        ("💻", "command", "執行 shell 命令", h.COLOR_BLUE, "預設選項", "lint、部署、阻擋"),
        ("🌐", "http", "POST 到 URL", h.COLOR_GREEN, "外部服務", "共用稽核、Slack 通知"),
        ("🔌", "mcp_tool", "呼叫 MCP 工具", h.COLOR_PRIMARY, "用 MCP 做事", "安全掃描、資料庫查詢"),
        ("🧠", "prompt", "LLM 評估（Haiku）", h.COLOR_RED, "需要判斷", "任務完成度檢查"),
        ("🤖", "agent", "Subagent 驗證器", RGBColor(0x7C, 0x3A, 0xED), "複雜驗證", "實驗性、檢查測試結果"),
    ]

    box_w = Inches(2.4)
    box_h = Inches(2.5)
    h_gap = Inches(0.15)
    total_w = box_w * 5 + h_gap * 4
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(1.8)

    for i, (icon, name, desc, color, label, use) in enumerate(types):
        x = start_x + i * (box_w + h_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        h.add_text_block(
            slide, icon,
            x, start_y + Inches(0.2), box_w, Inches(0.6),
            font_size=36, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, name,
            x, start_y + Inches(0.85), box_w, Inches(0.4),
            font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER, font="Consolas"
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), start_y + Inches(1.25), box_w - Inches(0.4), Inches(0.5),
            font_size=11, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, f"標籤：{label}",
            x + Inches(0.2), start_y + Inches(1.75), box_w - Inches(0.4), Inches(0.4),
            font_size=10, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, f"用途：{use}",
            x + Inches(0.2), start_y + Inches(2.1), box_w - Inches(0.4), Inches(0.3),
            font_size=9, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
        )

    h.add_text_block(
        slide, "預設逾時：command/http/mcp_tool = 10 分鐘 · prompt = 30 秒 · agent = 60 秒",
        Inches(0.5), Inches(4.7), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "⚠️ 並非所有事件都支援所有類型",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_RED, bold=True, align=PP_ALIGN.CENTER
    )

    h.add_comparison_table(
        slide,
        ["支援事件", "類型"],
        [
            ["所有 5 種工具事件（PreToolUse 等）", "command、http、mcp_tool、prompt、agent"],
            ["SessionStart、Setup", "僅 command、mcp_tool"],
            ["大多數其他事件", "command、http、mcp_tool（無 prompt/agent）"],
        ],
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(1.1),
        font_size=10
    )

    # ============================================================
    # Slide 13：Command Hook 詳細
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Command Hook 詳解",
        "最常用的類型",
        slide_num=13, total=TOTAL, source="06 § Command"
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Edit|Write",
        "hooks": [
          {
            "type": "command",
            "command": "/path/to/script.sh",
            "args": [],
            "timeout": 30,
            "if": "Edit(*.ts)"
          }
        ]
      }
    ]
  }
}""",
        Inches(0.5), Inches(1.7), Inches(7), Inches(3.0),
        font_size=12
    )

    h.add_text_block(
        slide, "🔧 Exec form vs Shell form",
        Inches(7.833), Inches(1.7), Inches(5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """// Exec form（有 args）
{
  "type": "command",
  "command": "node",
  "args": ["${CLAUDE_PLUGIN_ROOT}/script.js"]
}

// Shell form（無 args）
{
  "type": "command",
  "command": "node \"${CLAUDE_PLUGIN_ROOT}\"/script.js"
}""",
        Inches(7.833), Inches(2.1), Inches(5), Inches(2.6),
        font_size=10
    )

    h.add_text_block(
        slide, "⚠️ 引用路徑佔位符時用 exec form",
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_BLUE, bold=True
    )

    h.add_text_block(
        slide, "Windows 上的 .cmd/.bat shim 不是可執行檔 → 用 exec form 必須透過 node 呼叫腳本",
        Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.4),
        font_size=11, color=h.COLOR_GRAY_TXT
    )

    h.add_callout(
        slide, "Plugin 用 ${user_config.*} 必須用 exec form（shell form 會被拒絕）",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        icon="🔒", font_size=12
    )

    # ============================================================
    # Slide 14：HTTP / MCP / Prompt Hooks 簡介
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "HTTP、MCP Tool、Prompt Hooks",
        "另外 3 種類型快速了解",
        slide_num=14, total=TOTAL, source="06 § 其他類型"
    )

    h.add_text_block(
        slide, "🌐 HTTP Hook",
        Inches(0.5), Inches(1.7), Inches(4), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_GREEN
    )
    h.add_code_block(
        slide, """{
  "type": "http",
  "url": "http://localhost:8080/hooks",
  "headers": {
    "Authorization": "Bearer $MY_TOKEN"
  },
  "allowedEnvVars": ["MY_TOKEN"]
}

// POST JSON 輸入
// 回應 2xx + JSON 決策""",
        Inches(0.5), Inches(2.1), Inches(4), Inches(2.5),
        font_size=10
    )

    h.add_text_block(
        slide, "🔌 MCP Tool Hook",
        Inches(4.733), Inches(1.7), Inches(4), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )
    h.add_code_block(
        slide, """{
  "type": "mcp_tool",
  "server": "my_server",
  "tool": "security_scan",
  "input": {
    "file_path": "${tool_input.file_path}"
  }
}

// 呼叫已連接 MCP server 的工具""",
        Inches(4.733), Inches(2.1), Inches(4), Inches(2.5),
        font_size=10
    )

    h.add_text_block(
        slide, "🧠 Prompt Hook",
        Inches(8.967), Inches(1.7), Inches(4), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )
    h.add_code_block(
        slide, """{
  "type": "prompt",
  "prompt": "Evaluate if Claude should stop: $ARGUMENTS. Check if all tasks are complete.",
  "timeout": 30
}

// LLM 回應 {ok, reason, impossible}""",
        Inches(8.967), Inches(2.1), Inches(4), Inches(2.5),
        font_size=10
    )

    h.add_two_column_compare(
        slide,
        "Prompt 回應語意",
        [
            "ok: true → 允許",
            "ok: false + reason → 阻擋",
            "ok: false + impossible: true",
            "  → 條件不可能滿足，允許停止",
            "依事件而定的回應處理",
        ],
        "Agent Hook 差異",
        [
            "可使用工具（Read、Grep、Glob）",
            "多輪驗證，最多 50 輪",
            "預設 60 秒逾時",
            "實驗性，可能變動",
            "用於複雜驗證（檢查測試結果）",
        ],
        top=Inches(4.9), height=Inches(2.0)
    )

    # ============================================================
    # Slide 15：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 4", "Hook 輸入/輸出與退出碼", "資料格式與決策控制")

    # ============================================================
    # Slide 16：Hook 輸入
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 輸入：JSON 格式",
        "從 stdin 接收的資料",
        slide_num=16, total=TOTAL, source="06 § 輸入"
    )

    h.add_text_block(
        slide, "📥 範例：Bash 命令的 PreToolUse 輸入",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """{
  "session_id": "abc123",
  "prompt_id": "550e8400-e29b-41d4-a716-446655440000",
  "transcript_path": "/home/user/.claude/projects/.../transcript.jsonl",
  "cwd": "/home/user/my-project",
  "permission_mode": "default",
  "hook_event_name": "PreToolUse",
  "tool_name": "Bash",
  "tool_input": {
    "command": "npm test",
    "description": "Run test suite",
    "timeout": 120000
  },
  "tool_use_id": "toolu_01ABC123..."
}""",
        Inches(0.5), Inches(2.1), Inches(12.333), Inches(3.0),
        font_size=11
    )

    h.add_text_block(
        slide, "常見輸入欄位",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["欄位", "描述"],
        [
            ["session_id", "目前 session 識別碼"],
            ["cwd", "目前工作目錄"],
            ["permission_mode", "default/plan/acceptEdits/auto/dontAsk/bypassPermissions"],
            ["hook_event_name", "事件名稱"],
            ["tool_name + tool_input", "工具名稱與輸入（僅工具事件）"],
        ],
        Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.3),
        font_size=10
    )

    # ============================================================
    # Slide 17：退出碼語意
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "退出碼語意",
        "Exit code 決定 Claude 接下來做什麼",
        slide_num=17, total=TOTAL, source="06 § 退出碼"
    )

    h.add_comparison_table(
        slide,
        ["退出碼", "意義", "stdout 處理"],
        [
            ["0", "成功", "若有有效 JSON → 決策；純文字 → 寫入 debug log"],
            ["2", "阻塞錯誤", "stderr 作為原因，阻擋該動作（多數事件）"],
            ["其他", "取決於 stdout", "純文字 → 非阻塞錯誤；JSON → 套用 JSON 決策"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.0),
        font_size=12
    )

    h.add_code_block(
        slide, """#!/bin/bash
# 阻擋 rm 命令的 hook
input=$(cat)
command=$(echo "$input" | jq -r '.tool_input.command')

if [[ "$command" == rm* ]]; then
  echo "Blocked: rm commands are not allowed" >&2
  exit 2  # 阻塞錯誤
fi

exit 0  # 允許""",
        Inches(0.5), Inches(4.0), Inches(12.333), Inches(2.0),
        font_size=11
    )

    h.add_callout(
        slide, "exit 2 對大多數事件阻擋，但 PermissionRequest、PostToolUse 等不阻擋",
        Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.4),
        icon="⚠️", font_size=12
    )

    h.add_text_block(
        slide, "WorktreeCreate 例外：任何非零退出碼都會失敗",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        font_size=11, color=h.COLOR_RED, italic=True
    )

    # ============================================================
    # Slide 18：JSON 決策控制
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "JSON 決策控制",
        "比退出碼更精細的權限",
        slide_num=18, total=TOTAL, source="06 § JSON 輸出"
    )

    h.add_text_block(
        slide, "🎯 退出 0 + JSON = 結構化控制",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_code_block(
        slide, """{
  "continue": false,
  "stopReason": "Build failed, fix errors before continuing",
  "systemMessage": "⚠️ Lint errors detected",
  "hookSpecificOutput": {
    "hookEventName": "PreToolUse",
    "permissionDecision": "deny",
    "permissionDecisionReason": "Use rg instead of grep for better performance",
    "updatedInput": {
      "command": "rg pattern"
    },
    "additionalContext": "Current environment: production. Proceed with caution."
  }
}""",
        Inches(0.5), Inches(2.3), Inches(12.333), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "通用欄位 vs 事件特定欄位",
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "通用欄位",
        [
            "continue: 完全停止 Claude",
            "stopReason: 停止原因",
            "systemMessage: 警告訊息",
            "terminalSequence: 終端跳脫序列",
        ],
        "事件特定欄位",
        [
            "PreToolUse: permissionDecision",
            "PostToolUse: updatedToolOutput",
            "Stop: decision + reason",
            "UserPromptSubmit: additionalContext",
        ],
        top=Inches(5.4), height=Inches(1.6)
    )

    # ============================================================
    # Slide 19：決策優先級
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "決策優先級",
        "多個 hook 衝突時的規則",
        slide_num=19, total=TOTAL, source="06 § 決策"
    )

    h.add_text_block(
        slide, "🎯 PreToolUse 多個 hook 衝突時",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_RED,
        align=PP_ALIGN.CENTER
    )

    # 優先級橫向圖
    levels = [
        ("deny", "1. 最高", "阻擋", h.COLOR_RED),
        ("defer", "2. 次高", "延後", h.COLOR_PRIMARY),
        ("ask", "3. 確認", "詢問", h.COLOR_BLUE),
        ("allow", "4. 最低", "允許", h.COLOR_GREEN),
    ]

    box_w = Inches(2.8)
    box_h = Inches(2.0)
    h_gap = Inches(0.2)
    total_w = box_w * 4 + h_gap * 3
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.3)

    for i, (name, rank, action, color) in enumerate(levels):
        x = start_x + i * (box_w + h_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        h.add_text_block(
            slide, name,
            x, start_y + Inches(0.3), box_w, Inches(0.5),
            font_size=24, bold=True, color=color, align=PP_ALIGN.CENTER, font="Consolas"
        )
        h.add_text_block(
            slide, rank,
            x, start_y + Inches(0.9), box_w, Inches(0.4),
            font_size=12, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, action,
            x, start_y + Inches(1.4), box_w, Inches(0.4),
            font_size=14, bold=True, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )

        if i < len(levels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w + Inches(0.02), start_y + box_h/2 - Inches(0.15), h_gap - Inches(0.04), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = h.COLOR_GRAY_TXT
            arrow.line.fill.background()

    h.add_text_block(
        slide, "📌 其他事件（如 Stop）用 top-level decision + reason",
        Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
        font_size=13, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
    )

    h.add_code_block(
        slide, """{
  "decision": "block",
  "reason": "Tests must pass before stopping"
}""",
        Inches(2.0), Inches(5.4), Inches(9.333), Inches(1.3),
        font_size=12
    )

    # ============================================================
    # Slide 20：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 5", "重要事件詳解", "PreToolUse、PostToolUse、Stop")

    # ============================================================
    # Slide 21：PreToolUse
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "PreToolUse：工具執行前",
        "最常用的事件",
        slide_num=21, total=TOTAL, source="06 § PreToolUse"
    )

    h.add_text_block(
        slide, "🛡 用途：阻擋、修改、自動批准",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """{
  "hookSpecificOutput": {
    "hookEventName": "PreToolUse",
    "permissionDecision": "deny",
    "permissionDecisionReason": "Destructive command blocked"
  }
}""",
        Inches(0.5), Inches(2.2), Inches(6), Inches(1.8),
        font_size=11
    )

    h.add_code_block(
        slide, """{
  "hookSpecificOutput": {
    "hookEventName": "PreToolUse",
    "permissionDecision": "allow",
    "updatedInput": {
      "command": "rg pattern"
    }
  }
}""",
        Inches(6.833), Inches(2.2), Inches(6), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 注意：PreToolUse 只在 Claude 呼叫工具時執行",
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_BLUE
    )

    h.add_bullet_list(
        slide, [
            "檔案用 @ 引用 → 不觸發 PreToolUse（Read 也跳過）",
            "EndConversation → 不觸發",
            "若想阻擋 @ 引用：使用 Read deny 規則",
        ],
        Inches(0.5), Inches(4.9), Inches(12), Inches(1.5),
        font_size=11
    )

    h.add_callout(
        slide, "PreToolUse 不會被 hook 自動 deny 規則取代；deny 規則仍生效",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="🔒", font_size=12
    )

    # ============================================================
    # Slide 22：PostToolUse
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "PostToolUse：工具成功後",
        "自動 lint、格式化、通知",
        slide_num=22, total=TOTAL, source="06 § PostToolUse"
    )

    h.add_code_block(
        slide, """// 編輯後自動跑 prettier
{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Edit|Write",
        "hooks": [
          {
            "type": "command",
            "command": "jq -r '.tool_input.file_path' | xargs npx prettier --write"
          }
        ]
      }
    ]
  }
}""",
        Inches(0.5), Inches(1.7), Inches(7), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "📊 進階：替換工具輸出",
        Inches(7.833), Inches(1.7), Inches(5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_code_block(
        slide, """{
  "hookSpecificOutput": {
    "hookEventName": "PostToolUse",
    "additionalContext": "File is generated; edit src/schema.ts",
    "updatedToolOutput": {
      "stdout": "[redacted]",
      "stderr": "",
      "interrupted": false
    }
  }
}

// updatedToolOutput 必須匹配工具輸出形狀""",
        Inches(7.833), Inches(2.1), Inches(5), Inches(2.5),
        font_size=9
    )

    h.add_callout(
        slide, "⚠️ updatedToolOutput 只改 Claude 看到的內容；工具實際執行結果不變",
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.4),
        icon="⚠️", font_size=12
    )

    h.add_text_block(
        slide, "📋 寬鬆匹配技巧",
        Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "省略 matcher 或設 * → 任何工具後執行",
            "若想特定檔案變更時觸發 → 用 FileChanged 事件",
            "FileChanged 無決策控制，但有 watchPaths 動態管理",
        ],
        Inches(0.5), Inches(5.6), Inches(12), Inches(1.2),
        font_size=11
    )

    # ============================================================
    # Slide 23：Stop Hook
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Stop Hook：Claude 完成時",
        "確保任務完整結束",
        slide_num=23, total=TOTAL, source="06 § Stop"
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "Stop": [
      {
        "hooks": [
          {
            "type": "prompt",
            "prompt": "Check if all tasks are complete. If not, respond with {ok: false, reason: 'what remains'}."
          }
        ]
      }
    ]
  }
}

// 當 ok: false → reason 給 Claude 繼續
// 當 ok: true 或 ok: false + impossible: true → 允許停止""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "⚠️ Stop hook 連續 8 次阻擋會被覆寫",
        Inches(0.5), Inches(4.4), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_code_block(
        slide, """#!/bin/bash
# 避免無限迴圈
input=$(cat)
if [ "$(echo "$input" | jq -r '.stop_hook_active')" = "true" ]; then
  exit 0  # 已被觸發過，允許停止
fi
# ... 你的檢查邏輯""",
        Inches(0.5), Inches(4.8), Inches(12.333), Inches(1.5),
        font_size=10
    )

    h.add_callout(
        slide, "/goal 是 session-scoped Stop hook 的內建快捷方式",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="🎯", font_size=12
    )

    # ============================================================
    # Slide 24：SessionStart 與 SessionEnd
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "SessionStart 與 SessionEnd",
        "Session 生命週期",
        slide_num=24, total=TOTAL, source="06 § Session"
    )

    h.add_code_block(
        slide, """// 載入專案 context
{
  "hooks": {
    "SessionStart": [
      {
        "matcher": "compact",    // 只在 compaction 後
        "hooks": [
          {
            "type": "command",
            "command": "echo 'Reminder: use Bun, not npm. Run bun test before committing.'"
          }
        ]
      }
    ]
  }
}""",
        Inches(0.5), Inches(1.7), Inches(7), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "SessionStart matchers",
        Inches(7.833), Inches(1.7), Inches(5), Inches(0.4),
        font_size=13, bold=True, color=h.COLOR_BLUE
    )

    h.add_bullet_list(
        slide, [
            "startup — 新 session",
            "resume — /resume 或 --continue",
            "clear — /clear",
            "compact — auto/manual compaction",
            "fork — 從現有 session 分叉",
        ],
        Inches(7.833), Inches(2.1), Inches(5), Inches(2.5),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 SessionStart hooks 在每個 session 都執行 → 保持快速",
        Inches(0.5), Inches(4.4), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_BLUE, bold=True
    )

    h.add_code_block(
        slide, """// 持久化環境變數
#!/bin/bash
if [ -n "$CLAUDE_ENV_FILE" ]; then
  echo 'export NODE_ENV=production' >> "$CLAUDE_ENV_FILE"
  echo 'export PATH="$PATH:./node_modules/.bin"' >> "$CLAUDE_ENV_FILE"
fi""",
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(1.5),
        font_size=10
    )

    h.add_text_block(
        slide, "SessionEnd：清理任務、記錄 session 統計、儲存狀態",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_DARK
    )

    # ============================================================
    # Slide 25：Notification 事件
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Notification 事件：通知觸發",
        "桌面通知、Slack 整合",
        slide_num=25, total=TOTAL, source="06 § Notification"
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "Notification": [
      {
        "matcher": "permission_prompt",
        "hooks": [
          {
            "type": "command",
            "command": "osascript -e 'display notification \"Claude needs your attention\" with title \"Claude Code\"'"
          }
        ]
      }
    ]
  }
}""",
        Inches(0.5), Inches(1.7), Inches(7), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "跨平台桌面通知",
        Inches(7.833), Inches(1.7), Inches(5), Inches(0.4),
        font_size=13, bold=True, color=h.COLOR_BLUE
    )

    h.add_bullet_list(
        slide, [
            "macOS: osascript",
            "Linux: notify-send",
            "Windows: powershell MessageBox",
        ],
        Inches(7.833), Inches(2.1), Inches(5), Inches(1.5),
        font_size=11
    )

    h.add_text_block(
        slide, "Notification matchers",
        Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_comparison_table(
        slide,
        ["Matcher", "觸發時機"],
        [
            ["permission_prompt", "Claude 需要權限批准，等待約 6 秒"],
            ["idle_prompt", "Claude 完成回應約 60 秒，你未輸入"],
            ["auth_success", "認證完成"],
            ["elicitation_dialog", "MCP server 打開 elicitation 表格"],
            ["agent_needs_input", "背景 session 等待輸入"],
            ["agent_completed", "背景 session 完成或失敗"],
        ],
        Inches(0.5), Inches(4.4), Inches(12.333), Inches(2.5),
        font_size=11
    )

    # ============================================================
    # Slide 26：背景執行
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "背景執行 Hooks（async）",
        "不阻擋 Claude 的長時間任務",
        slide_num=26, total=TOTAL, source="06 § 背景執行"
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Write",
        "hooks": [
          {
            "type": "command",
            "command": "/path/to/run-tests.sh",
            "async": true,         ← 背景執行
            "timeout": 300         ← 最長 5 分鐘
          }
        ]
      }
    ]
  }
}""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.0),
        font_size=11
    )

    h.add_two_column_compare(
        slide,
        "✅ 背景 hook 可用於",
        [
            "跑測試套件",
            "部署到 staging",
            "上傳到 S3",
            "產生報告",
            "長時間 lint",
        ],
        "❌ 背景 hook 不能",
        [
            "阻擋 Claude 動作",
            "改變決策",
            "修改工具輸入/輸出",
            "被立刻讀取結果",
            "（下次 session 才知道）",
        ],
        top=Inches(4.0), height=Inches(2.5)
    )

    h.add_callout(
        slide, "async 結果在下個 session 開始時可見（--continue 時讀取）",
        Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
        icon="⏱", font_size=12
    )

    # ============================================================
    # Slide 27：為 Claude 新增上下文
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "為 Claude 新增上下文",
        "additionalContext 的使用",
        slide_num=27, total=TOTAL, source="06 § 上下文注入"
    )

    h.add_code_block(
        slide, """{
  "hookSpecificOutput": {
    "hookEventName": "PostToolUse",
    "additionalContext": "This file is generated. Edit src/schema.ts and run `bun generate` instead."
  }
}

// additionalContext 會包裝為系統提醒插入對話""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(1.7),
        font_size=11
    )

    h.add_text_block(
        slide, "📌 適合的內容",
        Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_GREEN
    )

    h.add_bullet_list(
        slide, [
            "環境狀態：目前 branch、deployment target",
            "條件式規則：哪些目錄是唯讀的",
            "外部資料：分配給你的 open issues、最近 CI 結果",
            "專案特定的測試指令",
        ],
        Inches(0.5), Inches(4.0), Inches(12), Inches(1.5),
        font_size=11
    )

    h.add_text_block(
        slide, "📝 寫成事實陳述而非系統指令",
        Inches(0.5), Inches(5.7), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "✅ 「部署目標是 production」",
            "❌ 「你必須部署到 production」",
            "事實讀起來像資訊；指令可能觸發 prompt-injection 防禦",
        ],
        Inches(0.5), Inches(6.1), Inches(12), Inches(0.9),
        font_size=11
    )

    # ============================================================
    # Slide 28：發送終端通知
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "發送終端通知",
        "桌面通知、視窗標題",
        slide_num=28, total=TOTAL, source="06 § 終端通知"
    )

    h.add_code_block(
        slide, """#!/bin/bash
# 從 Notification hook 觸發桌面通知
input=$(cat)
title="Claude Code"
body=$(jq -r '.message // "Needs your attention"' <<<"$input")
seq=$(printf '\\033]777;notify;%s;%s\\007' "$title" "$body")
jq -nc --arg seq "$seq" '{terminalSequence: $seq}'""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "🎯 terminalSequence 支援的跳脫序列",
        Inches(0.5), Inches(3.9), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "桌面通知",
        [
            "OSC 9：iTerm2、ConEmu、WezTerm",
            "OSC 99：Kitty 通知",
            "OSC 777：urxvt、Ghostty、Warp",
            "BEL：基本鈴聲",
        ],
        "視窗控制",
        [
            "OSC 0：視窗標題 + icon",
            "OSC 1：icon 標題",
            "OSC 2：視窗標題",
            "其他 OSC 會被忽略",
        ],
        top=Inches(4.3), height=Inches(2.0)
    )

    h.add_callout(
        slide, "terminalSequence 在非互動模式（-p）和 Agent SDK 中會被忽略",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    # ============================================================
    # Slide 29：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 6", "安全性與最佳實踐", "避免踩坑")

    # ============================================================
    # Slide 30：安全性警告
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "安全性：認真對待 Hooks",
        "⚠️ Hooks 以你的完整權限執行",
        slide_num=30, total=TOTAL, source="06 § 安全性"
    )

    h.add_callout(
        slide,
        "Command hooks 以你的完整使用者權限執行 shell 命令。它們可以修改、刪除或訪問你的使用者帳戶可訪問的任何檔案。",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(1.0),
        icon="⚠️", font_size=14
    )

    h.add_text_block(
        slide, "🔒 工作區信任機制",
        Inches(0.5), Inches(2.9), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "互動式 session：每個 hooks 都需工作區信任",
            "-p / SDK session：永遠不顯示對話框，直接執行（⚠️ 危險）",
            "外部 repo 用 --bare 或 --settings '{\"disableAllHooks\": true}'",
        ],
        Inches(0.5), Inches(3.3), Inches(12), Inches(1.8),
        font_size=12
    )

    h.add_text_block(
        slide, "🛡 安全最佳實踐",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_GREEN
    )

    h.add_two_column_compare(
        slide,
        "✅ 應該做",
        [
            "驗證所有輸入",
            "永遠用雙引號包變數",
            "檢查路徑遍歷（..）",
            "使用絕對路徑",
            "跳過敏感檔案（.env、.git/）",
        ],
        "❌ 不該做",
        [
            "盲目信任輸入",
            "不引號包 shell 變數",
            "從未驗證的 repo 跑 hooks",
            "在 plugin 內用 ../ 共享檔案",
        ],
        top=Inches(5.7), height=Inches(1.3)
    )

    # ============================================================
    # Slide 31：除錯技巧
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 除錯技巧",
        "找不到問題時的標準流程",
        slide_num=31, total=TOTAL, source="06 § 除錯"
    )

    h.add_text_block(
        slide, "🔍 標準除錯流程",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_bullet_list(
        slide, [
            "1. /hooks 確認 hook 已註冊",
            "2. claude --debug 啟動看完整日誌",
            "3. 手動測試：echo '{...}' | ./hook.sh",
            "4. 確認 chmod +x（最重要的常見錯誤！）",
            "5. 確認 jq 已安裝",
            "6. 看退出碼：echo $?",
        ],
        Inches(0.7), Inches(2.2), Inches(12), Inches(2.5),
        font_size=12
    )

    h.add_code_block(
        slide, """# 完整除錯命令
claude --debug-file /tmp/claude.log

# 在另一個 terminal
tail -f /tmp/claude.log""",
        Inches(0.5), Inches(4.9), Inches(12.333), Inches(1.2),
        font_size=12
    )

    h.add_text_block(
        slide, "📌 4 個最常見錯誤",
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "❌ matcher 拼錯 → 沉默失敗（用 /hooks 確認）",
            "❌ chmod 沒給 → 退出 126（檢查 echo $?）",
            "❌ jq 沒裝 → parse 錯誤（apt install jq）",
            "❌ hook JSON 沒生效 → 檢查 shell profile 沒有 echo 干擾",
        ],
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.9),
        font_size=10
    )

    # ============================================================
    # Slide 32：PowerShell on Windows
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "PowerShell on Windows",
        "在 Windows 上執行 hooks",
        slide_num=32, total=TOTAL, source="06 § PowerShell"
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Write",
        "hooks": [
          {
            "type": "command",
            "shell": "powershell",
            "command": "Write-Host 'File written'"
          }
        ]
      }
    ]
  }
}""",
        Inches(0.5), Inches(1.7), Inches(7), Inches(2.0),
        font_size=11
    )

    h.add_text_block(
        slide, "💡 PowerShell 注意事項",
        Inches(7.833), Inches(1.7), Inches(5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_bullet_list(
        slide, [
            "用 $env:CLAUDE_PROJECT_DIR 讀環境變數",
            "不要寫裸 $CLAUDE_PROJECT_DIR",
            "（會被解析為 null）",
            "Claude Code 自動重寫為 ${env:NAME}",
            "Windows .cmd/.bat shim 必須 shell form",
        ],
        Inches(7.833), Inches(2.1), Inches(5), Inches(2.0),
        font_size=11
    )

    h.add_code_block(
        slide, """// 安全寫法
{
  "type": "command",
  "shell": "powershell",
  "command": "& \"$env:CLAUDE_PROJECT_DIR\\.claude\\hooks\\check.ps1\""
}""",
        Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.5),
        font_size=11
    )

    h.add_callout(
        slide, "跨平台 plugin 建議：用 exec form 搭配 node 處理差異",
        Inches(0.5), Inches(5.7), Inches(12.333), Inches(0.4),
        icon="🌐", font_size=12
    )

    # ============================================================
    # Slide 33：常見 Hook 模式
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "5 個常見的 Hook 模式",
        "複製貼上就能用",
        slide_num=33, total=TOTAL, source="06 § 常用模式"
    )

    patterns = [
        ("1. 編輯後自動格式化", h.COLOR_BLUE,
         """{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Edit|Write",
        "hooks": [
          { "type": "command",
            "command": "jq -r '.tool_input.file_path' | xargs npx prettier --write" }
        ]
      }
    ]
  }
}"""),
        ("2. 阻擋危險命令", h.COLOR_RED,
         """{
  "hooks": {
    "PreToolUse": [
      {
        "matcher": "Bash",
        "hooks": [
          { "type": "command",
            "command": "./block-rm-rf.sh" }
        ]
      }
    ]
  }
}"""),
        ("3. 桌面通知", h.COLOR_GREEN,
         """{
  "hooks": {
    "Notification": [
      {
        "matcher": "",
        "hooks": [
          { "type": "command",
            "command": "osascript -e 'display notification ...'" }
        ]
      }
    ]
  }
}"""),
        ("4. 阻擋編輯敏感檔案", h.COLOR_PRIMARY,
         """{
  "hooks": {
    "PreToolUse": [
      {
        "matcher": "Edit|Write",
        "hooks": [
          { "type": "command",
            "command": "block-secrets.sh" }
        ]
      }
    ]
  }
}"""),
        ("5. Session 開始時載入 context", RGBColor(0x7C, 0x3A, 0xED),
         """{
  "hooks": {
    "SessionStart": [
      {
        "matcher": "startup",
        "hooks": [
          { "type": "command",
            "command": "echo 'Project: $(basename $PWD)'" }
        ]
      }
    ]
  }
}""")
    ]

    for i, (label, color, code) in enumerate(patterns):
        y = Inches(1.7) + i * Inches(1.0)

        # 標題列
        title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(3.5), Inches(0.85))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = color
        title_box.line.fill.background()
        h.add_text_block(
            slide, label,
            Inches(0.5), y, Inches(3.5), Inches(0.85),
            font_size=12, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 程式碼
        code_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), y, Inches(8.833), Inches(0.85))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = h.COLOR_CODE_BG
        code_box.line.fill.background()
        h.add_text_block(
            slide, code.replace("\n", " · ")[:100] + "...",
            Inches(4.1), y, Inches(8.7), Inches(0.85),
            font_size=8, color=h.COLOR_CODE_FG, font="Consolas"
        )

    h.add_callout(
        slide, "完整可運行的 hook 範例見 06-hooks.md § 背景執行",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="📖", font_size=12
    )

    # ============================================================
    # Slide 34：實戰：阻擋 rm -rf
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "實戰：阻擋 rm -rf 完整實作",
        "10 分鐘學會 hook 整合",
        slide_num=34, total=TOTAL, source="綜合實戰"
    )

    h.add_text_block(
        slide, "🎯 目標：阻擋危險的 rm -rf 命令",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_RED
    )

    h.add_text_block(
        slide, "Step 1：建立 hook 腳本",
        Inches(0.5), Inches(2.2), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """#!/bin/bash
# ~/.claude/hooks/block-rm-rf.sh

input=$(cat)
COMMAND=$(echo "$input" | jq -r '.tool_input.command // empty')

# 阻擋 rm -rf
if echo "$COMMAND" | grep -qE 'rm\\s+(-[a-zA-Z]*r[a-zA-Z]*f|-[a-zA-Z]*f[a-zA-Z]*r|-rf|-fr)'; then
  echo "Blocked: rm -rf is not allowed" >&2
  exit 2
fi

exit 0""",
        Inches(0.5), Inches(2.6), Inches(6), Inches(2.8),
        font_size=10
    )

    h.add_text_block(
        slide, "Step 2：設定 hook",
        Inches(6.833), Inches(2.2), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """// ~/.claude/settings.json
{
  "hooks": {
    "PreToolUse": [
      {
        "matcher": "Bash",
        "hooks": [
          {
            "type": "command",
            "command": "~/.claude/hooks/block-rm-rf.sh"
          }
        ]
      }
    ]
  }
}""",
        Inches(6.833), Inches(2.6), Inches(6), Inches(2.0),
        font_size=10
    )

    h.add_text_block(
        slide, "Step 3：賦予執行權限並測試",
        Inches(6.833), Inches(4.8), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """$ chmod +x ~/.claude/hooks/block-rm-rf.sh

# 測試
$ echo '{"tool_name":"Bash","tool_input":{"command":"rm -rf /"}}' \\
  | ~/.claude/hooks/block-rm-rf.sh
Blocked: rm -rf is not allowed
$ echo $?
2""",
        Inches(6.833), Inches(5.2), Inches(6), Inches(1.8),
        font_size=10
    )

    h.add_callout(
        slide, "完成！現在 Claude 嘗試 rm -rf 會被阻擋，stderr 訊息會給 Claude 看到",
        Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4),
        icon="🛡", font_size=13
    )

    # ============================================================
    # Slide 35：實戰：編輯後自動跑測試
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "實戰：編輯後背景跑測試",
        "async + JSON 回報",
        slide_num=35, total=TOTAL, source="綜合實戰"
    )

    h.add_text_block(
        slide, "🎯 目標：編輯 .ts/.js 後背景跑測試，結果回報給 Claude",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_GREEN
    )

    h.add_code_block(
        slide, """#!/bin/bash
# ~/.claude/hooks/run-tests-async.sh

input=$(cat)
FILE_PATH=$(echo "$input" | jq -r '.tool_input.file_path // empty')

# 只對原始碼跑
if [[ "$FILE_PATH" != *.ts && "$FILE_PATH" != *.js ]]; then
  exit 0
fi

# 跑測試
RESULT=$(npm test 2>&1)
EXIT_CODE=$?

if [ $EXIT_CODE -eq 0 ]; then
  MSG="Tests passed after editing $FILE_PATH"
else
  MSG="Tests failed after editing $FILE_PATH: $RESULT"
fi

# 透過 additionalContext 把結果給 Claude
jq -nc --arg msg "$MSG" '{
  hookSpecificOutput: {
    hookEventName: "PostToolUse",
    additionalContext: $msg
  }
}'""",
        Inches(0.5), Inches(2.2), Inches(7), Inches(4.0),
        font_size=10
    )

    h.add_text_block(
        slide, "設定",
        Inches(7.833), Inches(2.2), Inches(5), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Write|Edit",
        "hooks": [
          {
            "type": "command",
            "command": "~/.claude/hooks/run-tests-async.sh",
            "async": true,
            "timeout": 300
          }
        ]
      }
    ]
  }
}

// async: Claude 不等結果
// timeout: 5 分鐘
// additionalContext: 結果在下輪回報給 Claude""",
        Inches(7.833), Inches(2.6), Inches(5), Inches(3.5),
        font_size=10
    )

    h.add_callout(
        slide, "完整背景 hooks 機制：06-hooks.md § 背景執行",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        icon="📖", font_size=12
    )

    # ============================================================
    # Slide 36：Hooks 在 Skill/Subagent 中
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hooks 在 Skill 和 Subagent 中",
        "frontmatter 限定範圍的 hooks",
        slide_num=36, total=TOTAL, source="06 § Skills/Agents"
    )

    h.add_code_block(
        slide, """# Skill 中的 hook
---
name: secure-operations
description: Perform operations with security checks
hooks:
  PreToolUse:
    - matcher: "Bash"
      hooks:
        - type: command
          command: "./scripts/security-check.sh"
---

# Subagent 中的 hook
---
name: code-reviewer
description: Review code changes with automatic linting
hooks:
  PreToolUse:
    - matcher: "Bash"
      hooks:
        - type: command
          command: "./scripts/validate-command.sh"
  PostToolUse:
    - matcher: "Edit|Write"
      hooks:
        - type: command
          command: "./scripts/run-linter.sh"
---

// 當 agent 作為 subagent 觸發，Stop 自動轉為 SubagentStop""",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=10
    )

    h.add_text_block(
        slide, "📌 Hook 生命週期差異",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_PRIMARY
    )

    h.add_two_column_compare(
        slide,
        "Skill hooks",
        [
            "Skill 叫用時註冊",
            "整個 session 持續執行",
            "once: true 可限制為一次",
        ],
        "Subagent hooks",
        [
            "Subagent 啟動時註冊",
            "Subagent 完成時移除",
            "Stop 自動轉 SubagentStop",
        ],
        top=Inches(5.8), height=Inches(1.2)
    )

    # ============================================================
    # Slide 37：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 7", "常見問題與疑難排解", "實戰中會遇到的狀況")

    # ============================================================
    # Slide 38：疑難排解速查
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "疑難排解速查表",
        "症狀 → 原因 → 解決",
        slide_num=38, total=TOTAL, source="06 § 疑難排解"
    )

    issues = [
        ("🚫 Hook 未觸發", "matcher 拼錯、事件名錯、permission 未批准",
         "用 /hooks 確認；/debug 看日誌；檢查大小寫"),
        ("🚫 Hook 報錯", "腳本退出非零、jq 找不到、權限不足",
         "echo {...} | ./hook.sh 手動測試；echo $?"),
        ("🚫 /hooks 沒顯示", "JSON 格式錯、檔案權限、檔案監視器漏掉",
         "重啟 session；驗證 JSON 語法；檢查 settings 路徑"),
        ("🚫 Stop hook 卡住", "連續 8 次阻擋被覆寫；stop_hook_active 沒檢查",
         "在腳本中檢查 stop_hook_active"),
        ("🚫 JSON 沒生效", "shell profile 有 echo 干擾；解析失敗",
         "if [[ $- == *i* ]] 包裝 echo；檢查 stdout 第一字元"),
    ]

    for i, (symptom, cause, solution) in enumerate(issues):
        y = Inches(1.7) + i * Inches(1.0)

        # 症狀
        sym_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(2.8), Inches(0.85))
        sym_box.fill.solid()
        sym_box.fill.fore_color.rgb = h.COLOR_RED
        sym_box.line.fill.background()
        h.add_text_block(
            slide, symptom,
            Inches(0.5), y, Inches(2.8), Inches(0.85),
            font_size=12, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        # 原因 + 解決
        info_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.3), y, Inches(9.533), Inches(0.85))
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        info_box.line.color.rgb = h.COLOR_GRAY_TXT
        info_box.line.width = Pt(0.5)

        h.add_text_block(
            slide, f"原因：{cause}",
            Inches(3.5), y + Inches(0.05), Inches(9.3), Inches(0.35),
            font_size=11, color=h.COLOR_DARK
        )
        h.add_text_block(
            slide, f"解法：{solution}",
            Inches(3.5), y + Inches(0.45), Inches(9.3), Inches(0.4),
            font_size=11, color=h.COLOR_GREEN
        )

    h.add_callout(
        slide, "Ctrl+O 打開 transcript 看 hook 執行的詳細結果",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="🔍", font_size=13
    )

    # ============================================================
    # Slide 39：除錯工作流
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "除錯工作流",
        "6 步找到問題根源",
        slide_num=39, total=TOTAL, source="06 § Debug 流程"
    )

    steps = [
        ("1", "確認 hook 註冊", "/hooks", "沒看到 → 檢查 JSON 格式"),
        ("2", "確認事件觸發", "/debug 啟動", "看日誌哪個 hook 跑了"),
        ("3", "手動測試腳本", "echo JSON | ./hook.sh", "模擬 stdin 輸入"),
        ("4", "檢查退出碼", "echo $?", "0=成功、2=阻塞、其他=非阻塞錯誤"),
        ("5", "檢查 stdout 第一字元", "第一字元決定 JSON 或純文字", "若是 JSON 但解析失敗 → schema 問題"),
        ("6", "檢查 stderr", "若非 0，看 stderr", "exit 2 → 訊息給 Claude"),
    ]

    step_w = Inches(3.8)
    step_h = Inches(0.9)
    h_gap = Inches(0.2)
    grid_cols = 2
    grid_rows = 3
    grid_w = step_w * grid_cols + h_gap
    grid_h = step_h * grid_rows + Inches(0.4)
    start_x = (h.SLIDE_W - grid_w) / 2
    start_y = Inches(1.7)

    for i, (num, title, cmd, note) in enumerate(steps):
        row = i // grid_cols
        col = i % grid_cols
        x = start_x + col * (step_w + h_gap)
        y = start_y + row * (step_h + Inches(0.2))

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = h.COLOR_RED
        card.line.width = Pt(1.5)

        # 編號
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.15), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = h.COLOR_RED
        circle.line.fill.background()
        h.add_text_block(
            slide, num,
            x + Inches(0.1), y + Inches(0.15), Inches(0.6), Inches(0.6),
            font_size=20, bold=True, color=h.COLOR_WHITE, align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, title,
            x + Inches(0.8), y + Inches(0.05), step_w - Inches(0.9), Inches(0.4),
            font_size=14, bold=True, color=h.COLOR_DARK
        )
        h.add_code_block(
            slide, cmd,
            x + Inches(0.8), y + Inches(0.4), step_w - Inches(0.9), Inches(0.3),
            font_size=9
        )

    h.add_callout(
        slide, "99% 的 hook 問題都是：chmod 沒給、jq 沒裝、JSON 格式錯",
        Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4),
        icon="🎯", font_size=14
    )

    h.add_text_block(
        slide, "完整 6 步流程見 06-hooks.md § Debug 技巧",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "🔧 進階：CLAUDE_CODE_DEBUG_LOG_LEVEL=verbose 取得更詳細的 matcher 記錄",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GREEN
    )

    # ============================================================
    # Slide 40：退出碼速查
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "退出碼 + JSON 速查表",
        "實戰中隨手查",
        slide_num=40, total=TOTAL, source="06 § 速查"
    )

    h.add_comparison_table(
        slide,
        ["退出碼", "stdout 狀態", "效果"],
        [
            ["0", "純文字", "寫入 debug log；多數事件無效果"],
            ["0", "有效 JSON", "JSON 決策生效（permissionDecision 等）"],
            ["0", "無效 JSON", "非阻塞錯誤；hook error 通知"],
            ["2", "任何", "阻塞錯誤（多數事件）；stderr 給 Claude"],
            ["其他", "純文字", "非阻塞錯誤；hook error 通知"],
            ["其他", "有效 JSON", "忽略退出碼，JSON 決策生效"],
            ["其他", "無效 JSON", "非阻塞錯誤"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=11
    )

    h.add_text_block(
        slide, "⚠️ 重要例外",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_RED
    )

    h.add_bullet_list(
        slide, [
            "WorktreeCreate：任何非零退出碼都會失敗",
            "PermissionRequest：exit 2 不被尊重",
            "PostToolUse、PostToolUseFailure、Notification 等：exit 2 不阻塞",
            "PreToolUse、UserPromptSubmit、Stop：exit 2 阻塞",
        ],
        Inches(0.5), Inches(5.8), Inches(12), Inches(1.2),
        font_size=11
    )

    # ============================================================
    # Slide 41：章節分隔
    # ============================================================
    h.add_section_divider(prs, "Part 8", "練習題與實戰", "鞏固所學")

    # ============================================================
    # Slide 42：練習題
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "練習題：5 個實作挑戰",
        "由淺入深",
        slide_num=42, total=TOTAL, source="練習"
    )

    challenges = [
        ("🟢", "桌面通知", "Notification hook 用 osascript", "10 分鐘"),
        ("🟢", "阻擋 rm -rf", "PreToolUse + exit 2", "20 分鐘"),
        ("🟡", "自動跑 ESLint", "PostToolUse 編輯後跑 lint", "30 分鐘"),
        ("🟡", "背景跑測試", "async + additionalContext 回報", "40 分鐘"),
        ("🔴", "完整 plugin 整合", "hooks + skill + agent 組合", "1 小時+"),
    ]

    for i, (level, title, desc, time) in enumerate(challenges):
        y = Inches(1.8) + i * Inches(1.0)

        level_color = h.COLOR_GREEN if "🟢" in level else (h.COLOR_PRIMARY if "🟡" in level else h.COLOR_RED)
        level_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(1.2), Inches(0.85))
        level_box.fill.solid()
        level_box.fill.fore_color.rgb = level_color
        level_box.line.fill.background()
        h.add_text_block(
            slide, level,
            Inches(0.5), y, Inches(1.2), Inches(0.85),
            font_size=14, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, title,
            Inches(1.8), y + Inches(0.05), Inches(6.5), Inches(0.4),
            font_size=15, bold=True, color=h.COLOR_DARK
        )
        h.add_text_block(
            slide, desc,
            Inches(1.8), y + Inches(0.45), Inches(6.5), Inches(0.4),
            font_size=11, color=h.COLOR_GRAY_TXT
        )

        time_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y, Inches(4.333), Inches(0.85))
        time_box.fill.solid()
        time_box.fill.fore_color.rgb = h.COLOR_BG_GRAY
        time_box.line.color.rgb = h.COLOR_RED
        time_box.line.width = Pt(1.5)
        h.add_text_block(
            slide, f"⏱ {time}",
            Inches(8.5), y, Inches(4.333), Inches(0.85),
            font_size=15, bold=True, color=h.COLOR_RED,
            align=PP_ALIGN.CENTER
        )

    h.add_callout(
        slide, "完成後，你已掌握事件驅動自動化的核心能力",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="🎓", font_size=14
    )

    # ============================================================
    # Slide 43：Hooks + 其他元件
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hooks 與其他元件的組合",
        "建立完整自動化流程",
        slide_num=43, total=TOTAL, source="綜合應用"
    )

    h.add_comparison_table(
        slide,
        ["組合模式", "效果", "範例"],
        [
            ["Hook + Skill", "Skill 觸發後，Hook 自動跑後續", "Skill 跑測試，Hook 自動通知 Slack"],
            ["Hook + Subagent", "Hook 觸發 subagent 處理複雜任務", "PostToolUse 觸發 code-reviewer subagent"],
            ["Hook + MCP", "Hook 用 MCP 工具做外部操作", "PreToolUse 用 Slack MCP 通知團隊"],
            ["Skill + Subagent", "Skill 在 subagent 中執行（context: fork）", "deep-research skill 用 Explore agent"],
            ["完整整合", "Plugin 包含所有元件", "code-review plugin：skill + agent + hook + MCP"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(4.0),
        font_size=12
    )

    h.add_callout(
        slide, "完整 plugin = skills + agents + hooks + MCP servers + LSP servers + monitors + themes",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        icon="🧩", font_size=14
    )

    h.add_text_block(
        slide, "進階 plugin 設計：見 02-plugins.md § 開發更複雜的 Plugins",
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
    )

    # ============================================================
    # Slide 44：Hook 設計模式總結
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 設計模式總結",
        "什麼情況用什麼 hook",
        slide_num=44, total=TOTAL, source="06 § 設計模式"
    )

    h.add_two_column_compare(
        slide,
        "🔧 確定性動作 → Command",
        [
            "Lint 自動跑",
            "格式化自動套用",
            "阻擋危險命令",
            "桌面通知",
            "部署、構建",
            "檔案備份",
        ],
        "🧠 需要判斷 → Prompt/Agent",
        [
            "任務完成度檢查",
            "程式碼品質審查",
            "安全性評估",
            "「是否違反原則」",
            "需要讀檔才能判斷",
            "需要 LLM 推理",
        ]
    )

    h.add_callout(
        slide, "不確定時：先試 command（簡單、確定），不夠用再加 prompt/agent",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="💡", font_size=13
    )

    # ============================================================
    # Slide 45：Hook 對決策控制
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 對決策的精細控制",
        "不只是阻擋",
        slide_num=45, total=TOTAL, source="06 § 決策控制"
    )

    h_add_summary_slide_skipped = True  # placeholder
    h.add_text_block(
        slide, "🎯 Hook 不只能「yes/no」，還能「改寫」",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=18, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "❶ 改寫工具輸入（updatedInput）",
        Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """// grep → rg 自動改寫
{
  "hookSpecificOutput": {
    "hookEventName": "PreToolUse",
    "permissionDecision": "allow",
    "updatedInput": {
      "command": "rg pattern"
    }
  }
}""",
        Inches(0.5), Inches(2.7), Inches(6), Inches(1.8),
        font_size=10
    )

    h.add_text_block(
        slide, "❷ 改寫工具輸出（updatedToolOutput）",
        Inches(6.833), Inches(2.3), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """// 遮罩敏感輸出
{
  "hookSpecificOutput": {
    "hookEventName": "PostToolUse",
    "updatedToolOutput": {
      "stdout": "[redacted]",
      "stderr": "",
      "interrupted": false
    }
  }
}""",
        Inches(6.833), Inches(2.7), Inches(6), Inches(2.0),
        font_size=10
    )

    h.add_text_block(
        slide, "❸ 為 Claude 新增 context",
        Inches(0.5), Inches(4.9), Inches(6), Inches(0.4),
        font_size=14, bold=True, color=h.COLOR_BLUE
    )

    h.add_code_block(
        slide, """{
  "hookSpecificOutput": {
    "hookEventName": "UserPromptSubmit",
    "additionalContext": "Current branch: main. Run tests before committing."
  }
}""",
        Inches(0.5), Inches(5.3), Inches(6), Inches(1.5),
        font_size=10
    )

    h.add_callout(
        slide, "PreToolUse 改輸入、PostToolUse 改輸出、UserPromptSubmit 加 context、Stop 阻擋或加反饋",
        Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4),
        icon="🎨", font_size=12
    )

    # ============================================================
    # Slide 46：Hook 速查表（完整）
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "完整 Hook 速查表",
        "實戰中隨時參考",
        slide_num=46, total=TOTAL, source="06 § 速查表"
    )

    h.add_comparison_table(
        slide,
        ["事件", "可阻擋？", "通用 JSON 欄位", "特定欄位"],
        [
            ["PreToolUse", "✅", "continue、systemMessage", "permissionDecision、updatedInput"],
            ["PostToolUse", "❌", "continue、systemMessage", "updatedToolOutput、additionalContext"],
            ["Stop", "✅", "continue、systemMessage", "decision: block + reason"],
            ["UserPromptSubmit", "✅", "continue、systemMessage", "additionalContext"],
            ["SessionStart", "❌", "systemMessage", "additionalContext、initialUserMessage"],
            ["Notification", "❌", "terminalSequence", "（無特定欄位）"],
        ],
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(3.5),
        font_size=11
    )

    h.add_text_block(
        slide, "📚 完整 28 個事件見 06-hooks.md § 完整事件表",
        Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, italic=True, align=PP_ALIGN.CENTER
    )

    h.add_callout(
        slide, "通用欄位：continue、stopReason、systemMessage、terminalSequence",
        Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.4),
        icon="💡", font_size=12
    )

    h.add_text_block(
        slide, "所有 hooks 平行執行；多個結果合併：最嚴格的決策勝出",
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER
    )

    # ============================================================
    # Slide 47：Hook 最佳實踐
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "Hook 開發最佳實踐",
        "從一開始就做對",
        slide_num=47, total=TOTAL, source="06 § 最佳實踐"
    )

    h.add_two_column_compare(
        slide,
        "✅ 應該做",
        [
            "保持 hook 快速（10 分鐘預設）",
            "輸入驗證，不盲目信任",
            "用 exec form 引用路徑",
            "為錯誤日誌寫到 stderr",
            "為阻塞錯誤用 exit 2",
            "UserPromptSubmit 用 30s timeout",
            "所有 hook 寫到 plugin 的 hooks/",
        ],
        "❌ 不該做",
        [
            "hook 太慢（阻擋 Claude）",
            "未驗證的輸入（安全風險）",
            "用 shell form 引用路徑",
            "忘記 chmod +x",
            "在 hook 中執行用戶提供的代碼",
            "覆寫 manage settings 控制的 hook",
            "把 hooks 放在 .claude-plugin/",
        ]
    )

    h.add_callout(
        slide, "口訣：Hook 越簡單越好；能用 command 解決就不要 prompt/agent",
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4),
        icon="💎", font_size=14
    )

    # ============================================================
    # Slide 48：重點回顧
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_summary_slide(
        slide,
        title="重點回顧",
        key_points=[
            "Hook = 確定性觸發的腳本/HTTP/LLM（與 Skills/Subagents 互補）",
            "28 個事件分屬 3 個節奏：Session、Per-turn、Per-tool-call",
            "5 種 handler 類型：command、http、mcp_tool、prompt、agent",
            "退出碼 0/2 決定阻擋與否；JSON 提供更精細的決策控制",
            "決策優先級：deny > defer > ask > allow",
            "always-on、必 chmod、jq、路徑佔位符等注意事項",
        ],
        next_steps=[
            "🎯 立即：寫一個 Notification hook 桌面通知",
            "📚 進階：寫 PreToolUse hook 阻擋危險命令",
            "🛠 整合：把常用 hooks 打包到 plugin",
            "🚀 自動化：建立團隊的開發流程 hooks",
        ],
        source="06-hooks.md"
    )

    # ============================================================
    # Slide 49：系列回顧
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide)
    h.add_title_bar(
        slide, "系列回顧：你已完成的學習",
        "4 份簡報的累積",
        slide_num=49, total=TOTAL, source="00 § 全系列"
    )

    h.add_text_block(
        slide, "🎉 你已掌握 Claude Code Plugin 完整知識",
        Inches(0.5), Inches(1.7), Inches(12.333), Inches(0.4),
        font_size=20, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    ppts = [
        ("00", "系列總覽", "生態系全景、學習路徑", "30 張", h.COLOR_PRIMARY, "✅ 已完成"),
        ("02", "Plugin 開發", "5 步建立第一個 plugin", "25 張", h.COLOR_BLUE, "✅ 已完成"),
        ("04", "Skills", "Skills 設計、撰寫、評估", "40 張", h.COLOR_GREEN, "✅ 已完成"),
        ("06", "Hooks", "事件驅動自動化", "50 張", h.COLOR_RED, "✅ 已完成"),
    ]

    box_w = Inches(2.9)
    box_h = Inches(2.8)
    h_gap = Inches(0.25)
    total_w = box_w * 4 + h_gap * 3
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(2.4)

    for i, (num, title, desc, pages, color, status) in enumerate(ppts):
        x = start_x + i * (box_w + h_gap)
        y = start_y

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_BG_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.6))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()
        h.add_text_block(
            slide, f"#{num} {title}",
            x, y, box_w, Inches(0.6),
            font_size=14, bold=True, color=h.COLOR_WHITE,
            align=PP_ALIGN.CENTER
        )

        h.add_text_block(
            slide, f"{pages} | 簡報 #{num}",
            x + Inches(0.2), y + Inches(0.8), box_w - Inches(0.4), Inches(0.4),
            font_size=12, color=h.COLOR_DARK, font="Consolas",
            align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, desc,
            x + Inches(0.2), y + Inches(1.3), box_w - Inches(0.4), Inches(0.8),
            font_size=11, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, status,
            x + Inches(0.2), y + Inches(2.3), box_w - Inches(0.4), Inches(0.4),
            font_size=14, bold=True, color=h.COLOR_GREEN, align=PP_ALIGN.CENTER
        )

    h.add_text_block(
        slide, "🏆 總計 145 張投影片，覆蓋 4 個核心主題",
        Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.4),
        font_size=16, bold=True, color=h.COLOR_PRIMARY,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "未涵蓋（依需求再做）：#01 Marketplaces、#03 技術參考、#05 Subagents、#07 探索",
        Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.4),
        font_size=12, color=h.COLOR_GRAY_TXT, align=PP_ALIGN.CENTER, italic=True
    )

    h.add_callout(
        slide, "下一步：應用所學，建立你自己的 production plugin",
        Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
        icon="🚀", font_size=14
    )

    # ============================================================
    # Slide 50：結束頁
    # ============================================================
    slide = h.add_blank_slide(prs)
    h.set_slide_bg(slide, h.COLOR_BG_CREAM)

    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.0),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Hooks 大師之路 🎣"
    run.font.name = h.FONT_TITLE
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = h.COLOR_RED

    h.add_text_block(
        slide, "從一個簡單的 hook 開始，建立自動化流程",
        Inches(1), Inches(3.7), Inches(11.333), Inches(0.6),
        font_size=20, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    # 三個關鍵 takeaways
    boxes = [
        ("🎯 確定性", "必須做的事 → Hook", h.COLOR_RED),
        ("🛡 安全", "信任、chmod、路徑", h.COLOR_GREEN),
        ("🧩 組合", "Hook + Skill + Agent", h.COLOR_BLUE),
    ]

    box_w = Inches(3.5)
    box_h = Inches(1.2)
    h_gap = Inches(0.4)
    total_w = box_w * 3 + h_gap * 2
    start_x = (h.SLIDE_W - total_w) / 2
    start_y = Inches(4.8)

    for i, (label, text, color) in enumerate(boxes):
        x = start_x + i * (box_w + h_gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = h.COLOR_WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)

        h.add_text_block(
            slide, label,
            x, start_y + Inches(0.15), box_w, Inches(0.4),
            font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER
        )
        h.add_text_block(
            slide, text,
            x, start_y + Inches(0.6), box_w, Inches(0.4),
            font_size=12, color=h.COLOR_DARK, align=PP_ALIGN.CENTER
        )

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.666), Inches(6.4),
        Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = h.COLOR_RED
    bar.line.fill.background()

    h.add_text_block(
        slide, "Claude Code Plugin 完整學習系列 · #06 · 系列完結 🎉",
        Inches(1), Inches(6.7), Inches(11.333), Inches(0.5),
        font_size=14, bold=True, color=h.COLOR_DARK,
        align=PP_ALIGN.CENTER
    )

    h.add_text_block(
        slide, "整理自 code.claude.com/docs 官方文件",
        Inches(1), Inches(7.2), Inches(11.333), Inches(0.4),
        font_size=11, color=h.COLOR_GRAY_TXT,
        align=PP_ALIGN.CENTER
    )

    # 儲存
    output = "/home/elan/pi-proj/06-hooks.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build()
    print(f"✅ 簡報產生完成：{path}")
    import os
    size = os.path.getsize(path)
    print(f"   檔案大小：{size:,} bytes ({size/1024:.1f} KB)")
